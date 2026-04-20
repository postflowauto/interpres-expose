import os, io, json, base64, zipfile, requests, re, uuid, shutil
from copy import deepcopy
from flask import Flask, request, jsonify, send_file
from pptx import Presentation
from lxml import etree

app = Flask(__name__)
app.config['MAX_CONTENT_LENGTH'] = 100 * 1024 * 1024  # 100 MB pro Chunk (Render Proxy Limit)

# Chunk-Upload Verzeichnis
CHUNK_DIR = "/tmp/interpres_chunks"
os.makedirs(CHUNK_DIR, exist_ok=True)

from flask import make_response

@app.errorhandler(413)
def request_too_large(e):
    return jsonify({"error": "Datei zu groß (max. 500 MB)"}), 413

@app.after_request
def add_cors(response):
    response.headers["Access-Control-Allow-Origin"] = "*"
    response.headers["Access-Control-Allow-Headers"] = "Content-Type, X-API-Token"
    response.headers["Access-Control-Allow-Methods"] = "POST, GET, OPTIONS"
    return response

@app.route("/generate-expose", methods=["OPTIONS"])
@app.route("/fill-pptx", methods=["OPTIONS"])
@app.route("/health", methods=["OPTIONS"])
def options():
    return make_response("", 204)

API_TOKEN = os.environ.get("API_TOKEN", "interpres-secret-2026")
CLAUDE_API_KEY = os.environ.get("CLAUDE_API_KEY", "")
CLOUDCONVERT_KEY = os.environ.get("CLOUDCONVERT_KEY", "")
UNSPLASH_ACCESS_KEY = os.environ.get("UNSPLASH_ACCESS_KEY", "")
TEST_MODE = os.environ.get("TEST_MODE", "false").lower() == "true"
TEMPLATE_URL = "https://raw.githubusercontent.com/postflowauto/interpres-expose/main/urbanunits_Marketing_Expose_v3.pdf-7.pptx"

# Dummy-Daten für TEST_MODE (kein Claude-API-Call)
DUMMY_PROJEKTDATEN = {
    "projektname_roh": "Testprojekt Hannover", "adresse": "Lindener Marktplatz 5",
    "stadt": "Hannover", "stadtteil": "Linden", "plz": "30449",
    "bautraeger": "Urban Units GmbH", "anzahl_haeuser": "2",
    "we_pro_haus": "24", "anzahl_we_gesamt": "48", "kfw_standard": "KfW 55 EE",
    "energieversorgung": "Fernwärme + Photovoltaik", "stellplaetze": "24",
    "groesse_von": "28", "groesse_bis": "67", "kaufpreis_ab": "189.000",
    "besonderheiten": "Möbliert, Smart-Lock, Dachterrasse", "planungsphase": "Baugenehmigung erteilt"
}

DUMMY_EXPOSE_DATA = {
    "projekt_name": "Stadtquartier Linden", "projekt_titel": "Leben im Herzen Lindens",
    "entwickler_name": "Urban Units GmbH", "entwickler_name_gross": "URBAN UNITS GMBH",
    "stadt": "Hannover", "stadt_gross": "HANNOVER", "stadtteil": "Linden",
    "adresse_lang": "Lindener Marktplatz 5, 30449 Hannover", "plz": "30449",
    "quartier_name": "Linden-Mitte", "quartier_history": "Lebendiges Gründerzeitviertel mit urbanem Flair",
    "quartier_ref": "Hannover Linden", "stadt_bezeichnung": "Landeshauptstadt",
    "anzahl_we": "48", "anzahl_1zi": "12", "anzahl_2zi": "24", "anzahl_barrierefrei": "12",
    "groesse_von": "28", "groesse_bis": "67", "zimmer_typen": "1-Zimmer und 2-Zimmer",
    "produkt_beschreibung": "Vollmöblierte Mikro-Apartments mit Smart-Lock",
    "kaufpreis_ab": "189.000", "kfw_darlehen": "150.000", "stellplaetze": "24",
    "kfw_standard": "KfW 55 EE", "energieversorgung": "Fernwärme + Photovoltaik",
    "besonderheiten": "Möbliert, Smart-Lock, Dachterrasse, E-Bike-Sharing",
    "steuerliche_moeglichkeiten": "Dreifach AfA - 5% degressiv §7 Abs.5a EStG + 5% Sonder-AfA §7b EStG + 10% Möbel-AfA",
    "prospekt_datum": "April 2026",
    "text_kapitel_invest": "INVEST", "text_kapitel_live": "LIVE",
    "text_kapitel_stay": "STAY", "text_kapitel_know": "KNOW",
    "text_intro": "Hannover wächst – und Linden ist mittendrin.",
    "text_investment_pitch": "Solide Rendite in einem der dynamischsten Stadtteile Hannovers.",
    "text_hotel": "Möbliert, flexibel, sofort vermietbar – ideal für Kurzzeitvermietung.",
    "text_projekt_nachhaltig_1": "KfW 55 EE – höchster Förderstandard.",
    "text_projekt_nachhaltig_2": "Photovoltaik deckt 30% des Allgemeinstrombedarfs.",
    "text_greenliving_intro": "Nachhaltig wohnen in Hannover.",
    "text_greenliving_1": "Fernwärme aus regenerativen Quellen.",
    "text_greenliving_2": "E-Bike-Sharing für alle Bewohner.",
    "text_ausstattung_intro": "Hochwertig. Vollständig. Bezugsfertig.",
    "text_ausstattung_detail": "Designermöbel, Echtholzparkett, moderne Einbauküche.",
    "text_ausstattung_kurz": "Alles inklusive.", "text_ausstattung_lang": "Vom Bett bis zur Kaffeemaschine.",
    "text_grundriss_intro": "Clever geplante Grundrisse für maximale Nutzfläche.",
    "text_architektur": "Zeitloser Klinkerbau trifft moderne Glaselemente.",
    "text_nachhaltig_1": "KfW 55 EE", "text_nachhaltig_2": "Fernwärme",
    "text_nachhaltig_3": "Photovoltaik", "text_nachhaltig_4": "E-Mobilität",
    "text_standort_1": "Zentral in Hannover-Linden.", "text_standort_2": "Alles in Laufnähe.",
    "stadt_einwohner": "535.932", "stadt_bip": "38.500", "stadt_mietsteigerung": "+3,2%",
    "stadt_studierende": "48.000", "bundesland_bip": "310 Mrd. EUR",
    "text_einwohner_detail": "Hannover wächst kontinuierlich.",
    "text_bip_detail": "Niedersachsen – starke Industrie und Dienstleistungen.",
    "text_mietsteigerung_detail": "Stabile Mietsteigerungen über dem Bundesschnitt.",
    "text_studierende_detail": "Universitätsstadt mit hoher Nachfrage.",
    "text_stadt_wachstum_1": "Bevölkerungswachstum seit 2015 konstant.",
    "text_stadt_wachstum_2": "Zuzug aus Ballungsräumen verstärkt Nachfrage.",
    "text_stadt_intro": "Hannover – Niedersachsens Wirtschaftsmotor.",
    "text_stadt_wirtschaft_links": "Messe, Continental, TUI – globale Player vor Ort.",
    "text_stadt_wirtschaft_rechts": "Starker Mittelstand und wachsende Startup-Szene.",
    "stadt_invest_titel": "Investitionsstandort Hannover",
    "stadt_invest_label": "Rendite", "text_stadt_invest_detail": "Attraktive Nettomietrenditen von 4–5%.",
    "stadt_stat_1_zahl": "535.932", "stadt_stat_1_label": "Einwohner",
    "stadt_stat_2_zahl": "48.000", "stadt_stat_2_label": "Studierende",
    "stadt_stat_3_zahl": "+3,2%", "stadt_stat_3_label": "Mietsteigerung p.a.",
    "stadt_branche_titel": "Leitbranchen", "text_stadt_branche_1": "Messe & Kongress",
    "text_stadt_branche_2": "Automobil & Logistik",
    "quelle_1": "Statistik Hannover 2024", "quelle_2": "IHK Hannover 2024",
    "quelle_3": "Wohnmarktreport 2024", "quelle_4": "Bundesagentur für Arbeit 2024",
    "freizeit_1_name": "Maschsee", "freizeit_2_name": "Eilenriede",
    "freizeit_3_name": "Kröpcke", "freizeit_4_name": "Herrenhäuser Gärten",
    "min_freizeit_1": "8", "min_freizeit_2": "12", "min_freizeit_3": "15", "min_freizeit_4": "20",
    "min_uni": "18", "label_min_uni": "Leibniz Universität",
    "min_bahnhof": "12", "label_min_bahnhof": "Hannover Hbf",
    "min_altstadt": "14", "label_min_altstadt": "Altstadt",
    "feature_1_zahl": "48", "feature_1_label": "Wohneinheiten",
    "feature_2_zahl": "100", "feature_2_label": "Prozent möbliert",
    "feature_3_zahl": "24", "feature_3_label": "Stunden Zugang per Smart-Lock-System",
    "amenity_1": "E-Bike-Sharing", "amenity_2": "Solar-Carport",
    "amenity_3": "Fitnessstudio", "amenity_4": "Paketstation",
    "amenity_5": "Café im EG", "amenity_6": "Dachgarten",
    "amenity_7": "Fernwärme", "amenity_8": "Tiefgarage", "amenity_9": "Balkon",
    "grundriss_1_label": "Typ A – 28 m²", "grundriss_2_label": "Typ B – 42 m²",
    "grundriss_3_label": "Typ C – 55 m²", "grundriss_4_label": "Typ D – 67 m²",
    "ansicht_1_label": "Westfassade", "ansicht_2_label": "Südfassade",
    "we_bereich_1": "Wohnen & Schlafen", "we_bereich_2": "Bad & Küche",
    "we_beispiel_1": "Typ A", "we_beispiel_2": "Typ B",
    "we_typ_beschreibung": "Kompakte Grundrisse, maximale Funktionalität.",
    "we_flaeche_1": "28", "we_flaeche_2": "35", "we_flaeche_3": "42",
    "we_flaeche_4": "55", "we_flaeche_5": "67",
    "logo_initial": "S",
    # Kapitel-Seiten Subtext
    "text_kapitel_invest_1": "Nachhaltig investieren in Hannover.",
    "text_kapitel_invest_2": "Maximale Förderung, stabile Rendite.",
    "text_kapitel_live_1": "Die Stadt. Der Standort. Das Quartier.",
    "text_kapitel_live_2": "Hannover – Wirtschaftsmotor Niedersachsens.",
    "text_kapitel_stay_1": "Vollmöbliert. Nachhaltig. Bezugsfertig.",
    "text_kapitel_stay_2": "Design trifft Funktion in Hannover-Linden.",
    "text_kapitel_know_1": "Transparenz und Rechtssicherheit.",
    "text_kapitel_know_2": "Alle Fakten auf einen Blick.",
    # Stadtstatistik Details
    "text_stadt_stat_1_detail": "Hannover wächst kontinuierlich.",
    "text_stadt_stat_2_detail": "Universitätsstadt mit hoher Nachfrage.",
    "text_stadt_stat_3_detail": "Stabile Mietsteigerungen über dem Bundesschnitt.",
    # Grundriss-Bilder
    "bild_grundriss_1": "", "bild_grundriss_2": "", "bild_grundriss_3": "", "bild_grundriss_4": "",
    # Bundesland
    "bundesland": "Niedersachsen",
    "bild_titel": "", "bild_quartier": "",
    "bild_projekt_aussen": "", "bild_amenity_1": "", "bild_amenity_2": "",
    "bild_amenity_3": "", "bild_amenity_4": "", "bild_amenity_5": "",
    "bild_amenity_6": "", "bild_amenity_7": "", "bild_amenity_8": "",
    "bild_amenity_9": "", "bild_greenliving_1": "", "bild_greenliving_2": "",
    "bild_interior": "", "bild_ausstattung_1": "", "bild_ausstattung_2": "",
    "bild_ausstattung_3": "", "bild_ausstattung_4": "", "bild_ausstattung_5": "",
    "bild_ausstattung_6": "", "bild_grundriss_intro_1": "", "bild_grundriss_intro_2": "",
    "bild_ansicht_1": "", "bild_ansicht_2": "", "bild_we_1": "", "bild_we_2": "",
    "bild_stadt_presse": "", "bild_stadt_branche": "",
    "bild_rechtlich_1": "", "bild_rechtlich_2": "",
    "bild_collage_1": "", "bild_collage_2": "", "bild_collage_3": "",
    "bild_collage_4": "", "bild_collage_5": "",
    "bild_standort_innen": "", "bild_standort_aussen": "",
    "bild_hotel_1": "", "bild_hotel_2": "",
    "bild_stadt_gross": "", "bild_stadt_klein": "",
    "bild_lageplan": "", "bild_grundriss_intro_3": "",
    "bild_projekt": "",
    "produkt_beschreibung": "Vollmöblierte Mikro-Apartments mit Smart-Lock",
    "zitat_intro": "Wohnen neu gedacht.",
}

# Relevante PDF-Typen nach Priorität
PDF_PRIORITY = [
    (1, ["zusammenfassung", "summary"]),
    (1, ["berechnung-bri", "bri-berechnung"]),
    (2, ["grundriss", "floor", "lageplan"]),
    (2, ["wfl-berechnung", "wohnflaeche", "wfl_berechnung"]),
    (3, ["schnitt", "ansicht", "elevation"]),
]

UNSPLASH_QUERIES = {
    "BILD_TITEL": "modern luxury residential building exterior",
    "BILD_QUARTIER": "city neighborhood street urban architecture",
    "BILD_PROJEKT_AUSSEN": "modern apartment building exterior architecture",
    "BILD_AMENITY_1": "car sharing electric vehicle urban",
    "BILD_AMENITY_2": "solar panels rooftop renewable energy",
    "BILD_AMENITY_3": "gym fitness weights modern",
    "BILD_AMENITY_4": "parcel station locker delivery",
    "BILD_AMENITY_5": "cafe coffee interior modern",
    "BILD_AMENITY_6": "green roof garden urban",
    "BILD_AMENITY_7": "district heating pipes infrastructure",
    "BILD_AMENITY_8": "vintage classic car showroom",
    "BILD_AMENITY_9": "apartment balcony urban view",
    "BILD_GREENLIVING_1": "sustainable green building nature",
    "BILD_GREENLIVING_2": "modern residential building facade",
    "BILD_INTERIOR": "modern bedroom interior minimal design",
    "BILD_AUSSTATTUNG_1": "modern living room interior design",
    "BILD_AUSSTATTUNG_2": "hardwood floor interior",
    "BILD_AUSSTATTUNG_3": "modern bathroom tiles",
    "BILD_AUSSTATTUNG_4": "modern kitchen interior",
    "BILD_AUSSTATTUNG_5": "bedroom furniture design",
    "BILD_AUSSTATTUNG_6": "apartment interior detail",
    "BILD_GRUNDRISS_INTRO_1": "modern apartment living room",
    "BILD_GRUNDRISS_INTRO_2": "modern apartment bedroom",
    "BILD_ANSICHT_1": "apartment building exterior west",
    "BILD_ANSICHT_2": "modern residential building south",
    "BILD_WE_1": "apartment floor plan interior",
    "BILD_WE_2": "modern apartment bedroom interior",
    "BILD_STADT_PRESSE": "newspaper article table coffee",
    "BILD_STADT_BRANCHE": "scientist laboratory research modern",
    "BILD_RECHTLICH_1": "modern residential building exterior",
    "BILD_RECHTLICH_2": "apartment building facade evening",
    "BILD_COLLAGE_1": "modern apartment interior living room",
    "BILD_COLLAGE_2": "food lifestyle dinner table modern",
    "BILD_COLLAGE_3": "rooftop terrace modern apartment",
    "BILD_COLLAGE_4": "modern kitchen interior design",
    "BILD_COLLAGE_5": "apartment building exterior architecture",
    "BILD_STANDORT_INNEN": "modern bedroom interior minimal",
    "BILD_STANDORT_AUSSEN": "residential building exterior street",
    "BILD_HOTEL_1": "hotel bedroom luxury modern",
    "BILD_HOTEL_2": "hotel lobby modern interior",
    "BILD_STADT_GROSS": "city skyline aerial",
    "BILD_STADT_KLEIN": "city street urban",
    "BILD_LAGEPLAN": "city map urban district aerial overview",
    "BILD_GRUNDRISS_INTRO_3": "modern apartment interior living space",
    "BILD_PROJEKT": "modern luxury apartment building exterior night",
    "BILD_GRUNDRISS_1": "apartment floor plan architectural drawing",
    "BILD_GRUNDRISS_2": "apartment floor plan 2 room layout",
    "BILD_GRUNDRISS_3": "apartment floor plan 3 room layout",
    "BILD_GRUNDRISS_4": "apartment floor plan large layout",
}

PLATZHALTER = {
    "projekt_name": "", "projekt_titel": "", "entwickler_name": "", "entwickler_name_gross": "",
    "stadt": "", "stadt_gross": "", "stadtteil": "", "adresse_lang": "", "plz": "",
    "quartier_name": "", "quartier_history": "", "quartier_ref": "", "stadt_bezeichnung": "",
    "anzahl_we": "", "anzahl_1zi": "", "anzahl_2zi": "", "anzahl_barrierefrei": "",
    "groesse_von": "", "groesse_bis": "", "zimmer_typen": "", "produkt_beschreibung": "",
    "kaufpreis_ab": "", "kfw_darlehen": "150.000", "stellplaetze": "", "kfw_standard": "",
    "energieversorgung": "", "besonderheiten": "",
    "steuerliche_moeglichkeiten": "Dreifach AfA - 5% degressiv §7 Abs.5a EStG + 5% Sonder-AfA §7b EStG + 10% Möbel-AfA",
    "prospekt_datum": "", "text_kapitel_invest": "", "text_kapitel_live": "",
    "text_kapitel_stay": "", "text_kapitel_know": "", "text_intro": "",
    "text_investment_pitch": "", "text_hotel": "", "text_projekt_nachhaltig_1": "",
    "text_projekt_nachhaltig_2": "", "text_greenliving_intro": "", "text_greenliving_1": "",
    "text_greenliving_2": "", "text_ausstattung_intro": "", "text_ausstattung_detail": "",
    "text_ausstattung_kurz": "", "text_ausstattung_lang": "", "text_grundriss_intro": "",
    "text_architektur": "", "text_nachhaltig_1": "", "text_nachhaltig_2": "",
    "text_nachhaltig_3": "", "text_nachhaltig_4": "", "text_standort_1": "", "text_standort_2": "",
    "stadt_einwohner": "", "stadt_bip": "", "stadt_mietsteigerung": "", "stadt_studierende": "",
    "bundesland_bip": "", "text_einwohner_detail": "", "text_bip_detail": "",
    "text_mietsteigerung_detail": "", "text_studierende_detail": "",
    "text_stadt_wachstum_1": "", "text_stadt_wachstum_2": "", "text_stadt_intro": "",
    "text_stadt_wirtschaft_links": "", "text_stadt_wirtschaft_rechts": "",
    "stadt_invest_titel": "", "stadt_invest_label": "", "text_stadt_invest_detail": "",
    "stadt_stat_1_zahl": "", "stadt_stat_1_label": "", "stadt_stat_2_zahl": "",
    "stadt_stat_2_label": "", "stadt_stat_3_zahl": "", "stadt_stat_3_label": "",
    "stadt_branche_titel": "", "text_stadt_branche_1": "", "text_stadt_branche_2": "",
    "quelle_1": "", "quelle_2": "", "quelle_3": "", "quelle_4": "",
    "freizeit_1_name": "", "freizeit_2_name": "", "freizeit_3_name": "", "freizeit_4_name": "",
    "min_freizeit_1": "", "min_freizeit_2": "", "min_freizeit_3": "", "min_freizeit_4": "",
    "min_uni": "", "label_min_uni": "", "min_bahnhof": "", "label_min_bahnhof": "",
    "min_altstadt": "", "label_min_altstadt": "",
    "feature_1_zahl": "", "feature_1_label": "",
    "feature_2_zahl": "100", "feature_2_label": "Prozent möbliert",
    "feature_3_zahl": "24", "feature_3_label": "Stunden Zugang per Smart-Lock-System",
    "amenity_1": "", "amenity_2": "", "amenity_3": "", "amenity_4": "", "amenity_5": "",
    "amenity_6": "", "amenity_7": "", "amenity_8": "", "amenity_9": "",
    "grundriss_1_label": "", "grundriss_2_label": "", "grundriss_3_label": "", "grundriss_4_label": "",
    "ansicht_1_label": "", "ansicht_2_label": "",
    "we_bereich_1": "", "we_bereich_2": "", "we_beispiel_1": "", "we_beispiel_2": "",
    "we_typ_beschreibung": "", "we_flaeche_1": "", "we_flaeche_2": "",
    "we_flaeche_3": "", "we_flaeche_4": "", "we_flaeche_5": "",
    "logo_initial": "",
    # Kapitel-Seiten Subtext (2 Zeilen pro Kapitel)
    "text_kapitel_invest_1": "", "text_kapitel_invest_2": "",
    "text_kapitel_live_1": "", "text_kapitel_live_2": "",
    "text_kapitel_stay_1": "", "text_kapitel_stay_2": "",
    "text_kapitel_know_1": "", "text_kapitel_know_2": "",
    # Stadtstatistik Details
    "text_stadt_stat_1_detail": "", "text_stadt_stat_2_detail": "", "text_stadt_stat_3_detail": "",
    # Grundriss-Bilder (direkte Slots, nicht intro)
    "bild_grundriss_1": "", "bild_grundriss_2": "", "bild_grundriss_3": "", "bild_grundriss_4": "",
    # Bundesland (ohne BIP)
    "bundesland": "",
    "bild_titel": "", "bild_quartier": "",
    "bild_projekt_aussen": "", "bild_amenity_1": "", "bild_amenity_2": "", "bild_amenity_3": "",
    "bild_amenity_4": "", "bild_amenity_5": "", "bild_amenity_6": "", "bild_amenity_7": "",
    "bild_amenity_8": "", "bild_amenity_9": "", "bild_greenliving_1": "", "bild_greenliving_2": "",
    "bild_interior": "", "bild_ausstattung_1": "", "bild_ausstattung_2": "", "bild_ausstattung_3": "",
    "bild_ausstattung_4": "", "bild_ausstattung_5": "", "bild_ausstattung_6": "",
    "bild_grundriss_intro_1": "", "bild_grundriss_intro_2": "",
    "bild_ansicht_1": "", "bild_ansicht_2": "", "bild_we_1": "", "bild_we_2": "",
    "bild_stadt_presse": "", "bild_stadt_branche": "",
    "bild_rechtlich_1": "", "bild_rechtlich_2": "",
    "bild_collage_1": "", "bild_collage_2": "", "bild_collage_3": "",
    "bild_collage_4": "", "bild_collage_5": "",
    "bild_standort_innen": "", "bild_standort_aussen": "",
    "bild_hotel_1": "", "bild_hotel_2": "",
    "bild_stadt_gross": "", "bild_stadt_klein": "",
    "bild_lageplan": "", "bild_grundriss_intro_3": "",
    "bild_projekt": "",
    "zitat_intro": "",
}

def generate_logo_initial(projekt_name):
    """Nimmt den ersten markanten Buchstaben des Projektnamens als Logo-Initial."""
    if not projekt_name:
        return "P"
    skip_words = {"das", "der", "die", "ein", "eine", "am", "im", "an", "auf", "the", "a", "an"}
    words = re.split(r'[\s\-_\.]+', projekt_name)
    for word in words:
        cleaned = re.sub(r'[^a-zA-Z\u00c4\u00d6\u00dc\u00e4\u00f6\u00fc]', '', word)
        if cleaned and cleaned.lower() not in skip_words:
            return cleaned[0].upper()
    return projekt_name[0].upper()

def get_pdf_priority(filename):
    name_lower = filename.lower()
    for priority, keywords in PDF_PRIORITY:
        if any(kw in name_lower for kw in keywords):
            return priority
    return 99

def extract_pdfs_from_zip(zip_bytes):
    """
    Extrahiert PDFs aus ZIP inkl. verschachtelter Ordner (Haus A/, Haus B/ etc.)
    Gibt max 20 relevanteste PDFs zurück.
    """
    all_pdfs = []

    try:
        with zipfile.ZipFile(io.BytesIO(zip_bytes)) as zf:
            for name in zf.namelist():
                if not name.lower().endswith('.pdf'):
                    continue
                if '__MACOSX' in name or name.startswith('.'):
                    continue

                data = zf.read(name)
                if len(data) < 1000:
                    continue

                parts = name.split('/')
                folder = parts[-2] if len(parts) > 1 else "root"
                filename = parts[-1]
                priority = get_pdf_priority(filename)

                all_pdfs.append({
                    "name": filename,
                    "folder": folder,
                    "priority": priority,
                    "base64": base64.b64encode(data).decode(),
                })

    except Exception as e:
        print(f"ZIP Fehler: {e}")

    # Sortieren nach Priorität
    all_pdfs.sort(key=lambda x: (x["priority"], x["folder"]))

    # Auswahl: max 2 Prio-1 PDFs pro Ordner, max 1 Prio-2 pro Ordner, gesamt max 20
    selected = []
    folder_count = {}

    for pdf in all_pdfs:
        if len(selected) >= 20:
            break
        folder = pdf["folder"]
        prio = pdf["priority"]

        if prio == 99:
            continue

        key = f"{folder}_{prio}"
        folder_count[key] = folder_count.get(key, 0) + 1

        limit = 2 if prio == 1 else 1
        if folder_count[key] <= limit:
            selected.append(pdf)

    print(f"PDFs gesamt: {len(all_pdfs)}, ausgewählt: {len(selected)}")
    for p in selected:
        print(f"  [Prio {p['priority']}] {p['folder']} / {p['name']}")

    return selected

def fetch_unsplash_image(query):
    if not UNSPLASH_ACCESS_KEY:
        return ""
    try:
        resp = requests.get(
            "https://api.unsplash.com/photos/random",
            params={"query": query, "orientation": "landscape"},
            headers={"Authorization": f"Client-ID {UNSPLASH_ACCESS_KEY}"},
            timeout=10
        )
        if resp.status_code == 200:
            return resp.json()["urls"]["regular"]
    except Exception as e:
        print(f"Unsplash Fehler für '{query}': {e}")
    return ""

def fill_image_placeholders(data):
    stadt = data.get("stadt", "")
    queries = UNSPLASH_QUERIES.copy()
    if stadt:
        queries["BILD_TITEL"] = f"modern residential building {stadt} urban architecture"
        queries["BILD_QUARTIER"] = f"{stadt} city neighborhood street"
        queries["BILD_PROJEKT_AUSSEN"] = f"modern apartment building {stadt} exterior"
        queries["BILD_GREENLIVING_1"] = f"sustainable green building {stadt}"
        queries["BILD_GREENLIVING_2"] = f"modern residential {stadt} facade"
        queries["BILD_STADT_GROSS"] = f"city skyline aerial {stadt}"
        queries["BILD_STADT_KLEIN"] = f"city street urban {stadt}"
        queries["BILD_LAGEPLAN"] = f"city map district aerial {stadt}"
    for placeholder_key, query in queries.items():
        data_key = placeholder_key.lower()
        if data_key in data:
            url = fetch_unsplash_image(query)
            if url:
                data[data_key] = url
    return data

def analyze_pdfs_with_claude(pdfs):
    content = []
    for pdf in pdfs:
        content.append({
            "type": "document",
            "source": {"type": "base64", "media_type": "application/pdf", "data": pdf["base64"]},
            "title": f"{pdf['folder']} / {pdf['name']}"
        })
    content.append({
        "type": "text",
        "text": (
            "Analysiere diese Immobilien-Dokumente aus dem Projektdatenraum und extrahiere alle relevanten Projektdaten. "
            "Es kann sein dass du Dokumente von mehreren Häusern (Haus A, B, C...) siehst - das ist ein Gesamtprojekt. "
            "Antworte NUR mit einem JSON-Objekt mit diesen Feldern: "
            "projektname_roh, adresse, stadt, stadtteil, plz, bautraeger, anzahl_haeuser, "
            "we_pro_haus, anzahl_we_gesamt, kfw_standard, energieversorgung, stellplaetze, "
            "groesse_von, groesse_bis, kaufpreis_ab, besonderheiten, planungsphase. "
            "Kein Text davor oder danach, keine Markdown-Backticks."
        )
    })
    resp = requests.post(
        "https://api.anthropic.com/v1/messages",
        headers={"x-api-key": CLAUDE_API_KEY, "anthropic-version": "2023-06-01", "content-type": "application/json"},
        json={"model": "claude-haiku-4-5-20251001", "max_tokens": 4000,
              "messages": [{"role": "user", "content": content}]},
        timeout=120
    )
    resp.raise_for_status()
    text = resp.json()["content"][0]["text"].replace("```json", "").replace("```", "").strip()
    return json.loads(text)

def generate_expose_with_claude(projektdaten):
    stadt = projektdaten.get('stadt', 'der Stadt')
    prompt = (
        "Du bist ein Immobilien-Exposé-Spezialist bei INTERPRÉS GmbH. "
        "Antworte NUR mit einem validen JSON-Objekt. Kein Text davor oder danach. Keine Markdown-Backticks.\n\n"
        f"## PROJEKTDATEN\n{json.dumps(projektdaten, ensure_ascii=False)}\n\n"
        f"## AUFGABE\nFülle alle Felder aus. Nutze dein Wissen über {stadt} für Statistiken "
        "(Einwohner, BIP, Mietsteigerung, Studierende, Top-Arbeitgeber, Freizeiteinrichtungen).\n\n"
        f"## ALLE FELDER AUSFÜLLEN:\n{json.dumps(PLATZHALTER, ensure_ascii=False)}"
    )
    resp = requests.post(
        "https://api.anthropic.com/v1/messages",
        headers={"x-api-key": CLAUDE_API_KEY, "anthropic-version": "2023-06-01", "content-type": "application/json"},
        json={
            "model": "claude-sonnet-4-6", "max_tokens": 8000,
            "messages": [{"role": "user", "content": prompt}]
        },
        timeout=240
    )
    resp.raise_for_status()
    json_text = ""
    for block in resp.json()["content"]:
        if block.get("type") == "text":
            json_text = block["text"]
    json_text = json_text.replace("```json", "").replace("```", "").strip()
    if not json_text:
        raise ValueError("Claude hat keinen Text zurückgegeben. Stop-Reason: " +
                         str(resp.json().get("stop_reason")))
    return json.loads(json_text)

def _replace_placeholders(text, data):
    """Replace {{KEY}}, {{KEY|suffix}}, {{KEY | suffix}} in text using data dict.
    Case-insensitive key matching. Returns replaced string."""
    repl_map = {k.upper(): str(v or "") for k, v in data.items()}
    pattern = re.compile(r'\{\{\s*([A-Z0-9_]+)\s*(?:\|[^}]*)?\}\}', re.IGNORECASE)
    def _sub(m):
        return repl_map.get(m.group(1).upper().strip(), m.group(0))
    return pattern.sub(_sub, text)


def duplicate_slide(prs, slide_index):
    """Duplicates the slide at slide_index and inserts the copy at slide_index + 1."""
    template = prs.slides[slide_index]
    new_slide = prs.slides.add_slide(template.slide_layout)

    # Replace shape tree with a deep copy of the template's
    sp_tree = new_slide.shapes._spTree
    tmpl_sp_tree = template.shapes._spTree

    # Remove all children added by add_slide (keep only nvGrpSpPr + grpSpPr = first 2)
    for child in list(sp_tree)[2:]:
        sp_tree.remove(child)

    # Copy shapes from template (skip first 2 as well, copy the rest)
    for child in list(tmpl_sp_tree)[2:]:
        sp_tree.append(deepcopy(child))

    # Move new slide (currently last) to position slide_index + 1
    sldIdLst = prs.slides._sldIdLst
    moved_el = sldIdLst[-1]
    sldIdLst.remove(moved_el)
    sldIdLst.insert(slide_index + 1, moved_el)

    return prs.slides[slide_index + 1]


def duplicate_we_slides(prs, data):
    """
    Finds the WE template slide (contains WE_BEISPIEL_1 or BILD_WE_1),
    duplicates it for every additional Wohnungstyp found in data,
    and replaces _1 → _N and letter 'a' → 'b'/'c'/... in each duplicate.
    Called AFTER text replacement so the originals are already filled.
    """
    letters = ['a', 'b', 'c', 'd', 'e', 'f']

    # Count how many WE types are present
    we_count = 1
    for n in range(2, 7):
        if data.get(f"we_beispiel_{n}") or data.get(f"bild_we_{n}"):
            we_count = n

    if we_count <= 1:
        print("duplicate_we_slides: nur 1 WE-Typ, kein Duplizieren")
        return

    # Find WE template slide
    we_idx = None
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame:
                txt = shape.text_frame.text.upper()
                if "WE_BEISPIEL_1" in txt or "BILD_WE_1" in txt:
                    we_idx = i
                    break
        if we_idx is not None:
            break

    if we_idx is None:
        print("duplicate_we_slides: WE-Template-Slide nicht gefunden")
        return

    print(f"WE-Slide bei Index {we_idx}, {we_count - 1} Duplikate")

    for offset in range(1, we_count):
        new_slide = duplicate_slide(prs, we_idx + (offset - 1))
        n = offset + 1
        sp_tree = new_slide.shapes._spTree
        xml_str = etree.tostring(sp_tree, encoding="unicode")

        # Replace _1 field references with _N
        xml_str = xml_str.replace("WE_BEISPIEL_1", f"WE_BEISPIEL_{n}")
        xml_str = xml_str.replace("WE_BEREICH_1",  f"WE_BEREICH_{n}")
        xml_str = xml_str.replace("BILD_WE_1",     f"BILD_WE_{n}")
        xml_str = xml_str.replace("WE_FLAECHE_1",  f"WE_FLAECHE_{n}")
        # Replace stand-alone letter label 'a' → 'b'/'c'/... in XML text nodes
        old_letter = letters[0]        # 'a'
        new_letter = letters[offset]   # 'b', 'c', …
        xml_str = xml_str.replace(f">{old_letter}<", f">{new_letter}<")

        new_sp_tree = etree.fromstring(xml_str.encode("utf-8"))
        for child in list(sp_tree):
            sp_tree.remove(child)
        for child in list(new_sp_tree):
            sp_tree.append(child)

        # Now fill the duplicate's placeholders with actual data
        for shape in list(new_slide.shapes):
            try:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if not para.runs:
                            continue
                        full_text = "".join(r.text for r in para.runs)
                        modified = _replace_placeholders(full_text, data)
                        if modified != full_text:
                            para.runs[0].text = modified
                            for run in para.runs[1:]:
                                run.text = ""
            except Exception as e:
                print(f"WE-Duplikat Shape-Fehler: {e}")


def fill_pptx(template_bytes, data):
    """Fill PPTX template using python-pptx: embeds images and replaces text placeholders."""

    prs = Presentation(io.BytesIO(template_bytes))

    # Pre-download all images referenced by bild_* keys
    image_data = {}
    for key, value in data.items():
        if key.startswith("bild_") and isinstance(value, str) and value.startswith("http"):
            try:
                resp = requests.get(value, timeout=15)
                if resp.status_code == 200:
                    image_data[key] = resp.content
                    print(f"Bild geladen: {key} ({len(resp.content)//1024} KB)")
                else:
                    print(f"Bild HTTP-Fehler {key}: {resp.status_code}")
            except Exception as e:
                print(f"Bild Download-Fehler {key}: {e}")

    def make_replacement_map(data):
        """Build a case-insensitive lookup: UPPER_KEY -> value."""
        return {k.upper(): str(v or "") for k, v in data.items()}

    REPL_MAP = make_replacement_map(data)

    # Regex that matches {{KEY}}, {{KEY|suffix}}, {{KEY | suffix}} (spaces ok around pipe)
    PLACEHOLDER_RE = re.compile(r'\{\{\s*([A-Z0-9_]+)\s*(?:\|[^}]*)?\}\}', re.IGNORECASE)

    def replace_text(text):
        """Replace all placeholders in a string using REPL_MAP."""
        def _sub(m):
            key = m.group(1).upper().strip()
            return REPL_MAP.get(key, m.group(0))  # keep original if not found
        return PLACEHOLDER_RE.sub(_sub, text)

    def replace_in_paragraph(para):
        """Replace placeholders in a paragraph, handling split-run placeholders.

        Strategy: reassemble all runs into one string, replace, write back into
        runs[0] preserving its formatting, clear all other runs.
        """
        if not para.runs:
            return
        full_text = "".join(r.text for r in para.runs)
        if "{{" not in full_text:
            return
        modified = replace_text(full_text)
        if modified != full_text:
            para.runs[0].text = modified
            for run in para.runs[1:]:
                run.text = ""

    def replace_in_textframe(tf):
        """Replace placeholders across entire text frame, including cross-paragraph splits.

        Some placeholders span two paragraphs (e.g. '{{PRODUKT_BESCHREI-' / 'BUNG}}').
        We join the full text frame, detect these, and replace them in para[0] of the span.
        """
        # First do normal per-paragraph replacement
        for para in tf.paragraphs:
            replace_in_paragraph(para)

        # Then handle cross-paragraph splits: join all paragraph texts and check
        # Build list of (para_index, full_run_text) pairs
        para_texts = ["".join(r.text for r in p.runs) for p in tf.paragraphs]
        joined = "\n".join(para_texts)

        # Find any remaining {{...}} that survived (i.e. were split across paragraphs)
        # by scanning pairs of adjacent paragraphs
        for i in range(len(tf.paragraphs) - 1):
            combined = para_texts[i] + para_texts[i + 1]
            if "{{" in combined and "}}" in combined:
                replaced = replace_text(combined)
                if replaced != combined:
                    # Split the replacement back at the original boundary
                    split_at = len(para_texts[i])
                    new_p0 = replaced[:split_at] if len(replaced) >= split_at else replaced
                    new_p1 = replaced[split_at:] if len(replaced) >= split_at else ""
                    # Write to para i
                    p0 = tf.paragraphs[i]
                    if p0.runs:
                        p0.runs[0].text = new_p0
                        for r in p0.runs[1:]:
                            r.text = ""
                    # Write to para i+1
                    p1 = tf.paragraphs[i + 1]
                    if p1.runs:
                        p1.runs[0].text = new_p1
                        for r in p1.runs[1:]:
                            r.text = ""
                    # Update cache
                    para_texts[i] = new_p0
                    para_texts[i + 1] = new_p1

    def get_abs_coords(group_shape, child_shape):
        """Berechnet absolute Slide-Koordinaten eines Child-Shapes in einer Gruppe.
        Nutzt grpSpPr/xfrm aus dem p:-Namespace (nicht a:-Namespace!).
        """
        NS_A = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        NS_P = 'http://schemas.openxmlformats.org/presentationml/2006/main'
        grp_el  = group_shape._element
        grpSpPr = grp_el.find(f'{{{NS_P}}}grpSpPr')
        if grpSpPr is None:
            # fallback: simple addition
            return (
                (group_shape.left or 0) + (child_shape.left or 0),
                (group_shape.top  or 0) + (child_shape.top  or 0),
                child_shape.width  or 0,
                child_shape.height or 0,
            )
        xfrm  = grpSpPr.find(f'{{{NS_A}}}xfrm')
        if xfrm is None:
            return (
                (group_shape.left or 0) + (child_shape.left or 0),
                (group_shape.top  or 0) + (child_shape.top  or 0),
                child_shape.width  or 0,
                child_shape.height or 0,
            )
        off   = xfrm.find(f'{{{NS_A}}}off')
        ext   = xfrm.find(f'{{{NS_A}}}ext')
        chOff = xfrm.find(f'{{{NS_A}}}chOff')
        chExt = xfrm.find(f'{{{NS_A}}}chExt')
        if None in (off, ext, chOff, chExt):
            return (
                (group_shape.left or 0) + (child_shape.left or 0),
                (group_shape.top  or 0) + (child_shape.top  or 0),
                child_shape.width  or 0,
                child_shape.height or 0,
            )
        grp_x  = int(off.get('x',  0));  grp_y  = int(off.get('y',  0))
        grp_w  = int(ext.get('cx', 1));  grp_h  = int(ext.get('cy', 1))
        ch_x   = int(chOff.get('x', 0)); ch_y   = int(chOff.get('y', 0))
        ch_w   = int(chExt.get('cx', 1));ch_h   = int(chExt.get('cy', 1))
        scale_x = grp_w / ch_w if ch_w else 1
        scale_y = grp_h / ch_h if ch_h else 1
        abs_left = int(grp_x + ((child_shape.left or 0) - ch_x) * scale_x)
        abs_top  = int(grp_y + ((child_shape.top  or 0) - ch_y) * scale_y)
        abs_w    = int((child_shape.width  or 0) * scale_x)
        abs_h    = int((child_shape.height or 0) * scale_y)
        return abs_left, abs_top, abs_w, abs_h

    def process_shape(slide, shape, image_data):
        """Ersetzt Text oder bettet Bild ein. Gruppen werden korrekt mit absoluten Coords behandelt."""
        # Gruppe: Kinder einzeln verarbeiten
        if shape.shape_type == 6:
            for child in list(shape.shapes):
                try:
                    if child.has_text_frame:
                        txt = child.text_frame.text.strip()
                        m = PLACEHOLDER_RE.search(txt)  # search statt match — findet auch mit Whitespace davor
                        if m and len(txt) < 60:  # kurzer Text → reiner Placeholder
                            key = m.group(1).lower()
                            if key in image_data and image_data[key]:
                                # Bild mit korrekten absoluten Koordinaten einfügen
                                abs_left, abs_top, abs_w, abs_h = get_abs_coords(shape, child)
                                slide.shapes.add_picture(
                                    io.BytesIO(image_data[key]),
                                    abs_left, abs_top, abs_w, abs_h
                                )
                                # Placeholder-Text leeren
                                for para in child.text_frame.paragraphs:
                                    for r in para.runs:
                                        r.text = ""
                                print(f"  Bild eingesetzt: {key} @ {abs_left//914400*2.54:.1f},{abs_top//914400*2.54:.1f}cm {abs_w//914400*2.54:.1f}x{abs_h//914400*2.54:.1f}cm")
                                continue
                    # Kein Bild → Text ersetzen
                    if child.has_text_frame:
                        replace_in_textframe(child.text_frame)
                except Exception as e:
                    print(f"  Gruppen-Child Fehler {child.name}: {e}")
            # Auch Text direkt auf der Gruppe ersetzen (selten)
            if shape.has_text_frame:
                replace_in_textframe(shape.text_frame)
            return

        # Top-Level Non-Group Shape
        shape_name_lower = (shape.name or "").lower()

        # 1. Bild per Shape-Name
        if shape_name_lower in image_data and image_data[shape_name_lower]:
            try:
                left, top, width, height = shape.left, shape.top, shape.width, shape.height
                shape._element.getparent().remove(shape._element)
                slide.shapes.add_picture(
                    io.BytesIO(image_data[shape_name_lower]), left, top, width, height
                )
                return
            except Exception as e:
                print(f"Bild-Ersatz Fehler (name) {shape_name_lower}: {e}")

        # 2. Bild per Text-Inhalt
        if shape.has_text_frame:
            txt = shape.text_frame.text.strip()
            m = PLACEHOLDER_RE.match(txt)
            if m:
                key = m.group(1).lower()
                if key in image_data and image_data[key]:
                    try:
                        left, top, width, height = shape.left, shape.top, shape.width, shape.height
                        shape._element.getparent().remove(shape._element)
                        slide.shapes.add_picture(
                            io.BytesIO(image_data[key]), left, top, width, height
                        )
                        return
                    except Exception as e:
                        print(f"Bild-Ersatz Fehler (text) {key}: {e}")

        # 3. Text ersetzen
        if shape.has_text_frame:
            replace_in_textframe(shape.text_frame)

    for slide in prs.slides:
        for shape in list(slide.shapes):
            try:
                process_shape(slide, shape, image_data)
            except Exception as e:
                print(f"Shape-Fehler slide={slide.slide_id} shape={shape.name}: {e}")

    # Duplicate WE slides after all text/image replacement is done
    duplicate_we_slides(prs, data)

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()

def convert_to_pdf(pptx_bytes, filename):
    import time
    cc_headers = {"Authorization": f"Bearer {CLOUDCONVERT_KEY}", "Content-Type": "application/json"}

    # Job erstellen (async API)
    job_resp = requests.post(
        "https://api.cloudconvert.com/v2/jobs",
        headers=cc_headers,
        json={"tasks": {
            "upload":  {"operation": "import/upload"},
            "convert": {"operation": "convert", "input": "upload",
                        "input_format": "pptx", "output_format": "pdf", "engine": "office"},
            "export":  {"operation": "export/url", "input": "convert"}
        }}, timeout=30
    )
    job_resp.raise_for_status()
    job = job_resp.json()["data"]
    job_id = job["id"]

    # Datei hochladen
    upload_task = next(t for t in job["tasks"] if t["name"] == "upload")
    form = upload_task["result"]["form"]
    files = {"file": (filename, pptx_bytes,
                      "application/vnd.openxmlformats-officedocument.presentationml.presentation")}
    requests.post(form["url"], data=form.get("parameters", {}), files=files, timeout=60).raise_for_status()

    # Warten bis Job fertig (max 5 Minuten, alle 5s pollen)
    for _ in range(60):
        time.sleep(5)
        status_resp = requests.get(
            f"https://api.cloudconvert.com/v2/jobs/{job_id}",
            headers=cc_headers, timeout=30
        )
        status_resp.raise_for_status()
        job_status = status_resp.json()["data"]["status"]
        if job_status == "finished":
            tasks = status_resp.json()["data"]["tasks"]
            pdf_url = next(t for t in tasks if t["name"] == "export")["result"]["files"][0]["url"]
            return requests.get(pdf_url, timeout=60).content
        if job_status == "error":
            tasks = status_resp.json()["data"]["tasks"]
            err = next((t.get("message","") for t in tasks if t.get("status") == "error"), "Unbekannter Fehler")
            raise RuntimeError(f"CloudConvert Fehler: {err}")

    raise RuntimeError("CloudConvert Timeout nach 5 Minuten")

def assemble_session(session_id):
    """Liest alle Chunks einer Session von /tmp und gibt die assemblierten Bytes zurück."""
    session_dir = os.path.join(CHUNK_DIR, session_id)
    if not os.path.isdir(session_dir):
        raise ValueError(f"Session {session_id} nicht gefunden")
    meta_path = os.path.join(session_dir, "meta.json")
    with open(meta_path) as f:
        meta = json.load(f)
    total = meta["total_chunks"]
    assembled = b""
    for i in range(total):
        chunk_path = os.path.join(session_dir, f"chunk_{i:04d}")
        with open(chunk_path, "rb") as f:
            assembled += f.read()
    shutil.rmtree(session_dir, ignore_errors=True)
    return assembled


@app.route("/upload-chunk", methods=["POST", "OPTIONS"])
def upload_chunk():
    if request.method == "OPTIONS":
        return make_response("", 204)
    if request.headers.get("X-API-Token") != API_TOKEN:
        return jsonify({"error": "Unauthorized"}), 401
    try:
        session_id = request.form.get("session_id") or str(uuid.uuid4())
        chunk_index = int(request.form.get("chunk_index", 0))
        total_chunks = int(request.form.get("total_chunks", 1))
        filename = request.form.get("filename", "upload.zip")

        chunk_file = request.files.get("chunk")
        if not chunk_file:
            return jsonify({"error": "Kein 'chunk' im Request"}), 400

        session_dir = os.path.join(CHUNK_DIR, session_id)
        os.makedirs(session_dir, exist_ok=True)

        chunk_path = os.path.join(session_dir, f"chunk_{chunk_index:04d}")
        chunk_file.save(chunk_path)

        meta_path = os.path.join(session_dir, "meta.json")
        if not os.path.exists(meta_path):
            with open(meta_path, "w") as f:
                json.dump({"total_chunks": total_chunks, "filename": filename}, f)

        received = len([n for n in os.listdir(session_dir) if n.startswith("chunk_")])
        ready = received >= total_chunks
        return jsonify({"session_id": session_id, "chunk": chunk_index,
                        "received": received, "total": total_chunks, "ready": ready})
    except Exception as e:
        import traceback
        return jsonify({"error": str(e), "trace": traceback.format_exc()}), 500


@app.route("/")
def index():
    index_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "index.html")
    with open(index_path, encoding="utf-8") as f:
        return f.read(), 200, {"Content-Type": "text/html; charset=utf-8"}

@app.route("/health")
def health():
    return jsonify({"status": "ok", "service": "INTERPRES Full Pipeline v3",
                    "test_mode": TEST_MODE})

@app.route("/generate-expose", methods=["POST"])
def generate_expose():
    if request.headers.get("X-API-Token") != API_TOKEN:
        return jsonify({"error": "Unauthorized"}), 401
    try:
        pdfs = []

        # --- Chunked-Session Upload ---
        session_ids = request.form.getlist("session_ids")
        if session_ids:
            for sid in session_ids:
                zip_bytes = assemble_session(sid)
                pdfs.extend(extract_pdfs_from_zip(zip_bytes))

        elif request.content_type and "multipart" in request.content_type:
            uploaded = request.files.getlist("files") or request.files.getlist("file")
            if not uploaded:
                return jsonify({"error": "Keine Dateien im Request"}), 400
            for f in uploaded:
                pdfs.extend(extract_pdfs_from_zip(f.read()))
        else:
            body = request.get_json(force=True) or {}
            if "zip_base64_list" in body:
                for b64 in body["zip_base64_list"]:
                    pdfs.extend(extract_pdfs_from_zip(base64.b64decode(b64)))
            elif "zip_base64" in body:
                pdfs.extend(extract_pdfs_from_zip(base64.b64decode(body["zip_base64"])))
            else:
                return jsonify({"error": "zip_base64 oder zip_base64_list fehlt"}), 400

        if not pdfs:
            return jsonify({"error": "Keine relevanten PDFs gefunden"}), 400

        # Max. 3 PDFs senden (Kostenkontrolle)
        pdfs = sorted(pdfs, key=lambda x: x["priority"])[:3]

        if TEST_MODE:
            print("TEST_MODE aktiv – überspringe Claude API")
            expose_data = DUMMY_EXPOSE_DATA.copy()
            expose_data = fill_image_placeholders(expose_data)
        else:
            projektdaten = analyze_pdfs_with_claude(pdfs)
            expose_data = generate_expose_with_claude(projektdaten)
            expose_data["logo_initial"] = generate_logo_initial(expose_data.get("projekt_name", ""))
            expose_data = fill_image_placeholders(expose_data)

        tmpl_bytes = requests.get(TEMPLATE_URL, timeout=30).content
        pptx_bytes = fill_pptx(tmpl_bytes, expose_data)

        projekt_name = expose_data.get("projekt_name", "Expose").replace(" ", "_")
        pdf_bytes = convert_to_pdf(pptx_bytes, f"{projekt_name}.pptx")

        return send_file(io.BytesIO(pdf_bytes), mimetype="application/pdf",
                         as_attachment=True, download_name=f"{projekt_name}_Expose.pdf")

    except Exception as e:
        import traceback
        return jsonify({"error": str(e), "trace": traceback.format_exc()}), 500

if __name__ == "__main__":
    app.run(host="0.0.0.0", port=int(os.environ.get("PORT", 5000)))
