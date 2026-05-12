import os, io, json, base64, zipfile, requests, re, uuid, shutil
from copy import deepcopy
from flask import Flask, request, jsonify, send_file
from pptx import Presentation
from lxml import etree

app = Flask(__name__)
app.config['MAX_CONTENT_LENGTH'] = 100 * 1024 * 1024  # 100 MB pro Chunk (Render Proxy Limit)

# ── V2 (HTML/CSS-basierter Renderer) parallel registrieren ──────────────────
# /v2/* Routen sind unabhängig vom V1-Pipeline und liefern den neuen Look.
try:
    from v2.server import register_v2 as _register_v2
    _register_v2(app)
    print("[v2] HTML-Renderer-Routen unter /v2/* registriert")
except Exception as _v2e:
    print(f"[v2] WARN: V2-Registrierung fehlgeschlagen: {_v2e}")

# Chunk-Upload Verzeichnis
CHUNK_DIR = "/tmp/interpres_chunks"
os.makedirs(CHUNK_DIR, exist_ok=True)

from flask import make_response

@app.errorhandler(413)
def request_too_large(e):
    return jsonify({"error": "Datei zu groß (max. 500 MB)"}), 413

@app.errorhandler(500)
def internal_error(e):
    import traceback
    return jsonify({"error": str(e), "trace": traceback.format_exc()}), 500

@app.errorhandler(Exception)
def unhandled_exception(e):
    import traceback
    app.logger.error("Unhandled exception: %s\n%s", str(e), traceback.format_exc())
    return jsonify({"error": str(e), "trace": traceback.format_exc()}), 500

@app.after_request
def add_cors(response):
    response.headers["Access-Control-Allow-Origin"] = "*"
    response.headers["Access-Control-Allow-Headers"] = "Content-Type, X-API-Token"
    response.headers["Access-Control-Allow-Methods"] = "POST, GET, OPTIONS"
    return response

@app.route("/generate-expose", methods=["OPTIONS"])
@app.route("/fill-pptx", methods=["OPTIONS"])
@app.route("/debug-images", methods=["OPTIONS"])
@app.route("/health", methods=["OPTIONS"])
def options():
    return make_response("", 204)

API_TOKEN = os.environ.get("API_TOKEN", "interpres-secret-2026")
CLAUDE_API_KEY = os.environ.get("CLAUDE_API_KEY", "")
CLOUDCONVERT_KEY = os.environ.get("CLOUDCONVERT_KEY", "")
UNSPLASH_ACCESS_KEY = os.environ.get("UNSPLASH_ACCESS_KEY", "")
TAVILY_API_KEY = os.environ.get("TAVILY_API_KEY", "")
TEST_MODE = os.environ.get("TEST_MODE", "false").lower() == "true"
TEMPLATE_URL = "https://raw.githubusercontent.com/postflowauto/interpres-expose/main/urbanunits_Marketing_Expose_v3.pdf-25.pptx"

# Zusatz-Templates (Kurz-Exposé + Rechtliches). Werden nur im Multi-Template-
# Workflow geladen. Wenn URL leer/404 → der jeweilige Tab im Editor zeigt
# 'Template noch nicht hinterlegt' und blockiert nicht das Marketing-Expose.
KURZ_TEMPLATE_URL = os.environ.get(
    "KURZ_TEMPLATE_URL",
    "https://raw.githubusercontent.com/postflowauto/interpres-expose/main/urbanunits_Kurzexpose-4.pptx",
)
RECHTLICH_TEMPLATE_URL = os.environ.get(
    "RECHTLICH_TEMPLATE_URL",
    "https://raw.githubusercontent.com/postflowauto/interpres-expose/main/urbanunits_Rechtlich_v1.pptx",
)

# Mapping expose_typ → URL. Wird in fill_pptx + V2-Worker genutzt.
TEMPLATE_URLS = {
    "marketing": TEMPLATE_URL,
    "kurz":      KURZ_TEMPLATE_URL,
    "rechtlich": RECHTLICH_TEMPLATE_URL,
}

# Dummy-Daten für TEST_MODE (kein Claude-API-Call)
DUMMY_PROJEKTDATEN = {
    "projektname_roh": "Testprojekt Hannover", "adresse": "Lindener Marktplatz 5",
    "stadt": "Hannover", "stadtteil": "Linden", "plz": "30449",
    "bautraeger": "Urban Units GmbH", "anzahl_haeuser": "2",
    "we_pro_haus": "24", "anzahl_we_gesamt": "48", "kfw_standard": "KfW 55 EE",
    "energieversorgung": "Fernwärme + Photovoltaik", "stellplaetze": "24",
    "groesse_von": "28", "groesse_bis": "67", "kaufpreis_ab": "189.000",
    "besonderheiten": "Möbliert, Smart-Lock, Dachterrasse", "planungsphase": "Baugenehmigung erteilt"
}

DUMMY_EXPOSE_DATA = {
    "projekt_name": "Stadtquartier Linden", "projekt_titel": "Leben im Herzen Lindens",
    "entwickler_name": "Urban Units GmbH", "entwickler_name_gross": "URBAN UNITS GMBH",
    "stadt": "Hannover", "stadt_gross": "HANNOVER", "stadtteil": "Linden",
    "adresse_lang": "Lindener Marktplatz 5, 30449 Hannover", "plz": "30449",
    "quartier_name": "Linden-Mitte", "quartier_history": "Lebendiges Gründerzeitviertel mit urbanem Flair",
    "quartier_ref": "Hannover Linden", "stadt_bezeichnung": "Landeshauptstadt",
    "anzahl_we": "48", "anzahl_1zi": "12", "anzahl_2zi": "24", "anzahl_barrierefrei": "12",
    "groesse_von": "28", "groesse_bis": "67", "zimmer_typen": "1-Zimmer und 2-Zimmer",
    "produkt_beschreibung": "Vollmöblierte Mikro-Apartments mit Smart-Lock",
    "kaufpreis_ab": "189.000", "kfw_darlehen": "150.000", "stellplaetze": "24",
    "kfw_standard": "KfW 55 EE", "energieversorgung": "Fernwärme + Photovoltaik",
    "besonderheiten": "Möbliert, Smart-Lock, Dachterrasse, E-Bike-Sharing",
    "steuerliche_moeglichkeiten": "Dreifach AfA - 5% degressiv §7 Abs.5a EStG + 5% Sonder-AfA §7b EStG + 10% Möbel-AfA",
    "prospekt_datum": "April 2026",
    # Rechtlicher Teil B — Platzhalter aus rechtlich_template.pptx
    "entwickler_adresse":          "Lindener Marktplatz 5, 30449 Hannover",
    "entwickler_strasse":          "Lindener Marktplatz 5",
    "entwickler_plz_ort":          "30449 Hannover",
    "entwickler_handelsregister":  "Amtsgericht Hannover unter HRB 12345",
    "verwalter_name":              "Interpres Hausverwaltung GmbH",
    "verwalter_strasse":           "Lindener Marktplatz 5",
    "verwalter_geschaeftsfuehrer": "Max Mustermann und Anna Beispiel",
    "projekt_adresse":             "Lindener Marktplatz 5",
    "projekt_plz_stadt":           "30449 Hannover",
    "projekt_areal":               "Alten Stadtquartier Linden",
    "projekt_typ":                 "Neubau Mikroapartments",
    "fertig_geplant":              "30.06.2027",
    "fertig_verbindlich":          "31.12.2027",
    "verweisurkunde_datum":        "04.07.2025",
    "stand_steuer":                "Juli 2025",
    "kaufpreis_von":               "189.000,00",
    "kaufpreis_bis":               "325.000,00",
    "baujahr_range":               "2025-2027",
    "notar_name":                  "Peter Krolopp",
    "notar_adresse":               "Humboldtstr. 2, 39112 Magdeburg",
    "notar_stadt":                 "Magdeburg",
    "weg_bezeichnung":             "Stadtquartier Linden",
    "text_kapitel_invest": "INVEST", "text_kapitel_live": "LIVE",
    "text_kapitel_stay": "STAY", "text_kapitel_know": "KNOW",
    "text_kapitel_hotel": "HOTEL",
    "text_intro": "Hannover wächst – und Linden ist mittendrin.",
    "text_investment_pitch": "Solide Rendite in einem der dynamischsten Stadtteile Hannovers.",
    "text_hotel": "Möbliert, flexibel, sofort vermietbar – ideal für Kurzzeitvermietung.",
    "text_projekt_nachhaltig_1": "KfW 55 EE – höchster Förderstandard.",
    "text_projekt_nachhaltig_2": "Photovoltaik deckt 30% des Allgemeinstrombedarfs.",
    "text_greenliving_intro": "Nachhaltig wohnen in Hannover.",
    "text_greenliving_1": "Fernwärme aus regenerativen Quellen.",
    "text_greenliving_2": "E-Bike-Sharing für alle Bewohner.",
    "text_ausstattung_intro": "Hochwertig. Vollständig. Bezugsfertig.",
    "text_ausstattung_detail": "Designermöbel, Echtholzparkett, moderne Einbauküche.",
    "text_ausstattung_kurz": "Alles inklusive.", "text_ausstattung_lang": "Vom Bett bis zur Kaffeemaschine.",
    "text_grundriss_intro": "Clever geplante Grundrisse für maximale Nutzfläche.",
    "text_architektur": "Zeitloser Klinkerbau trifft moderne Glaselemente.",
    "text_nachhaltig_1": "KfW 55 EE", "text_nachhaltig_2": "Fernwärme",
    "text_nachhaltig_3": "Photovoltaik", "text_nachhaltig_4": "E-Mobilität",
    "text_standort_1": "Zentral in Hannover-Linden.", "text_standort_2": "Alles in Laufnähe.",
    "stadt_einwohner": "535.932", "stadt_bip": "38.500", "stadt_mietsteigerung": "+3,2%",
    "stadt_studierende": "48.000", "bundesland_bip": "310 Mrd.",
    "text_einwohner_detail": "Hannover wächst kontinuierlich.",
    "text_bip_detail": "Niedersachsen – starke Industrie und Dienstleistungen.",
    "text_mietsteigerung_detail": "Stabile Mietsteigerungen über dem Bundesschnitt.",
    "text_studierende_detail": "Universitätsstadt mit hoher Nachfrage.",
    "text_stadt_wachstum_1": "Bevölkerungswachstum seit 2015 konstant.",
    "text_stadt_wachstum_2": "Zuzug aus Ballungsräumen verstärkt Nachfrage.",
    "text_stadt_intro": "Hannover – Niedersachsens Wirtschaftsmotor.",
    "text_stadt_wirtschaft_links": "Messe, Continental, TUI – globale Player vor Ort.",
    "text_stadt_wirtschaft_rechts": "Starker Mittelstand und wachsende Startup-Szene.",
    "stadt_invest_titel": "Investitionsstandort Hannover",
    "stadt_invest_label": "Rendite", "text_stadt_invest_detail": "Attraktive Nettomietrenditen von 4–5%.",
    "stadt_stat_1_zahl": "535.932", "stadt_stat_1_label": "Einwohner",
    "stadt_stat_2_zahl": "48.000", "stadt_stat_2_label": "Studierende",
    "stadt_stat_3_zahl": "+3,2%", "stadt_stat_3_label": "Mietsteigerung p.a.",
    "stadt_branche_titel": "Leitbranchen", "text_stadt_branche_1": "Messe & Kongress",
    "text_stadt_branche_2": "Automobil & Logistik",
    "quelle_1": "Statistik Hannover 2024", "quelle_2": "IHK Hannover 2024",
    "quelle_3": "Wohnmarktreport 2024", "quelle_4": "Bundesagentur für Arbeit 2024",
    # ── Standort-Minuten (Slide 5) ────────────────────────────────────────────
    "min_uni": "18", "label_min_uni": "Leibniz Universität",
    "min_bahnhof": "12", "label_min_bahnhof": "Hauptbahnhof",
    "min_altstadt": "15", "label_min_altstadt": "Altstadt",
    # ── Alles ganz nah (Slide 14): 4 Freizeit-Einträge ───────────────────────
    "freizeit_1_name": "Maschsee", "min_freizeit_1": "8",
    "freizeit_2_name": "Eilenriede", "min_freizeit_2": "12",
    "freizeit_3_name": "Innenstadt", "min_freizeit_3": "15",
    "freizeit_4_name": "Herrenhäuser Gärten", "min_freizeit_4": "20",
    # ── WE-Typen: Original-Slide (Typen 1+2 nebeneinander) ───────────────────
    "we_beispiel_1": "WE 02", "we_nummern_1": "WE 01, WE 02, WE 05, WE 06",
    "we_beispiel_2": "WE 07", "we_nummern_2": "WE 07, WE 08, WE 11, WE 12",
    "we_raum_1_name_1": "Wohnen/Kochen", "we_raum_2_name_1": "Schlafen",
    "we_raum_3_name_1": "Bad",           "we_raum_4_name_1": "Abstellraum",
    "we_raum_5_name_1": "Balkon",
    "we_flaeche_1_1": "23,99 m²", "we_flaeche_2_1": "5,36 m²",
    "we_flaeche_3_1": "5,34 m²",  "we_flaeche_4_1": "2,33 m²", "we_flaeche_5_1": "32,02 m²",
    "we_typ_beschreibung_1": "1-Zimmer-Wohnung mit Balkon. Optimal für Studierende und Berufspendler.",
    "we_raum_1_name_2": "Wohnen/Kochen", "we_raum_2_name_2": "Schlafen",
    "we_raum_3_name_2": "Bad",           "we_raum_4_name_2": "Abstellraum",
    "we_raum_5_name_2": "Balkon",
    "we_flaeche_1_2": "28,45 m²", "we_flaeche_2_2": "6,10 m²",
    "we_flaeche_3_2": "5,80 m²",  "we_flaeche_4_2": "3,20 m²", "we_flaeche_5_2": "43,55 m²",
    "we_typ_beschreibung_2": "2-Zimmer-Wohnung mit großem Balkon. Ideal für Paare und Singles.",
    # Duplikat-Slide (Typen 3+4), leer = kein Duplikat
    "we_beispiel_3": "", "we_nummern_3": "",
    "we_beispiel_4": "", "we_nummern_4": "",
    # Duplikat-Slide 2 (Typen 5+6), leer = kein zweites Duplikat
    "we_beispiel_5": "", "we_nummern_5": "",
    "we_beispiel_6": "", "we_nummern_6": "",
    "feature_1_zahl": "48", "feature_1_label": "Wohneinheiten",
    "feature_2_zahl": "100", "feature_2_label": "Prozent möbliert",
    "feature_3_zahl": "24", "feature_3_label": "Stunden Zugang per Smart-Lock-System",
    "amenity_1": "E-Bike-Sharing", "amenity_2": "Solar-Carport",
    "amenity_3": "Fitnessstudio", "amenity_4": "Paketstation",
    "amenity_5": "Café im EG", "amenity_6": "Dachgarten",
    "amenity_7": "Fernwärme", "amenity_8": "Tiefgarage", "amenity_9": "Balkon",
    "grundriss_1_label": "Typ A – 28 m²", "grundriss_2_label": "Typ B – 42 m²",
    "grundriss_3_label": "Typ C – 55 m²", "grundriss_4_label": "Typ D – 67 m²",
    "ansicht_1_label": "Westfassade", "ansicht_2_label": "Südfassade",
    "logo_initial": "S",
    # Kapitel-Seiten Subtext
    "text_kapitel_invest_1": "Nachhaltig investieren in Hannover.",
    "text_kapitel_invest_2": "Maximale Förderung, stabile Rendite.",
    "text_kapitel_live_1": "Die Stadt. Der Standort. Das Quartier.",
    "text_kapitel_live_2": "Hannover – Wirtschaftsmotor Niedersachsens.",
    "text_kapitel_stay_1": "Vollmöbliert. Nachhaltig. Bezugsfertig.",
    "text_kapitel_stay_2": "Design trifft Funktion in Hannover-Linden.",
    "text_kapitel_know_1": "Transparenz und Rechtssicherheit.",
    "text_kapitel_know_2": "Alle Fakten auf einen Blick.",
    # Stadtstatistik Details
    "text_stadt_stat_1_detail": "Hannover wächst kontinuierlich.",
    "text_stadt_stat_2_detail": "Universitätsstadt mit hoher Nachfrage.",
    "text_stadt_stat_3_detail": "Stabile Mietsteigerungen über dem Bundesschnitt.",
    # Grundriss-Bilder
    "bild_grundriss_1": "", "bild_grundriss_2": "", "bild_grundriss_3": "", "bild_grundriss_4": "",
    # Bundesland
    "bundesland": "Niedersachsen",
    "bild_titel": "", "bild_quartier": "",
    "bild_projekt_aussen": "", "bild_amenity_1": "", "bild_amenity_2": "",
    "bild_amenity_3": "", "bild_amenity_4": "", "bild_amenity_5": "",
    "bild_amenity_6": "", "bild_amenity_7": "", "bild_amenity_8": "",
    "bild_amenity_9": "", "bild_greenliving_1": "", "bild_greenliving_2": "",
    "bild_interior": "", "bild_ausstattung_1": "", "bild_ausstattung_2": "",
    "bild_ausstattung_3": "", "bild_ausstattung_4": "", "bild_ausstattung_5": "",
    "bild_ausstattung_6": "", "bild_grundriss_intro_1": "", "bild_grundriss_intro_2": "",
    "bild_ansicht_1": "", "bild_ansicht_2": "",
    **{f"bild_we_{n}": "" for n in range(1, 21)},   # bild_we_1 … bild_we_20 für DUMMY
    "bild_stadt_presse": "", "bild_stadt_branche": "",
    "bild_rechtlich_1": "", "bild_rechtlich_2": "",
    "bild_collage_1": "", "bild_collage_2": "", "bild_collage_3": "",
    "bild_collage_4": "", "bild_collage_5": "",
    "bild_standort_innen": "", "bild_standort_aussen": "",
    "bild_hotel_1": "", "bild_hotel_2": "",
    "bild_stadt_gross": "", "bild_stadt_klein": "",
    "bild_lageplan": "", "bild_grundriss_intro_3": "",
    "bild_projekt": "",
    "produkt_beschreibung": "Vollmöblierte Mikro-Apartments mit Smart-Lock",
    "zitat_intro": "Wohnen neu gedacht.",
    # ── Kurz-Exposé Felder (TEST_MODE) ─────────────────────────────────────
    # Längen orientiert an den Prompt-Limits (max-Zeichen), damit das
    # Layout im TEST_MODE genauso aussieht wie mit echtem Claude-Output.
    "projekt_untertitel": "Designed to stay – urbanes Mikro-Living mit Konzept.",  # ~52 / 60
    "projekt_beschreibung": "In Hannover-Linden entsteht mit dem Stadtquartier Linden ein zukunftsweisendes Wohnprojekt mit 48 vollmöblierten 1- bis 2-Zimmer-Apartments. Geplant nach KfW-55-EE-Standard mit Photovoltaik und Fernwärme, ausgestattet mit Designermöbeln, Smart-Lock und Echtholzparkett. Das lebendige Gründerzeitviertel verbindet kurze Wege zur Leibniz Universität, zur Altstadt und zum Hauptbahnhof mit grünen Rückzugsorten und Quartiers-Flair. Ob als Eigennutzung oder Kapitalanlage – das Projekt überzeugt durch Lage, Effizienz und flexible Wohnungsgrößen.",  # ~540 / 560
    "projekt_beschreibung_kurz": "In Hannover-Linden entsteht mit dem Stadtquartier Linden ein zukunftsweisendes Wohnprojekt mit 48 vollmöblierten 1- bis 2-Zimmer-Apartments. Geplant nach KfW-55-EE-Standard mit Photovoltaik und Fernwärme, ausgestattet mit Designermöbeln, Smart-Lock und Echtholzparkett. Das lebendige Gründerzeitviertel verbindet kurze Wege zur Leibniz Universität, zur Altstadt und zum Hauptbahnhof mit grünen Rückzugsorten und Quartiers-Flair. Ob als Eigennutzung oder Kapitalanlage – das Projekt überzeugt durch Lage, Effizienz und flexible Wohnungsgrößen.",
    "text_relevanz":   "48.000 Studierende sorgen in Hannover für hohe Mikro-Wohnen-Nachfrage.",  # ~75 / 75 (2 Zeilen)
    "text_design":     "Designermöbel, Echtholzparkett und Walk-In-Dusche – bezugsfertig.",  # ~68 / 75
    "text_foerderung": "KfW-55-EE, Sonder-AfA §7b und degressive AfA – dreifacher Steuervorteil.",  # ~74 / 75
    "text_tech":       "Smart-Lock, Glasfaser, Photovoltaik und E-Mobility – moderne Infrastruktur.",  # ~75 / 75
    # Cover-Stichworte (Slide 1 Marketing): 4 kurze Stichworte, je max ~22 Zeichen.
    # User platziert die 4 Platzhalter selbst im Template (DQN-Stil: fette Caps).
    "cover_stichwort_1": "VOLLMÖBLIERT",
    "cover_stichwort_2": "FÖRDERFÄHIG",
    "cover_stichwort_3": "KFW-40 QNG PLUS",
    "cover_stichwort_4": "SMART HOME READY",
    "besonderheiten_liste": (
        ":  alle Apartments mit Balkon oder Dachterrasse\n"
        ":  barrierearm, Aufzug in alle Etagen\n"
        ":  E-Bike-Sharing für alle Bewohner\n"
        ":  Glasfaseranschluss und Smart-Lock-System\n"
        ":  nachhaltige Fernwärme für Heizung\n"
        ":  Eichenparkett, Fußbodenheizung, Designermöbel"
    ),  # 6 Punkte — Spacing wird per <a:spcBef>=1200 (12pt) in fill_pptx gesetzt
    "gesamtwohnflaeche": "2.142,40 m²",
    "zimmer_anzahl_min": "1",
    "zimmer_anzahl_max": "2",
}

# Relevante PDF-Typen nach Priorität
PDF_PRIORITY = [
    # Prio 1 – wichtigste Projektdokumente (bis 2 pro Ordner erlaubt)
    (1, ["zusammenfassung", "summary", "expose", "exposé", "verkauf", "vertrieb",
         "investment", "invest", "mieteinnahmen", "rendite", "broschüre", "brochure",
         "flyer", "projektbeschreibung", "projektinfo"]),
    (1, ["berechnung-bri", "bri-berechnung"]),
    # Prio 2 – Grundrisse & Flächenberechnungen (bis 3 pro Ordner erlaubt)
    (2, ["grundriss", "floor", "lageplan"]),
    (2, ["wfl-berechnung", "wohnflaeche", "wfl_berechnung", "flaeche", "fläche",
         "raumplan", "typ_a", "typ_b", "typ_c", "typ-a", "typ-b", "typ-c"]),
    (3, ["schnitt", "ansicht", "elevation"]),
    # Prio 4 – sonstige Projektdokumente (1 pro Ordner)
    (4, ["baugenehmigung", "genehmigung", "baubeschreibung", "leistungsverzeichnis",
         "ausstattung", "energieausweis"]),
]

UNSPLASH_QUERIES = {
    "BILD_TITEL": "modern luxury residential building exterior",
    "BILD_QUARTIER": "city neighborhood street urban architecture",
    "BILD_PROJEKT_AUSSEN": "modern apartment building exterior architecture",
    "BILD_AMENITY_1": "car sharing electric vehicle urban",
    "BILD_AMENITY_2": "solar panels rooftop renewable energy",
    "BILD_AMENITY_3": "gym fitness weights modern",
    "BILD_AMENITY_4": "parcel station locker delivery",
    "BILD_AMENITY_5": "cafe coffee interior modern",
    "BILD_AMENITY_6": "green roof garden urban",
    "BILD_AMENITY_7": "district heating pipes infrastructure",
    "BILD_AMENITY_8": "vintage classic car showroom",
    "BILD_AMENITY_9": "apartment balcony urban view",
    "BILD_GREENLIVING_1": "sustainable green building nature",
    "BILD_GREENLIVING_2": "modern residential building facade",
    "BILD_INTERIOR": "modern bedroom interior minimal design",
    "BILD_AUSSTATTUNG_1": "modern living room interior design",
    "BILD_AUSSTATTUNG_2": "hardwood floor interior",
    "BILD_AUSSTATTUNG_3": "modern bathroom tiles",
    "BILD_AUSSTATTUNG_4": "modern kitchen interior",
    "BILD_AUSSTATTUNG_5": "bedroom furniture design",
    "BILD_AUSSTATTUNG_6": "apartment interior detail",
    "BILD_GRUNDRISS_INTRO_1": "modern apartment living room",
    "BILD_GRUNDRISS_INTRO_2": "modern apartment bedroom",
    "BILD_ANSICHT_1": "apartment building exterior west",
    "BILD_ANSICHT_2": "modern residential building south",
    "BILD_WE_1":  "modern apartment interior design minimal furnished",
    "BILD_WE_2":  "studio apartment interior furnished modern design",
    "BILD_WE_3":  "cozy apartment living room interior design",
    "BILD_WE_4":  "modern apartment bedroom minimalist design",
    "BILD_WE_5":  "luxury apartment interior penthouse design",
    "BILD_WE_6":  "compact apartment smart interior design",
    "BILD_WE_7":  "modern apartment open plan living design",
    "BILD_WE_8":  "bright apartment interior contemporary",
    "BILD_WE_9":  "minimalist apartment interior Scandinavian",
    "BILD_WE_10": "modern apartment kitchen dining area",
    "BILD_WE_11": "studio loft apartment modern design",
    "BILD_WE_12": "apartment terrace balcony modern",
    "BILD_WE_13": "penthouse apartment interior luxury",
    "BILD_WE_14": "duplex apartment interior design",
    "BILD_WE_15": "modern apartment bathroom interior",
    "BILD_WE_16": "cozy apartment bedroom interior",
    "BILD_WE_17": "open plan apartment living dining",
    "BILD_WE_18": "modern apartment hallway entrance",
    "BILD_WE_19": "contemporary apartment interior style",
    "BILD_WE_20": "bright apartment modern interior design",
    "BILD_STADT_PRESSE": "newspaper article table coffee",
    "BILD_STADT_BRANCHE": "scientist laboratory research modern",
    "BILD_RECHTLICH_1": "modern residential building exterior",
    "BILD_RECHTLICH_2": "apartment building facade evening",
    "BILD_COLLAGE_1": "modern apartment interior living room",
    "BILD_COLLAGE_2": "food lifestyle dinner table modern",
    "BILD_COLLAGE_3": "rooftop terrace modern apartment",
    "BILD_COLLAGE_4": "modern kitchen interior design",
    "BILD_COLLAGE_5": "apartment building exterior architecture",
    "BILD_STANDORT_INNEN": "modern bedroom interior minimal",
    "BILD_STANDORT_AUSSEN": "residential building exterior street",
    "BILD_HOTEL_1": "hotel bedroom luxury modern",
    "BILD_HOTEL_2": "hotel lobby modern interior",
    "BILD_STADT_GROSS": "city skyline aerial",
    "BILD_STADT_KLEIN": "city street urban",
    "BILD_LAGEPLAN": "city map urban district aerial overview",
    "BILD_GRUNDRISS_INTRO_3": "modern apartment interior living space",
    "BILD_PROJEKT": "modern luxury apartment building exterior night",
    "BILD_GRUNDRISS_1": "apartment floor plan architectural drawing",
    "BILD_GRUNDRISS_2": "apartment floor plan 2 room layout",
    "BILD_GRUNDRISS_3": "apartment floor plan 3 room layout",
    "BILD_GRUNDRISS_4": "apartment floor plan large layout",
}

PLATZHALTER = {
    "projekt_name": "", "projekt_titel": "", "entwickler_name": "", "entwickler_name_gross": "",
    "stadt": "", "stadt_gross": "", "stadtteil": "", "adresse_lang": "", "plz": "",
    "quartier_name": "", "quartier_history": "", "quartier_ref": "", "stadt_bezeichnung": "",
    "anzahl_we": "", "anzahl_1zi": "", "anzahl_2zi": "", "anzahl_barrierefrei": "",
    "groesse_von": "", "groesse_bis": "", "zimmer_typen": "", "produkt_beschreibung": "",
    "kaufpreis_ab": "", "kfw_darlehen": "150.000", "stellplaetze": "", "kfw_standard": "",
    "energieversorgung": "", "besonderheiten": "",
    "steuerliche_moeglichkeiten": "Dreifach AfA - 5% degressiv §7 Abs.5a EStG + 5% Sonder-AfA §7b EStG + 10% Möbel-AfA",
    "prospekt_datum": "", "text_kapitel_invest": "", "text_kapitel_live": "",
    "text_kapitel_stay": "", "text_kapitel_know": "", "text_kapitel_hotel": "", "text_intro": "",
    "text_investment_pitch": "", "text_hotel": "", "text_projekt_nachhaltig_1": "",
    "text_projekt_nachhaltig_2": "", "text_greenliving_intro": "", "text_greenliving_1": "",
    "text_greenliving_2": "", "text_ausstattung_intro": "", "text_ausstattung_detail": "",
    "text_ausstattung_kurz": "", "text_ausstattung_lang": "", "text_grundriss_intro": "",
    "text_architektur": "", "text_nachhaltig_1": "", "text_nachhaltig_2": "",
    "text_nachhaltig_3": "", "text_nachhaltig_4": "", "text_standort_1": "", "text_standort_2": "",
    "stadt_einwohner": "", "stadt_bip": "", "stadt_mietsteigerung": "", "stadt_studierende": "",
    "bundesland_bip": "", "text_einwohner_detail": "", "text_bip_detail": "",
    "text_mietsteigerung_detail": "", "text_studierende_detail": "",
    "text_stadt_wachstum_1": "", "text_stadt_wachstum_2": "", "text_stadt_intro": "",
    "text_stadt_wirtschaft_links": "", "text_stadt_wirtschaft_rechts": "",
    "stadt_invest_titel": "", "stadt_invest_label": "", "text_stadt_invest_detail": "",
    "stadt_stat_1_zahl": "", "stadt_stat_1_label": "", "stadt_stat_2_zahl": "",
    "stadt_stat_2_label": "", "stadt_stat_3_zahl": "", "stadt_stat_3_label": "",
    "stadt_branche_titel": "", "text_stadt_branche_1": "", "text_stadt_branche_2": "",
    # ── 6 gute Gründe (Slide 9) — Titel + Text pro Punkt ─────────────────
    "text_grund_1_titel": "", "text_grund_1_text": "",
    "text_grund_2_titel": "", "text_grund_2_text": "",
    "text_grund_3_titel": "", "text_grund_3_text": "",
    "text_grund_4_titel": "", "text_grund_4_text": "",
    "text_grund_5_titel": "", "text_grund_5_text": "",
    "text_grund_6_titel": "", "text_grund_6_text": "",
    # ── Kurz-Exposé Felder (separates Template, eigenes PDF) ─────────────
    "projekt_untertitel": "",       # Tagline z.B. 'Urbaner Wohnkomfort in nachhaltigem Umfeld'
    "projekt_beschreibung": "",     # Cover-Untertitel-Block (links unten Kurz-S1)
    "projekt_beschreibung_kurz": "",  # grosser Pitch-Text auf Kurz-S2
    "text_relevanz":   "",  # max 75 Zeichen — 2-Zeiler USP Marktrelevanz
    "text_design":     "",  # max 75 Zeichen — 2-Zeiler USP Design/Wohnen
    "text_foerderung": "",  # max 75 Zeichen — 2-Zeiler USP Foerderung/Steuer
    "text_tech":       "",  # max 75 Zeichen — 2-Zeiler USP Technik/Smart-Home
    "besonderheiten_liste": "",  # max 6 Bullets im DQN-Stil
    # Cover-Stichworte (Slide 1 Marketing) — User platziert 4 Platzhalter im Template
    "cover_stichwort_1": "",
    "cover_stichwort_2": "",
    "cover_stichwort_3": "",
    "cover_stichwort_4": "",
    # ── Rechtlicher Teil B (eigenes PPTX-Template) ───────────────────────
    "entwickler_adresse":          "",
    "entwickler_strasse":          "",
    "entwickler_plz_ort":          "",
    "entwickler_handelsregister":  "",
    "verwalter_name":              "",
    "verwalter_strasse":           "",
    "verwalter_geschaeftsfuehrer": "",
    "projekt_adresse":             "",
    "projekt_plz_stadt":           "",
    "projekt_areal":               "",
    "projekt_typ":                 "",
    "fertig_geplant":              "",
    "fertig_verbindlich":          "",
    "verweisurkunde_datum":        "",
    "stand_steuer":                "",
    "kaufpreis_von":               "",
    "kaufpreis_bis":               "",
    "baujahr_range":               "",
    "notar_name":                  "",
    "notar_adresse":               "",
    "notar_stadt":                 "",
    "weg_bezeichnung":             "",
    # "Auf einen Blick" — rechte Spalte Kurz-S2
    "gesamtwohnflaeche": "",  # z.B. '3.741,58 m²' — Summe aller WE
    "zimmer_anzahl_min": "",  # '1' (kleinster Wohnungstyp)
    "zimmer_anzahl_max": "",  # '2' (groesster Wohnungstyp)
    # Bild-Slots Kurz-Cover (6) + Kurz-S2 (4) — werden vom Kunden hochgeladen
    **{f"bild_titel_{n}": "" for n in range(1, 7)},
    **{f"bild_kurz_{n}": ""  for n in range(1, 5)},
    "quelle_1": "", "quelle_2": "", "quelle_3": "", "quelle_4": "",
    # ── Standort-Minuten (Slide 5) ────────────────────────────────────────────
    "min_uni": "", "label_min_uni": "",
    "min_bahnhof": "", "label_min_bahnhof": "",
    "min_altstadt": "", "label_min_altstadt": "",
    # ── Alles ganz nah (Slide 14): Freizeit (4) ─────────────────────────────
    "freizeit_1_name": "", "min_freizeit_1": "",
    "freizeit_2_name": "", "min_freizeit_2": "",
    "freizeit_3_name": "", "min_freizeit_3": "",
    "freizeit_4_name": "", "min_freizeit_4": "",
    # ── Alles ganz nah (Slide 14): Einkaufen (4) – via Overpass ─────────────
    "einkaufen_1_name": "Bäckerei",   "min_einkaufen_1": "2",
    "einkaufen_2_name": "Supermarkt", "min_einkaufen_2": "2",
    "einkaufen_3_name": "Drogerie",   "min_einkaufen_3": "3",
    "einkaufen_4_name": "REWE",       "min_einkaufen_4": "4",
    # ── Alles ganz nah (Slide 14): Ärzte (4) – via Overpass ─────────────────
    "arzt_1_name": "Hausarzt",    "min_arzt_1": "5",
    "arzt_2_name": "Facharzt",    "min_arzt_2": "8",
    "arzt_3_name": "Apotheke",    "min_arzt_3": "3",
    "arzt_4_name": "Krankenhaus", "min_arzt_4": "12",
    # ── Alles ganz nah (Slide 14): Sport (4) – via Overpass ─────────────────
    "sport_1_name": "Fitnessstudio", "min_sport_1": "8",
    "sport_2_name": "Schwimmbad",    "min_sport_2": "10",
    "sport_3_name": "Sportanlage",   "min_sport_3": "6",
    "sport_4_name": "Sportpark",     "min_sport_4": "5",
    # ── Alles ganz nah (Slide 14): Bildung (4) – via Overpass ───────────────
    "bildung_1_name": "Kita",        "min_bildung_1": "5",
    "bildung_2_name": "Grundschule", "min_bildung_2": "8",
    "bildung_3_name": "Gymnasium",   "min_bildung_3": "10",
    "bildung_4_name": "Universität", "min_bildung_4": "15",
    # ── WE-Typen (v19-Schema: 1 Typ pro Slide, 2 Beispiel-WEs pro Slide) ─────
    # Geteilte Felder pro Typ (LINKE und RECHTE Spalte zeigen denselben Typ):
    "we_flaeche_1": "", "we_flaeche_2": "", "we_flaeche_3": "",
    "we_flaeche_4": "", "we_flaeche_5": "",
    "we_typ_beschreibung": "",
    # Per-Spalte Felder (LINKES und RECHTES Beispiel-WE):
    "we_beispiel_1": "", "we_bereich_1": "",
    "we_beispiel_2": "", "we_bereich_2": "",
    # Backwards-Compat (v20-Schema, falls altes Template gewählt wird):
    "we_nummern_1": "", "we_nummern_2": "",
    "we_raum_1_name_1": "", "we_raum_2_name_1": "", "we_raum_3_name_1": "",
    "we_raum_4_name_1": "", "we_raum_5_name_1": "",
    "we_flaeche_1_1": "", "we_flaeche_2_1": "", "we_flaeche_3_1": "",
    "we_flaeche_4_1": "", "we_flaeche_5_1": "",
    "we_typ_beschreibung_1": "",
    "we_raum_1_name_2": "", "we_raum_2_name_2": "", "we_raum_3_name_2": "",
    "we_raum_4_name_2": "", "we_raum_5_name_2": "",
    "we_flaeche_1_2": "", "we_flaeche_2_2": "", "we_flaeche_3_2": "",
    "we_flaeche_4_2": "", "we_flaeche_5_2": "",
    "we_typ_beschreibung_2": "",
    # v19: Beispiel-WEs für Typ 2..8 (jeder Typ = 1 Slide mit 2 Beispielen)
    **{f"we_flaeche_{n}_typ{t}": "" for t in range(2, 9) for n in range(1, 6)},
    **{f"we_typ_beschreibung_typ{t}": "" for t in range(2, 9)},
    "we_bereich_3": "", "we_bereich_4": "",
    "we_bereich_5": "", "we_bereich_6": "",
    "we_bereich_7": "", "we_bereich_8": "",
    # v20-Backwards
    "we_beispiel_3": "", "we_nummern_3": "",
    "we_beispiel_4": "", "we_nummern_4": "",
    "we_raum_1_name_3": "", "we_raum_2_name_3": "", "we_raum_3_name_3": "",
    "we_raum_4_name_3": "", "we_raum_5_name_3": "",
    "we_flaeche_1_3": "", "we_flaeche_2_3": "", "we_flaeche_3_3": "",
    "we_flaeche_4_3": "", "we_flaeche_5_3": "",
    "we_typ_beschreibung_3": "",
    "we_raum_1_name_4": "", "we_raum_2_name_4": "", "we_raum_3_name_4": "",
    "we_raum_4_name_4": "", "we_raum_5_name_4": "",
    "we_flaeche_1_4": "", "we_flaeche_2_4": "", "we_flaeche_3_4": "",
    "we_flaeche_4_4": "", "we_flaeche_5_4": "",
    "we_typ_beschreibung_4": "",
    "we_beispiel_5": "", "we_nummern_5": "",
    "we_beispiel_6": "", "we_nummern_6": "",
    "we_raum_1_name_5": "", "we_raum_2_name_5": "", "we_raum_3_name_5": "",
    "we_raum_4_name_5": "", "we_raum_5_name_5": "",
    "we_flaeche_1_5": "", "we_flaeche_2_5": "", "we_flaeche_3_5": "",
    "we_flaeche_4_5": "", "we_flaeche_5_5": "",
    "we_typ_beschreibung_5": "",
    "we_raum_1_name_6": "", "we_raum_2_name_6": "", "we_raum_3_name_6": "",
    "we_raum_4_name_6": "", "we_raum_5_name_6": "",
    "we_flaeche_1_6": "", "we_flaeche_2_6": "", "we_flaeche_3_6": "",
    "we_flaeche_4_6": "", "we_flaeche_5_6": "",
    "we_typ_beschreibung_6": "",
    "we_beispiel_7": "", "we_beispiel_8": "",
    "feature_1_zahl": "", "feature_1_label": "",
    "feature_2_zahl": "100", "feature_2_label": "Prozent möbliert",
    "feature_3_zahl": "24", "feature_3_label": "Stunden Zugang per Smart-Lock-System",
    "amenity_1": "", "amenity_2": "", "amenity_3": "", "amenity_4": "", "amenity_5": "",
    "amenity_6": "", "amenity_7": "", "amenity_8": "", "amenity_9": "",
    "grundriss_1_label": "", "grundriss_2_label": "", "grundriss_3_label": "", "grundriss_4_label": "",
    "ansicht_1_label": "", "ansicht_2_label": "",
    "logo_initial": "",
    # Kapitel-Seiten Subtext (2 Zeilen pro Kapitel)
    "text_kapitel_invest_1": "", "text_kapitel_invest_2": "",
    "text_kapitel_live_1": "", "text_kapitel_live_2": "",
    "text_kapitel_stay_1": "", "text_kapitel_stay_2": "",
    "text_kapitel_know_1": "", "text_kapitel_know_2": "",
    # Stadtstatistik Details
    "text_stadt_stat_1_detail": "", "text_stadt_stat_2_detail": "", "text_stadt_stat_3_detail": "",
    # Grundriss-Bilder (direkte Slots, nicht intro)
    "bild_grundriss_1": "", "bild_grundriss_2": "", "bild_grundriss_3": "", "bild_grundriss_4": "",
    # Bundesland (ohne BIP)
    "bundesland": "",
    "bild_titel": "", "bild_quartier": "",
    "bild_projekt_aussen": "", "bild_amenity_1": "", "bild_amenity_2": "", "bild_amenity_3": "",
    "bild_amenity_4": "", "bild_amenity_5": "", "bild_amenity_6": "", "bild_amenity_7": "",
    "bild_amenity_8": "", "bild_amenity_9": "", "bild_greenliving_1": "", "bild_greenliving_2": "",
    "bild_interior": "", "bild_ausstattung_1": "", "bild_ausstattung_2": "", "bild_ausstattung_3": "",
    "bild_ausstattung_4": "", "bild_ausstattung_5": "", "bild_ausstattung_6": "",
    "bild_grundriss_intro_1": "", "bild_grundriss_intro_2": "",
    "bild_ansicht_1": "", "bild_ansicht_2": "",
    **{f"bild_we_{n}": "" for n in range(1, 21)},   # bild_we_1 … bild_we_20
    "bild_stadt_presse": "", "bild_stadt_branche": "",
    "bild_rechtlich_1": "", "bild_rechtlich_2": "",
    "bild_collage_1": "", "bild_collage_2": "", "bild_collage_3": "",
    "bild_collage_4": "", "bild_collage_5": "",
    "bild_standort_innen": "", "bild_standort_aussen": "",
    "bild_hotel_1": "", "bild_hotel_2": "",
    "bild_stadt_gross": "", "bild_stadt_klein": "",
    "bild_lageplan": "", "bild_grundriss_intro_3": "",
    "bild_projekt": "",
    "zitat_intro": "",
}

def generate_logo_initial(projekt_name):
    """Nimmt den ersten markanten Buchstaben des Projektnamens als Logo-Initial."""
    if not projekt_name:
        return "P"
    skip_words = {"das", "der", "die", "ein", "eine", "am", "im", "an", "auf", "the", "a", "an"}
    words = re.split(r'[\s\-_\.]+', projekt_name)
    for word in words:
        cleaned = re.sub(r'[^a-zA-Z\u00c4\u00d6\u00dc\u00e4\u00f6\u00fc]', '', word)
        if cleaned and cleaned.lower() not in skip_words:
            return cleaned[0].upper()
    return projekt_name[0].upper()

def get_pdf_priority(filename):
    name_lower = filename.lower()
    for priority, keywords in PDF_PRIORITY:
        if any(kw in name_lower for kw in keywords):
            return priority
    return 99

def extract_pdfs_from_zip(zip_path, work_dir):
    """
    Extrahiert PDFs aus ZIP. Speichert auf Disk unter work_dir/pdfs/,
    behält nur Pfade in der Liste (NICHT bytes/base64) → spart RAM.
    base64 wird on-demand erzeugt (siehe analyze_pdfs_with_claude).
    Gibt max 20 relevanteste PDFs zurück.
    """
    all_pdfs = []
    pdfs_dir = os.path.join(work_dir, "pdfs")
    os.makedirs(pdfs_dir, exist_ok=True)
    seen_paths = set()

    try:
        with zipfile.ZipFile(zip_path) as zf:
            for name in zf.namelist():
                if not name.lower().endswith('.pdf'):
                    continue
                if '__MACOSX' in name or name.startswith('.'):
                    continue
                info = zf.getinfo(name)
                if info.file_size < 1000:
                    continue

                parts = name.split('/')
                folder = parts[-2] if len(parts) > 1 else "root"
                filename = parts[-1]
                priority = get_pdf_priority(filename)

                # Stream auf Disk, nicht in RAM
                disk_name = f"{len(all_pdfs)}_{filename}".replace('/', '_')[:120]
                disk_path = os.path.join(pdfs_dir, disk_name)
                if disk_path in seen_paths:
                    disk_path = os.path.join(pdfs_dir, f"{len(all_pdfs)}b_{filename}")
                seen_paths.add(disk_path)
                with zf.open(name) as src, open(disk_path, "wb") as dst:
                    while True:
                        chunk = src.read(64 * 1024)
                        if not chunk:
                            break
                        dst.write(chunk)

                all_pdfs.append({
                    "name": filename,
                    "folder": folder,
                    "priority": priority,
                    "path": disk_path,
                    "size": info.file_size,
                })

    except Exception as e:
        print(f"ZIP Fehler: {e}")

    # Sortieren nach Priorität
    all_pdfs.sort(key=lambda x: (x["priority"], x["folder"]))

    # Auswahl: Prio-1 max 2/Ordner, Prio-2 max 3/Ordner (WFL-Typen!), Prio-3/4 max 1/Ordner
    # Gesamt max 20
    _PRIO_LIMITS = {1: 2, 2: 3, 3: 1, 4: 1}
    selected = []
    folder_count = {}

    for pdf in all_pdfs:
        if len(selected) >= 20:
            break
        folder = pdf["folder"]
        prio = pdf["priority"]

        if prio == 99:
            continue

        key = f"{folder}_{prio}"
        folder_count[key] = folder_count.get(key, 0) + 1

        limit = _PRIO_LIMITS.get(prio, 1)
        if folder_count[key] <= limit:
            selected.append(pdf)

    print(f"PDFs gesamt: {len(all_pdfs)}, ausgewählt: {len(selected)}")
    for p in selected:
        print(f"  [Prio {p['priority']}] {p['folder']} / {p['name']}")

    return selected

def _extract_images_from_pdf_bytes(pdf_bytes, pdf_name, seen_hashes):
    """Extrahiert eingebettete Bilder aus PDF via PyMuPDF. Dedupliziert via Hash."""
    import hashlib
    images = []
    try:
        import fitz  # PyMuPDF
        doc = fitz.open(stream=pdf_bytes, filetype="pdf")
        n_pages = min(len(doc), 3)  # max 3 Seiten pro PDF → weniger RAM
        for page_num in range(n_pages):
            for img_idx, img in enumerate(doc[page_num].get_images(full=True)):
                xref = img[0]
                try:
                    base_image = doc.extract_image(xref)
                    raw = base_image["image"]
                    ext = "." + base_image["ext"]
                    if ext not in {'.jpg', '.jpeg', '.png', '.webp', '.bmp'}:
                        continue
                    if len(raw) < 20000:  # < 20 KB = Logo/Icon/Decoration
                        continue
                    # Deduplizieren (Logo erscheint auf jeder Seite!)
                    h = hashlib.md5(raw).hexdigest()
                    if h in seen_hashes:
                        continue
                    seen_hashes.add(h)
                    images.append({
                        'name': f"{pdf_name}_s{page_num+1}_b{img_idx+1}{ext}",
                        'bytes': raw,
                        'ext': ext,
                        'size': len(raw),
                    })
                except Exception:
                    continue
        doc.close()
    except ImportError:
        print("  PyMuPDF nicht verfügbar – keine PDF-Bildextraktion")
    except Exception as e:
        print(f"  PDF-Bildextraktion Fehler ({pdf_name}): {e}")
    return images


def extract_images_from_zip(zip_path, work_dir):
    """
    Extrahiert Bilder aus ZIP:
    1. Direkte Bilddateien (.jpg/.png/.webp)
    2. In PDFs eingebettete Bilder (via PyMuPDF, dedupliziert)

    Speichert ALLE Bilder auf Disk (work_dir/imgs/), behält nur Pfade in der
    Liste – kein bytes-Field! Bytes werden on-demand geladen wenn gebraucht.
    Gibt max. 15 Bilder zurück (sortiert nach Größe).
    """
    import hashlib
    IMAGE_EXTS = {'.jpg', '.jpeg', '.png', '.webp'}
    imgs_dir = os.path.join(work_dir, "imgs")
    os.makedirs(imgs_dir, exist_ok=True)
    images = []
    seen_hashes = set()
    img_counter = 0

    def _save(name, raw, ext):
        """Schreibt rohe Bild-Bytes auf Disk und gibt einen Eintrag zurück."""
        nonlocal img_counter
        img_counter += 1
        safe_name = f"{img_counter}_{name.replace('/', '_')[:80]}"
        if not safe_name.lower().endswith(ext):
            safe_name += ext
        path = os.path.join(imgs_dir, safe_name)
        with open(path, "wb") as fh:
            fh.write(raw)
        return {
            'name': safe_name,
            'path': path,
            'ext':  ext,
            'size': len(raw),
        }

    try:
        with zipfile.ZipFile(zip_path) as zf:
            names = zf.namelist()

            # 1. Direkte Bilddateien (Streaming)
            for name in names:
                if '__MACOSX' in name or name.startswith('.'):
                    continue
                ext = os.path.splitext(name.lower())[1]
                if ext not in IMAGE_EXTS:
                    continue
                info = zf.getinfo(name)
                if info.file_size < 8000:  # Skip Thumbnails
                    continue
                # Lese minimal in RAM für Hash + Save
                raw = zf.read(name)
                if len(raw) < 8000:
                    continue
                h = hashlib.md5(raw).hexdigest()
                if h in seen_hashes:
                    del raw
                    continue
                seen_hashes.add(h)
                images.append(_save(name.split('/')[-1], raw, ext))
                del raw

            # 2. Aus PDFs extrahieren – Streaming pro PDF, max 6 PDFs
            pdf_entries = [
                (n, zf.getinfo(n).file_size)
                for n in names
                if n.lower().endswith('.pdf')
                and '__MACOSX' not in n and not n.startswith('.')
            ]
            pdf_entries.sort(key=lambda x: x[1], reverse=True)
            for pdf_name_in_zip, _ in pdf_entries[:6]:
                if len(images) >= 15:
                    break
                try:
                    # Stream PDF auf Disk, dann mit fitz.open(path) öffnen
                    tmp_pdf = os.path.join(imgs_dir, f"_tmp_{img_counter}.pdf")
                    with zf.open(pdf_name_in_zip) as src, open(tmp_pdf, "wb") as dst:
                        while True:
                            chunk = src.read(64 * 1024)
                            if not chunk:
                                break
                            dst.write(chunk)
                    pdf_imgs = _extract_images_from_pdf_path(
                        tmp_pdf, pdf_name_in_zip.split('/')[-1],
                        seen_hashes, _save
                    )
                    images.extend(pdf_imgs)
                    print(f"  PDF {pdf_name_in_zip.split('/')[-1]}: {len(pdf_imgs)} Bilder")
                    try: os.unlink(tmp_pdf)
                    except OSError: pass
                except Exception as e:
                    print(f"  PDF lesen Fehler ({pdf_name_in_zip}): {e}")

    except Exception as e:
        print(f"extract_images_from_zip Fehler: {e}")

    images.sort(key=lambda x: x.get('size', 0), reverse=True)
    images = images[:15]
    print(f"Bilder gesamt: {len(images)} (auf Disk unter {imgs_dir})")
    for i, img in enumerate(images[:8]):
        print(f"    {i+1}. {img['name']} ({img['size']//1024} KB)")
    return images


def _extract_images_from_pdf_path(pdf_path, pdf_name, seen_hashes, save_fn):
    """Extrahiert Bilder aus PDF (auf Disk) via PyMuPDF. save_fn() schreibt sie weiter."""
    import hashlib
    images = []
    try:
        import fitz
        doc = fitz.open(pdf_path)
        n_pages = min(len(doc), 3)
        for page_num in range(n_pages):
            for img_idx, img in enumerate(doc[page_num].get_images(full=True)):
                xref = img[0]
                try:
                    base = doc.extract_image(xref)
                    raw  = base["image"]
                    ext  = "." + base["ext"]
                    if ext not in {'.jpg', '.jpeg', '.png', '.webp', '.bmp'}:
                        continue
                    if len(raw) < 20000:
                        continue
                    h = hashlib.md5(raw).hexdigest()
                    if h in seen_hashes:
                        continue
                    seen_hashes.add(h)
                    images.append(save_fn(
                        f"{pdf_name}_s{page_num+1}_b{img_idx+1}{ext}", raw, ext
                    ))
                except Exception:
                    continue
        doc.close()
    except ImportError:
        print("  PyMuPDF nicht verfügbar")
    except Exception as e:
        print(f"  PDF-Bild-Extract Fehler: {e}")
    return images


# Mapping: welche bild_* Slots können Kundenbilder aufnehmen (in Prioritätsreihenfolge)
# Nur diese Kategorien werden automatisch klassifiziert. Alle Projekt-Fotos
# (Außenansichten, Innenraum, Wohnungstypen, Grundrisse, Hotel-Feeling, Collagen)
# sowie Amenity-Fotos bleiben LEER und werden vom Kunden via Preview-UI
# hochgeladen – das verhindert falsche Auto-Zuweisungen wie Fahrradbilder
# in einem Dachterrasse-Slot.
CUSTOMER_IMAGE_SLOTS = {
    "lageplan":       ["bild_lageplan"],
    "quartier":       ["bild_quartier", "bild_stadt_gross", "bild_stadt_klein"],
}


def classify_and_assign_customer_images(images):
    """
    Sendet Kundenbilder an Claude Vision und lässt es sie den richtigen bild_*-Slots zuweisen.
    Gibt {bild_key: image_bytes} zurück.
    Fällt auf regelbasierte Zuweisung zurück wenn kein Claude-Key vorhanden.
    """
    if not images:
        return {}

    # Helfer: Bytes für ein Bild lesen (von Disk wenn 'path', sonst aus 'bytes')
    def _img_bytes(img):
        if img.get('path') and os.path.exists(img['path']):
            with open(img['path'], 'rb') as fh:
                return fh.read()
        return img.get('bytes', b'')

    # Regelbasierter Fallback (anhand Dateiname). Nutzt nur Kategorien, die
    # tatsächlich in CUSTOMER_IMAGE_SLOTS existieren – alles andere ignoriert.
    def _rule_based(images):
        slot_counters = {cat: 0 for cat in CUSTOMER_IMAGE_SLOTS}
        result = {}
        for img in images:
            name_lower = img['name'].lower()
            cat = None
            if any(k in name_lower for k in ('lageplan', 'lage', 'map', 'karte', 'site-plan')):
                cat = 'lageplan'
            elif any(k in name_lower for k in ('quartier', 'strasse', 'straße', 'stadt', 'neighborhood', 'umgebung')):
                cat = 'quartier'
            if cat is None or cat not in CUSTOMER_IMAGE_SLOTS:
                continue
            slots = CUSTOMER_IMAGE_SLOTS[cat]
            idx = slot_counters[cat]
            if idx < len(slots):
                result[slots[idx]] = _img_bytes(img)
                slot_counters[cat] += 1
        return result

    if not CLAUDE_API_KEY:
        print("classify_customer_images: kein Claude-Key → regelbasierte Zuweisung")
        result = _rule_based(images)
        print(f"  {len(result)} Bilder zugewiesen (regelbasiert)")
        return result

    # Claude Vision: bis zu 25 Bilder, in Batches à 10
    all_images = images[:25]
    slot_list = "\n".join(
        f"- {cat}: {', '.join(slots[:4])}" + (" …" if len(slots) > 4 else "")
        for cat, slots in CUSTOMER_IMAGE_SLOTS.items()
    )
    classify_prompt = (
        "Analysiere diese Immobilien-Bilder und weise jedem Bild EINE passende Kategorie zu, "
        "ODER 'skip' wenn keine Kategorie passt.\n\n"
        "Kategorien:\n" + slot_list + "\n\n"
        "STRENGE REGELN:\n"
        "- lageplan: NUR Karten/Stadtpläne/Lagepläne (Vogelperspektive auf Stadt/Quartier mit Markierungen)\n"
        "- quartier: Straßen, Stadtteile, Nachbarschaft, urbane Umgebung mit Stadtcharakter\n"
        "- skip: Alles andere (Renderings, Innenräume, Grundrisse, Fahrradbilder, "
        "Logos, etc.) – diese Bilder werden vom Nutzer manuell zugewiesen.\n\n"
        "Antworte NUR mit JSON: {\"1\": \"kategorie\", \"2\": \"kategorie\", ...}\n"
        "Jedes Bild bekommt genau EINE Kategorie oder 'skip'. Keine Erklärungen."
    )

    all_assignments = {}  # global_idx → category
    BATCH_SIZE = 10

    for batch_start in range(0, len(all_images), BATCH_SIZE):
        batch = all_images[batch_start:batch_start + BATCH_SIZE]
        content = []
        for i, img in enumerate(batch):
            mt = ('image/jpeg' if img['ext'] in ('.jpg', '.jpeg')
                  else 'image/png' if img['ext'] == '.png'
                  else 'image/webp')
            # Bytes on-demand von Disk (oder bytes-Field falls direkt gesetzt)
            raw = _img_bytes(img)
            content.append({"type": "image", "source": {
                "type": "base64", "media_type": mt,
                "data": base64.b64encode(raw).decode()
            }})
            content.append({"type": "text", "text": f"Bild {i+1}: {img['name']}"})
            del raw
        content.append({"type": "text", "text": classify_prompt})

        try:
            resp = requests.post(
                "https://api.anthropic.com/v1/messages",
                headers={"x-api-key": CLAUDE_API_KEY, "anthropic-version": "2023-06-01",
                         "content-type": "application/json"},
                json={"model": "claude-haiku-4-5-20251001", "max_tokens": 500,
                      "messages": [{"role": "user", "content": content}]},
                timeout=90
            )
            resp.raise_for_status()
            json_text = resp.json()["content"][0]["text"]
            json_text = json_text.replace("```json", "").replace("```", "").strip()
            batch_assignments = json.loads(json_text)
            # Mapping local_idx (1-basiert) → global_idx
            for local_str, cat in batch_assignments.items():
                global_idx = batch_start + int(local_str) - 1
                if 0 <= global_idx < len(all_images):
                    all_assignments[global_idx] = cat
            print(f"  Batch {batch_start//BATCH_SIZE + 1}: {len(batch_assignments)} Bilder klassifiziert")
        except Exception as e:
            print(f"  Batch {batch_start//BATCH_SIZE + 1} Claude-Fehler: {e}")
            # Fallback: keine Auto-Zuweisung für diesen Batch
            for i, img in enumerate(batch):
                if (batch_start + i) not in all_assignments:
                    all_assignments[batch_start + i] = 'skip'

    if not all_assignments:
        print("classify_customer_images: keine Klassifizierung → regelbasierter Fallback")
        return _rule_based(all_images)

    # Sortiere Bilder pro Kategorie nach Größe (größte zuerst → wichtigste Slots)
    by_category = {cat: [] for cat in CUSTOMER_IMAGE_SLOTS}
    for global_idx, cat in all_assignments.items():
        cat = cat.strip().lower() if isinstance(cat, str) else ''
        # Unbekannte oder 'skip'-Kategorien werden nicht zugewiesen
        if cat not in CUSTOMER_IMAGE_SLOTS:
            continue
        by_category[cat].append((all_images[global_idx]['size'], global_idx))

    # Innerhalb jeder Kategorie: größte zuerst
    for cat in by_category:
        by_category[cat].sort(reverse=True)

    # Slots befüllen – Bytes erst hier on-demand von Disk lesen
    result = {}
    for cat, items in by_category.items():
        slots = CUSTOMER_IMAGE_SLOTS[cat]
        for slot_idx, (size, global_idx) in enumerate(items):
            if slot_idx >= len(slots):
                break
            img = all_images[global_idx]
            result[slots[slot_idx]] = _img_bytes(img)
            print(f"  {img['name']} ({size//1024} KB) → {slots[slot_idx]} [{cat}]")

    print(f"classify_customer_images: {len(result)} Bilder zugewiesen (Claude Vision)")
    return result


def _fetch_wikipedia_city_image(city, lang="de"):
    """Holt das Hauptbild des Wikipedia-Artikels über eine Stadt.
    Zuverlässiger als Wikimedia Commons Suche – gibt immer das offizielle Stadtfoto zurück.
    """
    title = city.replace(" ", "_")
    for wiki_lang in [lang, "en"]:
        try:
            resp = requests.get(
                f"https://{wiki_lang}.wikipedia.org/api/rest_v1/page/summary/{requests.utils.quote(title)}",
                headers={"User-Agent": "interpres-expose/1.0"},
                timeout=10,
            )
            if resp.status_code == 200:
                d = resp.json()
                img = ((d.get("originalimage") or {}).get("source") or
                       (d.get("thumbnail") or {}).get("source"))
                if not img:
                    continue
                # Filter: SVG (Wappen, Logos, rasterisierte Vektoren) ablehnen.
                # Wikipedia rendert SVGs zu PNG, Pfad enthält dann '.svg/' im Mittelteil.
                low = img.lower()
                if ".svg/" in low or low.split("?")[0].endswith(".svg"):
                    print(f"  Wikipedia ({wiki_lang}) '{city}' → SVG/Wappen abgelehnt")
                    continue
                # Auch Wappen explizit blocken (URL enthält 'wappen', 'coa', 'flag')
                if any(b in low for b in ('wappen', 'coa_', 'coat_of_arms', '_flag', 'logo_')):
                    print(f"  Wikipedia ({wiki_lang}) '{city}' → Wappen/Logo abgelehnt")
                    continue
                if low.split("?")[0].rsplit(".", 1)[-1] in ("jpg", "jpeg", "png"):
                    print(f"  Wikipedia ({wiki_lang}) '{city}': {img[:80]}")
                    return img
        except Exception as e:
            print(f"  Wikipedia {wiki_lang} Fehler für '{city}': {e}")
    return None


def _fetch_city_landmark_image(stadt, landmark_type="dom"):
    """Lädt das Wikipedia-Hauptbild eines Wahrzeichens der Stadt.
    Garantiert ein echtes Foto der Stadt (nicht generisch wie Unsplash).

    landmark_type: 'dom' | 'hbf' | 'park' | 'altstadt' | 'rathaus' | 'markt'
    """
    if not stadt:
        return None
    s = stadt.strip()
    # Pattern pro landmark_type, in Reihenfolge (erstes Bild gewinnt)
    patterns = {
        "dom":      [f"{s}er Dom", f"Dom zu {s}", f"Kathedrale {s}", f"{s} Dom"],
        "hbf":      [f"Hauptbahnhof {s}", f"{s} Hauptbahnhof", f"Bahnhof {s}"],
        "park":     [f"Stadtpark {s}", f"Bürgerpark {s}", f"{s} Stadtpark", f"{s}er Stadtpark"],
        "altstadt": [f"Altstadt {s}", f"{s}er Altstadt", f"Marktplatz {s}", f"Rathausplatz {s}"],
        "rathaus":  [f"Rathaus {s}", f"{s}er Rathaus", f"Altes Rathaus {s}"],
        "markt":    [f"Marktplatz {s}", f"{s}er Markt", f"Hauptmarkt {s}"],
        "skyline":  [f"Skyline {s}", f"Panorama {s}", f"{s} Skyline", f"{s} Stadtansicht"],
        "fluss":    [f"Elbe {s}", f"Donau {s}", f"Rhein {s}", f"Main {s}", f"Oder {s}"],
        "uni":      [f"Universität {s}", f"{s}er Universität", f"Hochschule {s}"],
    }
    candidates = patterns.get(landmark_type, [s])

    for query in candidates:
        title = query.replace(" ", "_")
        for wiki_lang in ["de", "en"]:
            try:
                resp = requests.get(
                    f"https://{wiki_lang}.wikipedia.org/api/rest_v1/page/summary/{requests.utils.quote(title)}",
                    headers={"User-Agent": "interpres-expose/1.0"},
                    timeout=10,
                )
                if resp.status_code != 200:
                    continue
                d = resp.json()
                img = ((d.get("originalimage") or {}).get("source") or
                       (d.get("thumbnail") or {}).get("source"))
                if not img:
                    continue
                low = img.lower()
                if ".svg/" in low or low.split("?")[0].endswith(".svg"):
                    continue
                if any(b in low for b in ('wappen', 'coa_', 'coat_of_arms', '_flag', 'logo_')):
                    continue
                if low.split("?")[0].rsplit(".", 1)[-1] in ("jpg", "jpeg", "png"):
                    print(f"  Wahrzeichen ({landmark_type}/{query}): {img[:80]}")
                    return img
            except Exception as e:
                print(f"  Wahrzeichen Fehler '{query}' ({wiki_lang}): {e}")
    return None


def _fetch_wikimedia_image(query, top_n=1):
    """Sucht thematisch passende Bilder via Wikimedia Commons.
    top_n=1 → eine URL (str|None), top_n>1 → Liste der besten N URLs.
    """
    try:
        resp = requests.get(
            "https://commons.wikimedia.org/w/api.php",
            params={
                "action": "query",
                "generator": "search",
                "gsrnamespace": 6,
                "gsrsearch": query,
                "gsrlimit": 12,
                "prop": "imageinfo",
                "iiprop": "url|size",
                "iiurlwidth": 1200,
                "format": "json",
            },
            headers={"User-Agent": "interpres-expose/1.0"},
            timeout=10,
        )
        if resp.status_code != 200:
            return None if top_n == 1 else []
        pages = resp.json().get("query", {}).get("pages", {})
        candidates = []
        for page in pages.values():
            for ii in page.get("imageinfo", []):
                url = ii.get("thumburl") or ii.get("url", "")
                w, h = ii.get("thumbwidth", 0), ii.get("thumbheight", 0)
                if not url:
                    continue
                ext = url.lower().split("?")[0].rsplit(".", 1)[-1]
                if ext not in ("jpg", "jpeg", "png"):
                    continue
                if w > 0 and h > 0 and w < h:
                    continue
                candidates.append((w * h, url))
        candidates.sort(reverse=True)
        urls = [u for _, u in candidates]
        if top_n == 1:
            return urls[0] if urls else None
        return urls[:top_n]
    except Exception as e:
        print(f"  Wikimedia Fehler für '{query}': {e}")
    return None if top_n == 1 else []


def _validate_image_with_claude_vision(url, expected_subject_de):
    """Prüft via Claude Vision ob das Bild wirklich `expected_subject_de` zeigt.
    Gibt True/False zurück. Bei Fehler (kein Key, Timeout, etc.) → False (lieber kein Bild).
    """
    if not CLAUDE_API_KEY or not url:
        return False
    try:
        # Bild herunterladen
        img_resp = requests.get(url, timeout=15, headers={"User-Agent": "interpres-expose/1.0"})
        if img_resp.status_code != 200 or len(img_resp.content) < 1000:
            return False
        ext = url.lower().split("?")[0].rsplit(".", 1)[-1]
        media_type = {"jpg": "image/jpeg", "jpeg": "image/jpeg", "png": "image/png"}.get(ext, "image/jpeg")
        b64 = base64.b64encode(img_resp.content).decode()
        # Sehr knapper JA/NEIN-Prompt für Geschwindigkeit
        resp = requests.post(
            "https://api.anthropic.com/v1/messages",
            headers={"x-api-key": CLAUDE_API_KEY, "anthropic-version": "2023-06-01", "content-type": "application/json"},
            json={
                "model": "claude-haiku-4-5-20251001",
                "max_tokens": 10,
                "messages": [{
                    "role": "user",
                    "content": [
                        {"type": "image", "source": {"type": "base64", "media_type": media_type, "data": b64}},
                        {"type": "text", "text": f"Zeigt dieses Bild eindeutig: {expected_subject_de}? Antworte NUR mit 'ja' oder 'nein'."},
                    ],
                }],
            },
            timeout=20,
        )
        if resp.status_code != 200:
            return False
        text = "".join(b.get("text", "") for b in resp.json().get("content", [])).strip().lower()
        return text.startswith("ja")
    except Exception as e:
        print(f"  Vision-Validation Fehler: {e}")
        return False


def _find_validated_amenity_image(query, expected_subject_de, max_tries=4):
    """Holt mehrere Wikimedia-Treffer, validiert sie der Reihe nach mit Claude Vision.
    Gibt erste passende URL zurück oder None.
    """
    urls = _fetch_wikimedia_image(query, top_n=max_tries) or []
    for url in urls:
        if _validate_image_with_claude_vision(url, expected_subject_de):
            print(f"    ✓ Validiert: {expected_subject_de} → {url[:80]}")
            return url
        else:
            print(f"    ✗ Verworfen: {expected_subject_de} ({url[:80]})")
    return None


def fetch_unsplash_image(query):
    """Holt Bild-URL von Unsplash via /search/photos (relevanteste Treffer).
    Vorher: /photos/random → unzuverlässig (lieferte oft unrelated photos).
    Bei fehlendem Key oder 0 Treffern → Picsum-Fallback.
    """
    if UNSPLASH_ACCESS_KEY:
        try:
            resp = requests.get(
                "https://api.unsplash.com/search/photos",
                params={"query": query, "orientation": "landscape", "per_page": 5},
                headers={"Authorization": f"Client-ID {UNSPLASH_ACCESS_KEY}"},
                timeout=10
            )
            print(f"  Unsplash [{resp.status_code}] search={query!r}")
            if resp.status_code == 200:
                results = resp.json().get("results", [])
                if not results:
                    print(f"    → 0 Treffer")
                else:
                    # Top-Treffer (relevanteste Match)
                    url = results[0]["urls"]["regular"]
                    alt = (results[0].get("alt_description") or "")[:60]
                    print(f"    → '{alt}' / {url[:60]}")
                    return url
            else:
                print(f"    Fehler-Body: {resp.text[:120]}")
        except Exception as e:
            print(f"  Unsplash Exception für '{query}': {e}")
    else:
        print(f"  UNSPLASH_ACCESS_KEY fehlt – Picsum-Fallback für '{query}'")

    # Picsum-Fallback: deterministisches Bild anhand des Query-Hashes
    seed = abs(hash(query)) % 1000
    url = f"https://picsum.photos/seed/{seed}/1200/800"
    print(f"  Picsum-Fallback → {url}")
    return url

def fill_image_placeholders(data):
    """Füllt bild_*-Slots die noch leer sind.
    Stadtbilder: Wikimedia Commons (echte Fotos der Stadt).
    Alle anderen: Picsum-Fallback (deterministisch, kein Rate-Limit).
    """
    stadt = data.get("stadt", "")

    # Stadtspezifische Slots: Wikipedia REST API (zuverlässig) → Wikimedia-Fallback
    stadtteil = data.get("stadtteil", "")
    # Stadt-Bilder: Wikipedia/Commons – NUR für echte Stadt-Slots, NICHT für
    # Projekt-Cover (bild_titel) und nicht für Standort-Außenfoto (das soll
    # das Kundenbild bleiben, sonst läuft der Eigentum-Renderer rein).
    # (slot, landmark_type, unsplash_query)
    # landmark_type wählt den Wikipedia-Wahrzeichen-Artikel (Dom, Hbf, Park ...)
    # → garantiert echte Stadt-Fotos statt zufälliger Stock-Photos.
    _CITY_SLOT_SEARCHES = [
        ("bild_quartier",      "altstadt", f"{stadtteil or stadt} {stadt} historic architecture"),
        ("bild_stadt_gross",   "skyline",  f"{stadt} skyline panorama"),
        ("bild_stadt_klein",   "markt",    f"{stadt} marketplace downtown"),
        ("bild_stadt_presse",  "dom",      f"{stadt} cathedral landmark"),
        ("bild_stadt_branche", "hbf",      f"{stadt} train station modern"),
    ]

    def _norm_url(u):
        """Normalisiert Wikimedia-Thumb-URLs: '/120px-' und '/800px-' zeigen
        auf dasselbe Bild — beim Dedup gleich behandeln."""
        if not u:
            return u
        return re.sub(r'/\d+px-', '/', u.split('?', 1)[0]).lower()

    used_city_urls = set()  # Dedup: dieselbe Datei nie zweimal in zwei Slots

    def _try(url):
        """Akzeptiert URL nur wenn Inhalt nicht schon belegt ist."""
        if not url:
            return None
        n = _norm_url(url)
        if n in used_city_urls:
            return None
        return url

    for slot, landmark_type, unsplash_q in _CITY_SLOT_SEARCHES:
        if slot not in data:
            continue
        if data.get(slot) and str(data[slot]).startswith("http"):
            used_city_urls.add(_norm_url(data[slot]))
            continue
        if not stadt:
            continue
        # 1. Wikipedia-Wahrzeichen-Artikel (Dom, Hbf, etc.) – ECHTE Stadt-Fotos
        url = _try(_fetch_city_landmark_image(stadt, landmark_type))
        # 2. Wikipedia REST mit Stadt selbst (filtert Wappen/SVG raus)
        if not url:
            url = _try(_fetch_wikipedia_city_image(stadt))
        # 3. Unsplash search mit englischem Begriff
        if not url:
            cand = fetch_unsplash_image(unsplash_q)
            if cand and "picsum.photos" in cand:
                cand = None
            url = _try(cand)
        # 4. Wikimedia Commons Suche — top_n=4 → bei Dedup-Treffer naechsten nehmen
        if not url:
            for cand in (_fetch_wikimedia_image(f"{stadt} {landmark_type}", top_n=4) or []):
                url = _try(cand)
                if url:
                    break
        # 5. Wikimedia Commons mit alternativem Begriff (mehr Variation pro Stadt)
        if not url:
            for cand in (_fetch_wikimedia_image(f"{stadt} {unsplash_q}", top_n=4) or []):
                url = _try(cand)
                if url:
                    break
        # 6. Allerletzter Notfall: Picsum (mit Slot-spezifischem Seed → unterschiedlich)
        if not url:
            seed = abs(hash(f"{stadt}_{slot}_{landmark_type}")) % 10000
            url = f"https://picsum.photos/seed/{seed}/1200/800"
        data[slot] = url
        used_city_urls.add(_norm_url(url))
        print(f"  {slot} ({landmark_type}) → {url[:70]}")

    # Slots die NUR mit echten Projektfotos vom Kunden befüllt werden.
    # Alles andere kriegt automatisch ein passendes Stock-/Wikimedia-Bild.
    # Ziel: Kunde sieht ein vollständiges Exposé und ersetzt nur die paar
    # echten Projekt-spezifischen Bilder (Cover, Außenansichten, Grundrisse,
    # konkrete Wohnungen). Generische Themenbilder (Greenliving, Interior,
    # Ausstattung, Rechtliches, etc.) liefert die Pipeline.
    _NO_FALLBACK_SLOTS = {
        # ── Echte Projekt-Visualisierungen / Renderings vom Bauträger ──
        "bild_titel", "bild_projekt", "bild_projekt_aussen",
        "bild_ansicht_1", "bild_ansicht_2",
        # ── Außenansichten Standort (echtes Quartier-Foto wenn möglich) ──
        "bild_standort_aussen",
        # ── Konkrete Grundriss-Pläne (DIN-277 vom Architekt) ──
        "bild_grundriss_1", "bild_grundriss_2", "bild_grundriss_3", "bild_grundriss_4",
        # ── Amenity-Bilder werden separat via Wikimedia + Vision validiert ──
        *{f"bild_amenity_{n}" for n in range(1, 10)},
        # ── Wohnungstyp-Bilder: Kunde lädt seine Grundriss-/Innenraum-Fotos selbst.
        # Vorher kamen Berge/Auto/Wölfe via Picsum-Fallback rein — sah unprofessionell.
        *{f"bild_we_{n}" for n in range(1, 21)},
        # ── Kurz-Exposé Cover-Collage (6 Slots) + Seite-2-Bilder (4 Slots).
        # Kunde lädt seine echten Projekt-/Lifestyle-Fotos hoch.
        *{f"bild_titel_{n}" for n in range(1, 7)},
        *{f"bild_kurz_{n}"  for n in range(1, 5)},
    }

    queries = UNSPLASH_QUERIES.copy()
    # Lageplan-URL kommt aus OSM, nicht aus Picsum
    queries.pop("BILD_LAGEPLAN", None)

    # Themen-Slots: hier bringt Unsplash deutlich passendere Bilder als Picsum
    # (Architektur, Interior, Nachhaltigkeit, Hotel-Look, Magdeburg-Branche).
    _UNSPLASH_THEMED_PREFIXES = (
        "bild_greenliving", "bild_interior",
        "bild_ausstattung", "bild_hotel", "bild_collage",
        "bild_rechtlich", "bild_grundriss_intro",
        "bild_standort_innen",
        "bild_stadt_presse", "bild_stadt_branche",
    )

    filled = 0
    for placeholder_key, query in queries.items():
        data_key = placeholder_key.lower()
        if data_key not in data:
            continue
        # Skip wenn Slot bereits durch Kundenbild oder Wikimedia belegt
        if data.get(data_key) and str(data[data_key]).startswith("http"):
            continue
        # Kein Fallback für projekt-spezifische Slots ohne Kundenbild
        if data_key in _NO_FALLBACK_SLOTS:
            continue
        # Skip bild_we_N für nicht vorhandene WE-Typen
        _m = re.match(r'^bild_we_(\d+)$', data_key)
        if _m:
            n = int(_m.group(1))
            if n > 2:
                pair_k   = (n + 1) // 2
                left_n   = pair_k * 2 - 1
                right_n  = pair_k * 2
                has_text = (data.get(f"we_beispiel_{left_n}") or data.get(f"we_nummern_{left_n}")
                            or data.get(f"we_beispiel_{right_n}") or data.get(f"we_nummern_{right_n}"))
                if not has_text:
                    continue
        # Themen-Slots: erst Unsplash mit echtem Suchbegriff, sonst Picsum-Fallback
        url = None
        if any(data_key.startswith(p) for p in _UNSPLASH_THEMED_PREFIXES):
            try:
                url = fetch_unsplash_image(query)
                if url and "picsum.photos" in url:
                    url = None
            except Exception:
                url = None
        if not url:
            seed = abs(hash(query)) % 1000
            url = f"https://picsum.photos/seed/{seed}/1200/800"
        data[data_key] = url
        filled += 1

    # Amenity-Bilder: Unsplash als PRIMÄR (zuverlässiger als Wikimedia für
    # spezifische Begriffe wie Fahrradabstellplatz, Solaranlage etc.).
    # Mapping: Keyword → englischer Unsplash-Suchbegriff (englisch funktioniert
    # auf Unsplash deutlich besser).
    _AMENITY_UNSPLASH = [
        ("dachterras",   "rooftop terrace garden modern"),
        ("balkon",       "modern apartment balcony"),
        ("terras",       "outdoor terrace residential modern"),
        ("fahrrad",      "bicycle parking storage room"),
        ("e-bike",       "electric bike charging station"),
        ("bike",         "bicycle parking modern"),
        ("spindel",      "storage room basement organized"),
        ("gemeinschaft", "community lounge modern interior"),
        ("lounge",       "lounge interior modern minimal"),
        ("smart-lock",   "smart door lock keyless entry"),
        ("smart",        "smart home tablet control"),
        ("sanitär",      "modern bathroom shower minimal"),
        ("bad",          "modern bathroom luxury"),
        ("boden",        "wooden floor parquet apartment"),
        ("außenanlage",  "modern building landscaping garden"),
        ("aufzug",       "modern elevator interior"),
        ("gym",          "modern gym equipment fitness"),
        ("fitness",      "fitness studio modern equipment"),
        ("pool",         "modern swimming pool architecture"),
        ("küche",        "modern kitchen apartment"),
        ("tiefgarage",   "underground garage parking"),
        ("parken",       "parking garage modern"),
        ("stellplat",    "parking spot modern building"),
        ("solar",        "solar panels rooftop building"),
        ("photovoltaik", "photovoltaic solar panels modern"),
        ("fernwärme",    "modern heating radiator interior"),
        ("concierge",    "hotel reception lobby modern"),
        ("paket",        "parcel locker delivery box"),
        ("post",         "mailbox modern apartment building"),
        ("möbli",        "furnished modern apartment interior"),
        ("möbel",        "modern furniture interior design"),
        ("dach",         "green roof rooftop garden"),
        ("garten",       "modern garden landscaping residential"),
        ("ladestat",     "electric vehicle charging station home"),
        ("café",         "modern cafe interior"),
        ("internet",     "fiber internet router home"),
        ("glasfaser",    "fiber optic cable network"),
        ("barriere",     "accessible apartment ramp design"),
        ("aufzüge",      "modern elevator architecture"),
    ]
    amenity_filled = 0
    amenity_skipped = 0
    for n in range(1, 10):
        amenity_val = str(data.get(f"amenity_{n}", "")).strip().lower()
        bild_key = f"bild_amenity_{n}"
        if bild_key not in data or data.get(bild_key):
            continue
        if not amenity_val:
            continue
        # Mehrere Suchbegriffe versuchen: erst gemapptes Keyword, dann amenity_val
        # selbst, dann generisch. Erstes Treffer gewinnt.
        candidates = []
        for kw, q in _AMENITY_UNSPLASH:
            if kw in amenity_val:
                candidates.append(q)
                break
        # IMMER auch amenity_val selbst probieren (oft schon English-tauglich)
        candidates.append(f"{amenity_val} modern interior")
        candidates.append(f"modern apartment {amenity_val.split()[0] if amenity_val else 'amenity'}")

        url = None
        for q in candidates:
            url = fetch_unsplash_image(q)
            if url and "picsum.photos" in url:
                url = None
                continue
            if url:
                break
        if not url:
            url = _fetch_wikimedia_image(amenity_val + " modern")
        if url:
            data[bild_key] = url
            amenity_filled += 1
            print(f"  ✓ {bild_key} ({amenity_val[:30]}) → {url[:60]}")
        else:
            amenity_skipped += 1
            print(f"  ✗ {bild_key} ({amenity_val[:30]}) → kein Bild gefunden")

    print(f"fill_image_placeholders: {filled} Picsum-Slots, {amenity_filled} Amenities befüllt, "
          f"{amenity_skipped} leer (kein Bild gefunden)")
    return data


# ── Geocoding + Proximity Calculation via Nominatim + Overpass ─────────────

def _geocode_address(adresse, stadt):
    """Geocodiert eine Adresse via Nominatim.
    Gibt (lat, lon, official_city) zurück oder None.
    official_city = offizieller Gemeinde-/Stadtname laut OSM (kann von 'stadt' abweichen).
    """
    try:
        resp = requests.get(
            "https://nominatim.openstreetmap.org/search",
            params={
                "q": f"{adresse}, {stadt}, Deutschland",
                "format": "json",
                "addressdetails": 1,
                "limit": 1,
            },
            headers={"User-Agent": "interpres-expose/1.0 (contact@interpres.de)"},
            timeout=10,
        )
        if resp.status_code == 200 and resp.json():
            loc = resp.json()[0]
            addr = loc.get("address", {})
            # Offizielle Stadt/Gemeinde aus OSM-Adressdaten (Priorität: city > town > village > municipality)
            official_city = (
                addr.get("city")
                or addr.get("town")
                or addr.get("village")
                or addr.get("municipality")
                or stadt  # Fallback auf User-Angabe
            )
            print(f"  Nominatim: stadt_input='{stadt}' → official_city='{official_city}'")
            return float(loc['lat']), float(loc['lon']), official_city
    except Exception as e:
        print(f"  Geocoding Fehler für '{adresse}, {stadt}': {e}")
    return None


def _osm_lageplan_url(lat, lon, zoom=15):
    """Generiert eine OpenStreetMap Static Map URL für den Projektstandort."""
    # Wikimedia static map: free, no API key needed, uses OSM tiles
    return f"https://maps.wikimedia.org/img/osm-intl,{zoom},{lat:.5f},{lon:.5f},800x600.png"


def _search_city_info(stadt, stadtteil=""):
    """Sucht aktuelle Infos zur Stadt via Tavily API.
    Gibt strukturierten Text-Block mit Quellen-URLs zurück.
    """
    if not TAVILY_API_KEY:
        print(f"  TAVILY_API_KEY fehlt – kein Web-Search für '{stadt}'")
        return ""

    # Mehrere fokussierte Queries für unterschiedliche Themenbereiche
    queries = [
        (f"{stadt} Einwohner aktuell Statistik Bevölkerung", "demografie"),
        (f"{stadt} Mietpreise Wohnungsmarkt Mietsteigerung Neubau 2024 2025", "mietmarkt"),
        (f"{stadt} BIP Wirtschaftsleistung Bundesland Wachstum", "wirtschaftskraft"),
        (f"{stadt} Universität Hochschule Studierende Forschung", "bildung"),
        (f"{stadt} größte Arbeitgeber Unternehmen Industriepark Investitionen", "arbeitgeber"),
        (f"{stadt} Infrastrukturprojekte Bahnhof Verkehr Stadtentwicklung 2024 2025", "infrastruktur"),
        (f"{stadt} Logistik Hafen Industrie Ansiedlungen", "industrie"),
    ]
    if stadtteil and stadtteil.lower() != stadt.lower():
        queries.append((f"{stadtteil} {stadt} Stadtteil Entwicklung Projekte", "stadtteil"))

    sections = []
    all_sources = []  # für quelle_1-4
    for q, topic in queries:
        try:
            resp = requests.post(
                "https://api.tavily.com/search",
                json={
                    "api_key": TAVILY_API_KEY,
                    "query": q,
                    "search_depth": "advanced",  # tiefer für bessere Snippets
                    "max_results": 4,
                    "include_answer": True,
                    "include_raw_content": False,
                },
                timeout=25,
            )
            if resp.status_code != 200:
                print(f"  Tavily Fehler {resp.status_code} für '{q}': {resp.text[:100]}")
                continue
            d = resp.json()
            block = [f"### {topic.upper()}: {q}"]
            if d.get("answer"):
                block.append(f"ZUSAMMENFASSUNG: {d['answer']}")
            for r in d.get("results", [])[:3]:
                title = (r.get("title") or "").strip()
                url   = (r.get("url") or "").strip()
                cnt   = (r.get("content") or r.get("snippet") or "")[:500]
                if cnt:
                    block.append(f"- {title} [{url}]\n  {cnt}")
                if url:
                    all_sources.append({"topic": topic, "title": title, "url": url})
            sections.append("\n".join(block))
        except Exception as e:
            print(f"  Tavily Exception für '{q}': {e}")

    combined = "\n\n".join(sections)
    if all_sources:
        combined += "\n\n### VERWENDBARE QUELLEN (für quelle_1-4):\n"
        for s in all_sources[:8]:
            combined += f"- [{s['topic']}] {s['title']} — {s['url']}\n"
    print(f"  Tavily: {len(sections)} Themen-Sektionen, {len(all_sources)} Quellen, {len(combined)} Zeichen")
    return combined


def _calculate_proximity_data(adresse, stadt, lat, lon):
    """
    Berechnet Entfernungen zu nahegelegenen POIs via Overpass API.
    Gibt dict mit einkaufen_N_name, min_einkaufen_N etc. zurück.
    Walking speed: 80 m/min | Cycling: 250 m/min
    """
    import math

    def _haversine_m(lat1, lon1, lat2, lon2):
        R = 6371000
        p1, p2 = math.radians(lat1), math.radians(lat2)
        dp, dl = math.radians(lat2 - lat1), math.radians(lon2 - lon1)
        a = math.sin(dp / 2) ** 2 + math.cos(p1) * math.cos(p2) * math.sin(dl / 2) ** 2
        return R * 2 * math.atan2(math.sqrt(a), math.sqrt(1 - a))

    def _walk(m): return str(max(1, round(m / 80)))
    def _bike(m): return str(max(1, round(m / 250)))

    # Defaults (match template DQN reference values)
    result = {
        "einkaufen_1_name": "Bäckerei",    "min_einkaufen_1": "2",
        "einkaufen_2_name": "Supermarkt",  "min_einkaufen_2": "2",
        "einkaufen_3_name": "Drogerie",    "min_einkaufen_3": "3",
        "einkaufen_4_name": "Getränke",    "min_einkaufen_4": "4",
        "arzt_1_name": "Hausarzt",         "min_arzt_1": "5",
        "arzt_2_name": "Facharzt",         "min_arzt_2": "8",
        "arzt_3_name": "Apotheke",         "min_arzt_3": "3",
        "arzt_4_name": "Krankenhaus",      "min_arzt_4": "12",
        "sport_1_name": "Fitnessstudio",   "min_sport_1": "8",
        "sport_2_name": "Schwimmbad",      "min_sport_2": "10",
        "sport_3_name": "Sportanlage",     "min_sport_3": "6",
        "sport_4_name": "Sportpark",       "min_sport_4": "5",
        "bildung_1_name": "Kita",          "min_bildung_1": "5",
        "bildung_2_name": "Grundschule",   "min_bildung_2": "8",
        "bildung_3_name": "Gymnasium",     "min_bildung_3": "10",
        "bildung_4_name": "Universität",   "min_bildung_4": "15",
    }

    try:
        overpass_q = f"""[out:json][timeout:15];
(
  node(around:2000,{lat},{lon})[amenity=bakery];
  node(around:2000,{lat},{lon})[shop=supermarket];
  node(around:2000,{lat},{lon})[shop=convenience];
  node(around:2000,{lat},{lon})[shop=beverages];
  node(around:2000,{lat},{lon})[shop=chemist];
  node(around:2000,{lat},{lon})[shop=drugstore];
  node(around:2000,{lat},{lon})[amenity=doctors];
  node(around:2000,{lat},{lon})[amenity=clinic];
  node(around:2000,{lat},{lon})[amenity=pharmacy];
  node(around:2000,{lat},{lon})[amenity=hospital];
  node(around:2000,{lat},{lon})[leisure=fitness_centre];
  node(around:2000,{lat},{lon})[leisure=sports_centre];
  node(around:2000,{lat},{lon})[leisure=swimming_pool];
  node(around:2000,{lat},{lon})[leisure=pitch];
  node(around:2000,{lat},{lon})[amenity=kindergarten];
  node(around:2000,{lat},{lon})[amenity=school];
  node(around:2000,{lat},{lon})[amenity=university];
  node(around:2000,{lat},{lon})[amenity=college];
);out body;"""

        op = requests.post(
            "https://overpass-api.de/api/interpreter",
            data={"data": overpass_q},
            timeout=25,
            headers={"User-Agent": "interpres-expose/1.0"}
        )
        if op.status_code != 200:
            print(f"  Overpass HTTP {op.status_code}")
            return result

        cats = {k: [] for k in (
            "bakery", "supermarket", "beverages", "chemist",
            "doctors", "pharmacy", "hospital",
            "fitness", "swimming", "pitch",
            "kindergarten", "school", "university"
        )}
        for el in op.json().get("elements", []):
            tags = el.get("tags", {})
            el_lat, el_lon = float(el.get("lat", lat)), float(el.get("lon", lon))
            name = (tags.get("name") or "").strip()
            dist = _haversine_m(lat, lon, el_lat, el_lon)
            amenity = tags.get("amenity", "")
            shop    = tags.get("shop", "")
            leisure = tags.get("leisure", "")

            if amenity == "bakery" or shop == "bakery":
                cats["bakery"].append((dist, name or "Bäckerei"))
            elif shop == "supermarket":
                cats["supermarket"].append((dist, name or "Supermarkt"))
            elif shop in ("beverages",):
                cats["beverages"].append((dist, name or "Getränkemarkt"))
            elif shop in ("chemist", "drugstore"):
                cats["chemist"].append((dist, name or "Drogerie"))
            elif amenity in ("doctors", "clinic"):
                cats["doctors"].append((dist, name or "Arzt"))
            elif amenity == "pharmacy":
                cats["pharmacy"].append((dist, name or "Apotheke"))
            elif amenity == "hospital":
                cats["hospital"].append((dist, name or "Krankenhaus"))
            elif leisure in ("fitness_centre", "sports_centre"):
                cats["fitness"].append((dist, name or "Fitnessstudio"))
            elif leisure == "swimming_pool":
                cats["swimming"].append((dist, name or "Schwimmbad"))
            elif leisure == "pitch":
                cats["pitch"].append((dist, name or "Sportanlage"))
            elif amenity == "kindergarten":
                cats["kindergarten"].append((dist, name or "Kita"))
            elif amenity == "school":
                cats["school"].append((dist, name or "Schule"))
            elif amenity in ("university", "college"):
                cats["university"].append((dist, name or "Universität"))

        for k in cats:
            cats[k].sort()

        def _truncate(s, maxlen=18):
            """Kürzt Namen auf maxlen Zeichen mit … wenn nötig (verhindert
            Layoutbruch in den schmalen Template-Spalten)."""
            s = (s or "").strip()
            if len(s) <= maxlen:
                return s
            return s[:maxlen - 1].rstrip() + "…"

        def _near(cat, default, default_min, mode="walk"):
            pois = cats.get(cat, [])
            if pois:
                d, n = pois[0]
                return _truncate(n), (_walk(d) if mode == "walk" else _bike(d))
            return default, default_min

        # Einkaufen (walking)
        n, m = _near("bakery", "Bäckerei", "2")
        result["einkaufen_1_name"] = n; result["min_einkaufen_1"] = m
        n, m = _near("supermarket", "Supermarkt", "2")
        result["einkaufen_2_name"] = n; result["min_einkaufen_2"] = m
        n, m = _near("chemist", "Drogerie", "3")
        result["einkaufen_3_name"] = n; result["min_einkaufen_3"] = m
        n, m = _near("beverages", "Getränkemarkt", "4")
        result["einkaufen_4_name"] = n; result["min_einkaufen_4"] = m

        # Ärzte (walking)
        docs = cats.get("doctors", [])
        if len(docs) >= 1:
            d, n = docs[0]; result["arzt_1_name"] = _truncate(n); result["min_arzt_1"] = _walk(d)
        if len(docs) >= 2:
            d, n = docs[1]; result["arzt_2_name"] = _truncate(n); result["min_arzt_2"] = _walk(d)
        n, m = _near("pharmacy", "Apotheke", "3")
        result["arzt_3_name"] = n; result["min_arzt_3"] = m
        n, m = _near("hospital", "Krankenhaus", "12")
        result["arzt_4_name"] = n; result["min_arzt_4"] = m

        # Sport (biking)
        n, m = _near("fitness", "Fitnessstudio", "8", "bike")
        result["sport_1_name"] = n; result["min_sport_1"] = m
        n, m = _near("swimming", "Schwimmbad", "10", "bike")
        result["sport_2_name"] = n; result["min_sport_2"] = m
        n, m = _near("pitch", "Sportanlage", "6", "bike")
        result["sport_3_name"] = n; result["min_sport_3"] = m

        # Bildung (walking/biking)
        n, m = _near("kindergarten", "Kita", "5")
        result["bildung_1_name"] = n; result["min_bildung_1"] = m

        schools = cats.get("school", [])
        grund = [(d, n) for d, n in schools if "gymnasium" not in n.lower()]
        gyms  = [(d, n) for d, n in schools if "gymnasium" in n.lower()]
        if grund:
            d, n = grund[0]; result["bildung_2_name"] = _truncate(n); result["min_bildung_2"] = _walk(d)
        if gyms:
            d, n = gyms[0]; result["bildung_3_name"] = _truncate(n); result["min_bildung_3"] = _walk(d)
        elif len(schools) >= 2:
            d, n = schools[1]; result["bildung_3_name"] = _truncate(n); result["min_bildung_3"] = _walk(d)

        n, m = _near("university", "Universität", "15", "bike")
        result["bildung_4_name"] = n; result["min_bildung_4"] = m

        print(f"  Proximity: {sum(1 for k in result if not k.startswith('min_'))} POI-Namen berechnet")
    except Exception as e:
        print(f"  Proximity Fehler: {e}")

    return result


def analyze_pdfs_with_claude(pdfs):
    content = []
    # PDFs werden hier on-demand von Disk gelesen (statt vorher base64-Bytes
    # im RAM zu halten). Nach dem API-Call wird content[] gedroppt.
    for pdf in pdfs:
        b64 = pdf.get("base64")  # Backwards-compat
        if not b64 and pdf.get("path"):
            with open(pdf["path"], "rb") as fh:
                b64 = base64.b64encode(fh.read()).decode()
        if not b64:
            continue
        content.append({
            "type": "document",
            "source": {"type": "base64", "media_type": "application/pdf", "data": b64},
            "title": f"{pdf['folder']} / {pdf['name']}"
        })
    content.append({
        "type": "text",
        "text": (
            "Analysiere diese Immobilien-Dokumente aus dem Projektdatenraum und extrahiere alle relevanten Projektdaten. "
            "Es kann sein dass du Dokumente von mehreren Häusern (Haus A, B, C...) siehst - das ist ein Gesamtprojekt. "
            "Antworte NUR mit einem JSON-Objekt mit diesen Feldern: "
            "projektname_roh, adresse, stadt, stadtteil, plz, bautraeger, anzahl_haeuser, "
            "we_pro_haus, anzahl_we_gesamt, kfw_standard, energieversorgung, stellplaetze, "
            "groesse_von, groesse_bis, kaufpreis_ab, besonderheiten, planungsphase, "
            "we_typen_count, we_typen_liste, "
            "anzahl_1zi (Anzahl 1-Zimmer-Wohnungen), anzahl_2zi (Anzahl 2-Zimmer), anzahl_barrierefrei. "
            "\n\n"
            "⚠️ DATENQUALITÄT IST KRITISCH – durchsuche ALLE PDFs gründlich:\n"
            "- stellplaetze: Suche in Bauunterlagen/Beschreibungen nach Stellplatz-Anzahl. "
            "Zähle Außenstellplätze + Tiefgarage-Plätze + Garagen separat. "
            "Format: '13 Außenstellplätze' oder '24 TG + 6 außen' oder einfach Zahl '13'. "
            "NIE 'auf Anfrage' wenn Zahl irgendwo im Datenraum steht.\n"
            "- kfw_standard: Suche in Energieausweis, Baubeschreibung, Förderübersicht. "
            "Format: 'KfW-Effizienzhaus 40 QNG-PLUS' oder 'KfW 55' oder 'KfW-40 NH'. "
            "NIE 'auf Anfrage' wenn Standard erkennbar.\n"
            "- energieversorgung: Suche in TGA/Heizungs-PDFs. Format: "
            "'Fernwärme + Photovoltaik' oder 'Wärmepumpe + Solar' oder 'Gas-Brennwert'.\n"
            "- besonderheiten: Konkrete Highlights wie 'Alle Einheiten mit Balkon/Terrasse, "
            "vollmöbliert, Aufzug, barrierearm, Fernwärme'. KEINE Floskeln.\n"
            "- kaufpreis_ab: Aus Preisliste, Format als Zahl '189000' (ohne €/Punkte/Komma).\n"
            "- anzahl_we_gesamt: Aus WFL-Übersicht oder Vertragsstand.\n"
            "- anzahl_1zi / anzahl_2zi: Anzahl der jeweiligen Wohnungstypen (zähle aus WFL-Übersicht).\n"
            "- anzahl_barrierefrei: WEs mit Barrierefrei-Markierung im WFL-Plan.\n"
            "\n"
            "WICHTIG für bautraeger: Nur den exakten Firmennamen, OHNE Fußnotenzahlen oder Sonderzeichen. "
            "Beispiel: 'SBB Bauträgergesellschaft mbH' (nicht 'SBB Bauträgergesellschaft1 mbH').\n"
            "WICHTIG für projektname_roh: Nur der vermarktbare Projektname (1-3 Wörter). "
            "Falls kein expliziter Projektname in den Dokumenten: ERFINDE einen kurzen, "
            "kreativen Markennamen passend zu Lage + Produktart. "
            "Beispiele: 'The Rothenseer', 'New Living 72', 'Central Magdeburg', 'Die 72'. "
            "NIEMALS den Firmennamen als Projektname verwenden!\n"
            "we_typen_liste: Array von Objekten mit {bezeichnung, typ, wohnflaeche_qm} für jeden WE-Typ "
            "wenn aus WFL-PDFs erkennbar.\n"
            "Kein Text davor oder danach, keine Markdown-Backticks."
        )
    })
    resp = requests.post(
        "https://api.anthropic.com/v1/messages",
        headers={"x-api-key": CLAUDE_API_KEY, "anthropic-version": "2023-06-01", "content-type": "application/json"},
        json={"model": "claude-haiku-4-5-20251001", "max_tokens": 4000,
              "messages": [{"role": "user", "content": content}]},
        timeout=120
    )
    resp.raise_for_status()
    print(f"analyze_pdfs_with_claude: HTTP {resp.status_code}, stop_reason={resp.json().get('stop_reason')}")
    text = resp.json()["content"][0]["text"].replace("```json", "").replace("```", "").strip()
    try:
        result = json.loads(text)
        print(f"  Projektdaten extrahiert: {list(result.keys())}")
        return result
    except json.JSONDecodeError as e:
        print(f"  JSON-Fehler in analyze_pdfs: {e} – letzten 200 Zeichen: ...{text[-200:]}")
        raise

def generate_expose_with_claude(projektdaten, city_context=""):
    stadt = projektdaten.get('stadt', 'der Stadt')
    projekt = projektdaten.get('projektname_roh', 'das Projekt')
    bautraeger = projektdaten.get('bautraeger', 'urbanunits')
    we_typen_liste = projektdaten.get('we_typen_liste', [])

    # WE-Typen Info für Prompt aufbereiten
    we_typen_hint = ""
    if we_typen_liste:
        lines = []
        for i, t in enumerate(we_typen_liste[:6]):
            bez = t.get('bezeichnung', f'Typ {i+1}')
            typ = t.get('typ', '')
            wfl = t.get('wohnflaeche_qm', '')
            lines.append(f"  - {bez}{': ' + typ if typ else ''}{', ' + str(wfl) + ' m²' if wfl else ''}")
        we_typen_hint = (
            f"\nEXTRAHIERTE WE-TYPEN aus WFL-PDFs ({len(we_typen_liste)} Typen):\n"
            + "\n".join(lines)
            + "\nNutze diese Daten für we_beispiel_N, we_nummern_N, we_raum_*_name_N und we_flaeche_*_N!\n"
        )

    prompt = (
        "Du bist ein erfahrener Immobilien-Exposé-Texter bei INTERPRÉS GmbH. "
        "Antworte NUR mit einem validen JSON-Objekt. Kein Text davor oder danach. Keine Markdown-Backticks.\n"
        "⚠️ WICHTIG: Verwende '€' NUR für Geldbeträge (z.B. '189.000 €' oder '17 Mrd. €'). "
        "NIEMALS '€' als Silbe in deutschen Wörtern wie 'Europas', 'Europa', 'europäisch' – "
        "schreibe diese Wörter immer vollständig aus!\n\n"

        f"⚠️ ANTI-HALLUZINATION: Das Projekt liegt in **{stadt}**"
        + (f" (Stadtteil: {projektdaten.get('stadtteil')})" if projektdaten.get('stadtteil') else "")
        + ". "
        "Schreibe ALLE Stadt-Texte (Wirtschaft, Wachstum, Standort, Branchen) AUSSCHLIESSLICH ueber "
        f"{stadt} selbst. NIEMALS ueber andere Staedte oder Nachbarorte schreiben (z.B. nicht "
        "'Wolmirstedt liegt vor den Toren von ...'). Auch die Statistiken (Einwohner, BIP, "
        f"Mietsteigerung, Studierende) muessen sich auf {stadt} bzw. dessen Bundesland beziehen.\n\n"

        f"## PROJEKTDATEN\n{json.dumps(projektdaten, ensure_ascii=False)}\n"
        f"{we_typen_hint}\n"

        "## SCHREIBSTIL – DQN-STIL als VERBINDLICHE REFERENZ\n"
        "Das Exposé soll im EXAKT gleichen Stil wie ein Premium-DQN-Exposé wirken.\n\n"
        "🎯 ZIELGRUPPE — gleichgewichtig auf BEIDE Seiten ansprechen:\n"
        "  1. KAUF-INTERESSENT (Eigennutzer): Person, die einziehen will. Ansprechen mit:\n"
        "     Wohnqualität, Lifestyle, Stadtteil-Charakter, kurze Wege, Nachbarschaft, Komfort,\n"
        "     Atmosphäre, Tagesabläufe, Gefühl. Beispiel-Worte: 'wohnen', 'leben', 'genießen',\n"
        "     'Rückzug', 'Alltag', 'Zuhause', 'Atmosphäre', 'Stadtteil'.\n"
        "  2. KAPITALANLEGER: Person, die kaufen + vermieten will. Ansprechen mit:\n"
        "     Förderung, Rendite, AfA, Stabilität, Vermietbarkeit, Wertentwicklung. Beispiel-\n"
        "     Worte: 'Investment', 'Förderung', 'Rendite', 'Vermietung', 'Wertstabilität'.\n\n"
        "🚫 NICHT nur Kapitalanlage-Sprech! Das war ein Vorgängerproblem. Texte sollen sich\n"
        "   AUCH an Menschen richten die selbst einziehen wollen — egal auf welcher Seite.\n"
        "   Nur die expliziten Investment-Felder (text_investment_pitch, text_kapitel_invest_1/2,\n"
        "   text_grund_*) duerfen vorrangig anlegerorientiert sein. Alle anderen Texte balanciert.\n\n"
        "📏 LÄNGE: Die DQN-Texte sind LANG und narrativ (300-600 Zeichen pro Block, mehrere\n"
        "   Sätze, mit konkreten Fakten, Firmennamen, Zahlen, Quellen). NICHT 1-2 dünne Sätze.\n"
        "   Lieber detailliert + konkret als kurz + generisch.\n\n"
        "WICHTIG: Diese Struktur funktioniert für JEDE Stadt – nicht nur Magdeburg.\n"
        "Schreibe ECHTE Fakten zur jeweiligen Stadt, generischen Inhalten ist eine Absage zu erteilen.\n\n"

        "### DQN-Strukturreferenzen (übernehme Stil & Länge, ersetze nur die Stadt-Fakten):\n\n"

        "REFERENZ text_intro (Slide 3 'Moderne Apartments in {stadt}', ~200 Zeichen):\n"
        "  'Im »{projekt_titel}« stehen alle Zeichen auf zukunftsfähiges Bauen. Das nachhaltige "
        "Gesamtkonzept schafft eine grüne Oase inmitten des urbanen Umfelds der Landeshauptstadt – "
        "durchdacht, effizient und umweltbewusst.'\n\n"

        "REFERENZ text_investment_pitch (DQN: 'kleine Einstiegspreise + KfW + AfA-Vorteil'):\n"
        "  'Kleine Einstiegspreise, attraktive KfW-Förderung und dreifach-AfA bieten ideale "
        "Voraussetzungen für Kapitalanleger, die Wert auf Effizienz und Stabilität legen.'\n\n"

        "REFERENZ text_kapitel_invest_2 (Verbindung Investment + Lage + Effizienz, ~300 Zeichen):\n"
        "  'Die aufstrebende Lage in {stadt}, die energieeffiziente Bauweise (KfW-Effizienzhaus-Stufe 40 "
        "QNG PLUS) sowie die durchdachte Möblierung machen {projekt_titel} zu einem Investment, "
        "das heute überzeugt – und morgen relevant bleibt.'\n\n"

        "REFERENZ text_stadt_intro (Slide 'Magdeburg/Stadt – Die Lage'):\n"
        "  'Die Landeshauptstadt wächst. Der Stadtteil {stadtteil} ist heute einer der spannendsten "
        "Orte {stadt}s: gewachsen, urban, im Wandel. {projekt_titel} entsteht genau hier – zwischen "
        "{landmark_1} und {landmark_2}, zwischen Universität und Einkaufszentren. Eine Lage, die "
        "vieles verbindet: Nähe zur City, kurze Wege, grüne Rückzugsorte und ein hohes Maß an "
        "Lebensqualität.'\n\n"

        "REFERENZ Stadt-wirtschaft Subsektionen (DQN Slide 14 hat 3 Sektionen):\n"
        "Sektion 1 'Technologie und Industrie' (~280 Zeichen mit Firmennamen):\n"
        "  'Im {industriepark} entstehen neue Flächen für internationale Unternehmen aus den "
        "Bereichen Halbleiter, Batterietechnik, Rechenzentren und Pharma. Namen wie {firma_1} – "
        "{firma_1_zusatz} – oder {firma_2} unterstreichen die internationale Relevanz des Standorts. "
        "Insgesamt entstehen hier mehrere tausend neue Arbeitsplätze.'\n\n"
        "Sektion 2 'Infrastruktur \"Knoten {stadt}\"' (~280 Zeichen, Verkehr/Bahn):\n"
        "  'Parallel dazu investiert die Stadt massiv in ihre Infrastruktur. Der Bahnhof "
        "{stadt}-{stadtteil} wird im Rahmen des Großprojekts \"Knoten {stadt}\" umfassend modernisiert "
        "– mit einem Volumen von rund {summe} Millionen Euro. Neue S-Bahn-Anbindungen, bessere "
        "Taktungen und optimierte Logistikwege stärken die Verbindung zu regionalen und überregionalen "
        "Wirtschaftsräumen.'\n\n"
        "Sektion 3 'Logistiksektor / Hafen' (~280 Zeichen):\n"
        "  'Der {hafen_name} – {größenbeschreibung} – wird strategisch ausgebaut und entwickelt sich "
        "mit Neuansiedlungen wie {firma_3} und {firma_4} zu einem zentralen Logistikknoten im "
        "europäischen Netzwerk.'\n\n"

        "REFERENZ Stadt-Stat-Cards (4 Zahlen mit Detail-Texten):\n"
        "  '245.279 Einwohner' → '{stadt} wächst kontinuierlich und gewinnt als Wohn- und "
        "Wirtschaftsstandort zunehmend an Bedeutung.'\n"
        "  '78,4 Mrd. BIP' → 'Das BIP des Bundeslandes lag 2023 bei etwa 78,4 Milliarden Euro, "
        "und ist seit 2013 über 40 % gestiegen.'\n"
        "  '+31% Mietsteigerung' → 'Für Neubauwohnungen sind die Mietpreise in {stadt} seit 2017 "
        "um rund 31 % gestiegen.'\n"
        "  '21.000 Studierende' → 'sind an den Hochschulen und Universitäten {stadt}s eingeschrieben.'\n\n"

        "REFERENZ MedTech-Hotspot / Branchen-Detail-Slide (~340 Zeichen Hauptabsatz):\n"
        "  '{stadt} entwickelt sich zu einem führenden Standort für {hauptbranche}. Über {anzahl} "
        "Unternehmen der Branche, zahlreiche Start-ups sowie international renommierte "
        "Forschungseinrichtungen wie {institut_1}, {institut_2} und {institut_3} bündeln hier ihre "
        "Kompetenzen. Zukunftsweisende Projekte wie {projekt_1} zeigen: {stadt} ist nicht nur "
        "Standort, sondern Treiber medizinischer Innovationen.'\n\n"

        "REFERENZ text_kapitel_stay (Stay/Projekt-Eröffnung, DQN p.18):\n"
        "  'Ein Neubauprojekt nach KfW-Effizienzhaus-Stufe 40 mit QNG-Zertifizierung, das "
        "Nachhaltigkeit und urbane Lebensqualität vereint. Energieeffiziente Bauweise, Fernwärme, "
        "Photovoltaik und Gründach sorgen für geringe Betriebskosten und ein besseres Mikroklima. "
        "Kompakte, durchdacht geschnittene Mikroapartments und moderne Technik schaffen Wohnraum, "
        "der ressourcenschonend, zukunftssicher und attraktiv für Mieter ist – mitten in {stadt}.'\n\n"

        "### Slogans (text_kapitel_invest/live/stay/know/hotel – NUR die Hauptüberschrift):\n"
        "Kurze Phrasen mit Punkt. Maximal 3-4 Wörter. Beispiele:\n"
        "'feels like a hotel.'  'think green. live smart.'  'naturban.'  'work, life balance.'\n"
        "'designed to stay.'  'stilvoll. durchdacht.'  'simply more.'\n\n"
        "### Kapitelseiten-Bodytexte (text_kapitel_invest_1, text_kapitel_invest_2 etc.):\n"
        "Das sind die Fließtexte auf den Kapitel-Trennseiten (links neben dem Slogan).\n"
        "Stil: 2-3 prägnante Sätze, emotional, projekt-spezifisch. Kein Bullet-Point-Stil.\n"
        "Referenz text_kapitel_invest_1: 'Kleine Einstiegspreise, attraktive KfW-Förderung und "
        "dreifach-AfA bieten ideale Voraussetzungen für Kapitalanleger, die Wert auf Effizienz "
        "und Stabilität legen.'\n"
        "Referenz text_kapitel_invest_2: 'Die aufstrebende Lage, die energieeffiziente Bauweise "
        "sowie die durchdachte Möblierung machen [Projekt] zu einem Investment, das heute "
        "überzeugt – und morgen relevant bleibt.'\n"
        "Referenz text_kapitel_live_1: 'Ein Ort, an dem man das Leben in der Stadt in vollen Zügen "
        "genießen kann – ohne auf die Schönheit der Natur zu verzichten.'\n"
        "Referenz text_kapitel_stay_1: '[Projekt] steht für eine Wohnform, die den Alltag neu denkt: "
        "kompakt, hochwertig, durchdacht.'\n\n"

        "### Fließtexte (text_intro, text_investment_pitch, text_greenliving_*, text_ausstattung_detail):\n"
        "Maximal 2-3 Sätze. Prägnant, emotional, auf den Punkt. Kein Fließtext-Aufsatz.\n"
        "Beispiel text_intro: 'Mitten in [Stadtteil] entsteht [Projektname] – [WE-Anzahl] möblierte Apartments "
        "für Studierende, Berufstätige und Investoren. Modern ausgestattet, smart vernetzt und sofort bezugsfertig.'\n\n"

        "### Key-Facts (feature_N_label, amenity_N):\n"
        "amenity_1 bis amenity_9: Kurz, max 22 Zeichen, 1-3 Wörter. WICHTIG: Begriffe müssen "
        "BILDBAR sein – also Dinge die als Foto erkennbar sind.\n"
        "GUTE Beispiele: 'Dachterrasse', 'E-Bike-Station', 'Fahrradabstellplatz', "
        "'Solaranlage', 'Fernwärme', 'Smart-Lock', 'Concierge', 'Paketstation', "
        "'Tiefgarage', 'Aufzug', 'Gemeinschaftsraum', 'Innenhof', 'Spielplatz'\n"
        "SCHLECHT (zu abstrakt für Foto-Suche): 'Effiziente Grundrisse', "
        "'4 identische Häuser', 'Hochwertige Ausstattung', 'Moderne Bauweise', "
        "'Innovative Technik' – diese werden NICHT als Amenity verwendet, "
        "sondern in feature_N_label oder text_ausstattung.\n"
        "Wähle für die 9 Amenities AUSSCHLIESSLICH visuell darstellbare Features.\n\n"

        "### Entwicklername (entwickler_name, entwickler_name_gross):\n"
        f"entwickler_name: Kurzer Markenname des Entwicklers (max 15 Zeichen, kein 'GmbH/mbH/AG').\n"
        f"Falls der Bauträger '{bautraeger}' heißt und das lang/korporativ klingt:\n"
        "→ Nimm nur den prägnanten Teil (z.B. 'SBB' statt 'SBB Bauträgergesellschaft mbH')\n"
        "→ oder den Markennamen falls bekannt (z.B. 'Urban Units' statt 'Urbanunits GmbH')\n"
        "entwickler_name_gross: Großbuchstaben-Version davon (z.B. 'SBB' → 'SBB')\n\n"
        "### Zahlen (feature_N_zahl, min_*, stadtstatistiken):\n"
        "Nur die Zahl, kein Text. Fahrrad-/Gehminuten realistisch für die Stadt.\n\n"

        "### Stadttext (text_stadt_*, text_einwohner_detail etc.):\n"
        "Nutze ECHTES Wissen über die Stadt. Max. 1-2 Sätze pro Feld.\n"
        "Beispiel: 'Magdeburg wächst: 245.000 Einwohner, über 21.000 Studierende, +34% Mieten seit 2017.'\n\n"

        "### Investmenttext:\n"
        "Konkret in 2 Sätzen: Preis, KfW-Darlehen, AfA-Vorteil. Kein Fließtext.\n\n"

        "### Nachhaltigkeit:\n"
        "text_greenliving_1: max 2 Sätze über Fernwärme/Photovoltaik/Energiekonzept.\n"
        "text_greenliving_2: max 2 Sätze über Außenbereiche/Mobilität/Lebensqualität.\n\n"

        "### Ausstattung:\n"
        "text_ausstattung_detail: max 2 Sätze. Bodenbelag, Heizung, Balkone. Konkret.\n\n"

        "## ⚠️ ZEICHENLIMITS – Templatefelder haben feste Größen.\n"
        "Bleibe IM Limit, aber NUTZE den Platz aus. Premium-Exposé heißt: AUSFÜHRLICH und SPEZIFISCH.\n"
        "Generische 1-Satz-Antworten = sofort verworfen. Liefere DQN-Niveau (mit Firmennamen, Zahlen, Quellen).\n\n"
        "produkt_beschreibung: ZIEL 35-55 Zeichen — Cover-Untertitel mit Stadt-Bezug.\n"
        "  Steht direkt nach {{ANZAHL_WE}} auf der Cover-Slide. DQN-Beispiel:\n"
        "  'vollmoeblierte Design-Apartments im Herzen Magdeburgs' (53 chars).\n"
        "  Bei 32 WE wird daraus: '32 vollmoeblierte Design-Apartments im Herzen Magdeburgs · Prospektteil A'.\n"
        "  Format: 'TYP-Beschreibung (vollmoebliert/Mikro/Design-) Apartments im Herzen/Zentrum/Norden {stadt}s'\n"
        "  oder analog. NIE nur '1-2 Zi.' — das ist zu duenn fuer das Cover.\n"
        "text_kapitel_invest/live/stay/know/hotel (Slogan): max 40 Zeichen\n"
        "text_kapitel_invest_1/2, text_kapitel_live_1/2, text_kapitel_stay_1/2, text_kapitel_know_1/2:\n"
        "  ZIEL 200-280 Zeichen – 2-3 Sätze, projekt-spezifisch, narrativ (DQN-Niveau)\n"
        "text_hotel: max 40 Zeichen\n"
        "text_intro: ZIEL 500-700 Zeichen – mehrere Sätze, emotional + konkret, BEWOHNER+INVESTOR\n"
        "  DQN-Beispiel hat ~720 Zeichen Fließtext: 'Ein modernes, nachhaltiges Wohnquartier...\n"
        "  Hier wird nicht nur gewohnt – hier wird gelebt: Ideal für Studierende, Berufstätige,\n"
        "  Zweitwohnsitznutzer – und jeden, der flexibel wohnen möchte...'\n"
        "text_investment_pitch: ZIEL 480-580 Zeichen – DQN-Beispiel: 'Weil gute Investments heute\n"
        "  anders aussehen. urbanunits verbindet, was Kapitalanleger heute suchen: niedrige\n"
        "  Einstiegspreise, maximale staatliche Förderung, steuerliche Vorteile und eine klare\n"
        "  Vermietungsperspektive. Das Projekt wurde konsequent vom Nutzer her gedacht – aber\n"
        "  mit dem Fokus auf nachhaltigen Ertrag und langfristige Wertentwicklung. Die\n"
        "  Apartments sind kompakt, designorientiert und energetisch optimiert. Die Lage in\n"
        "  {stadt} sorgt für konstant hohe Nachfrage – bei Mietern wie bei Investoren.'\n"
        "text_greenliving_intro: max 90 Zeichen — DQN: 'Ein modernes, nachhaltiges Wohnquartier,\n"
        "  realisiert mit zukunftsweisender Planung, effizienten Strukturen und einem klaren\n"
        "  Fokus auf Qualität und Umweltbewusstsein.' (~150 chars). Wenn etwas länger OK.\n"
        "text_greenliving_1: ZIEL 380-500 Zeichen — LINKE Body-Spalte. DQN-Beispiel:\n"
        "  'Ein neuer Ort zum Leben und Wohlfühlen. Eine grüne Oase inmitten der Landeshauptstadt\n"
        "  Sachsen-Anhalts – gebaut für die Anforderungen von heute, mit Blick auf das Morgen.\n"
        "  Inmitten der Landeshauptstadt, auf dem südlichen Areal der ehemaligen Magdeburger\n"
        "  Diamant Brauerei, entsteht das Quartier. In begehrter Wohnlage im\n"
        "  aufstrebenden Stadtteil {stadtteil} gehen urbane Lebendigkeit und Rückzug ins Grüne\n"
        "  eine besondere Verbindung ein.'  (~470 chars).\n"
        "  HINWEIS: nutze ECHTE Stadtteil-/Areal-/Quartier-Geschichte aus den PDFs.\n"
        "text_greenliving_2: ZIEL 380-500 Zeichen — RECHTE Body-Spalte. DQN-Beispiel:\n"
        "  'Wo früher einmal Pferdekutschen hielten und das beliebte Diamant Bier gebraut wurde,\n"
        "  entstehen heute 104 moderne Mikroapartments – kompakt, hochwertig und zukunfts-\n"
        "  orientiert. Die nachhaltige Bauweise nach KfW-Effizienzhaus-Stufe 40 mit QNG PLUS,\n"
        "  ergänzt durch Photovoltaik, Fernwärme und ein Gründachkonzept, sorgt für einen\n"
        "  reduzierten ökologischen Fußabdruck. So entsteht zeitgemäßer Wohnraum auf höchstem\n"
        "  Niveau – {stadt}s grüner Diamant.'  (~430 chars).\n"
        "  HINWEIS: NUR Fakten aus den PDFs verwenden — keine erfundenen Geschichten.\n"
        "text_ausstattung_intro: max 90 Zeichen\n"
        "text_ausstattung_detail: ZIEL 280-380 Zeichen – Materialien, Boden, Bad, Smart-Home, narrativ\n"
        "text_ausstattung_kurz: max 80 Zeichen\n"
        "text_ausstattung_lang: ZIEL 280-380 Zeichen\n"
        "text_grundriss_intro: ZIEL 280-380 Zeichen, DQN-Stil mit konkreten Größen + Zielgruppen.\n"
        "text_architektur: max 130 Zeichen\n"
        "text_nachhaltig_1/2/3/4: FESTE Themen-Zuordnung (rendert in 4 'Einfach …'-Spalten):\n"
        "  text_nachhaltig_1 → 'Einfach NACHHALTIGER' — Dämmung/Effizienzhaus/Zertifizierung.\n"
        "    DQN: 'Gebaut nach KfW-Effizienzhaus-Stufe-40, QNG PLUS zertifiziert.' (60 chars)\n"
        "  text_nachhaltig_2 → 'Einfach GRÜNER' — Photovoltaik/Gründach/Bepflanzung/Mikroklima.\n"
        "    DQN: 'Photovoltaik, Gründach und klimaaktive Bepflanzung sorgen für ein besseres\n"
        "    Mikroklima und urbane Ruhe.' (115 chars)\n"
        "  text_nachhaltig_3 → 'Einfach EFFIZIENTER' — Energieversorgung/Smart-Home/E-Mobility.\n"
        "    DQN: 'Fernwärme, Smart-Home-Ready und vorbereitete Ladeinfrastruktur für\n"
        "    E-Mobility.' (90 chars)\n"
        "  text_nachhaltig_4 → 'Einfach STADTNÄHER' — Lage/Zielgruppen/Kompaktheit.\n"
        "    DQN: 'Kompakte Mikroapartments in zentraler Lage – ideal für Pendler, Studierende,\n"
        "    Senioren und urbane Singles.' (130 chars)\n"
        "  Limit pro Eintrag: max 130 Zeichen.\n"
        "  REIHENFOLGE STRENG einhalten — Slot-Position bestimmt das Thema!\n"
        "text_standort_1/2: ZIEL 280-380 Zeichen pro Eintrag\n"
        "text_projekt_nachhaltig_1/2: ZIEL 280-380 Zeichen\n"
        "text_stadt_intro: ZIEL 100-160 Zeichen — KURZER Headline-Satz, Lead-In für die Slide.\n"
        "  DQN: '{stadt} ist einer der Vorzeige-Standorte, der mitten im Wandel steckt, aber\n"
        "  einen Großteil erfolgreich hinter sich gebracht hat.' (~135 chars)\n"
        "  NICHT der lange beschreibende Stadt-Text — der kommt in die Wirtschaft-Felder.\n"
        "stadt_bezeichnung: 1 Wort, z.B. 'Landeshauptstadt' / 'Hansestadt' / 'Großstadt' /\n"
        "  'Universitätsstadt' — passend zur jeweiligen Stadt.\n"
        "text_stadt_wachstum_1: ZIEL 280-380 Zeichen – Branchenüberblick mit Firmennamen + Investitionssummen\n"
        "text_stadt_wachstum_2: max 240 Zeichen – konkrete Projekte/Investitionssummen\n"
        "text_stadt_wirtschaft_links: ZIEL 350-450 Zeichen — SUB-SEKTION 1 mit Mini-Header.\n"
        "  Format: '<HEADER>\\n\\n<BODY>'. HEADER: 2-4 Wörter, kein Punkt, z.B. 'Technologie und\n"
        "  Industrie'. BODY: 320-410 chars mit echten Firmen + Volumen.\n"
        "  DQN-Beispiel:\n"
        "  'Technologie und Industrie\\n\\nIm Industriepark Eulenberg entstehen neue Flächen für\n"
        "  internationale Unternehmen aus den Bereichen Halbleiter, Batterietechnik, Rechen-\n"
        "  zentren und Pharma. Namen wie CATL – der weltweit größte Hersteller von Lithium-\n"
        "  Ionen-Batterien – oder Mercury unterstreichen die internationale Relevanz des\n"
        "  Standorts. Insgesamt entstehen hier mehrere tausend neue Arbeitsplätze.'\n"
        "text_stadt_wirtschaft_rechts: ZIEL 350-450 Zeichen — SUB-SEKTION 2 mit Mini-Header.\n"
        "  Gleiche Struktur. Thema: Logistik/Hafen/Verkehr ODER Infrastruktur.\n"
        "  DQN-Beispiel:\n"
        "  'Logistiksektor\\n\\nDer Magdeburger Hafen – größter Binnenhafen Ostdeutschlands –\n"
        "  wird strategisch ausgebaut und entwickelt sich mit Neuansiedlungen wie Amazon und\n"
        "  HelloFresh zu einem zentralen Logistikknoten im europäischen Netzwerk. Die direkte\n"
        "  Anbindung an Wasser, Schiene und Straße macht ihn zu einem gefragten Standort für\n"
        "  moderne, multimodale Warenströme.'\n"
        "stadt_invest_titel: max 18 Zeichen — Highlight-Box-Titel mit Volumen.\n"
        "  DQN: 'Invest 3 Mrd.'\n"
        "stadt_invest_label: max 60 Zeichen — Untertitel mit Projekt + Volumen-Detail.\n"
        "  DQN: 'FMC-Chipfabrik | 3 Mrd. Euro Investition geplant'\n"
        "text_stadt_invest_detail: ZIEL 280-380 Zeichen — Highlight-Box-Body mit Investor + Was\n"
        "  + Wo + Wann.\n"
        "  DQN: 'Die Ferroelectric Memory Company (FMC), ein Dresdner Hightech-Unternehmen im\n"
        "  Bereich innovativer Speichertechnologien, plant einen Produktionsstandort auf rund\n"
        "  35 Hektar im Süden Magdeburgs. Ziel ist es, dort energieeffiziente Speicherchips für\n"
        "  den Einsatz in Künstlicher Intelligenz und Rechenzentren zu fertigen – mit\n"
        "  Produktionsstart bis Ende 2025.' (~360 chars)\n"
        "text_einwohner_detail: max 130 Zeichen – konkrete Einwohner-Entwicklung\n"
        "text_bip_detail: max 130 Zeichen – BIP-Entwicklung mit % über mehrere Jahre\n"
        "text_mietsteigerung_detail: max 110 Zeichen – Mietpreis-Tendenz mit €/m²\n"
        "text_studierende_detail: max 110 Zeichen – Hochschulen + Fachgebiete\n"
        "stadt_branche_titel: 'MedTech-Hotspot' / 'Tech-Hub' / 'Logistik-Drehkreuz' / 'Forschungs-\n"
        "  standort' — 1-2 Wörter, kategorisiert die wichtigste 2. Branche der Stadt.\n"
        "text_stadt_branche_1: ZIEL 380-500 Zeichen – HAUPTFLIESSTEXT der Branche-Slide. Mit\n"
        "  konkreten Forschungsinstituten + Firmen + Projekten. DQN-Beispiel:\n"
        "  '{stadt} entwickelt sich zu einem führenden Standort für Medizintechnik und\n"
        "  Neurowissenschaften. Über 30 Unternehmen der Branche, zahlreiche Start-ups sowie\n"
        "  international renommierte Forschungseinrichtungen wie das Leibniz-Institut für\n"
        "  Neurobiologie (LIN), das Deutsche Zentrum für Neurodegenerative Erkrankungen (DZNE),\n"
        "  das Fraunhofer IBMT und das Max-Planck-Institut bündeln hier ihre Kompetenzen.\n"
        "  Zukunftsweisende Projekte wie STIMULATE – ein vom Bund geförderter Forschungscampus\n"
        "  für bildgestützte Minimalinvasive Medizin – sowie das Zentrum für Neurowissen-\n"
        "  schaftliche Innovation und Technologie (ZENIT) zeigen: {stadt} ist nicht nur\n"
        "  Standort, sondern Treiber medizinischer Innovationen.' (~640 chars)\n"
        "text_stadt_branche_2: ZIEL 200-300 Zeichen – Ergänzungs-Block z.B. zu Klinikum/Uni.\n"
        "  DQN: 'Das Universitätsklinikum {stadt} ist mit 4.100 Mitarbeitern der größte\n"
        "  Arbeitgeber im Gesundheitswesen Nord-Sachsen-Anhalts und bildet über 1.000\n"
        "  Medizinstudierende aus.' (~225 chars)\n"
        "text_stadt_stat_N_detail: max 180 Zeichen pro Stat — DQN-Beispiele:\n"
        "  '1,06 Mrd. Investition in Uniklinikum' → 'Mit einem Investitionsvolumen von 1,06\n"
        "  Milliarden Euro markiert der Spatenstich im Oktober 2024 den Beginn eines der\n"
        "  bedeutendsten Bauprojekte für die Gesundheitsversorgung in der Region.' (~190 chars)\n"
        "  'Top 10 Krankenhäuser' → 'Das Klinikum {stadt} zählt laut Focus Klinikliste 2025 zu\n"
        "  den besten Krankenhäusern in Sachsen-Anhalt – bewertet unter 1.629 Kliniken mit\n"
        "  Ihren 14.383 Fachkliniken bundesweit.' (~205 chars)\n"
        "  '1.486 Studenten Humanmedizin' → 'Über 1.000 Medizinstudierende, mehr als 20\n"
        "  Forschungseinrichtungen – {stadt} gehört zu den Städten seiner Größenordnung mit\n"
        "  einer deutlich überdurchschnittlichen MedTech-Dichte.' (~190 chars)\n"
        "\n"
        "### 6 GUTE GRÜNDE (Slide 'gute Gründe' — 6 Punkte mit Titel + Beschreibung):\n"
        "Pro Grund EIN Titel (kurz, prägnant, EINZEILIG!) + EINE Beschreibung (1-2 Sätze).\n"
        "🚨 LAYOUT-KRITISCH: Titel MUSS auf eine Zeile passen, Text MUSS in 2 Zeilen passen.\n"
        "Sonst überlappt die Beschreibung mit dem nächsten Punkt darunter.\n"
        "Stil DQN: Konkrete Zahlen, Förderdetails, klarer Vorteil. Mix Käufer + Anleger.\n"
        "Felder:\n"
        "  text_grund_1_titel + text_grund_1_text  → Lage/Erreichbarkeit. DQN-Beispiel:\n"
        "    Titel: 'Zentrale Lage in {stadt}'  Text: '3 Minuten zur Universität, 8 Minuten in\n"
        "    die City, 3 Minuten ins Grüne – perfekte Voraussetzungen für dauerhafte Vermietbarkeit.'\n"
        "    Nutze die ECHTEN Min-Werte aus min_uni/min_bahnhof/min_altstadt.\n"
        "  text_grund_2_titel + text_grund_2_text  → Kaufpreis/Einstieg. DQN-Beispiel:\n"
        "    Titel: 'Kaufpreise ab {kaufpreis_ab} €'  Text: 'Voll förderfähiges Neubauprojekt\n"
        "    mit geringer Einstiegsschwelle. Ab {kaufpreis_ab} € investieren Sie in nachhaltige,\n"
        "    möblierte Apartments.'\n"
        "  text_grund_3_titel + text_grund_3_text  → KfW-Förderung. DQN-Beispiel:\n"
        "    Titel: 'KfW-Niedrigzins-Darlehen'  Text: 'von bis zu {kfw_darlehen} € je Einheit.\n"
        "    Kapitalanleger profitieren von attraktiven Zinssätzen, und tilgungsfreien Anlaufjahren.'\n"
        "  text_grund_4_titel + text_grund_4_text  → AfA/Steuern. DQN-Beispiel:\n"
        "    Titel: '3-fach Abschreibung'  Text: '5 % degressive AfA + 5 % Sonder-AfA + separate\n"
        "    Abschreibung der Möbel. Maximale Steuerwirkung für Kapitalanleger.'\n"
        "  text_grund_5_titel + text_grund_5_text  → Möblierung/Smart-Home. DQN-Beispiel:\n"
        "    Titel: 'Möblierungskonzept'  Text: 'Individuell gestaltete Apartments im Westwing-\n"
        "    Design, kombiniert mit modernster Smart-Home-Technologie.'\n"
        "  text_grund_6_titel + text_grund_6_text  → Mietgarantie/Vermietbarkeit. DQN-Beispiel:\n"
        "    Titel: 'Mietgarantie für 3 Monate'  Text: 'Nach Fertigstellung garantiert –\n"
        "    20 €/m² Kaltmiete, für drei Monate.'\n"
        "Wenn ein Punkt fuer das konkrete Projekt nicht passt (z.B. keine Mietgarantie):\n"
        "Titel + Text trotzdem ausfuellen mit projekt-spezifischer Alternative (z.B. 'Premium-\n"
        "Lage', 'Fertigstellungstermin', 'Energiekonzept').\n"
        "📏 HARTE LIMITS:\n"
        "  Titel: max 32 Zeichen — MUSS einzeilig sein. DQN-Beispiele:\n"
        "    'Zentrale Lage in Magdeburg' (26)\n"
        "    'Kaufpreise ab 185.000 €'    (23)\n"
        "    'KfW-Niedrigzins-Darlehen'   (24)\n"
        "    '3-fach Abschreibung'        (19)\n"
        "    'Möblierungskonzept'         (18)\n"
        "    'Mietgarantie für 3 Monate'  (25)\n"
        "  Text:  100-150 Zeichen — MUSS in 2 Zeilen passen.\n"
        "  Wenn dein Text laenger wird: kuerzen, nicht aufweiten!\n\n"

        "### KURZ-EXPOSÉ FELDER (separates 2-Seiten-PDF, ZUSAETZLICH zum Marketing-Expose).\n"
        "Diese Felder bedienen ein Kompakt-Format — aehnliche Tonalitaet wie Marketing,\n"
        "ABER deutlich knapper + zugespitzter. Pflicht zu liefern.\n"
        "⚠️ STILREGEL: Die unten gezeigten DQN-Saetze sind REINE STIL-INSPIRATION —\n"
        "NICHT abschreiben, NICHT umformulieren, NICHT als Vorlage benutzen. Eigener\n"
        "Wortlaut, eigene Satzstruktur, eigene Bilder. Wenn der DQN-Satz 'Urbaner Wohnkomfort\n"
        "in nachhaltigem Umfeld' lautet, schreib NICHT 'Urbaner Lebenskomfort in nachhaltigem\n"
        "Quartier' — das ist 1:1 mit Synonymen. Schreib stattdessen etwas frisches wie\n"
        "'Modern wohnen, smart vermieten' oder 'Nah am Puls, fernab vom Trubel'.\n\n"
        "  projekt_untertitel: max 60 Zeichen — Tagline. Format frei: Halbsatz, Slogan oder\n"
        "    zwei Begriffe, die das Projekt-Gefuehl tragen. Keine Zahlen, keine Floskel-\n"
        "    Adjektive ('hochwertig', 'modern', 'attraktiv'). Beispiele fuer Tonalitaet\n"
        "    (NICHT kopieren!): 'Designed to stay.', 'Wohnen wie geplant.', 'Smart. Klein. Komplett.'\n"
        "  projekt_beschreibung: ZIEL 420-560 Zeichen — DER Hauptpitch-Text auf Kurz-S2,\n"
        "    grosser Block links neben dem Hauptbild. 3-4 ausfuehrliche Saetze. Struktur:\n"
        "      Satz 1: WO + WAS — Stadt, Lage/Stadtteil, Projektname, Anzahl WE,\n"
        "              Kernfeature (z.B. 'Mikro-Living mit durchdachter Ausstattung').\n"
        "      Satz 2: Architektur + Nachhaltigkeit — KfW-Standard (z.B. 'KfW-Effizienz-\n"
        "              haus-Stufe 40 Klimafreundlicher Neubau'), QNG-Plus, Materialien.\n"
        "      Satz 3: Quartier/Umfeld-Mehrwert — was macht die Lage besonders (Spiel-\n"
        "              platz, Gruenanlagen, Gastro, Nachbarschaft).\n"
        "      Satz 4 (optional): Investitions-/Nutzungs-Aspekt — Eigennutzung vs.\n"
        "              Kapitalanlage, Vermietungsperspektive.\n"
        "    Tonalitaet: faktenstark mit Lebensgefuehl. Konkrete Zahlen einbauen.\n"
        "    NICHT mit 'Im Herzen des neuen ... entsteht ...' starten — das ist die DQN-Phrase.\n"
        "    Variiere den Satzanfang: 'In {stadt} entsteht ...', 'Mit [anzahl_we] Apartments\n"
        "    setzt [projekt_titel] ...', '[projekt_titel] bringt ... nach {stadt}.'.\n"
        "  projekt_beschreibung_kurz: gleicher Inhalt wie projekt_beschreibung — Backwards-\n"
        "    Compat-Alias. Nimm einfach denselben Text.\n"
        "  text_relevanz: ZIEL 60-75 Zeichen, ABSOLUTE OBERGRENZE 75 — passt damit sauber auf\n"
        "    EXAKT 2 Zeilen in der Icon-Spalte. Konkreter Markt-Bezug, keine Buzzwords.\n"
        "    Statt 'im Trend von Mikroapartments': nimm Markt-Treiber wie Studierende,\n"
        "    Pendler, Single-Haushalte, Mietnachfrage in {stadt}.\n"
        "  text_design: ZIEL 60-75 Zeichen, max 75 — 2-Zeiler USP DESIGN/WOHNEN. Nimm konkrete\n"
        "    Materialien aus dem Datenraum (Eiche, Fliesen, Walk-in-Dusche, Kueche-Marke).\n"
        "  text_foerderung: ZIEL 60-75 Zeichen, max 75 — 2-Zeiler USP FOERDERUNG/STEUER. Liste\n"
        "    die echten Foerderbausteine (KfW-Stufe, QNG, Sonder-AfA, etc).\n"
        "  text_tech: ZIEL 60-75 Zeichen, max 75 — 2-Zeiler USP TECH/SMART-HOME. Konkrete Tech:\n"
        "    Smart-Lock, Glasfaser, PV, Waermepumpe, E-Ladesaeulen.\n"
        "  cover_stichwort_1..4: 4 kurze Stichworte fuer die Cover-Box (Slide 1 Marketing),\n"
        "    je max 22 Zeichen, GROSSBUCHSTABEN. DQN-Stil: 1-3 Worte pro Stichwort.\n"
        "    Reihenfolge: 1) Wohnform/Ausstattung (z.B. 'VOLLMOEBLIERT', 'KOMPAKT-WOHNEN'),\n"
        "    2) Foerder-/Steuer-Status (z.B. 'FOERDERFAEHIG', 'SONDER-AFA'),\n"
        "    3) Energie-/KfW-Stufe (z.B. 'KFW-40 QNG PLUS', 'KFW-55 EE'),\n"
        "    4) Tech-Highlight (z.B. 'SMART HOME READY', 'GLASFASER & PV').\n"
        "    KEIN Lauftext, KEINE Saetze — pure Schlagworte.\n"
        "  besonderheiten_liste: GENAU 6 Stichpunkte im DQN-Stil (NICHT mehr als 6 — sonst\n"
        "    laeuft die Box ueber). JEDER Punkt auf eigener\n"
        "    Zeile, beginnend mit ':  ' (Doppelpunkt + 2 Leerzeichen) als Marker.\n"
        "    Trenner zwischen Zeilen: EINFACHER Zeilenumbruch (\\n). Spacing wird\n"
        "    automatisch zwischen den Zeilen gesetzt — KEINE doppelten Newlines.\n"
        "    \n"
        "    ⚠️ DQN-STRUKTUR (GENAU 6 Punkte, je kurz, max ~50 Zeichen pro Punkt):\n"
        "      Punkt 1: Aussenbereich (Balkon, Terrasse, Dachterrasse)\n"
        "      Punkt 2: Erschliessung (barrierearm, Aufzug)\n"
        "      Punkt 3: Quartier-Feature (Spielplatz, Gemeinschaftsbereich, Gruenanlage)\n"
        "      Punkt 4: Anschluss/Tech (Glasfaser, Smart-Lock)\n"
        "      Punkt 5: Energie (Fernwaerme, Photovoltaik, Waermepumpe)\n"
        "      Punkt 6: Innenausstattung (Eichenparkett, Fussbodenheizung, Designermoebel)\n"
        "    \n"
        "    Beispiel-Format (Format zeigen, Inhalt aus Projekt-Daten!):\n"
        "      ':  alle Wohnungen mit Balkon oder Dachterrasse\\n"
        ":  barrierearm, Aufzug in alle Etagen\\n"
        ":  quartierseigener Spielplatz\\n"
        ":  Glasfaseranschluss und Smart-Lock\\n"
        ":  nachhaltige Fernwaerme fuer Heizung\\n"
        ":  Eichenparkett, Fussbodenheizung, Designermoebel'\n"
        "    \n"
        "    Konkrete Features aus dem Datenraum nehmen — KEINE Buzzwords ohne Inhalt.\n"
        "    Gesamtlaenge inkl. Marker + Newlines: 250-380 Zeichen.\n"
        "  gesamtwohnflaeche: berechne anzahl_we * Mittel(groesse_von, groesse_bis) und\n"
        "    formatiere mit Tausender-Punkt + Komma + ' m²', z.B. '3.741,58 m²'. Wenn die\n"
        "    Summe explizit im Datenraum steht (WFL-Berechnung), die echte Zahl nehmen.\n"
        "  zimmer_anzahl_min / zimmer_anzahl_max: Zahlen aus zimmer_typen ableiten\n"
        "    (z.B. '1-Zimmer und 2-Zimmer' → min='1', max='2'). Reine Ziffern, keine\n"
        "    Einheit, keine Worte.\n"
        "  bild_titel_1..6: leer lassen — Kunde laedt Cover-Collage-Bilder selbst hoch.\n"
        "  bild_kurz_1..4: leer lassen — Kunde laedt Seite-2-Bilder selbst hoch.\n\n"

        "WICHTIG: alle obigen Zeichen-Limits sind HARTE Obergrenzen. Lieber 1-2 starke Sätze\n"
        "als ein dritter Halbsatz der das Layout sprengt.\n"
        "feature_N_label: max 28 Zeichen\n"
        "amenity_N: max 22 Zeichen, 1-3 Wörter, BILDBAR (siehe Key-Facts-Liste oben)\n"
        "we_typ_beschreibung_N: max 50 Zeichen, 1 Zeile (DQN-Stil: '1-Zi mit Balkon, Barrierefrei')\n"
        "besonderheiten: max 80 Zeichen\n"
        "steuerliche_moeglichkeiten: max 110 Zeichen\n"
        "quartier_history: ZIEL 260-340 Zeichen – Geschichte/Charakter des Stadtteils\n"
        "quartier_ref: max 130 Zeichen\n"
        "zitat_intro: max 160 Zeichen\n"
        "quelle_1/2/3/4: kurz wie 'statistik.sachsen-anhalt.de 2024' (max 80 Zeichen)\n\n"

        f"## STANDORT-MINUTEN ({stadt} – Slide 5):\n"
        f"min_uni / label_min_uni: Fahrradminuten + Name der nächsten Uni/FH in {stadt}\n"
        f"min_bahnhof / label_min_bahnhof: Fahrradminuten + Hauptbahnhof\n"
        f"min_altstadt / label_min_altstadt: Fahrradminuten + Altstadt/Innenstadt\n"
        f"WICHTIG: 'min_*'-Felder nur die Zahl, z.B. '3'. 'label_min_*' nur den Namen, z.B. 'Leibniz Universität'.\n\n"

        f"## FREIZEIT NAH ({stadt} – Slide 14, 4 Einträge):\n"
        f"freizeit_N_name: Die WICHTIGSTEN/BEKANNTESTEN Freizeit-Highlights in {stadt} – "
        f"keine obskuren Lokal-Tipps, sondern Sehenswürdigkeiten/Attraktionen die JEDER kennt:\n"
        f"  - Magdeburg → Elbauenpark, Magdeburger Dom, Hundertwasserhaus, Elbufer-Promenade, "
        f"Allee-Center, Jahrtausendturm\n"
        f"  - Andere Städte: Hauptpark, Hauptdom/Kirche, Hauptmuseum, Hauptpromenade, "
        f"Hauptstadion, beliebter See/Strand\n"
        f"min_freizeit_N: REALISTISCHE Minutenangabe – wenn ein Highlight zu weit weg ist (>30 Min "
        f"zu Fuß), nimm es NICHT (lieber näheres Highlight wählen). Format: nur Zahl, z.B. '8'\n"
        f"  Faustregel: Fuß-Minuten = (km × 12), Fahrrad = (km × 4)\n\n"

        f"## WOHNUNGSTYPEN (Template v19 – 1 Typ pro Slide, 2 Beispiel-WEs nebeneinander):\n"
        f"Analysiere alle WFL-Berechnungs-PDFs und Grundrisse.\n"
        f"WICHTIG: Pro Wohnungstyp wird EIN Slide erzeugt. Die zwei Spalten links/rechts auf jedem\n"
        f"Slide zeigen ZWEI Beispiel-WEs DESSELBEN Typs (z.B. links 'WE 02', rechts 'WE 09' –\n"
        f"beide Studios mit gleicher Aufteilung). Beide Seiten teilen sich Raumflächen + Beschreibung.\n\n"
        f"Flächen-Format: '23,99 m²' (Komma, immer mit ' m²'). Lies echte Werte aus den WFL-PDFs!\n"
        f"we_flaeche_5 = Gesamtfläche der Wohnung (Wohnfläche total)\n\n"
        f"⚠️ we_typ_beschreibung: SEHR KURZ im DQN-Stil!\n"
        f"NUR 1 Zeile, max 50 Zeichen. Stil:\n"
        f"  '1-Zimmer-Wohnung mit Balkon, Barrierefrei'\n"
        f"  '2-Zimmer-Wohnung mit Terrasse, Bad rechts'\n"
        f"  '1-Zimmer-Studio mit großem Balkon'\n"
        f"NICHT: 'Kompaktes EG-Apartment mit 23,34 m² – klarer Grundriss mit offenem...'\n"
        f"Nur Zimmeranzahl + Hauptfeature + Besonderheit. KEINE Quadratmeter im Text "
        f"(stehen schon in der Tabelle), KEINE Werbephrasen wie 'renditestark', 'modern'.\n\n"

        f"### TYP 1 (immer befüllt, Slide 1):\n"
        f"- we_beispiel_1: WE-Nummer des LINKEN Beispiels, z.B. 'WE 02'\n"
        f"- we_beispiel_2: WE-Nummer des RECHTEN Beispiels (gleicher Typ!), z.B. 'WE 09'\n"
        f"- we_bereich_1: Beschreibung links, z.B. 'EG, 23,34 m²' oder '1-Zi., Balkon'\n"
        f"- we_bereich_2: Beschreibung rechts, z.B. 'OG, 23,34 m²'\n"
        f"- we_flaeche_1 bis we_flaeche_5: Raumflächen (geteilt – beide Spalten zeigen dasselbe!)\n"
        f"  Reihenfolge entspricht Template: Wohnen/Kochen, Schlafen, Bad, Abstellraum, Balkon (oder Gesamt für 5)\n"
        f"- we_typ_beschreibung: KURZ, max 50 Zeichen (siehe Beispiele oben)\n\n"

        f"### TYP 2 (wenn ≥2 Wohnungstypen) – wird auf neuem dupliziertem Slide angezeigt:\n"
        f"- we_beispiel_3 / we_beispiel_4: Zwei Beispiel-WEs des Typ 2\n"
        f"- we_bereich_3 / we_bereich_4\n"
        f"- we_flaeche_1_typ2 bis we_flaeche_5_typ2: Raumflächen für Typ 2\n"
        f"- we_typ_beschreibung_typ2\n\n"

        f"### TYP 3 (wenn ≥3 Typen): Suffix _typ3 analog. TYP 4: _typ4. Etc bis TYP 8.\n"
        f"Felder leer lassen ('') wenn dieser Typ nicht existiert.\n\n"

        f"Beispiel für 4 WE-Typen:\n"
        f"- Typ 1 (Studio EG): we_beispiel_1='WE 01', we_beispiel_2='WE 09'\n"
        f"- Typ 2 (Studio OG): we_beispiel_3='WE 02', we_beispiel_4='WE 10'\n"
        f"- Typ 3 (2-Zi EG):   we_beispiel_5='WE 03', we_beispiel_6='WE 11'\n"
        f"- Typ 4 (2-Zi OG):   we_beispiel_7='WE 04', we_beispiel_8='WE 12'\n\n"

        + (
        f"## AKTUELLE RECHERCHE FÜR {stadt.upper()} (Web-Suche heute, Stand 2024/2025):\n"
        f"Diese Recherche enthält ZUSAMMENFASSUNGEN, ARTIKEL-SNIPPETS und QUELLEN-URLs.\n"
        f"DEINE PFLICHT:\n"
        f"1. Lies ALLE Sektionen.\n"
        f"2. Extrahiere für die Stadttexte KONKRETE FAKTEN: echte Firmennamen (z.B. 'Intel', "
        f"'CATL', 'Amazon', 'HelloFresh'), benannte Projekte (z.B. 'Knoten Magdeburg 400 Mio €', "
        f"'Industriepark Eulenberg'), exakte Zahlen mit Jahresangabe ('245.279 Einwohner 2024').\n"
        f"3. KEINE generischen Floskeln wie 'wachsende Wirtschaft' oder 'attraktive Lage' - immer mit "
        f"einem konkreten Beleg (Firma, Projekt, Zahl, Investitionssumme).\n"
        f"4. Fülle quelle_1, quelle_2, quelle_3, quelle_4 mit den TATSÄCHLICHEN URLs aus der Recherche "
        f"(NICHT erfundene oder Beispiel-URLs!). Format: kurz wie 'statistik.sachsen-anhalt.de 2024' "
        f"oder 'IHK Magdeburg 2024' – aber die URL aus der Recherche muss dahinterstehen können.\n"
        f"5. text_stadt_branche_1/2 und text_einwohner_detail/text_bip_detail/text_mietsteigerung_detail/"
        f"text_studierende_detail MÜSSEN spezifisch sein – mindestens EINE genannte Firma, EIN benanntes "
        f"Projekt oder EINE konkrete Zahl pro Feld.\n"
        f"6. WIEDERHOLE NICHTS – jeder Abschnitt nennt andere Fakten. Nicht überall 'Intel investiert "
        f"17 Mrd €' wiederholen.\n\n"
        f"--- RECHERCHE-DATEN ---\n"
        f"{city_context}\n"
        f"--- ENDE RECHERCHE ---\n\n"
        if city_context else ""
        )
        + f"## STADTSTATISTIKEN ({stadt}):\n"
        f"Verwende echte, aktuelle Zahlen für {stadt}:\n"
        f"stadt_einwohner: Einwohnerzahl als formatierte Zahl, z.B. '245.279'\n"
        f"bundesland_bip: BIP des Bundeslandes NUR als Zahl+Einheit OHNE 'EUR'/'Euro', z.B. '310 Mrd.' oder '78,4 Mrd.'\n"
        f"  (Das Template-Label schreibt 'in €' bereits dahinter – niemals doppelt!)\n"
        f"stadt_mietsteigerung: Mietsteigerung des Mietniveaus seit 2017/2018, z.B. '+31%'\n"
        f"stadt_studierende: Studierende an Hochschulen, z.B. '21.000'\n"
        f"stadt_bip: BIP der Stadt/Region als formatierte Zahl (optional, falls vorhanden)\n\n"
        + "\n"
        f"## REFERENZ-BEISPIELE (Premium-Niveau wie DQN-Exposé):\n"
        f"So sieht GUTER Stadttext aus – konkrete Firmen, Projekte, Zahlen, Quellen:\n\n"
        f"text_stadt_intro (Magdeburg-Beispiel, ~330 Zeichen):\n"
        f"  'Die Landeshauptstadt wächst. Der Stadtteil Neue Neustadt ist heute einer der spannendsten "
        f"Orte Magdeburgs: gewachsen, urban, im Wandel. {{projekt_name}} entsteht genau hier – zwischen "
        f"Elbufer und Altstadt, zwischen Universität und Einkaufszentren. Eine Lage, die vieles verbindet: "
        f"Nähe zur City, kurze Wege, grüne Rückzugsorte und ein hohes Maß an Lebensqualität.'\n\n"
        f"text_stadt_wirtschaft_links (~280 Zeichen, Sektor mit Firmennamen):\n"
        f"  'Im Industriepark Eulenberg entstehen neue Flächen für internationale Unternehmen aus den "
        f"Bereichen Halbleiter, Batterietechnik, Rechenzentren und Pharma. Namen wie CATL – der weltweit "
        f"größte Hersteller von Lithium-Ionen-Batterien – oder Mercury unterstreichen die Relevanz des Standorts.'\n\n"
        f"text_stadt_wirtschaft_rechts (~280 Zeichen, anderer Sektor):\n"
        f"  'Der Magdeburger Hafen – größter Binnenhafen Ostdeutschlands – wird strategisch ausgebaut "
        f"und entwickelt sich mit Neuansiedlungen wie Amazon und HelloFresh zu einem zentralen Logistikknoten "
        f"im europäischen Netzwerk.'\n\n"
        f"text_stadt_invest_detail (~260 Zeichen, eine konkrete Großinvestition):\n"
        f"  'Die Ferroelectric Memory Company (FMC), ein Dresdner Hightech-Unternehmen, plant einen "
        f"Produktionsstandort auf rund 35 Hektar im Süden Magdeburgs – mit 3 Mrd. € Investition für "
        f"energieeffiziente Speicherchips für KI und Rechenzentren.'\n\n"
        f"text_einwohner_detail (~150 Zeichen):\n"
        f"  'Magdeburg wächst kontinuierlich und gewinnt als Wohn- und Wirtschaftsstandort zunehmend an Bedeutung.'\n\n"
        f"DAS IST DAS NIVEAU. Wende es auf {stadt} an mit den ECHTEN Daten aus der Recherche oben.\n"
        f"Erfinde NICHTS – nutze nur, was wirklich in der Recherche steht.\n\n"
        + (
        f"## BEKANNTE FAKTEN FÜR MAGDEBURG (zusätzlich, falls Recherche unvollständig):\n"
        f"  stadt_einwohner: '245.278'\n"
        f"  bundesland_bip: '73,4 Mrd.' (Sachsen-Anhalt)\n"
        f"  stadt_mietsteigerung: '+28%' (seit 2017)\n"
        f"  stadt_studierende: '21.000'\n"
        f"  Intel-Chip-Werksinvestition (~17 Mrd. €) sowie der Industriepark Eulenberg mit CATL/Mercury sind\n"
        f"  Schlüsselfakten. Nutze sie konkret in den Stadttexten.\n\n"
        if 'magdeburg' in stadt.lower() else ""
        )
        + f"## ⚠️ KEIN 'auf Anfrage' / KEINE Floskeln:\n"
        f"Felder wie kaufpreis_ab, kfw_darlehen, kfw_standard, stellplaetze, energieversorgung,\n"
        f"steuerliche_moeglichkeiten, besonderheiten MÜSSEN aus den PDFs extrahiert werden.\n"
        f"Wenn ein Wert wirklich nicht im Datenraum steht, lass das Feld LEER ('') – schreibe\n"
        f"NIE 'auf Anfrage' oder 'tba' oder Platzhalter wie 'XXX'.\n"
        f"❌ EXPLIZIT VERBOTEN als Wert (egal in welchem Feld):\n"
        f"   'nicht erkennbar', 'nicht ermittelbar', 'nicht angegeben', 'nicht bekannt',\n"
        f"   'nicht verfügbar', 'k.A.', 'k. A.', 'n/a', 'N/A', '-', '—', 'unbekannt',\n"
        f"   'wird nachgereicht', 'siehe Anlage', 'fehlt im Datenraum'.\n"
        f"Wenn nicht klar extrahierbar: LEEREN String '' liefern. NIE einen Pseudo-Wert.\n"
        f"Bei kfw_standard, energieversorgung, stellplaetze: lieber konservativ aus dem\n"
        f"Datenraum kombinieren (z.B. wenn QNG erwähnt → 'KfW-Effizienzhaus 40 QNG PLUS')\n"
        f"als leer lassen — aber NIE 'nicht erkennbar' o.ae. schreiben.\n\n"

        f"## ALLE FELDER – PFLICHT:\n"
        f"Jedes Feld MUSS befüllt werden. Leere Strings sind nicht akzeptabel außer bei\n"
        f"we_beispiel_N/we_nummern_N/we_raum_*_N/we_flaeche_*_N/we_typ_beschreibung_N für nicht vorhandene WE-Typen.\n\n"
        f"{json.dumps(PLATZHALTER, ensure_ascii=False)}"
    )
    resp = requests.post(
        "https://api.anthropic.com/v1/messages",
        headers={"x-api-key": CLAUDE_API_KEY, "anthropic-version": "2023-06-01", "content-type": "application/json"},
        json={
            "model": "claude-sonnet-4-6", "max_tokens": 16000,
            "messages": [{"role": "user", "content": prompt}]
        },
        timeout=300
    )
    resp.raise_for_status()
    resp_data = resp.json()
    stop_reason = resp_data.get("stop_reason", "unknown")
    print(f"generate_expose_with_claude: stop_reason={stop_reason}")

    json_text = ""
    for block in resp_data.get("content", []):
        if block.get("type") == "text":
            json_text = block["text"]

    json_text = json_text.replace("```json", "").replace("```", "").strip()
    if not json_text:
        raise ValueError(f"Claude hat keinen Text zurückgegeben. stop_reason={stop_reason}")

    # ── Text-Korrekturen: Claude-Fehler bei der Verwendung von Sonderzeichen ──
    # Claude schreibt manchmal "€pas" statt "Europas" (verwechselt € mit "Euro-").
    # Auch IIntel/IIm-Doppelbuchstaben sind regelmäßige Halluzinationen.
    _TEXT_FIXES = [
        ("€pas",       "Europas"),
        ("€pa ",       "Europa "),
        ("€päisch",    "europäisch"),
        ("€pä",        "europä"),
        ("€pe",        "Europe"),
        ("IIntel",     "Intel"),
        ("IIm ",       "Im "),
        ("IIn ",       "In "),
        ("IIst ",      "Ist "),
        ("  ",         " "),
    ]
    for wrong, right in _TEXT_FIXES:
        json_text = json_text.replace(wrong, right)

    # Wenn Claude durch max_tokens abgeschnitten wurde → JSON reparieren
    if stop_reason == "max_tokens":
        print("  WARNUNG: Antwort durch max_tokens abgeschnitten – versuche JSON zu reparieren")
        # Finde letztes vollständiges Key-Value-Paar und schließe JSON
        last_comma = json_text.rfind('",')
        last_quote = json_text.rfind('"')
        cut = max(last_comma, last_quote)
        if cut > 0:
            json_text = json_text[:cut + 1]
            # Schließe alle offenen Objekte/Arrays
            depth = json_text.count('{') - json_text.count('}')
            json_text += '}' * max(depth, 1)
        print(f"  Repariertes JSON (letzte 100 Zeichen): ...{json_text[-100:]}")

    try:
        result = json.loads(json_text)
    except json.JSONDecodeError as e:
        print(f"  JSON-Parse-Fehler: {e}")
        print(f"  Letzten 500 Zeichen: ...{json_text[-500:]}")
        # Fallback: gib leeres Dict zurück, PLATZHALTER-Defaults werden benutzt
        print("  Fallback: verwende leere Felder (PLATZHALTER-Defaults)")
        return {}

    # Sicherheitsnetz: dieselben Text-Fixes auch auf jedem geparsten String-Wert.
    # Falls Claude die Halluzination erst beim Schreiben ins JSON-Feld macht,
    # wird sie hier nochmal abgefangen.
    fixed_count = 0
    for k, v in list(result.items()):
        if not isinstance(v, str):
            continue
        new_v = v
        for wrong, right in _TEXT_FIXES:
            new_v = new_v.replace(wrong, right)
        if new_v != v:
            result[k] = new_v
            fixed_count += 1
    if fixed_count:
        print(f"  Text-Fix nach Parse: {fixed_count} Felder bereinigt (€pas, IIm, etc.)")

    # Anti-Pseudo-Wert: Claude schreibt manchmal trotz Prompt-Anweisung
    # 'nicht erkennbar', 'k.A.', 'n/a' o.ae. — das landet sichtbar im PDF
    # ('Stellplaetze: nicht erkennbar'). Lieber LEER → Slot zeigt '{{...}}'
    # was das v2-Editor-UI dann manuell befuellbar macht, oder wird via
    # leere-Platzhalter-Cleanup entfernt.
    _PSEUDO_VALUES = {
        "nicht erkennbar", "nicht ermittelbar", "nicht angegeben",
        "nicht bekannt", "nicht verfügbar", "nicht verfuegbar",
        "k.a.", "k. a.", "n/a", "n.a.", "unbekannt",
        "wird nachgereicht", "siehe anlage", "fehlt im datenraum",
        "-", "—", "tba", "tbd", "xxx",
    }
    pseudo_cleared = 0
    for k, v in list(result.items()):
        if not isinstance(v, str):
            continue
        v_norm = v.strip().lower()
        if v_norm in _PSEUDO_VALUES:
            result[k] = ""
            pseudo_cleared += 1
    if pseudo_cleared:
        print(f"  Pseudo-Wert-Cleanup: {pseudo_cleared} Felder geleert "
              f"('nicht erkennbar' etc. → '')")
    return result

# Regex: matcht {{KEY}}, {{KEY-SPLIT}}, {{KEY|suffix}}, {{KEY | suffix}}
# Bindestriche im Key werden beim Lookup entfernt (z.B. {{PRODUKT_BESCHREI-BUNG}})
_PH_RE = re.compile(r'\{\{\s*([A-Z0-9_-]+)\s*(?:\|[^}]*)?\}\}', re.IGNORECASE)

def _replace_placeholders(text, data):
    """Ersetzt alle {{KEY}} und {{KEY|suffix}} Platzhalter. Case-insensitiv.
    Bindestriche im Key werden vor dem Lookup entfernt (für Template-Zeilenumbrüche).
    """
    repl_map = {k.upper(): str(v or "") for k, v in data.items()}
    def _sub(m):
        key = m.group(1).upper().strip().replace('-', '')
        return repl_map.get(key, m.group(0))
    return _PH_RE.sub(_sub, text)

# Regex für Cross-Paragraph-Splits mit optionalem Bindestrich am Zeilenende
# z.B. "{{PRODUKT_BESCHREI-\nBUNG}}" → key = "PRODUKT_BESCHREIBUNG"
_PH_SPLIT_RE = re.compile(
    r'\{\{\s*([A-Z0-9_]+-?)\n([A-Z0-9_]*)\s*(?:\|[^}]*)?\}\}',
    re.IGNORECASE
)

def _replace_split_placeholder(text, data):
    """Behandelt Platzhalter die mit Bindestrich über zwei Zeilen gesplittet sind."""
    repl_map = {k.upper(): str(v or "") for k, v in data.items()}
    def _sub(m):
        part1 = m.group(1).rstrip('-')
        part2 = m.group(2)
        key = (part1 + part2).upper().strip()
        return repl_map.get(key, m.group(0))
    return _PH_SPLIT_RE.sub(_sub, text)


def duplicate_slide(prs, slide_index):
    """Duplicates the slide at slide_index and inserts the copy at slide_index + 1."""
    template = prs.slides[slide_index]
    new_slide = prs.slides.add_slide(template.slide_layout)

    # Replace shape tree with a deep copy of the template's
    sp_tree = new_slide.shapes._spTree
    tmpl_sp_tree = template.shapes._spTree

    # Remove all children added by add_slide (keep only nvGrpSpPr + grpSpPr = first 2)
    for child in list(sp_tree)[2:]:
        sp_tree.remove(child)

    # Copy shapes from template (skip first 2 as well, copy the rest)
    for child in list(tmpl_sp_tree)[2:]:
        sp_tree.append(deepcopy(child))

    # Move new slide (currently last) to position slide_index + 1
    sldIdLst = prs.slides._sldIdLst
    moved_el = sldIdLst[-1]
    sldIdLst.remove(moved_el)
    sldIdLst.insert(slide_index + 1, moved_el)

    return prs.slides[slide_index + 1]


def duplicate_we_slides(prs, data):
    """
    v19-Modell: 1 Wohnungstyp pro Slide (mit 2 Beispiel-WEs nebeneinander).
    Original-Slide zeigt Typ 1 (we_beispiel_1, we_beispiel_2, we_flaeche_*, we_typ_beschreibung).
    Für jeden weiteren Typ (typ2..typ8) wird ein Slide-Duplikat erstellt.

    Dupliziert nur wenn TEXT-Keys für diesen Typ befüllt sind
    (we_beispiel_3/_5/_7/... oder we_flaeche_*_typ2/typ3/...).
    """
    from pptx.oxml import parse_xml
    import string

    # Snapshot der Original-Seitenzahlen + TOC-Folie BEVOR irgendetwas dupliziert
    # oder umnummeriert wird. Wird nach allen Modifikationen für den finalen
    # Page-Resync gebraucht.
    template_n_to_slide, toc_slide = _build_template_page_map(prs)

    # Ermittle alle extra Typen mit Daten — vorher: while+break, brach beim
    # ersten leeren Typ ab → Lücken (z.B. typ2 leer, typ3 voll) fielen weg.
    # Jetzt: alle 8 prüfen, Liste der gefundenen typ_nums sammeln.
    extra_typs = []
    for typ in range(2, 9):
        left_n  = typ * 2 - 1
        right_n = typ * 2
        has_data = (
            data.get(f"we_beispiel_{left_n}") or data.get(f"we_beispiel_{right_n}")
            or data.get(f"we_typ_beschreibung_typ{typ}")
            or any(data.get(f"we_flaeche_{n}_typ{typ}") for n in range(1, 6))
        )
        if has_data:
            extra_typs.append(typ)
    extra_slides = len(extra_typs)

    # Find original WE slide
    we_idx = None
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            txt = ""
            if shape.has_text_frame:
                txt = shape.text_frame.text.upper()
            elif shape.shape_type == 6:
                for c in shape.shapes:
                    if c.has_text_frame:
                        txt += c.text_frame.text.upper()
            if "WE_BEISPIEL_1" in txt or "WE_BEREICH_1" in txt or "BILD_WE_1" in txt:
                we_idx = i
                break
        if we_idx is not None:
            break

    if we_idx is None:
        # Try also matching new placeholder names
        for i, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                txt = ""
                if shape.has_text_frame:
                    txt = shape.text_frame.text.upper()
                elif shape.shape_type == 6:
                    for c in shape.shapes:
                        if c.has_text_frame:
                            txt += c.text_frame.text.upper()
                if "WE_NUMMERN_1" in txt or "WE_RAUM_1_NAME_1" in txt:
                    we_idx = i
                    break
            if we_idx is not None:
                break

    if we_idx is None:
        print("duplicate_we_slides: WE-Slide nicht gefunden – überspringe Buchstaben-Fix")
        return

    print(f"WE-Slide bei Index {we_idx} → {extra_slides} extra Slide(s) "
          f"({2 + extra_slides * 2} WE-Typen insgesamt)")

    # Buchstaben für Dekorativ-Beschriftung: a, b, c, d, e, f, g, h, ...
    # Für Pair k: linker Buchstabe = letters[2*(k-1)], rechter = letters[2*(k-1)+1]
    all_letters = string.ascii_lowercase  # a-z (genug für 13 WE-Paare)

    def _fix_we_letters(xml_str, left_letter, right_letter):
        """Ersetzt die zwei dekorativen 'a'-Buchstaben (links/rechts) im XML."""
        tag = '<a:t>a</a:t>'
        first_pos = xml_str.find(tag)
        if first_pos < 0:
            return xml_str
        last_pos = xml_str.rfind(tag)
        rep_left  = f'<a:t>{left_letter}</a:t>'
        rep_right = f'<a:t>{right_letter}</a:t>'
        if first_pos == last_pos:
            return xml_str[:first_pos] + rep_left + xml_str[first_pos + len(tag):]
        # Letzten (rechts) zuerst ersetzen, damit first_pos valide bleibt
        result = xml_str[:last_pos] + rep_right + xml_str[last_pos + len(tag):]
        fp = result.find(tag)
        if fp >= 0:
            result = result[:fp] + rep_left + result[fp + len(tag):]
        return result

    # Original-XML VOR jeglicher Modifikation speichern
    orig_xml = etree.tostring(prs.slides[we_idx].shapes._spTree, encoding="unicode")

    def _fix_we_letters_single(xml_str, letter):
        """v19: BEIDE 'a'-Buchstaben (links + rechts) zeigen DENSELBEN Typ-Letter."""
        return xml_str.replace('<a:t>a</a:t>', f'<a:t>{letter}</a:t>')

    # Duplikate in UMGEKEHRTER Reihenfolge erstellen:
    # Jedes neue Slide wird bei we_idx+1 eingefügt → vorherige rücken vor
    # → Endergebnis: extra_typs[0] bei we_idx+1, extra_typs[1] bei we_idx+2, ...
    # Buchstabe = laufende Position in der Sequenz (b, c, d, ...) — auch wenn
    # extra_typs Lücken hat (z.B. [3, 5] → b, c). Das DQN-Layout zeigt
    # konsekutive Buchstaben unabhängig von der Typ-Nummer in den Daten.
    for slot_idx in range(extra_slides - 1, -1, -1):
        typ_num     = extra_typs[slot_idx]              # echte Typ-Nummer aus den Daten
        left_n      = typ_num * 2 - 1
        right_n     = typ_num * 2
        type_letter = all_letters[slot_idx + 1]         # Position +1 (a=Typ1, b=erstes extra, ...)

        new_slide = duplicate_slide(prs, we_idx)
        sp_tree   = new_slide.shapes._spTree

        xml_str = orig_xml

        # Per-Spalte (Beispiel-WEs): links=we_beispiel_(2N-1), rechts=we_beispiel_(2N)
        xml_str = xml_str.replace("WE_BEISPIEL_2", f"WE_BEISPIEL_{right_n}")
        xml_str = xml_str.replace("WE_BEISPIEL_1", f"WE_BEISPIEL_{left_n}")
        xml_str = xml_str.replace("WE_BEREICH_2", f"WE_BEREICH_{right_n}")
        xml_str = xml_str.replace("WE_BEREICH_1", f"WE_BEREICH_{left_n}")
        xml_str = xml_str.replace("BILD_WE_2",     f"BILD_WE_{right_n}")
        xml_str = xml_str.replace("BILD_WE_1",     f"BILD_WE_{left_n}")

        # Geteilte Felder pro TYP: WE_FLAECHE_1..5 → WE_FLAECHE_1_typN
        # WE_TYP_BESCHREIBUNG → WE_TYP_BESCHREIBUNG_typN
        xml_str = xml_str.replace("WE_TYP_BESCHREIBUNG", f"WE_TYP_BESCHREIBUNG_typ{typ_num}")
        for fn in range(5, 0, -1):
            xml_str = xml_str.replace(f"WE_FLAECHE_{fn}", f"WE_FLAECHE_{fn}_typ{typ_num}")

        # v20-Backwards: WE_NUMMERN, WE_RAUM (auch ersetzen, falls noch Platzhalter da)
        xml_str = xml_str.replace("WE_NUMMERN_2", f"WE_NUMMERN_{right_n}")
        xml_str = xml_str.replace("WE_NUMMERN_1", f"WE_NUMMERN_{left_n}")
        for rn in range(5, 0, -1):
            xml_str = xml_str.replace(f"WE_RAUM_{rn}_NAME_2", f"WE_RAUM_{rn}_NAME_{right_n}")
            xml_str = xml_str.replace(f"WE_RAUM_{rn}_NAME_1", f"WE_RAUM_{rn}_NAME_{left_n}")

        # Beide Buchstaben-Positionen → derselbe Typ-Letter (v19: 1 Typ pro Slide)
        xml_str = _fix_we_letters_single(xml_str, type_letter)

        new_sp_tree = parse_xml(xml_str.encode("utf-8"))
        for child in list(sp_tree):
            sp_tree.remove(child)
        for child in list(new_sp_tree):
            sp_tree.append(child)

    # Original-Slide: bleibt bei 'a' (Typ 1) – beide Spalten haben schon 'a'

    # ── Page-Resync: Bottom-Pages + TOC auf echte PDF-Seitenzahl mappen ────
    # Ersetzt _renumber_bottom_pages + _renumber_toc_pages. Funktioniert auch
    # wenn extra_slides=0 (Template-Defaults stimmen sonst nicht zur Konvertier-
    # Pipeline 1 PPTX-Folie = 1 PDF-Seite).
    _resync_pages_to_actual(prs, template_n_to_slide, toc_slide)


# Page-number-like text matcher: 1-3 digit number, possibly with whitespace
import re as _re_pg
_PAGE_NUM_RE = _re_pg.compile(r'^\s*(\d{1,3})\s*$')


def _renumber_bottom_pages(prs, we_idx, extra_slides):
    """Verschiebt Seitenzahlen unten auf Slides nach we_idx.
    - Duplizierte WE-Slides (we_idx+1 .. we_idx+extra_slides): shift = 2 * Position
    - Alle Slides DANACH: shift = 2 * extra_slides
    Heuristik für 'page number': Text-Frame mit isolierter 1-3-stelliger Zahl,
    in Position unten auf der Slide (top > 80% der Slide-Höhe).
    """
    sh = int(prs.slide_height or 6858000)
    # Nur Text-Frames im unteren Drittel als Page-Number-Kandidaten
    bottom_threshold = int(sh * 0.85)

    def _shift_in_slide(slide, shift):
        for shape in slide.shapes:
            try:
                top = shape.top or 0
            except Exception:
                top = 0
            # Group-Shape: Children prüfen
            if shape.shape_type == 6:
                gx = shape.top or 0
                for child in shape.shapes:
                    ctop = (child.top or 0) + gx
                    if ctop < bottom_threshold:
                        continue
                    if child.has_text_frame:
                        _shift_text_frame(child.text_frame, shift)
            else:
                if top < bottom_threshold:
                    continue
                if shape.has_text_frame:
                    _shift_text_frame(shape.text_frame, shift)

    def _shift_text_frame(tf, shift):
        for p in tf.paragraphs:
            full = "".join(r.text for r in p.runs).strip()
            m = _PAGE_NUM_RE.match(full)
            if not m:
                continue
            n = int(m.group(1))
            # Nur 1-2 stellige Zahlen oder bis 200 (echte Seitenzahlen)
            if n < 1 or n > 200:
                continue
            new_n = n + shift
            # Schreibe neuen Wert in den ersten Run, alle anderen leer
            if p.runs:
                p.runs[0].text = str(new_n)
                for r in p.runs[1:]:
                    r.text = ""

    # Duplizierte WE-Slides: jeder Slide kriegt seinen eigenen Shift
    for i in range(1, extra_slides + 1):
        slide_idx = we_idx + i
        if slide_idx < len(prs.slides):
            _shift_in_slide(prs.slides[slide_idx], 2 * i)
            print(f"  Page-Renumber: WE-Slide {slide_idx} → +{2*i}")

    # Alle Slides nach dem WE-Block: shift um 2 * extra_slides
    total_shift = 2 * extra_slides
    for i in range(we_idx + extra_slides + 1, len(prs.slides)):
        _shift_in_slide(prs.slides[i], total_shift)
    print(f"  Page-Renumber: {len(prs.slides) - we_idx - extra_slides - 1} "
          f"Slides nach WE-Block um +{total_shift} verschoben")


def _read_max_bottom_page(prs, slide):
    """Liest die größte Bottom-Page-Nummer (1-3 stellig, untere 15% des Slides).
    Genutzt als Schwelle: alles > diesem Wert ist 'nach dem WE-Block'.
    Gibt None zurück wenn keine Page-Nummer gefunden."""
    sh_h = int(prs.slide_height or 6858000)
    bottom_threshold = int(sh_h * 0.85)
    found = []

    def _scan_tf(tf, top_abs):
        if top_abs < bottom_threshold:
            return
        for p in tf.paragraphs:
            full = "".join(r.text for r in p.runs).strip()
            m = _PAGE_NUM_RE.match(full)
            if m:
                n = int(m.group(1))
                if 1 <= n <= 200:
                    found.append(n)

    for shape in slide.shapes:
        try:
            top = shape.top or 0
        except Exception:
            top = 0
        if shape.shape_type == 6:
            for child in shape.shapes:
                ctop = (child.top or 0) + (shape.top or 0)
                if child.has_text_frame:
                    _scan_tf(child.text_frame, ctop)
        else:
            if shape.has_text_frame:
                _scan_tf(shape.text_frame, top)

    return max(found) if found else None


def _scan_bottom_page_paragraphs(prs, slide):
    """Findet Paragraphen im unteren 15% des Slides, die nur eine 1-3-stellige
    Zahl enthalten (1..200). Gibt Liste [(paragraph, n)] zurück."""
    sh_h = int(prs.slide_height or 6858000)
    bottom_threshold = int(sh_h * 0.85)
    out = []
    def _go(tf, abs_top):
        if abs_top < bottom_threshold:
            return
        for p in tf.paragraphs:
            full = "".join(r.text for r in p.runs).strip()
            m = _PAGE_NUM_RE.match(full)
            if not m:
                continue
            n = int(m.group(1))
            if 1 <= n <= 200:
                out.append((p, n))
    for shape in slide.shapes:
        try: top = shape.top or 0
        except Exception: top = 0
        if shape.shape_type == 6:
            for child in shape.shapes:
                if child.has_text_frame:
                    _go(child.text_frame, (child.top or 0) + (shape.top or 0))
        else:
            if shape.has_text_frame:
                _go(shape.text_frame, top)
    return out


def _scan_isolated_numbers(slide):
    """Alle Paragraphen mit isolierter Zahl 1..200, egal wo auf der Folie."""
    out = []
    def _go(tf):
        for p in tf.paragraphs:
            full = "".join(r.text for r in p.runs).strip()
            m = _PAGE_NUM_RE.match(full)
            if not m:
                continue
            n = int(m.group(1))
            if 1 <= n <= 200:
                out.append((p, n))
    for shape in slide.shapes:
        if shape.shape_type == 6:
            for child in shape.shapes:
                if child.has_text_frame:
                    _go(child.text_frame)
        else:
            if shape.has_text_frame:
                _go(shape.text_frame)
    return out


def _build_template_page_map(prs):
    """Snapshot VOR jeglicher Modifikation: Template-Spread-Page → Slide-Objekt.
    Wird nach WE-Duplikation genutzt, um TOC-Einträge auf neue Slide-Positionen
    zu mappen. Slide-Objekte überleben die Duplikation/Reorder, Indizes nicht.

    Mapping basiert auf Spread-Konvention: Slide N (1-basiert) repräsentiert
    PDF-Seiten (2N-1) und (2N). Vorher: nur Slides mit Bottom-Pages → TOC-
    Einträge die auf Header-Slides zeigen (z.B. 'u invest'-Header) landeten
    in unsicherer Interpolation. Jetzt: jede mögliche Spread-Seite eindeutig.

    Gibt ({n: slide}, toc_slide_or_None) zurück.
    """
    out = {}
    toc_slide = None
    for idx, slide in enumerate(prs.slides):
        spread = idx + 1
        # Beide Seiten des Spreads → diese Slide
        out.setdefault(2 * spread - 1, slide)  # linke Seite
        out.setdefault(2 * spread,     slide)  # rechte Seite
        # TOC erkennen: viele isolierte Zahlen, keine Bottom-Pages
        if toc_slide is None:
            bottoms = _scan_bottom_page_paragraphs(prs, slide)
            all_nums = _scan_isolated_numbers(slide)
            unique = {n for (_, n) in all_nums}
            if len(unique) >= 8 and len(bottoms) == 0:
                toc_slide = slide
    return out, toc_slide


def _resync_pages_to_actual(prs, template_n_to_slide, toc_slide):
    """Spread-Numbering: jede Folie repräsentiert 2 Seiten (links + rechts).
    Nur Folien die im Template ursprünglich Bottom-Pages hatten werden nummeriert
    (Cover + TOC bleiben ohne Nummer). TOC-Einträge werden via Mapping korrigiert.

    Schema:
      Erste numerierte Folie  (spread_idx=1):  links=1, rechts=2
      Zweite numerierte Folie (spread_idx=2):  links=3, rechts=4
      ...
    """
    # Slide-Objekte sind nicht hashbar → Lookup via id().
    id_to_idx = {id(s): i for i, s in enumerate(prs.slides)}

    # Helper: Bottom-Page-Paragraphen mit X-Position (für links/rechts-Sortierung)
    def _scan_bottom_with_x(slide):
        sh_h = int(prs.slide_height or 6858000)
        bottom_threshold = int(sh_h * 0.85)
        out = []  # [(paragraph, n, abs_x)]
        for shape in slide.shapes:
            try:
                top = shape.top or 0
            except Exception:
                top = 0
            if shape.shape_type == 6:
                for child in shape.shapes:
                    ctop = (child.top or 0) + (shape.top or 0)
                    if ctop < bottom_threshold:
                        continue
                    if child.has_text_frame:
                        cx = (child.left or 0) + (shape.left or 0)
                        for p in child.text_frame.paragraphs:
                            full = "".join(r.text for r in p.runs).strip()
                            m = _PAGE_NUM_RE.match(full)
                            if m and 1 <= int(m.group(1)) <= 999:
                                out.append((p, int(m.group(1)), cx))
            else:
                if top < bottom_threshold:
                    continue
                if shape.has_text_frame:
                    sx = shape.left or 0
                    for p in shape.text_frame.paragraphs:
                        full = "".join(r.text for r in p.runs).strip()
                        m = _PAGE_NUM_RE.match(full)
                        if m and 1 <= int(m.group(1)) <= 999:
                            out.append((p, int(m.group(1)), sx))
        return out

    # Spread-Index = PPTX-Slide-Index + 1 (1-basiert).
    # ALLE Folien zählen mit – auch Cover/TOC (die haben aber keine Bottom-Pages
    # und zeigen ihre Spread-Nummer deshalb nicht). Beispiel:
    #   Folie 0 (Cover): spread=1 → Seiten 1+2 (unsichtbar)
    #   Folie 1 (TOC):   spread=2 → Seiten 3+4 (unsichtbar)
    #   Folie 2 (Inhalt): spread=3 → Seiten 5+6 (sichtbar)
    slide_to_spread = {}
    for idx, slide in enumerate(prs.slides):
        slide_to_spread[id(slide)] = idx + 1

    # Bottom-Pages: links = 2*spread-1, rechts = 2*spread (X-sortiert)
    bottom_changes = 0
    for idx, slide in enumerate(prs.slides):
        si = slide_to_spread.get(id(slide))
        if si is None:
            continue
        bottoms = sorted(_scan_bottom_with_x(slide), key=lambda t: t[2])  # by X
        for i, (p, old_n, _x) in enumerate(bottoms):
            new_n = str(2 * si - 1 + i)  # 0 → links, 1 → rechts
            if str(old_n) == new_n:
                continue
            if p.runs:
                p.runs[0].text = new_n
                for r in p.runs[1:]:
                    r.text = ""
            bottom_changes += 1

    # TOC: für jede Template-Zahl → Slide finden → neue Spread-Seite setzen
    toc_changes = 0
    toc_idx = id_to_idx.get(id(toc_slide)) if toc_slide is not None else None
    if toc_slide is not None and toc_idx is not None:
        # Mapping Template-N → neue PDF-Seite. Parität bewahren:
        # alt ungerade → linke Seite (2*si-1), alt gerade → rechte Seite (2*si).
        template_n_to_new_page = {}
        for n, s in template_n_to_slide.items():
            si = slide_to_spread.get(id(s))
            if si is not None:
                template_n_to_new_page[n] = 2 * si if n % 2 == 0 else 2 * si - 1

        # Sortierte Liste für Fallback
        sorted_keys = sorted(template_n_to_new_page.keys())
        max_known   = sorted_keys[-1] if sorted_keys else 0
        last_spread = len(prs.slides)
        last_page   = 2 * last_spread  # rechte Seite des letzten Spreads (war: len(slides))

        def _resolve(old_n):
            """Liefert die neue Seitenzahl für eine TOC-Template-Zahl.
            Mit Spread-Vollmapping + Paritaets-Bewahrung im Mapping ist der
            Treffer in 99% der Faelle exakt. Stub-Werte oberhalb (z.B.
            Template-Original '68' fuer Rechtlich) → letzte Spread-Seite,
            Paritaet bewahren. Sonst None → Original behalten.
            """
            exact = template_n_to_new_page.get(old_n)
            if exact is not None:
                return exact
            if old_n > max_known:
                return last_page if old_n % 2 == 0 else last_page - 1
            return None

        toc_idx_skip = {n for _, n in _scan_bottom_page_paragraphs(prs, toc_slide)}

        # Hilfs-Funktion: Parent-Shape eines Paragraphen finden + TextBox so
        # konfigurieren dass CloudConvert keine 2-stelligen Zahlen umbricht
        # (das war der Output-Bug 'Magdeburg → 2 / 7' usw. — Boxen waren
        # 0.18" schmal mit auto_size=SHAPE_TO_FIT_TEXT, was CloudConvert nicht
        # respektiert; wraps zweistellige zu zwei Zeilen).
        from pptx.enum.text import MSO_AUTO_SIZE as _MSO_AUTO_SIZE
        _MIN_TOC_BOX_W = int(0.6 * 914400)  # 0.6" — sicher fuer 3-stellige Zahlen

        def _find_shape_for_paragraph(slide, paragraph):
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for sp in shape.text_frame.paragraphs:
                        if sp._p is paragraph._p:
                            return shape
                if shape.shape_type == 6:
                    for c in shape.shapes:
                        if c.has_text_frame:
                            for sp in c.text_frame.paragraphs:
                                if sp._p is paragraph._p:
                                    return c
            return None

        def _harden_toc_box(shape):
            """Verhindert Word-Wrap im PDF-Render: Wrap aus, Auto-Size aus,
            Mindest-Breite. Idempotent."""
            if shape is None:
                return
            try:
                shape.text_frame.word_wrap = False
                shape.text_frame.auto_size = _MSO_AUTO_SIZE.NONE
            except Exception:
                pass
            try:
                if (shape.width or 0) < _MIN_TOC_BOX_W:
                    shape.width = _MIN_TOC_BOX_W
            except Exception:
                pass

        # Heuristik-Whitelist: Kapitel-Marker (1..4) auf TOC bleiben, also nicht remappen
        # wenn 1..len(chapter_count) und sehr klein. Solche Zahlen sollten unverändert bleiben.
        for p, old_n in _scan_isolated_numbers(toc_slide):
            # Kapitel-Markierungen 1,2,3,4 unverändert lassen — sie sind Layout-Elemente,
            # keine Seitenreferenzen.
            if old_n <= 4:
                continue
            # Box auch bei unveraenderten Werten haerten — verhindert Render-Wrap
            shape = _find_shape_for_paragraph(toc_slide, p)
            _harden_toc_box(shape)
            new_n = _resolve(old_n)
            if new_n is None or new_n == old_n:
                continue
            if p.runs:
                p.runs[0].text = str(new_n)
                for r in p.runs[1:]:
                    r.text = ""
            toc_changes += 1

    print(f"  Page-Resync: {bottom_changes} Bottom-Pages → echte PDF-Seite, "
          f"{toc_changes} TOC-Einträge gemappt (TOC-Slide={toc_idx})")


def _crop_image_to_aspect(img_bytes, target_w_emu, target_h_emu, max_side_px=1400):
    """Schneidet ein Bild auf das Aspect-Ratio des Ziel-Shapes (Center-Crop) und
    skaliert auf eine vernünftige Pixelgröße. Verhindert Verzerrungen und Blur
    durch Stretch-Modus im PPTX-blipFill.

    Args:
        img_bytes: Original-Bilddaten (JPEG/PNG/WebP).
        target_w_emu, target_h_emu: Shape-Dimensionen in EMU (1 inch = 914400).
        max_side_px: Maximale Pixel-Kantenlänge (Standardwert 1400 ≈ 3 MB JPEG).
            Reduziert von 1800 → 1400 wegen 512 MB Render-Plan.
    """
    if not target_w_emu or not target_h_emu:
        return img_bytes
    try:
        from PIL import Image
        src = Image.open(io.BytesIO(img_bytes))
        # JPEG draft() = DCT-Skalierung: dekodiert direkt 1/2, 1/4 oder 1/8.
        # Spart 75-87% RAM beim Decode (4K-JPEG: ~48 MB → ~6 MB).
        if (src.format or "").upper() == "JPEG":
            try:
                src.draft("RGB", (max_side_px * 2, max_side_px * 2))
            except Exception:
                pass
        # Auto-rotate basierend auf EXIF-Orientation
        try:
            from PIL import ImageOps
            src = ImageOps.exif_transpose(src)
        except Exception:
            pass
        src_w, src_h = src.size
        if src_w <= 0 or src_h <= 0:
            return img_bytes

        target_aspect = target_w_emu / target_h_emu
        src_aspect    = src_w / src_h

        # Center-Crop zur Anpassung an Ziel-Aspect
        if abs(src_aspect - target_aspect) > 0.02:  # ≥2% Unterschied → croppen
            if src_aspect > target_aspect:
                # Quelle ist breiter → links/rechts beschneiden
                new_w = int(src_h * target_aspect)
                left  = (src_w - new_w) // 2
                src   = src.crop((left, 0, left + new_w, src_h))
            else:
                # Quelle ist höher → oben/unten beschneiden
                new_h = int(src_w / target_aspect)
                top   = (src_h - new_h) // 2
                src   = src.crop((0, top, src_w, top + new_h))

        # Skalieren auf vernünftige Größe – nie hoch, nur runter
        cw, ch = src.size
        long_side = max(cw, ch)
        if long_side > max_side_px:
            scale = max_side_px / long_side
            src = src.resize((int(cw * scale), int(ch * scale)), Image.LANCZOS)

        # Konvertieren zu RGB falls nötig (PNG mit Alpha → JPEG)
        if src.mode in ("RGBA", "LA", "P"):
            bg = Image.new("RGB", src.size, (255, 255, 255))
            if src.mode == "P":
                src = src.convert("RGBA")
            bg.paste(src, mask=src.split()[-1] if src.mode in ("RGBA", "LA") else None)
            src = bg
        elif src.mode != "RGB":
            src = src.convert("RGB")

        out = io.BytesIO()
        src.save(out, format="JPEG", quality=85, optimize=True)
        result = out.getvalue()
        src.close()
        del src, out
        return result
    except Exception as e:
        print(f"  _crop_image_to_aspect Fehler: {e} – verwende Original")
        return img_bytes


def _shrink_for_storage(raw, max_px=1600):
    """Pre-shrink: Bild SOFORT beim Laden in image_data verkleinern.
    Spart RAM-Footprint des Dicts (60 Bilder × 1 MB → ~60 × 200 KB = 12 MB).
    Nutzt JPEG draft() für minimalen Decode-Spike.
    Bei Fehler: Original zurückgeben (defensive)."""
    if not raw or len(raw) < 250_000:
        return raw  # bereits klein genug
    try:
        from PIL import Image
        img = Image.open(io.BytesIO(raw))
        if (img.format or "").upper() == "JPEG":
            try:
                img.draft("RGB", (max_px, max_px))
            except Exception:
                pass
        img.load()
        w, h = img.size
        long_side = max(w, h)
        if long_side <= max_px and (img.format or "").upper() == "JPEG":
            img.close()
            return raw
        if long_side > max_px:
            scale = max_px / long_side
            img = img.resize((max(1, int(w * scale)), max(1, int(h * scale))),
                             Image.LANCZOS)
        if img.mode in ("RGBA", "LA", "P"):
            bg = Image.new("RGB", img.size, (255, 255, 255))
            if img.mode == "P":
                img = img.convert("RGBA")
            bg.paste(img, mask=img.split()[-1] if img.mode in ("RGBA", "LA") else None)
            img = bg
        elif img.mode != "RGB":
            img = img.convert("RGB")
        out = io.BytesIO()
        img.save(out, format="JPEG", quality=85, optimize=True)
        result = out.getvalue()
        img.close()
        del img, out
        return result if len(result) < len(raw) else raw
    except Exception as e:
        print(f"  _shrink_for_storage Fehler: {e}")
        return raw


def _override_grund_texts(prs, data):
    """Ueberschreibt die fest hinterlegten Titel + Texte der '6 gute Gruende'-
    Slide mit projekt-spezifischen Werten aus expose_data.

    Template-Layout (rechte Spalte): 6 Punkte je 2 Zeilen — Titel oben, Text unten.
    Nur 3 von 12 Texten sind Placeholder-basiert ({{STADT}}, {{KAUFPREIS_AB}},
    {{KFW_DARLEHEN}}); 9 sind hartkodiert. Diese Funktion ersetzt die hart-
    kodierten via Substring-Match gegen den Original-Template-Text.

    Wenn das jeweilige expose_data-Feld leer ist → Default belassen (Template-
    Texte sind brauchbare Fallbacks).

    Idempotent — Substring-Match findet nach erstem Override nichts mehr.
    """
    # (Original-Template-Substring, expose_data-Key, max_len fuer Layout-Schutz)
    OVERRIDES = [
        # TITEL-Boxen — max 32 Zeichen einzeilig (Box-Breite begrenzt)
        ("Zentrale Lage in",            "text_grund_1_titel", 32),
        ("{{KAUFPREIS_AB",              "text_grund_2_titel", 32),
        ("{{KFW_DARLEHEN",              "text_grund_3_titel", 32),
        ("3-fach Abschreibung",         "text_grund_4_titel", 32),
        ("Möblierungskonzept",          "text_grund_5_titel", 32),
        ("Mietgarantie für 3 Monate",   "text_grund_6_titel", 32),
        # TEXT-Boxen — max 160 Zeichen (2 Zeilen + Sicherheits-Puffer)
        ("3 Minuten zur Universität",   "text_grund_1_text", 160),
        ("Voll förderfähiges Neubauprojekt", "text_grund_2_text", 160),
        ("von bis zu 150.000",          "text_grund_3_text", 160),
        ("5 % degressive AfA",          "text_grund_4_text", 160),
        ("Individuell gestaltete Apartments", "text_grund_5_text", 160),
        ("Nach Fertigstellung garantiert", "text_grund_6_text", 160),
    ]

    def _truncate_smart(text, max_len):
        """Kuerzt auf max_len Zeichen, abschneiden auf naechstem Wort.
        Wenn das letzte Wort >half max_len ist → harter Cut. Sonst Wort-Boundary."""
        if len(text) <= max_len:
            return text
        cut = text[:max_len].rstrip()
        # Auf Satzende oder Wortgrenze zurueck
        for sep in (". ", "! ", "? ", "; ", ", ", " "):
            idx = cut.rfind(sep)
            if idx > max_len * 0.6:  # mind. 60% des Limits behalten
                return cut[:idx + len(sep.rstrip())].rstrip()
        return cut

    overrides_applied = 0
    for slide in prs.slides:
        # Heuristisch '6 gute Gruende' Slide erkennen
        all_text = ""
        for sh in slide.shapes:
            if sh.has_text_frame:
                all_text += " " + sh.text_frame.text
            if sh.shape_type == 6:
                for c in sh.shapes:
                    if c.has_text_frame:
                        all_text += " " + c.text_frame.text
        if "6 gute" not in all_text.lower():
            continue
        # Slide gefunden — alle Shapes scannen + ggf. ersetzen
        def _walk(shape):
            yield shape
            if shape.shape_type == 6:
                for c in shape.shapes:
                    yield from _walk(c)
        for top_shape in list(slide.shapes):
            for shape in _walk(top_shape):
                if not shape.has_text_frame:
                    continue
                tf_text = shape.text_frame.text
                for substr, key, max_len in OVERRIDES:
                    new_text = (data or {}).get(key, "")
                    if not new_text or substr not in tf_text:
                        continue
                    # Layout-Schutz: kuerzen wenn Claude das Limit gerissen hat
                    new_text = _truncate_smart(str(new_text).strip(), max_len)
                    # Im ersten Paragraphen den Run-Text ersetzen,
                    # weitere Paragraphen leeren (1-Zeilen-Boxen).
                    if shape.text_frame.paragraphs:
                        first_p = shape.text_frame.paragraphs[0]
                        if first_p.runs:
                            first_p.runs[0].text = str(new_text)
                            for r in first_p.runs[1:]:
                                r.text = ""
                        else:
                            # Kein Run vorhanden → einen erzeugen
                            from pptx.util import Pt
                            run = first_p.add_run()
                            run.text = str(new_text)
                        # Weitere Paragraphen leeren — die festen Zwei-Zeilen-
                        # Texte (z.B. '5 % degressive AfA + ...') stehen oft
                        # in mehreren Paragraphen; alle nachfolgenden auf
                        # leeren-Run setzen damit nur new_text zu sehen ist.
                        for p in shape.text_frame.paragraphs[1:]:
                            for r in p.runs:
                                r.text = ""
                    overrides_applied += 1
                    break  # pro Shape nur EIN Override
    if overrides_applied:
        print(f"  6-gute-Gruende: {overrides_applied} Texte projekt-spezifisch ueberschrieben")


def _override_legal_text(prs, data):
    """Ersetzt magdeburg-spezifische Hardcoded-Texte auf der Rechtliche-Hinweise-
    Slide (S26) mit projekt-/stadt-dynamischen Werten.

    Aktuell: 'auf dem Areal der ehemaligen Diamant Brauerei' →
    aus expose_data.quartier_name oder stadtteil oder stadt erzeugen.

    Damit funktioniert das Template auch fuer Projekte ausserhalb von Magdeburg
    ohne dass der User das PPTX manuell editieren muss.
    """
    quartier = (data or {}).get("quartier_name") or ""
    stadtteil = (data or {}).get("stadtteil") or ""
    stadt = (data or {}).get("stadt") or ""
    plz = (data or {}).get("plz") or ""

    # Dynamischer Ersatz-Text fuer 'auf dem Areal der ehemaligen Diamant Brauerei'.
    # Falls quartier_name etwas wie 'Diamant Quartier' enthaelt → 'auf dem Areal {quartier}'.
    # Sonst stadtteil-basiert: 'im Stadtteil {stadtteil}'.
    # Sonst Fallback nur Stadt.
    if quartier and quartier.strip():
        replacement = f"auf dem Areal {quartier.strip()}"
    elif stadtteil and stadtteil.strip():
        replacement = f"im Stadtteil {stadtteil.strip()}"
    elif stadt:
        replacement = f"in {stadt.strip()}"
    else:
        replacement = ""

    if not replacement:
        return  # nichts zu ersetzen

    # Nur bei Stadt != Magdeburg ersetzen — Magdeburg-Projekte sollen den
    # Original-Text behalten weil dort tatsaechlich die Brauerei ist.
    if "magdeburg" in stadt.lower():
        return

    # Pattern: 'Diamant Brau' / 'Diamant Brauerei' (auch mit Bindestrich) ueber
    # Paragraph-Grenzen hinweg. Wir behandeln den TextFrame als Ganzes.
    import re as _re
    DIAMANT_RE = _re.compile(
        r'auf\s+dem\s+Areal\s+der\s+ehemaligen\s+Diamant\s+Brau-?\s*erei',
        _re.IGNORECASE | _re.DOTALL
    )

    overrides = 0
    for slide in prs.slides:
        # Heuristisch S26 erkennen
        all_text = ""
        for sh in slide.shapes:
            if sh.has_text_frame:
                all_text += " " + sh.text_frame.text
        if "Verkaufsprospekt" not in all_text and "Diamant" not in all_text:
            continue
        # TextBoxen scannen, die den Diamant-Text enthalten
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            full = shape.text_frame.text
            if "Diamant" not in full:
                continue
            new_full = DIAMANT_RE.sub(replacement, full)
            if new_full == full:
                continue
            # Den ganzen TextFrame neu aufbauen, dabei Linebreaks erhalten
            # (per Newline → eigene Paragraphen).
            # Stil: erster Paragraph behaelt sein Run-Format (font/size/color).
            tf = shape.text_frame
            paragraphs = new_full.split("\n")
            # Ersten Paragraph: ersten Run mit erster Zeile setzen
            first_p = tf.paragraphs[0]
            if first_p.runs:
                first_p.runs[0].text = paragraphs[0]
                for r in first_p.runs[1:]:
                    r.text = ""
            else:
                first_p.add_run().text = paragraphs[0]
            # Restliche Paragraphen: vorhandene leeren + neue Zeilen reinpacken
            for i, line in enumerate(paragraphs[1:], start=1):
                if i < len(tf.paragraphs):
                    p = tf.paragraphs[i]
                    if p.runs:
                        p.runs[0].text = line
                        for r in p.runs[1:]:
                            r.text = ""
                    else:
                        p.add_run().text = line
                else:
                    # Mehr Zeilen als Paragraphen → neuen Paragraph anhaengen
                    new_p = tf.add_paragraph()
                    new_p.text = line
            # Falls ueberzaehlige alte Paragraphen → leeren
            for j in range(len(paragraphs), len(tf.paragraphs)):
                p = tf.paragraphs[j]
                for r in p.runs:
                    r.text = ""
            overrides += 1
    if overrides:
        print(f"  Legal-Text: {overrides} Magdeburg-spezifische Stelle(n) → {replacement!r}")


def fill_pptx(template_bytes, data, customer_images=None):
    """Fill PPTX template using python-pptx: embeds images and replaces text placeholders.
    customer_images: optional dict {bild_key: bytes} – Kundenbilder haben Vorrang vor URLs."""

    prs = Presentation(io.BytesIO(template_bytes))

    # Pre-load images: Kundenbilder zuerst, dann URLs für fehlende Slots
    image_data = {}

    # 1. Kundenbilder direkt einsetzen (höchste Priorität)
    # Vor-Verkleinerung: spart RAM permanent während fill_pptx läuft
    if customer_images:
        for key, raw in customer_images.items():
            shrunk = _shrink_for_storage(raw)
            image_data[key] = shrunk
            print(f"  ✓ Kundenbild: {key} ({len(raw)//1024} KB → {len(shrunk)//1024} KB)")
        print(f"  Kundenbilder: {len(customer_images)} eingeladen")
        import gc; gc.collect()

    # 2. Unsplash/Picsum-URLs für alle noch fehlenden bild_* Keys herunterladen
    bild_keys_total = [k for k in data if k.startswith("bild_")]
    bild_keys_with_url = [k for k in bild_keys_total
                          if isinstance(data[k], str) and data[k].startswith("http")
                          and k not in image_data]
    # Priorität: Stadt-Bilder zuerst (sonst gehen sie verloren), dann Amenities,
    # dann Picsum-Fallbacks. Mit Disk-Streaming + 1MB-Limit pro Bild ist 25 OK.
    _PRIO_PREFIXES = ["bild_stadt", "bild_quartier", "bild_lageplan",
                      "bild_amenity",
                      "bild_titel", "bild_projekt", "bild_greenliving",
                      "bild_interior", "bild_ausstattung",
                      "bild_ansicht", "bild_standort", "bild_grundriss_intro",
                      "bild_we_1", "bild_we_2", "bild_collage",
                      "bild_hotel", "bild_rechtlich"]
    def _sort_key(k):
        for i, prefix in enumerate(_PRIO_PREFIXES):
            if k.startswith(prefix):
                return i
        return 99
    # Cap auf 50 — mit pre-shrink (~200 KB pro Bild) ergibt das ~10 MB statt 60 MB Dauer-RAM.
    bild_keys_with_url_sorted = sorted(bild_keys_with_url, key=_sort_key)[:50]
    print(f"fill_pptx: {len(bild_keys_total)} bild_* Keys, "
          f"{len(bild_keys_with_url)} URLs, lade {len(bild_keys_with_url_sorted)} (RAM-Limit)")

    # OSM/Wikipedia/Wikimedia blocken Requests die nicht wie ein Browser aussehen.
    # Echte Chrome-UA + komplette Browser-Header → besteht Anti-Bot-Checks.
    _DL_HEADERS = {
        "User-Agent": "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) "
                      "AppleWebKit/537.36 (KHTML, like Gecko) "
                      "Chrome/121.0.0.0 Safari/537.36",
        "Accept": "image/avif,image/webp,image/apng,image/svg+xml,image/*,*/*;q=0.8",
        "Accept-Language": "de-DE,de;q=0.9,en;q=0.8",
        "Sec-Fetch-Dest": "image",
        "Sec-Fetch-Mode": "no-cors",
        "Sec-Fetch-Site": "cross-site",
    }
    import gc as _gc
    for idx, key in enumerate(bild_keys_with_url_sorted):
        value = data[key]
        try:
            resp = requests.get(value, timeout=15, headers=_DL_HEADERS, allow_redirects=True)
            if resp.status_code == 200:
                raw = resp.content
                # Cap raw download auf 3 MB; alles darüber wird verworfen statt geshrinkt
                # (sonst Decode-Spike auf einem 8K-Foto = 100+ MB RAM).
                if len(raw) >= 3_000_000:
                    print(f"  ⚠️  Bild zu groß übersprungen: {key} ({len(raw)//1024} KB)")
                else:
                    shrunk = _shrink_for_storage(raw)
                    image_data[key] = shrunk
                    print(f"  ✓ Bild geladen: {key} ({len(raw)//1024} KB → {len(shrunk)//1024} KB)")
                del raw
            else:
                print(f"  ✗ Bild HTTP-Fehler {key}: {resp.status_code} ({value[:80]})")
            del resp
        except Exception as e:
            print(f"  ✗ Bild Download-Fehler {key}: {e}")
        # Alle 8 Bilder gc — spart Pillow-Restpuffer
        if idx % 8 == 7:
            _gc.collect()
    _gc.collect()

    print(f"  image_data gesamt: {len(image_data)} Bilder "
          f"(~{sum(len(v) for v in image_data.values())//1024} KB)")

    def make_replacement_map(data):
        """Build a case-insensitive lookup: UPPER_KEY -> value.
        EXCLUDES bild_* keys — image slots must never be text-substituted.
        If an image fails to embed, the {{BILD_X}} placeholder stays visible
        rather than becoming a URL string in the slide text."""
        return {k.upper(): str(v or "") for k, v in data.items()
                if not k.lower().startswith('bild_')}

    REPL_MAP = make_replacement_map(data)

    # Regex that matches {{KEY}}, {{KEY-SPLIT}}, {{KEY|suffix}}, {{KEY | suffix}}
    # Includes '-' in key chars so split-across-linebreak placeholders like
    # {{PRODUKT_BESCHREI-BUNG}} (after joining runs) are captured.
    PLACEHOLDER_RE = re.compile(r'\{\{\s*([A-Z0-9_-]+)\s*(?:\|[^}]*)?\}\}', re.IGNORECASE)
    # Matches the |Xpt font-size hint inside a placeholder, e.g. {{MIN_UNI|50pt}}
    _SIZE_HINT_RE = re.compile(r'\|\s*(\d+)\s*pt\b', re.IGNORECASE)

    def replace_text(text):
        """Replace all placeholders in a string using REPL_MAP.
        Strips hyphens from keys before lookup so {{PRODUKT_BESCHREI-BUNG}}
        resolves to PRODUKT_BESCHREIBUNG (PowerPoint line-break artefact).
        """
        def _sub(m):
            key = m.group(1).upper().strip().replace('-', '')
            return REPL_MAP.get(key, m.group(0))  # keep original if not found
        return PLACEHOLDER_RE.sub(_sub, text)

    # Invisible/formatting chars that PowerPoint inserts inside placeholders
    _INVIS_RE = re.compile(r'[\u00AD\u200B\u200C\u200D\uFEFF\u00A0]')

    def replace_in_paragraph(para):
        """Replace placeholders in a paragraph, handling split-run placeholders.

        Strategy: reassemble all runs into one string, replace, write back into
        runs[0] preserving its formatting, clear all other runs.
        If the placeholder contains a |Xpt font-size hint (e.g. {{MIN_UNI|50pt}}),
        apply that size to the replacement run — Canva PPTX exports lose font sizes.
        If the replacement contains \\n, ECHTE PowerPoint-Paragraphen werden
        nach diesem Paragraph eingefuegt (mit kopierter Formatierung — Bullet,
        Font, Alignment werden vererbt).
        """
        if not para.runs:
            return
        full_text = "".join(r.text for r in para.runs)
        if "{{" not in full_text and "{{" not in _INVIS_RE.sub("", full_text):
            return
        # Strip invisible chars (soft-hyphens, zero-width spaces, etc.) that
        # PowerPoint inserts inside placeholder names (breaks regex matching)
        clean_text = _INVIS_RE.sub("", full_text)
        if "{{" not in clean_text:
            return
        # Extract font-size hint BEFORE stripping the suffix
        size_hint = None
        sh = _SIZE_HINT_RE.search(clean_text)
        if sh:
            size_hint = int(sh.group(1))
        # Use clean_text (soft-hyphens removed) for replacement
        modified = replace_text(clean_text)
        if modified != clean_text:
            # Multi-Line: erste Zeile in den aktuellen Paragraph, weitere
            # als neue Paragraphen direkt danach (mit kopierter Formatierung).
            if "\n" in modified:
                lines = modified.split("\n")
                para.runs[0].text = lines[0]
                for run in para.runs[1:]:
                    run.text = ""
                if size_hint is not None:
                    from pptx.util import Pt
                    para.runs[0].font.size = Pt(size_hint)
                # Weitere Zeilen als neue <a:p>-Elemente nach diesem einfuegen.
                # Jeder neue Paragraph bekommt zusaetzlich <a:spcBef> = 1200 (12pt),
                # damit zwischen den Bullet-Points luftiges Spacing entsteht (DQN-Look)
                # ohne leere Paragraphen einzufuegen.
                from copy import deepcopy
                from pptx.oxml.ns import qn
                from lxml import etree as _etree
                para_xml = para._p
                parent = para_xml.getparent()
                if parent is not None:
                    insert_after = para_xml
                    for line in lines[1:]:
                        new_p = deepcopy(para_xml)
                        runs = new_p.findall(qn('a:r'))
                        if runs:
                            first_t = runs[0].find(qn('a:t'))
                            if first_t is not None:
                                first_t.text = line
                            for r_extra in runs[1:]:
                                t = r_extra.find(qn('a:t'))
                                if t is not None:
                                    t.text = ""
                        # Space-Before setzen: 12pt = 1200 (val ist 100stel pt)
                        pPr = new_p.find(qn('a:pPr'))
                        if pPr is None:
                            pPr = _etree.SubElement(new_p, qn('a:pPr'))
                            new_p.insert(0, pPr)
                        # Existierendes spcBef entfernen
                        for old in pPr.findall(qn('a:spcBef')):
                            pPr.remove(old)
                        spc_before = _etree.SubElement(pPr, qn('a:spcBef'))
                        spc_pts = _etree.SubElement(spc_before, qn('a:spcPts'))
                        spc_pts.set('val', '1200')
                        # spcBef muss als erstes Kind von pPr stehen
                        pPr.insert(0, spc_before)
                        insert_after.addnext(new_p)
                        insert_after = new_p
            else:
                para.runs[0].text = modified
                for run in para.runs[1:]:
                    run.text = ""
                # Apply explicit font size if hint was present
                if size_hint is not None:
                    from pptx.util import Pt
                    para.runs[0].font.size = Pt(size_hint)

    def replace_in_textframe(tf):
        """Replace placeholders across entire text frame, including cross-paragraph splits.

        Some placeholders span two paragraphs (e.g. '{{PRODUKT_BESCHREI-' / 'BUNG}}').
        We join the full text frame, detect these, and replace them in para[0] of the span.
        """
        # First do normal per-paragraph replacement
        for para in tf.paragraphs:
            replace_in_paragraph(para)

        # Then handle cross-paragraph splits: join all paragraph texts and check
        # Build list of (para_index, full_run_text) pairs, stripping invisible chars
        para_texts = [_INVIS_RE.sub("", "".join(r.text for r in p.runs)) for p in tf.paragraphs]

        # Find any remaining {{...}} that survived (i.e. were split across paragraphs)
        # by scanning pairs of adjacent paragraphs
        for i in range(len(tf.paragraphs) - 1):
            combined = para_texts[i] + para_texts[i + 1]
            if "{{" in combined and "}}" in combined:
                replaced = replace_text(combined)
                if replaced != combined:
                    # Split the replacement back at the original boundary
                    split_at = len(para_texts[i])
                    new_p0 = replaced[:split_at] if len(replaced) >= split_at else replaced
                    new_p1 = replaced[split_at:] if len(replaced) >= split_at else ""
                    # Write to para i
                    p0 = tf.paragraphs[i]
                    if p0.runs:
                        p0.runs[0].text = new_p0
                        for r in p0.runs[1:]:
                            r.text = ""
                    # Write to para i+1
                    p1 = tf.paragraphs[i + 1]
                    if p1.runs:
                        p1.runs[0].text = new_p1
                        for r in p1.runs[1:]:
                            r.text = ""
                    # Update cache
                    para_texts[i] = new_p0
                    para_texts[i + 1] = new_p1

    def get_abs_coords(group_shape, child_shape):
        """Berechnet absolute Slide-Koordinaten eines Child-Shapes in einer Gruppe.
        Nutzt grpSpPr/xfrm aus dem p:-Namespace (nicht a:-Namespace!).
        """
        NS_A = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        NS_P = 'http://schemas.openxmlformats.org/presentationml/2006/main'
        grp_el  = group_shape._element
        grpSpPr = grp_el.find(f'{{{NS_P}}}grpSpPr')
        if grpSpPr is None:
            # fallback: simple addition
            return (
                (group_shape.left or 0) + (child_shape.left or 0),
                (group_shape.top  or 0) + (child_shape.top  or 0),
                child_shape.width  or 0,
                child_shape.height or 0,
            )
        xfrm  = grpSpPr.find(f'{{{NS_A}}}xfrm')
        if xfrm is None:
            return (
                (group_shape.left or 0) + (child_shape.left or 0),
                (group_shape.top  or 0) + (child_shape.top  or 0),
                child_shape.width  or 0,
                child_shape.height or 0,
            )
        off   = xfrm.find(f'{{{NS_A}}}off')
        ext   = xfrm.find(f'{{{NS_A}}}ext')
        chOff = xfrm.find(f'{{{NS_A}}}chOff')
        chExt = xfrm.find(f'{{{NS_A}}}chExt')
        if None in (off, ext, chOff, chExt):
            return (
                (group_shape.left or 0) + (child_shape.left or 0),
                (group_shape.top  or 0) + (child_shape.top  or 0),
                child_shape.width  or 0,
                child_shape.height or 0,
            )
        grp_x  = int(off.get('x',  0));  grp_y  = int(off.get('y',  0))
        grp_w  = int(ext.get('cx', 1));  grp_h  = int(ext.get('cy', 1))
        ch_x   = int(chOff.get('x', 0)); ch_y   = int(chOff.get('y', 0))
        ch_w   = int(chExt.get('cx', 1));ch_h   = int(chExt.get('cy', 1))
        scale_x = grp_w / ch_w if ch_w else 1
        scale_y = grp_h / ch_h if ch_h else 1
        abs_left = int(grp_x + ((child_shape.left or 0) - ch_x) * scale_x)
        abs_top  = int(grp_y + ((child_shape.top  or 0) - ch_y) * scale_y)
        abs_w    = int((child_shape.width  or 0) * scale_x)
        abs_h    = int((child_shape.height or 0) * scale_y)
        return abs_left, abs_top, abs_w, abs_h

    # Mindestgröße in EMU: nur für echte Null-/Winzig-Shapes als Fallback
    # (10 000 EMU ≈ 0,028 cm – alles darunter gilt als fehlende Dimension)
    MIN_IMG_SIZE = 10_000

    def insert_at_z(slide, pic_element, z_index):
        """Setzt pic_element an Position z_index im spTree (preserviert Z-Order des Originals).
        Indices 0+1 sind nvGrpSpPr/grpSpPr → min. Index 2.
        """
        sp_tree = slide.shapes._spTree
        sp_tree.remove(pic_element)
        sp_tree.insert(max(2, z_index), pic_element)

    def process_shape(slide, shape, image_data):
        """Ersetzt Text oder bettet Bild via blipFill ein.
        Das Template nutzt zwei Gruppen pro Bildslot:
          - Placeholder-Gruppe: solidFill Freeform + TextBox {{BILD_X}} (sichtbar, oben)
          - Target-Gruppe:      blipFill Freeform, kein Text (dahinter, wartet auf rId)
        Strategie:
          Case A (Target existiert): rId in Target-blipFill eintragen → Placeholder-Gruppe entfernen
          Case B (kein Target):      solidFill der Freeform im Placeholder durch blipFill ersetzen,
                                     TextBox-Kind entfernen
        """
        from lxml import etree

        _NS_A = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        _NS_P = 'http://schemas.openxmlformats.org/presentationml/2006/main'
        _NS_R = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'

        # ── Gruppe ──────────────────────────────────────────────────────────────
        if shape.shape_type == 6:
            # Suche TextBox-Kind mit BILD_-Placeholder
            bild_child = None
            for child in shape.shapes:
                if child.has_text_frame:
                    # Invisible chars (soft-hyphen, NBSP etc.) strippen → sonst schlägt Regex fehl!
                    txt = _INVIS_RE.sub("", child.text_frame.text).strip()
                    m = PLACEHOLDER_RE.match(txt)
                    if m:
                        key = m.group(1).lower()
                        if key in image_data and image_data[key]:
                            bild_child = (child, key)
                            break

            if bild_child is not None:
                child, key = bild_child
                try:
                    print(f"  → blipFill: key={key!r} in Gruppe {shape.name!r}")

                    # Aspect-Ratio-Crop basierend auf Placeholder-Group-Dims
                    _ph_grpSpPr = shape._element.find(f'{{{_NS_P}}}grpSpPr')
                    _ph_xfrm    = _ph_grpSpPr.find(f'{{{_NS_A}}}xfrm') if _ph_grpSpPr is not None else None
                    _ph_ext     = _ph_xfrm.find(f'{{{_NS_A}}}ext')     if _ph_xfrm    is not None else None
                    _tw = int(_ph_ext.get('cx', 0)) if _ph_ext is not None else (shape.width  or 0)
                    _th = int(_ph_ext.get('cy', 0)) if _ph_ext is not None else (shape.height or 0)
                    img_for_embed = _crop_image_to_aspect(image_data[key], _tw, _th)

                    # ── Schritt 1: Bild einbetten → rId holen ─────────────────
                    temp_pic = slide.shapes.add_picture(
                        io.BytesIO(img_for_embed), 0, 0, 914400, 914400
                    )
                    temp_el = temp_pic._element
                    blip_in_temp = temp_el.find(f'.//{{{_NS_A}}}blip')
                    new_rid = blip_in_temp.get(f'{{{_NS_R}}}embed') if blip_in_temp is not None else None
                    temp_el.getparent().remove(temp_el)   # p:pic entfernen; Relation bleibt

                    if not new_rid:
                        print(f"  FEHLER: kein rId für {key!r}")
                        return

                    print(f"    rId={new_rid!r}")

                    # ── Schritt 2: Position der Placeholder-Gruppe ─────────────
                    ph_grpSpPr = shape._element.find(f'{{{_NS_P}}}grpSpPr')
                    ph_xfrm    = ph_grpSpPr.find(f'{{{_NS_A}}}xfrm') if ph_grpSpPr is not None else None
                    ph_off     = ph_xfrm.find(f'{{{_NS_A}}}off')     if ph_xfrm    is not None else None
                    ph_x = ph_off.get('x', '0') if ph_off is not None else '0'
                    ph_y = ph_off.get('y', '0') if ph_off is not None else '0'

                    # ── Schritt 3: Target-Gruppe (blipFill, gleiche Position) suchen ──
                    target_info = None   # (target_group_shape, freeform_spPr, blipFill_el)
                    for other in slide.shapes:
                        if other.shape_id == shape.shape_id or other.shape_type != 6:
                            continue
                        o_grpSpPr = other._element.find(f'{{{_NS_P}}}grpSpPr')
                        o_xfrm    = o_grpSpPr.find(f'{{{_NS_A}}}xfrm') if o_grpSpPr is not None else None
                        o_off     = o_xfrm.find(f'{{{_NS_A}}}off')     if o_xfrm    is not None else None
                        if o_off is None:
                            continue
                        if (abs(int(o_off.get('x', 0)) - int(ph_x)) > 50000 or
                                abs(int(o_off.get('y', 0)) - int(ph_y)) > 50000):
                            continue
                        # Gleiche Position → hat diese Gruppe eine blipFill-Freeform?
                        for grp_child_o in other.shapes:
                            sp_pr_o = grp_child_o._element.find(f'{{{_NS_P}}}spPr')
                            if sp_pr_o is None:
                                continue
                            bf = sp_pr_o.find(f'{{{_NS_A}}}blipFill')
                            if bf is not None:
                                target_info = (other, sp_pr_o, bf)
                                break
                        if target_info is not None:
                            break

                    if target_info is not None:
                        # ── Case A: Target-Gruppe aktualisieren ────────────────
                        tgt_grp, tgt_spPr, tgt_blipFill = target_info
                        existing_blip = tgt_blipFill.find(f'{{{_NS_A}}}blip')
                        if existing_blip is not None:
                            existing_blip.set(f'{{{_NS_R}}}embed', new_rid)
                            print(f"  Case A ✓ blip.r:embed={new_rid!r} in {tgt_grp.name!r}")
                        else:
                            new_blip = etree.SubElement(tgt_blipFill, f'{{{_NS_A}}}blip')
                            new_blip.set(f'{{{_NS_R}}}embed', new_rid)
                            print(f"  Case A ✓ neues a:blip in {tgt_grp.name!r}")
                        # Reset fillRect so new image fills the shape without template crop offsets
                        stretch = tgt_blipFill.find(f'{{{_NS_A}}}stretch')
                        if stretch is None:
                            stretch = etree.SubElement(tgt_blipFill, f'{{{_NS_A}}}stretch')
                        fr = stretch.find(f'{{{_NS_A}}}fillRect')
                        if fr is None:
                            etree.SubElement(stretch, f'{{{_NS_A}}}fillRect')
                        else:
                            fr.attrib.clear()   # remove l/t/r/b crop offsets → full fill
                        # Placeholder-Gruppe entfernen (solidFill+TextBox waren oben drüber)
                        shape._element.getparent().remove(shape._element)

                    else:
                        # ── Case B: solidFill in eigener Freeform → blipFill ──
                        print(f"  Case B: kein Target für {key!r} @ ({ph_x},{ph_y})")
                        # TextBox-Placeholder-Kind entfernen
                        for grp_child_b in list(shape.shapes):
                            if grp_child_b.has_text_frame:
                                txt_b = grp_child_b.text_frame.text.strip()
                                if PLACEHOLDER_RE.match(txt_b):
                                    grp_child_b._element.getparent().remove(grp_child_b._element)
                                    break
                        # solidFill in der Freeform durch blipFill ersetzen
                        for grp_child_b in list(shape.shapes):
                            sp_pr_b = grp_child_b._element.find(f'{{{_NS_P}}}spPr')
                            if sp_pr_b is None:
                                continue
                            solid = sp_pr_b.find(f'{{{_NS_A}}}solidFill')
                            if solid is None:
                                continue
                            idx = list(sp_pr_b).index(solid)
                            sp_pr_b.remove(solid)
                            bf_el = etree.Element(f'{{{_NS_A}}}blipFill')
                            bl_el = etree.SubElement(bf_el, f'{{{_NS_A}}}blip')
                            bl_el.set(f'{{{_NS_R}}}embed', new_rid)
                            st_el = etree.SubElement(bf_el, f'{{{_NS_A}}}stretch')
                            etree.SubElement(st_el, f'{{{_NS_A}}}fillRect')
                            sp_pr_b.insert(idx, bf_el)
                            print(f"  Case B ✓ solidFill→blipFill in {grp_child_b.name!r}")
                            break

                except Exception as e:
                    print(f"  blipFill Fehler {shape.name!r}/{key!r}: {e}")
                    import traceback; traceback.print_exc()
                return

            # Keine BILD_-Gruppe → Text in allen Kindern ersetzen
            for child in list(shape.shapes):
                try:
                    if child.has_text_frame:
                        replace_in_textframe(child.text_frame)
                except Exception as e:
                    print(f"  Gruppen-Child Fehler {child.name}: {e}")
            if shape.has_text_frame:
                replace_in_textframe(shape.text_frame)
            return

        # ── Top-Level Non-Group Shape ────────────────────────────────────────
        sp_tree = slide.shapes._spTree
        shape_name_lower = (shape.name or "").lower()

        # 1. Bild per Shape-Name
        if shape_name_lower in image_data and image_data[shape_name_lower]:
            try:
                left   = shape.left   or 0
                top    = shape.top    or 0
                width  = shape.width  or 0
                height = shape.height or 0
                if width  < MIN_IMG_SIZE: width  = 3_000_000
                if height < MIN_IMG_SIZE: height = 2_400_000
                shape_z = list(sp_tree).index(shape._element)
                shape._element.getparent().remove(shape._element)
                _img = _crop_image_to_aspect(image_data[shape_name_lower], width, height)
                _pic = slide.shapes.add_picture(
                    io.BytesIO(_img), left, top, width, height
                )
                insert_at_z(slide, _pic._element, shape_z)
                return
            except Exception as e:
                print(f"Bild-Ersatz Fehler (name) {shape_name_lower}: {e}")

        # 2. Bild per Text-Inhalt
        if shape.has_text_frame:
            txt = _INVIS_RE.sub("", shape.text_frame.text).strip()
            m = PLACEHOLDER_RE.match(txt)
            if m:
                key = m.group(1).lower()
                if key in image_data and image_data[key]:
                    try:
                        # Check if this placeholder TextBox sits inside a solidFill group
                        # (e.g. BILD_LAGEPLAN inside dark left-panel Group) → replace that
                        # group's solidFill with the image instead of inserting a tiny picture.
                        ph_cx = (shape.left or 0) + (shape.width or 0) // 2
                        ph_cy = (shape.top  or 0) + (shape.height or 0) // 2
                        covering_target = None
                        for other in slide.shapes:
                            if other.shape_id == shape.shape_id or other.shape_type != 6:
                                continue
                            g_left  = other.left  or 0
                            g_top   = other.top   or 0
                            g_right = g_left + (other.width  or 0)
                            g_bot   = g_top  + (other.height or 0)
                            if not (g_left <= ph_cx <= g_right and g_top <= ph_cy <= g_bot):
                                continue
                            for gc in other.shapes:
                                sp_pr_gc = gc._element.find(f'{{{_NS_P}}}spPr')
                                if sp_pr_gc is None:
                                    continue
                                solid_gc = sp_pr_gc.find(f'{{{_NS_A}}}solidFill')
                                if solid_gc is None:
                                    continue
                                covering_target = (other, gc, sp_pr_gc, solid_gc)
                                break
                            if covering_target:
                                break

                        if covering_target:
                            # Embed image → get rId
                            grp, grp_child, sp_pr_gc, solid_gc = covering_target
                            _gw = grp.width  or 0
                            _gh = grp.height or 0
                            _img = _crop_image_to_aspect(image_data[key], _gw, _gh)
                            temp_pic = slide.shapes.add_picture(
                                io.BytesIO(_img), 0, 0, 914400, 914400
                            )
                            temp_el = temp_pic._element
                            blip_in_temp = temp_el.find(f'.//{{{_NS_A}}}blip')
                            new_rid = blip_in_temp.get(f'{{{_NS_R}}}embed') if blip_in_temp is not None else None
                            temp_el.getparent().remove(temp_el)
                            if new_rid:
                                idx = list(sp_pr_gc).index(solid_gc)
                                sp_pr_gc.remove(solid_gc)
                                bf_el = etree.Element(f'{{{_NS_A}}}blipFill')
                                bl_el = etree.SubElement(bf_el, f'{{{_NS_A}}}blip')
                                bl_el.set(f'{{{_NS_R}}}embed', new_rid)
                                st_el = etree.SubElement(bf_el, f'{{{_NS_A}}}stretch')
                                etree.SubElement(st_el, f'{{{_NS_A}}}fillRect')
                                sp_pr_gc.insert(idx, bf_el)
                                shape._element.getparent().remove(shape._element)
                                print(f"  Panel fill: {key!r} → group {grp.name!r}")
                                return

                        # Fallback: insert picture at TextBox dimensions
                        left   = shape.left   or 0
                        top    = shape.top    or 0
                        width  = shape.width  or 0
                        height = shape.height or 0
                        if width  < MIN_IMG_SIZE: width  = 3_000_000
                        if height < MIN_IMG_SIZE: height = 2_400_000
                        shape_z = list(sp_tree).index(shape._element)
                        shape._element.getparent().remove(shape._element)
                        _img = _crop_image_to_aspect(image_data[key], width, height)
                        _pic = slide.shapes.add_picture(
                            io.BytesIO(_img), left, top, width, height
                        )
                        insert_at_z(slide, _pic._element, shape_z)
                        return
                    except Exception as e:
                        print(f"Bild-Ersatz Fehler (text) {key}: {e}")
                        import traceback; traceback.print_exc()
                        return  # IMPORTANT: don't text-substitute on failure → leave {{BILD_X}}

        # 3. Text ersetzen – nur für Nicht-Bild-Shapes
        # (bild_* keys sind aus REPL_MAP ausgeschlossen, also kann diese Zeile
        #  keine Bild-URL als Text einsetzen – aber return oben verhindert auch die
        #  Fallthrough-Situation bei Shape-Typ 2 = image slot)
        if shape.has_text_frame:
            replace_in_textframe(shape.text_frame)

    # Duplicate WE slides BEFORE text/image replacement so placeholders are still intact
    duplicate_we_slides(prs, data)

    # ── 6-gute-Gruende-Override: Template hat 6 fest hinterlegte Titel + Texte.
    # Wenn Claude dynamische Werte (text_grund_X_titel/_text) liefert, ueber-
    # schreiben wir die Defaults BEVOR der Standard-Replace-Loop laeuft.
    # Substring-Match gegen die ORIGINAL-Default-Texte des Templates.
    _override_grund_texts(prs, data)

    # ── Rechtliche-Hinweise: hartkodierte Magdeburg-Referenzen ('Diamant
    # Brauerei') durch projekt-spezifische Texte ersetzen. Nur aktiv wenn
    # Stadt != Magdeburg.
    _override_legal_text(prs, data)

    # ── Zahlen-Spalten-Geometrie normalisieren VOR dem Replace,
    # solange die Platzhalter-Strings noch identifizierbar sind.
    _normalize_number_columns(prs)

    for slide in prs.slides:
        for shape in list(slide.shapes):
            try:
                process_shape(slide, shape, image_data)
            except Exception as e:
                print(f"Shape-Fehler slide={slide.slide_id} shape={shape.name}: {e}")

    # ── Brute-Force-Fallback: jeder verbliebene {{BILD_X}} wird per add_picture
    # auf seine TextBox-Position gesetzt, sofern image_data[X] vorhanden ist.
    # Greift auch bei tief verschachtelten Groups die process_shape uebersprungen
    # hat. Wichtig: nur ausfuehren wenn der Key wirklich Bilddaten hat — sonst
    # den Platzhalter sichtbar lassen.
    BILD_INLINE_RE = re.compile(r'\{\{\s*(BILD_[A-Z0-9_-]+)\s*(?:\|[^}]*)?\}\}', re.IGNORECASE)

    def _walk_with_offset(top_shape, off_x=0, off_y=0):
        """Yield (shape, abs_left, abs_top) für top_shape und alle (rekursiven)
        children in Groups."""
        try:
            sl = (top_shape.left or 0) + off_x
            st = (top_shape.top  or 0) + off_y
        except Exception:
            sl, st = off_x, off_y
        yield top_shape, sl, st
        if top_shape.shape_type == 6:  # Group
            for c in top_shape.shapes:
                yield from _walk_with_offset(c, sl, st)

    bf_embedded = 0
    bf_errors   = 0
    try:
        for slide in prs.slides:
            try:
                top_list = list(slide.shapes)
            except Exception:
                continue
            for top_shape in top_list:
                try:
                    walked = list(_walk_with_offset(top_shape))
                except Exception:
                    continue
                for shape, abs_l, abs_t in walked:
                    try:
                        if not shape.has_text_frame:
                            continue
                        txt = _INVIS_RE.sub("", shape.text_frame.text)
                        m = BILD_INLINE_RE.search(txt)
                        if not m:
                            continue
                        key = m.group(1).lower()
                        if key not in image_data or not image_data[key]:
                            continue
                        w = shape.width  or 0
                        h = shape.height or 0
                        if w < MIN_IMG_SIZE: w = 3_000_000
                        if h < MIN_IMG_SIZE: h = 2_400_000
                        img = _crop_image_to_aspect(image_data[key], w, h)
                        slide.shapes.add_picture(io.BytesIO(img), abs_l, abs_t, w, h)
                        parent = shape._element.getparent()
                        if parent is not None:
                            parent.remove(shape._element)
                        bf_embedded += 1
                        print(f"  [BF-Embed] {key} → slide_id={slide.slide_id} "
                              f"@ ({abs_l},{abs_t}) {w}x{h}")
                    except Exception as e:
                        bf_errors += 1
                        print(f"  [BF-Embed] Fehler in shape: {e}")
    except Exception as e:
        print(f"  [BF-Embed] Top-Level-Fehler abgefangen: {e}")
    if bf_embedded or bf_errors:
        print(f"  Brute-Force-Embed: {bf_embedded} Bild(er) eingesetzt, {bf_errors} Fehler")

    # Cleanup-Pass: Template-Texte die "in Euro" statt "in €" enthalten korrigieren
    _euro_fixes = [("in Euro", "in €"), (" Euro", " €"), ("in EUR", "in €"), (" EUR", " €")]
    for slide in prs.slides:
        for shape in slide.shapes:
            shapes_to_check = [shape]
            if shape.shape_type == 6:
                shapes_to_check += list(shape.shapes)
            for s in shapes_to_check:
                if not s.has_text_frame:
                    continue
                for para in s.text_frame.paragraphs:
                    for run in para.runs:
                        for old, new in _euro_fixes:
                            if old in run.text:
                                run.text = run.text.replace(old, new)

    # Final-Cleanup: Übrig gebliebene {{...}}-Platzhalter entfernen.
    # Tritt auf, wenn (a) Claude einen Key gar nicht zurückgegeben hat, (b) Template-
    # Key nicht in PLATZHALTER, (c) Whitespace im Key durch Run-Splitting. Statt
    # "{{TEXT_X}}" im finalen Dokument lieber leere Stelle. {{BILD_*}} bleibt erhalten,
    # damit die Placeholder-Shapes (auf die der Kunde klickt zum Upload) sichtbar bleiben.
    _LEFTOVER_RE = re.compile(r'\{\{\s*(?!BILD_)[A-Z0-9_\s\-|.]+\}\}', re.IGNORECASE)
    leftover_count = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            shapes_to_check = [shape]
            if shape.shape_type == 6:
                shapes_to_check += list(shape.shapes)
            for s in shapes_to_check:
                if not s.has_text_frame:
                    continue
                for para in s.text_frame.paragraphs:
                    if not para.runs:
                        continue
                    full = "".join(r.text for r in para.runs)
                    if "{{" not in full:
                        continue
                    cleaned = _LEFTOVER_RE.sub("", full)
                    if cleaned != full:
                        leftover_count += full.count("{{") - cleaned.count("{{")
                        para.runs[0].text = cleaned
                        for r in para.runs[1:]:
                            r.text = ""
    if leftover_count:
        print(f"  Cleanup: {leftover_count} unbefüllte Text-Platzhalter entfernt")

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()


# Gruppen von Platzhaltern, deren Textfelder geometrisch normalisiert werden.
# Innerhalb einer Gruppe: alle Felder bekommen die gleiche X-Position, Width und
# Height (jeweils das Median der Gruppe). Vertikale Y-Position bleibt unverändert.
_NUMBER_COLUMN_GROUPS = [
    # Slide 5 — Standort-Minuten (Zahl-Bereich)
    ("min_uni", "min_bahnhof", "min_altstadt"),
    ("label_min_uni", "label_min_bahnhof", "label_min_altstadt"),
    # Slide 6 — Hotel-Features
    ("feature_1_zahl", "feature_2_zahl", "feature_3_zahl"),
    ("feature_1_label", "feature_2_label", "feature_3_label"),
    # Slide 21 — Grundriss-Anzahl (Zahlen)
    ("anzahl_1zi", "anzahl_2zi", "anzahl_barrierefrei"),
]


def _normalize_number_columns(prs):
    """Macht Zahlen-Textfelder einer Gruppe geometrisch konsistent.

    Hintergrund: Im Canva-Marketing-Template wurden die Zahlen-Felder
    (z.B. {{MIN_UNI}}, {{MIN_BAHNHOF}}, {{MIN_ALTSTADT}}) leicht unterschiedlich
    positioniert + verschieden breit/hoch. Beim Render wirken die Zahlen dann
    nicht sauber untereinander zentriert. Diese Funktion sucht alle Shapes mit
    einem Platzhalter der zur selben Gruppe gehört und setzt deren left, width,
    height auf den Median-Wert der Gruppe. top bleibt unverändert (vertikale
    Stapelung beibehalten).
    """
    import re as _re
    PH = _re.compile(r'\{\{\s*([a-z][a-z0-9_]*)', _re.IGNORECASE)

    for slide in prs.slides:
        # Sammle alle Shapes pro Platzhalter-Key auf diesem Slide
        shapes_by_key = {}
        for shape in slide.shapes:
            try:
                if not shape.has_text_frame:
                    continue
            except Exception:
                continue
            text = shape.text_frame.text
            m = PH.search(text)
            if not m:
                continue
            key = m.group(1).lower()
            shapes_by_key[key] = shape

        # Font-Size-Hint Pattern: {{KEY|50pt}}
        SIZE_HINT = _re.compile(r'\|\s*(\d+)\s*pt\b', _re.IGNORECASE)

        for group in _NUMBER_COLUMN_GROUPS:
            present = [shapes_by_key[k] for k in group if k in shapes_by_key]
            if len(present) < 2:
                continue
            xs = sorted(s.left for s in present)
            ws = sorted(s.width for s in present)
            hs = sorted(s.height for s in present)
            median_x = xs[len(xs) // 2]
            median_w = ws[len(ws) // 2]
            target_h = max(hs)

            # Mindesthoehe basierend auf Font-Size-Hint sicherstellen.
            # Canva-Export hat die Textbox-Hoehe oft nur ~95885 EMU (0.1cm)
            # gesetzt — viel zu klein fuer 33-50pt Schrift. Daraus folgt:
            # Text wird oben/unten geclippt, scheint nicht zentriert.
            # 1 pt = 12700 EMU; wir geben 1.8x Schrifthoehe als Mindestbox.
            for s in present:
                m = SIZE_HINT.search(s.text_frame.text)
                if m:
                    font_pt = int(m.group(1))
                    min_h = int(font_pt * 12700 * 1.8)
                    if min_h > target_h:
                        target_h = min_h

            from pptx.enum.text import MSO_ANCHOR
            for s in present:
                # Box-Mitte beim Vergroessern erhalten: top wird nach oben
                # verschoben, damit der visuelle Mittelpunkt (= Zahl-Position
                # mit vertical_anchor=MIDDLE) gleich bleibt wie im Template.
                orig_center = s.top + s.height // 2
                s.left   = median_x
                s.width  = median_w
                s.height = target_h
                s.top    = orig_center - target_h // 2
                try:
                    s.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                except Exception:
                    pass

            # Vertikale Symmetrie: erste + letzte top bleiben, mittlere werden
            # linear zwischen ihnen interpoliert. Zugehoerige Labels (festes
            # Text-Shape direkt unter der Zahl, horizontal ueberlappend) werden
            # mit dem gleichen delta_y mitbewegt — sonst trennt sich die Zahl
            # vom Label-Text. Platzhalter-Shapes ({{...}}) werden NICHT als
            # Label genommen, weil sie evtl. in einer eigenen Gruppe sitzen
            # und dort selbst normalisiert werden (sonst doppelte Verschiebung).
            if len(present) >= 3:
                num_h = target_h
                labels_for = {}
                for num_shape in present:
                    n_left   = num_shape.left
                    n_right  = num_shape.left + num_shape.width
                    n_bottom = num_shape.top  + num_shape.height
                    candidates = []
                    for other in slide.shapes:
                        if other is num_shape:
                            continue
                        try:
                            if not other.has_text_frame:
                                continue
                        except Exception:
                            continue
                        other_text = other.text_frame.text.strip()
                        if not other_text or "{{" in other_text:
                            continue
                        o_center_x = other.left + other.width // 2
                        # Horizontal: Label-Mittelpunkt innerhalb erweiterter Zahl-Box
                        if not (n_left - num_shape.width <= o_center_x <= n_right + num_shape.width):
                            continue
                        # Vertikal: Label muss unter der Zahl-Box sein, max ~4 Box-Hoehen weg
                        if other.top < num_shape.top:
                            continue
                        dist = other.top - n_bottom
                        if dist > num_h * 4:
                            continue
                        candidates.append((dist, other))
                    if candidates:
                        candidates.sort(key=lambda x: x[0])
                        labels_for[id(num_shape)] = candidates[0][1]

                sorted_shapes  = sorted(present, key=lambda s: s.top)
                tops_sorted    = [s.top for s in sorted_shapes]
                y_first, y_last = tops_sorted[0], tops_sorted[-1]
                n              = len(sorted_shapes)
                even_tops      = [y_first + (y_last - y_first) * i // (n - 1) for i in range(n)]
                for shape, new_top in zip(sorted_shapes, even_tops):
                    delta = new_top - shape.top
                    if delta == 0:
                        continue
                    shape.top = new_top
                    lbl = labels_for.get(id(shape))
                    if lbl is not None:
                        lbl.top = lbl.top + delta

            print(f"  ↔ Zahlen-Spalte normalisiert: {[k for k in group if k in shapes_by_key]} "
                  f"(x={median_x}, w={median_w}, h={target_h})")

# ── Slot-Labels für Preview-UI: gibt {bild_key: human-readable label} zurück ──
_SLOT_LABELS = {
    "bild_titel":              "Titelbild (Außenansicht)",
    "bild_projekt_aussen":     "Projekt – Außenansicht",
    "bild_projekt":            "Projekt – Bild",
    "bild_quartier":           "Quartier / Umgebung",
    "bild_greenliving_1":      "Green Living – Bild 1",
    "bild_greenliving_2":      "Green Living – Bild 2",
    "bild_interior":           "Innenraum / Interior",
    "bild_ausstattung_1":      "Ausstattung 1",
    "bild_ausstattung_2":      "Ausstattung 2",
    "bild_ausstattung_3":      "Ausstattung 3",
    "bild_ausstattung_4":      "Ausstattung 4",
    "bild_ausstattung_5":      "Ausstattung 5",
    "bild_ausstattung_6":      "Ausstattung 6",
    "bild_ansicht_1":          "Außenansicht 1",
    "bild_ansicht_2":          "Außenansicht 2",
    "bild_standort_innen":     "Standort innen",
    "bild_standort_aussen":    "Standort außen",
    "bild_lageplan":           "Lageplan",
    "bild_stadt_gross":        "Stadtbild groß",
    "bild_stadt_klein":        "Stadtbild klein",
    "bild_grundriss_intro_1":  "Grundriss-Intro 1",
    "bild_grundriss_intro_2":  "Grundriss-Intro 2",
    "bild_grundriss_intro_3":  "Grundriss-Intro 3",
    "bild_grundriss_1":        "Grundriss 1",
    "bild_grundriss_2":        "Grundriss 2",
    "bild_grundriss_3":        "Grundriss 3",
    "bild_grundriss_4":        "Grundriss 4",
    "bild_collage_1":          "Collage 1",
    "bild_collage_2":          "Collage 2",
    "bild_collage_3":          "Collage 3",
    "bild_collage_4":          "Collage 4",
    "bild_collage_5":          "Collage 5",
    "bild_hotel_1":            "Hotel-Feeling 1",
    "bild_hotel_2":            "Hotel-Feeling 2",
    "bild_rechtlich_1":        "Rechtliches – Bild 1",
    "bild_rechtlich_2":        "Rechtliches – Bild 2",
    "bild_stadt_presse":       "Stadtbild – Presse",
    "bild_stadt_branche":      "Stadtbild – Branche",
}
def _slot_label(key):
    if key in _SLOT_LABELS:
        return _SLOT_LABELS[key]
    m = re.match(r"^bild_amenity_(\d+)$", key)
    if m:
        return f"Ausstattung Amenity {m.group(1)}"
    m = re.match(r"^bild_we_(\d+)$", key)
    if m:
        return f"WE-Grundriss {m.group(1)}"
    return key.replace("_", " ").title()


def extract_bild_placeholders(pptx_bytes):
    """Scannt das (gefüllte oder leere) PPTX und extrahiert für jeden {{BILD_*}}
    Platzhalter die absolute Slide-Position in EMU sowie die Slide-Größe.

    Gibt zurück:
      {
        "slide_w_emu": int, "slide_h_emu": int,
        "slides": [
          {"index": 0, "placeholders": [
              {"slot": "bild_titel", "x": 12345, "y": 12345, "w": 1234, "h": 1234,
               "label": "Titelbild"}
          ]}, ...
        ]
      }
    Positionen kommen aus der Gruppe (grpSpPr/xfrm) bzw. dem Shape selbst.
    """
    NS_A = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    NS_P = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    BILD_RE = re.compile(r'\{\{\s*(BILD_[A-Z0-9_]+)\s*(?:\|[^}]*)?\}\}', re.IGNORECASE)

    prs = Presentation(io.BytesIO(pptx_bytes))
    sw  = int(prs.slide_width  or 12192000)
    sh  = int(prs.slide_height or 6858000)

    out_slides = []
    for s_idx, slide in enumerate(prs.slides):
        placeholders = []
        seen_slots = set()

        def _record(slot, x, y, w, h):
            slot_l = slot.lower()
            if slot_l in seen_slots:
                return
            seen_slots.add(slot_l)
            placeholders.append({
                "slot":  slot_l,
                "x":     int(x or 0),
                "y":     int(y or 0),
                "w":     int(w or 0),
                "h":     int(h or 0),
                "label": _slot_label(slot_l),
            })

        for shape in slide.shapes:
            # Group shapes: scan children for BILD_-Text-Platzhalter
            if shape.shape_type == 6:
                grpSpPr = shape._element.find(f'{{{NS_P}}}grpSpPr')
                xfrm    = grpSpPr.find(f'{{{NS_A}}}xfrm') if grpSpPr is not None else None
                off     = xfrm.find(f'{{{NS_A}}}off')     if xfrm    is not None else None
                ext     = xfrm.find(f'{{{NS_A}}}ext')     if xfrm    is not None else None
                gx = int(off.get('x', 0)) if off is not None else 0
                gy = int(off.get('y', 0)) if off is not None else 0
                gw = int(ext.get('cx', 0)) if ext is not None else 0
                gh = int(ext.get('cy', 0)) if ext is not None else 0
                for child in shape.shapes:
                    if not child.has_text_frame:
                        continue
                    txt = child.text_frame.text
                    m = BILD_RE.search(txt)
                    if m:
                        _record(m.group(1), gx, gy, gw, gh)
                        break
            elif shape.has_text_frame:
                m = BILD_RE.search(shape.text_frame.text)
                if m:
                    _record(m.group(1), shape.left, shape.top, shape.width, shape.height)

        if placeholders:
            out_slides.append({"index": s_idx, "placeholders": placeholders})

    return {
        "slide_w_emu": sw,
        "slide_h_emu": sh,
        "slides":      out_slides,
    }


def _find_pdftoppm():
    """Sucht den pdftoppm Binary (poppler-utils)."""
    import shutil as _shutil
    p = _shutil.which("pdftoppm")
    if p:
        return p
    for c in ("/usr/bin/pdftoppm", "/usr/local/bin/pdftoppm", "/opt/homebrew/bin/pdftoppm"):
        if os.path.isfile(c):
            return c
    return None


def _render_pdf_to_jpgs_pymupdf(pdf_bytes, out_dir, dpi=110):
    """Fallback-Renderer mit PyMuPDF (fitz). Funktioniert ohne System-Binaries.
    Liefert dieselbe Liste an slide_<n>.jpg-Pfaden wie der pdftoppm-Pfad."""
    import fitz  # PyMuPDF
    os.makedirs(out_dir, exist_ok=True)
    doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    try:
        zoom = dpi / 72.0
        mat  = fitz.Matrix(zoom, zoom)
        out  = []
        for i, page in enumerate(doc, start=1):
            pix = page.get_pixmap(matrix=mat, alpha=False)
            path = os.path.join(out_dir, f"slide_{i}.jpg")
            pix.pil_save(path, format="JPEG", quality=82)
            out.append(path)
            del pix
            if i % 5 == 0:
                import gc as _gc
                _gc.collect()
        for p in out:
            _trim_white_borders(p)
        return out
    finally:
        doc.close()


def _add_hyperlinks_to_pdf(pdf_bytes):
    """Findet alle URLs (http/https) im PDF-Text und macht sie zu klickbaren
    Hyperlinks (PyMuPDF link annotations). DQN macht das auch — Quellen sind
    so direkt aus dem PDF aufrufbar.
    Gibt das modifizierte PDF zurück. Bei Fehler: Original unverändert.
    """
    try:
        import fitz  # PyMuPDF
        import re as _re
        # URL-Regex: http(s):// + alles bis Whitespace/typische End-Zeichen
        url_re = _re.compile(r'https?://[^\s)\]>"\']+', _re.IGNORECASE)
        doc = fitz.open(stream=pdf_bytes, filetype="pdf")
        added = 0
        for page in doc:
            text = page.get_text("text")
            seen = set()
            for m in url_re.finditer(text):
                url = m.group(0).rstrip('.,;:')
                if url in seen:
                    continue
                seen.add(url)
                # Suche alle Vorkommen im Page-Layout (mit Bounding-Boxen)
                rects = page.search_for(url)
                if not rects:
                    # Vielleicht nur Anfang des Strings (URLs wickeln sich oft um)
                    short = url[:60]
                    rects = page.search_for(short)
                for r in rects:
                    try:
                        page.insert_link({
                            "kind": fitz.LINK_URI,
                            "from": r,
                            "uri":  url,
                        })
                        added += 1
                    except Exception:
                        pass
        out = doc.tobytes(garbage=3, deflate=True)
        doc.close()
        if added:
            print(f"  Hyperlinks: {added} URL-Anker im PDF hinzugefuegt")
        return out
    except Exception as e:
        print(f"  Hyperlinks Fehler (Original-PDF unveraendert): {e}")
        return pdf_bytes


def _trim_white_borders(jpg_path):
    """Entfernt einheitliche weiße Margins (CloudConvert padded 16:9-Slides
    in A4-Landscape → weiße Streifen oben/unten). Erkennt nur reine 255-er
    Bereiche, lässt cremefarbene/dunkle Slide-Inhalte unangetastet."""
    try:
        from PIL import Image
        img = Image.open(jpg_path).convert("RGB")
        w, h = img.size
        if w < 100 or h < 100:
            return
        # Sample 50 columns evenly distributed
        sample_x = list(range(0, w, max(1, w // 50)))
        pixels = img.load()

        def _is_white_row(y):
            return all(pixels[x, y] == (255, 255, 255) for x in sample_x)

        # Top boundary
        top = 0
        while top < h and _is_white_row(top):
            top += 1
        # Bottom boundary
        bottom = h
        while bottom > top + 1 and _is_white_row(bottom - 1):
            bottom -= 1

        # Nur croppen wenn deutliche Margins (mind. 2% des Bildes)
        margin_threshold = max(int(h * 0.02), 10)
        if top < margin_threshold and (h - bottom) < margin_threshold:
            return  # nichts zu tun
        cropped = img.crop((0, top, w, bottom))
        cropped.save(jpg_path, "JPEG", quality=85, optimize=True)
    except Exception as e:
        print(f"  _trim_white_borders Fehler {jpg_path}: {e}")


def render_pdf_to_jpgs(pdf_bytes, out_dir, dpi=110):
    """Konvertiert PDF zu einzelnen JPG-Slide-Bildern.
    Bevorzugt pdftoppm (poppler-utils, schnell), Fallback PyMuPDF (pure Python).
    Schreibt slide_1.jpg, slide_2.jpg, ... ins out_dir.
    """
    import subprocess, tempfile
    bin_path = _find_pdftoppm()
    if not bin_path:
        # Fallback: PyMuPDF (in requirements.txt) – funktioniert ohne System-Binaries
        try:
            print("  pdftoppm fehlt → Fallback PyMuPDF")
            return _render_pdf_to_jpgs_pymupdf(pdf_bytes, out_dir, dpi=dpi)
        except Exception as e:
            raise RuntimeError(
                f"pdftoppm nicht gefunden UND PyMuPDF-Fallback gescheitert: {e}"
            )

    os.makedirs(out_dir, exist_ok=True)
    with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as f:
        f.write(pdf_bytes)
        pdf_path = f.name

    try:
        # Erst Seitenanzahl feststellen via pdfinfo (oder fallback: pdftoppm einzeln)
        page_count = None
        try:
            pdfinfo = subprocess.run(["pdfinfo", pdf_path],
                                     capture_output=True, text=True, timeout=30)
            for line in pdfinfo.stdout.splitlines():
                if line.startswith("Pages:"):
                    page_count = int(line.split()[1])
                    break
        except Exception:
            pass
        if not page_count:
            page_count = 50  # Defensive Obergrenze

        print(f"  pdftoppm: rendere {page_count} Seiten einzeln @ {dpi} DPI …")
        prefix = os.path.join(out_dir, "slide")
        for n in range(1, page_count + 1):
            result = subprocess.run(
                [bin_path, "-jpeg", "-r", str(dpi), "-f", str(n), "-l", str(n),
                 pdf_path, prefix],
                capture_output=True, text=True, timeout=60
            )
            if result.returncode != 0:
                # Wenn keine Datei rauskam → wahrscheinlich Ende des PDFs
                print(f"  Seite {n}: rc={result.returncode} → stop")
                break
            # Prüfen ob die Seite tatsächlich erzeugt wurde
            import glob as _glob
            if not _glob.glob(os.path.join(out_dir, f"slide-{n}.jpg")) and \
               not _glob.glob(os.path.join(out_dir, f"slide-{n:02d}.jpg")) and \
               not _glob.glob(os.path.join(out_dir, f"slide-{n:03d}.jpg")):
                # Keine Datei für diese Seite → wahrscheinlich Ende
                break
            if n % 5 == 0:
                # gc nach jedem 5er-Block
                import gc as _gc
                _gc.collect()
    finally:
        try: os.unlink(pdf_path)
        except OSError: pass

    import glob as _glob
    files = sorted(_glob.glob(os.path.join(out_dir, "slide-*.jpg")))
    # Normalize to slide_<n>.jpg (1-based, no zero-pad)
    normalized = []
    for p in files:
        m = re.search(r'slide-(\d+)\.jpg$', p)
        if not m:
            continue
        n = int(m.group(1))
        new_path = os.path.join(out_dir, f"slide_{n}.jpg")
        if new_path != p:
            os.rename(p, new_path)
        normalized.append(new_path)
    normalized.sort(key=lambda p: int(re.search(r'slide_(\d+)\.jpg$', p).group(1)))
    # Weiße Margins (CloudConvert-Padding) wegcroppen
    for p in normalized:
        _trim_white_borders(p)
    return normalized


def _convert_to_pdf_cloudconvert(pptx_bytes, filename, max_wait_s=300):
    """Konvertiert PPTX → PDF über CloudConvert API."""
    if not CLOUDCONVERT_KEY:
        raise RuntimeError("CLOUDCONVERT_KEY nicht gesetzt")

    base = "https://api.cloudconvert.com/v2"
    headers = {"Authorization": f"Bearer {CLOUDCONVERT_KEY}"}
    print(f"convert_to_pdf (CloudConvert): {filename} ({len(pptx_bytes)//1024} KB)")

    # 1) Job mit drei Tasks anlegen: upload → convert → export
    job_resp = requests.post(
        f"{base}/jobs",
        headers={**headers, "Content-Type": "application/json"},
        json={
            "tasks": {
                "import-1": {"operation": "import/upload"},
                "convert-1": {
                    "operation":     "convert",
                    "input":         "import-1",
                    "input_format":  "pptx",
                    "output_format": "pdf",
                },
                "export-1": {"operation": "export/url", "input": "convert-1"},
            }
        },
        timeout=30,
    )
    job_resp.raise_for_status()
    job   = job_resp.json()["data"]
    jobid = job["id"]

    # 2) Upload-URL aus import-Task holen + Datei dorthin POSTen
    import_task = next(t for t in job["tasks"] if t["name"] == "import-1")
    form        = import_task["result"]["form"]
    up_resp = requests.post(
        form["url"],
        data=form["parameters"],
        files={"file": (filename, pptx_bytes,
                        "application/vnd.openxmlformats-officedocument.presentationml.presentation")},
        timeout=300,
    )
    up_resp.raise_for_status()
    print(f"  PPTX hochgeladen, warte auf Konvertierung …")

    # 3) Job-Status pollen bis "finished" / "error"
    deadline = _time.time() + max_wait_s
    pdf_url  = None
    while _time.time() < deadline:
        _time.sleep(3)
        st_resp = requests.get(f"{base}/jobs/{jobid}", headers=headers, timeout=20)
        st_resp.raise_for_status()
        jdata = st_resp.json()["data"]
        status = jdata.get("status")
        if status == "finished":
            export_task = next(t for t in jdata["tasks"] if t["name"] == "export-1")
            files = (export_task.get("result") or {}).get("files") or []
            if not files:
                raise RuntimeError("CloudConvert export-1: keine Dateien")
            pdf_url = files[0]["url"]
            break
        if status == "error":
            err = ""
            for t in jdata["tasks"]:
                if t.get("status") == "error":
                    err = (t.get("message") or t.get("code") or "unbekannt")
                    break
            raise RuntimeError(f"CloudConvert Fehler: {err}")

    if not pdf_url:
        raise RuntimeError("CloudConvert Timeout nach 5 Minuten")

    # 4) PDF herunterladen
    pdf_resp = requests.get(pdf_url, timeout=120)
    pdf_resp.raise_for_status()
    print(f"  PDF erzeugt (CloudConvert): {len(pdf_resp.content)//1024} KB")
    return pdf_resp.content


def convert_to_pdf(pptx_bytes, filename):
    """Konvertiert PPTX-Bytes zu PDF via CloudConvert."""
    if not CLOUDCONVERT_KEY:
        raise RuntimeError("CLOUDCONVERT_KEY nicht gesetzt")
    return _convert_to_pdf_cloudconvert(pptx_bytes, filename)


def _can_convert_to_pdf():
    return bool(CLOUDCONVERT_KEY)

def assemble_session(session_id):
    """Liest alle Chunks einer Session von /tmp und gibt die assemblierten Bytes zurück."""
    session_dir = os.path.join(CHUNK_DIR, session_id)
    if not os.path.isdir(session_dir):
        raise ValueError(f"Session {session_id} nicht gefunden")
    meta_path = os.path.join(session_dir, "meta.json")
    with open(meta_path) as f:
        meta = json.load(f)
    total = meta["total_chunks"]
    assembled = b""
    for i in range(total):
        chunk_path = os.path.join(session_dir, f"chunk_{i:04d}")
        with open(chunk_path, "rb") as f:
            assembled += f.read()
    shutil.rmtree(session_dir, ignore_errors=True)
    return assembled


@app.route("/upload-chunk", methods=["POST", "OPTIONS"])
def upload_chunk():
    if request.method == "OPTIONS":
        return make_response("", 204)
    if request.headers.get("X-API-Token") != API_TOKEN:
        return jsonify({"error": "Unauthorized"}), 401
    try:
        session_id = request.form.get("session_id") or str(uuid.uuid4())
        chunk_index = int(request.form.get("chunk_index", 0))
        total_chunks = int(request.form.get("total_chunks", 1))
        filename = request.form.get("filename", "upload.zip")

        chunk_file = request.files.get("chunk")
        if not chunk_file:
            return jsonify({"error": "Kein 'chunk' im Request"}), 400

        session_dir = os.path.join(CHUNK_DIR, session_id)
        os.makedirs(session_dir, exist_ok=True)

        chunk_path = os.path.join(session_dir, f"chunk_{chunk_index:04d}")
        chunk_file.save(chunk_path)

        meta_path = os.path.join(session_dir, "meta.json")
        if not os.path.exists(meta_path):
            with open(meta_path, "w") as f:
                json.dump({"total_chunks": total_chunks, "filename": filename}, f)

        received = len([n for n in os.listdir(session_dir) if n.startswith("chunk_")])
        ready = received >= total_chunks
        return jsonify({"session_id": session_id, "chunk": chunk_index,
                        "received": received, "total": total_chunks, "ready": ready})
    except Exception as e:
        import traceback
        return jsonify({"error": str(e), "trace": traceback.format_exc()}), 500


@app.route("/")
def index():
    index_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "index.html")
    with open(index_path, encoding="utf-8") as f:
        return f.read(), 200, {"Content-Type": "text/html; charset=utf-8"}


@app.route("/health")
def health():
    """Status + verfügbare Tools (Diagnose für Vorschau-Probleme)."""
    info = {
        "status":         "ok",
        "service":        "INTERPRES Full Pipeline v3",
        "test_mode":      TEST_MODE,
        "cloudconvert":   bool(CLOUDCONVERT_KEY),
        "pdftoppm":       _find_pdftoppm() is not None,
    }
    # PyMuPDF-Verfügbarkeit (Fallback-Renderer)
    try:
        import fitz  # noqa
        info["pymupdf"] = True
    except Exception:
        info["pymupdf"] = False
    info["preview_ready"] = (
        info["cloudconvert"]
        and (info["pdftoppm"] or info["pymupdf"])
    )
    return jsonify(info)

# ── Async Job Store ────────────────────────────────────────────────────────────
import threading as _threading
import time as _time

# ── Dateisystem-basierter Job-Store (/tmp) ────────────────────────────────────
# Funktioniert mit mehreren Prozessen/Workern, da alle dasselbe /tmp teilen.
_JOB_DIR = "/tmp/interpres_jobs"
os.makedirs(_JOB_DIR, exist_ok=True)

def _job_meta_path(job_id):
    return os.path.join(_JOB_DIR, f"{job_id}.json")

def _job_pdf_path(job_id):
    return os.path.join(_JOB_DIR, f"{job_id}.pdf")

def _job_dir(job_id):
    """Per-Job Working-Dir für State + Slide-JPGs + Uploads."""
    d = os.path.join(_JOB_DIR, f"work_{job_id}")
    os.makedirs(d, exist_ok=True)
    return d

def _job_state_path(job_id):
    return os.path.join(_job_dir(job_id), "state.json")

def _job_pptx_path(job_id):
    return os.path.join(_job_dir(job_id), "first_pass.pptx")

def _job_slides_dir(job_id, typ="marketing"):
    """Slide-JPGs pro Expose-Typ in separatem Unterordner.
    marketing → ../slides   (Default, Backwards-Compat)
    kurz/rechtlich → ../slides_<typ>
    """
    base = _job_dir(job_id)
    if typ and typ != "marketing":
        d = os.path.join(base, f"slides_{typ}")
    else:
        d = os.path.join(base, "slides")
    os.makedirs(d, exist_ok=True)
    return d

def _job_uploads_dir(job_id):
    d = os.path.join(_job_dir(job_id), "uploads")
    os.makedirs(d, exist_ok=True)
    return d

def _write_job(job_id, **fields):
    """Schreibt/aktualisiert Job-Metadaten atomar."""
    path = _job_meta_path(job_id)
    # Lese bestehende Daten falls vorhanden
    try:
        with open(path) as f:
            data = json.load(f)
    except Exception:
        data = {}
    data.update(fields)
    tmp = path + ".tmp"
    with open(tmp, "w") as f:
        json.dump(data, f)
    os.replace(tmp, path)  # atomarer Schreibvorgang

def _read_job(job_id):
    """Liest Job-Metadaten. Gibt None zurück wenn nicht gefunden."""
    try:
        with open(_job_meta_path(job_id)) as f:
            return json.load(f)
    except Exception:
        return None


def _authorize_job(job_id):
    """Auth für /job/<job_id>/...-Endpoints. Akzeptiert X-API-Token-Header oder
    ?token=-Query (Letzteres für <img src=…>)."""
    token = request.headers.get("X-API-Token") or request.args.get("token")
    return token == API_TOKEN

def _cleanup_old_jobs():
    """Entfernt Jobs und PDFs älter als 30 Minuten."""
    cutoff = _time.time() - 1800
    try:
        for fname in os.listdir(_JOB_DIR):
            fpath = os.path.join(_JOB_DIR, fname)
            try:
                if os.path.getmtime(fpath) < cutoff:
                    os.remove(fpath)
            except Exception:
                pass
    except Exception:
        pass

def _run_expose_job(job_id, zip_paths):
    """Läuft in einem Background-Thread. Liest ZIPs von Disk, schreibt Status+PDF nach /tmp."""
    def _set(**kw):
        _write_job(job_id, **kw)

    input_dir = os.path.join(_JOB_DIR, f"input_{job_id}")

    try:
        # ── Schritt 0: ZIP-Dateien lesen + Extraktion (im Thread!) ────────────
        _set(status="processing", phase="Dateien werden verarbeitet …")
        pdfs = []
        customer_image_list = []

        # ── Direkt hochgeladene Bilder zuerst laden (höchste Priorität) ──────
        _DIRECT_EXTS = {".jpg", ".jpeg", ".png", ".webp"}
        if os.path.isdir(input_dir):
            for fname in sorted(os.listdir(input_dir)):
                if not fname.startswith("_direct_img_"):
                    continue
                ext = os.path.splitext(fname)[1].lower()
                if ext not in _DIRECT_EXTS:
                    continue
                fpath = os.path.join(input_dir, fname)
                try:
                    # Direktbilder: Pfad speichern, NICHT bytes (RAM sparen)
                    customer_image_list.append({
                        "name": fname,
                        "ext":  ext,
                        "path": fpath,
                        "size": os.path.getsize(fpath),
                    })
                    print(f"[{job_id}] Direktbild: {fname} ({os.path.getsize(fpath)//1024} KB)")
                except Exception as e:
                    print(f"[{job_id}] Direktbild Fehler {fname}: {e}")
        direct_image_count = len(customer_image_list)

        # Working-Dir für extrahierte Files (PDFs + Bilder)
        work_dir = _job_dir(job_id)
        for zip_path in zip_paths:
            try:
                # NEU: ZIP wird per Path geöffnet (kein zip_bytes im RAM!)
                pdfs.extend(extract_pdfs_from_zip(zip_path, work_dir))
                customer_image_list.extend(extract_images_from_zip(zip_path, work_dir))
                print(f"[{job_id}] ZIP {os.path.basename(zip_path)}: "
                      f"{len(pdfs)} PDFs, {len(customer_image_list)} Bilder")
                import gc; gc.collect()
            except Exception as e:
                print(f"[{job_id}] ZIP-Fehler {zip_path}: {e}")

        # Input-Verzeichnis aufräumen
        shutil.rmtree(input_dir, ignore_errors=True)

        if not pdfs:
            _set(status="error", phase="Fehler",
                 error="Keine relevanten PDFs im Datenraum gefunden.")
            return

        _set(status="processing", phase="Dokumente werden analysiert …")

        # Bis zu 8 PDFs senden (mehr Kontext = bessere Daten- & WE-Extraktion)
        pdfs = sorted(pdfs, key=lambda x: x["priority"])[:8]
        print(f"[{job_id}] sende {len(pdfs)} PDFs an Claude:")
        for p in pdfs:
            print(f"    [Prio {p['priority']}] {p['folder']} / {p['name']}")

        if TEST_MODE:
            print(f"[{job_id}] TEST_MODE – überspringe Claude API")
            expose_data = DUMMY_EXPOSE_DATA.copy()
            geo_coords = None
        else:
            print(f"[{job_id}] Schritt 1/5: analyze_pdfs_with_claude…")
            projektdaten = analyze_pdfs_with_claude(pdfs)

            # Schritt 1b: Geocoding + offizielle Stadt + Web-Suche
            adresse   = projektdaten.get("adresse", "")
            stadt     = projektdaten.get("stadt", "")
            stadtteil = projektdaten.get("stadtteil", "")
            geo_result = None
            official_city = stadt  # Fallback
            if adresse and stadt:
                print(f"[{job_id}] Schritt 1b: Geocoding {adresse!r}, {stadt!r} …")
                geo_result = _geocode_address(adresse, stadt)
                if geo_result:
                    lat, lon, official_city = geo_result
                    print(f"[{job_id}]   → ({lat:.4f}, {lon:.4f}), official_city='{official_city}'")
                    # Offizielle Stadt in projektdaten setzen (für Stadtbilder + Stats)
                    if official_city and official_city.lower() != stadt.lower():
                        print(f"[{job_id}]   Stadt korrigiert: '{stadt}' → '{official_city}'")
                        projektdaten["stadt_offiziell"] = official_city
                        # Für Bildsuche und Expose-Generierung die offizielle Stadt nutzen
                        projektdaten["stadt"] = official_city

            # Web-Suche: aktuelle Stadtinfos via Tavily
            city_context = ""
            search_city = projektdaten.get("stadt", stadt)
            if search_city:
                print(f"[{job_id}] Schritt 1c: Tavily Web-Suche für '{search_city}' …")
                city_context = _search_city_info(search_city, stadtteil)

            _set(status="processing", phase="Exposé-Texte werden generiert …")
            print(f"[{job_id}] Schritt 2/5: generate_expose_with_claude…")
            raw_expose = generate_expose_with_claude(projektdaten, city_context=city_context)
            print(f"[{job_id}]   Claude-Ausgabe: {len(raw_expose)} Felder")
            expose_data = {**PLATZHALTER, **raw_expose}
            expose_data["logo_initial"] = generate_logo_initial(expose_data.get("projekt_name", ""))

            # Hard-Cap besonderheiten_liste auf 6 Bullets (Box-Hoehe ist fest).
            _bl = expose_data.get("besonderheiten_liste", "") or ""
            if _bl:
                _lines = [ln for ln in _bl.split("\n") if ln.strip()]
                if len(_lines) > 6:
                    print(f"[{job_id}]   besonderheiten_liste: {len(_lines)} Bullets → auf 6 gekappt")
                    expose_data["besonderheiten_liste"] = "\n".join(_lines[:6])

            # Direktes Flow von projektdaten-Fakten in expose_data (Daten aus PDF-Analyse
            # haben Vorrang vor evtl. leeren Generator-Outputs für diese Felder)
            for k in ("stellplaetze", "kfw_standard", "energieversorgung", "besonderheiten",
                      "kaufpreis_ab", "anzahl_we", "anzahl_1zi", "anzahl_2zi", "anzahl_barrierefrei"):
                pd_val = projektdaten.get(k, "")
                if pd_val and not expose_data.get(k):
                    expose_data[k] = str(pd_val)
                    print(f"[{job_id}]   Faktenfeld {k} aus PDFs: {str(pd_val)[:60]!r}")
            # Auch anzahl_we als Synonym (älterer Name)
            if not expose_data.get("anzahl_we") and projektdaten.get("anzahl_we_gesamt"):
                expose_data["anzahl_we"] = str(projektdaten["anzahl_we_gesamt"])

            # Schritt 2b: Proximity-Daten + Lageplan
            lat = lon = None
            if geo_result:
                lat, lon, _ = geo_result
                print(f"[{job_id}] Schritt 2b: Proximity via Overpass …")
                proximity = _calculate_proximity_data(adresse, stadt, lat, lon)
                expose_data.update(proximity)
            else:
                # Fallback: Stadt allein geocodieren (ohne Adresse)
                # → Lageplan zeigt zumindest die Stadt-Mitte, nicht leer
                if stadt:
                    fallback_geo = _geocode_address("", stadt)
                    if fallback_geo:
                        lat, lon, _ = fallback_geo
                        print(f"[{job_id}]   Stadt-Fallback-Geocoding: ({lat:.4f}, {lon:.4f})")

            # Lageplan setzen wenn wir Koordinaten haben
            if lat is not None and lon is not None:
                expose_data["bild_lageplan"] = _osm_lageplan_url(lat, lon)
                print(f"[{job_id}]   Lageplan URL: {expose_data['bild_lageplan']}")

        _set(status="processing", phase="Bilder werden ausgewählt …")
        print(f"[{job_id}] Schritt 3/5: Bilder …")

        customer_images = {}
        if customer_image_list:
            customer_images = classify_and_assign_customer_images(customer_image_list)

        # ── customer_image_list ist jetzt überflüssig (zugewiesene Bilder sind in
        #     customer_images, der Rest wird nicht mehr gebraucht). Sofort freigeben.
        import gc
        del customer_image_list
        gc.collect()

        expose_data = fill_image_placeholders(expose_data)
        gc.collect()

        # ── Leere Raumflächen-Felder mit "—" füllen ─────────────────────────
        # Hintergrund: Bei 1-Zimmer-Wohnungen gibt es kein separates Schlafzimmer
        # → Claude liefert we_flaeche_2 leer. Im Template würde dann nur ein
        # leerer Zahlenslot stehen. Statt blank lieber ein Bindestrich.
        # Suffixe: "" = Typ 1, "_typN" = Typ 2..8 (v19), "_N" = Backwards (v20)
        _flaeche_suffixes = [""] + [f"_typ{t}" for t in range(2, 9)] \
                                + [f"_{n}" for n in range(1, 7)]
        for suf in _flaeche_suffixes:
            # Typ ist "befüllt" wenn we_typ_beschreibung ODER irgendein
            # we_flaeche_*{suf} einen echten Wert hat
            typ_active = bool(
                str(expose_data.get(f"we_typ_beschreibung{suf}", "")).strip()
                or any(str(expose_data.get(f"we_flaeche_{fn}{suf}", "")).strip()
                       for fn in range(1, 6))
            )
            if not typ_active:
                continue
            for fn in range(1, 6):
                key = f"we_flaeche_{fn}{suf}"
                if key in expose_data and not str(expose_data.get(key, "")).strip():
                    expose_data[key] = "—"

        # ── Kundennamen (Entwickler + Projekttitel) overriden Claude-Output ──
        # Beide werden unabhängig voneinander angewendet:
        #   entwickler_name  → ENTWICKLER_NAME (z.B. "SBB")
        #   entwickler_name_gross → großgeschrieben
        #   logo_initial     → erster Buchstabe des Entwicklernamens
        #   projekt_titel    → PROJEKT_TITEL (z.B. "The Rothenseer – Modern …")
        #   projekt_name     → kurzer Name aus Projekttitel (für Dateinamen etc.)
        try:
            job_meta = _read_job(job_id) or {}
            user_entw = (job_meta.get("user_entwicklername") or "").strip()
            user_titel = (job_meta.get("user_projekttitel")   or "").strip()
            claude_entw  = expose_data.get("entwickler_name", "")
            claude_titel = expose_data.get("projekt_titel",   "")
            claude_pname = expose_data.get("projekt_name",    "")
            print(f"[{job_id}] Namen: entw_claude={claude_entw!r} → user={user_entw!r} | "
                  f"titel_claude={claude_titel!r} → user={user_titel!r}")

            if user_entw:
                expose_data["entwickler_name"]       = user_entw
                expose_data["entwickler_name_gross"] = user_entw.upper()
                expose_data["logo_initial"]          = generate_logo_initial(user_entw)
                # Replace alle Vorkommen des alten Namens
                if claude_entw and claude_entw != user_entw:
                    n = 0
                    for k, v in list(expose_data.items()):
                        if isinstance(v, str) and claude_entw in v and k not in (
                                "entwickler_name", "entwickler_name_gross"):
                            expose_data[k] = v.replace(claude_entw, user_entw)
                            n += 1
                    if n: print(f"[{job_id}]   Entwicklername in {n} Feldern ersetzt")

            if user_titel:
                expose_data["projekt_titel"] = user_titel
                # projekt_name = erstes Komma/Bindestrich-Segment des Titels (für Dateinamen)
                short = re.split(r'[–\-,:]', user_titel)[0].strip() or user_titel
                expose_data["projekt_name"] = short
                if claude_titel and claude_titel != user_titel:
                    n = 0
                    for k, v in list(expose_data.items()):
                        if isinstance(v, str) and claude_titel in v and k != "projekt_titel":
                            expose_data[k] = v.replace(claude_titel, user_titel)
                            n += 1
                    if n: print(f"[{job_id}]   Projekttitel in {n} Feldern ersetzt")
                if claude_pname and claude_pname != short:
                    n = 0
                    for k, v in list(expose_data.items()):
                        if isinstance(v, str) and claude_pname in v and k not in (
                                "projekt_name", "projekt_titel"):
                            expose_data[k] = v.replace(claude_pname, short)
                            n += 1
                    if n: print(f"[{job_id}]   Projektname-Kurz in {n} Feldern ersetzt")

            print(f"[{job_id}] ✓ Final: entwickler={expose_data.get('entwickler_name')!r}, "
                  f"projekt_titel={expose_data.get('projekt_titel')!r}, "
                  f"projekt_name={expose_data.get('projekt_name')!r}")
        except Exception as e:
            import traceback as _tbb
            print(f"[{job_id}] Namen-Override Fehler: {e}\n{_tbb.format_exc()[:300]}")

        _set(status="processing", phase="Präsentation wird erstellt …")
        print(f"[{job_id}] Schritt 4/6: fill_pptx …")
        tmpl_bytes = requests.get(TEMPLATE_URL, timeout=30).content
        print(f"[{job_id}] Template geladen: {len(tmpl_bytes)//1024} KB")

        pptx_bytes = fill_pptx(tmpl_bytes, expose_data, customer_images=customer_images)
        del tmpl_bytes
        gc.collect()

        # ── Bbox-Map aus der FINALEN PPTX extrahieren ────────────────────────
        # {{BILD_*}}-Platzhalter, die fill_pptx NICHT mit Bildern befüllt hat,
        # bleiben als Text-Platzhalter erhalten – genau die Slots, die der
        # Kunde via Preview-UI mit eigenen Fotos befüllen soll.
        # Slide-Indices stimmen 1:1 mit den gerenderten JPGs überein.
        try:
            bbox_map = extract_bild_placeholders(pptx_bytes)
            n_placeholders = sum(len(s['placeholders']) for s in bbox_map['slides'])
            print(f"[{job_id}] Bbox-Map (finale PPTX): {n_placeholders} freie Slots "
                  f"auf {len(bbox_map['slides'])} Slides "
                  f"(slide_dim={bbox_map.get('slide_w_emu')}x{bbox_map.get('slide_h_emu')})")
        except Exception as be:
            import traceback as _tb
            print(f"[{job_id}] Bbox-Extraktion Fehler: {be}\n{_tb.format_exc()[:500]}")
            bbox_map = {"slide_w_emu": 12192000, "slide_h_emu": 6858000, "slides": []}

        projekt_name = expose_data.get("projekt_name", "Expose").replace(" ", "_")

        # ── State persistieren (für Second-Pass / Finalize) ───────────────────
        first_pass_pptx = _job_pptx_path(job_id)
        with open(first_pass_pptx, "wb") as fh:
            fh.write(pptx_bytes)

        # Auto-zugeordnete Kundenbilder speichern (damit Second-Pass sie behalten kann)
        customer_images_files = {}
        for slot, raw in (customer_images or {}).items():
            ext = ".jpg"
            try:
                if raw[:4] == b"\x89PNG":
                    ext = ".png"
            except Exception:
                pass
            cpath = os.path.join(_job_dir(job_id), f"auto_{slot}{ext}")
            with open(cpath, "wb") as fh:
                fh.write(raw)
            customer_images_files[slot] = cpath

        with open(_job_state_path(job_id), "w") as fh:
            json.dump({
                "expose_data": expose_data,
                "customer_images_files": customer_images_files,
                "projekt_name": projekt_name,
            }, fh, ensure_ascii=False)

        # ── Schritt 5: PPTX → PDF → JPGs für Preview ─────────────────────────
        _set(status="processing", phase="Vorschau wird erstellt …")
        print(f"[{job_id}] Schritt 5/6: PDF + Slide-Bilder für Vorschau …")

        slide_jpgs = []
        # bbox_map wurde oben aus pptx_bytes extrahiert – NICHT überschreiben!
        # Diagnose: warum versagt Preview ggf.?
        _can_pdf  = _can_convert_to_pdf()
        _ppm_path = _find_pdftoppm()
        # PyMuPDF-Verfügbarkeit (Fallback wenn pdftoppm fehlt)
        try:
            import fitz as _fitz  # noqa
            _has_pymupdf = True
        except Exception:
            _has_pymupdf = False
        _can_render_jpgs = _ppm_path or _has_pymupdf
        print(f"[{job_id}] Preview-Tools: PDF-Konverter={'CloudConvert' if CLOUDCONVERT_KEY else 'KEINER'}, "
              f"pdftoppm={_ppm_path or 'FEHLT'}, pymupdf={_has_pymupdf}")
        try:
            if _can_pdf and _can_render_jpgs:
                pdf_bytes = convert_to_pdf(pptx_bytes, f"{projekt_name}.pptx")
                # PDF mit klickbaren Hyperlinks anreichern
                pdf_bytes = _add_hyperlinks_to_pdf(pdf_bytes)
                # PPTX wurde an CloudConvert gesendet → kann jetzt aus dem RAM
                # (ist als Datei unter first_pass_pptx auf Disk)
                del pptx_bytes
                gc.collect()
                _set(phase="Slide-Bilder werden gerendert …")
                # PDF auf Disk schreiben statt in RAM behalten
                pdf_tmp = os.path.join(_job_dir(job_id), "_first_pass.pdf")
                with open(pdf_tmp, "wb") as fh:
                    fh.write(pdf_bytes)
                pdf_size_kb = len(pdf_bytes) // 1024
                del pdf_bytes
                gc.collect()
                # render_pdf_to_jpgs liest das PDF wieder von Disk – stabil
                with open(pdf_tmp, "rb") as fh:
                    # DPI=150 → scharf auch auf Retina-Display
                    jpg_paths = render_pdf_to_jpgs(fh.read(), _job_slides_dir(job_id), dpi=150)
                try: os.unlink(pdf_tmp)
                except OSError: pass
                slide_jpgs = [os.path.basename(p) for p in jpg_paths]
                print(f"[{job_id}] ✓ {len(slide_jpgs)} Slide-JPGs erstellt (PDF war {pdf_size_kb} KB) "
                      f"– Pfad: {_job_slides_dir(job_id)}")
            else:
                missing = []
                if not _can_pdf:        missing.append("CLOUDCONVERT_KEY")
                if not _can_render_jpgs: missing.append("pdftoppm UND pymupdf")
                print(f"[{job_id}] ⚠️ Preview übersprungen – fehlt: {', '.join(missing)}")
        except Exception as pe:
            import traceback as _tb
            print(f"[{job_id}] ✗ Preview-Render Fehler: {pe}\n{_tb.format_exc()[:600]}")

        # Welche Slots wurden bereits durch automatische Pipeline befüllt?
        already_filled = []
        for k in expose_data:
            if k.startswith("bild_") and isinstance(expose_data[k], str) and expose_data[k]:
                already_filled.append(k.lower())
        for slot in (customer_images or {}):
            if slot.lower() not in already_filled:
                already_filled.append(slot.lower())

        # ── Slot-Liste: nur ECHT im Template vorhandene Slots aus bbox_map.
        # Der frühere PLATZHALTER-Fallback hat alle 20 bild_we_* aus dem
        # Defaults-Dict eingefügt (auch für nicht existierende WE-Typen) und
        # zusätzlich phantom-Slots wie bild_grundriss_2 die im Template fehlen.
        # Wir nehmen jetzt ausschließlich was die bbox_map liefert; Fallback
        # nur wenn die bbox_map komplett leer ist (Extraktion versagt).
        slot_list = []
        seen_slots = set()
        for s in bbox_map.get("slides", []):
            for ph in s.get("placeholders", []):
                if ph["slot"] in seen_slots:
                    continue
                seen_slots.add(ph["slot"])
                if ph["slot"] in already_filled:
                    continue
                slot_list.append({
                    "slot":  ph["slot"],
                    "label": ph.get("label", ph["slot"]),
                    "slide": s["index"] + 1,
                })

        # Notfall-Fallback: bbox_map versagt → konservative Liste mit nur den
        # WE-Slots die zu aktiven Typen gehören + Standard-Projekt-Slots.
        if not slot_list and not bbox_map.get("slides"):
            print(f"[{job_id}] ⚠️ bbox_map leer → konservativer Slot-Fallback")
            # Aktive WE-Paare ermitteln
            active_we_n = set()
            for typ in range(1, 9):
                left_n  = typ * 2 - 1
                right_n = typ * 2
                if (expose_data.get(f"we_beispiel_{left_n}") or expose_data.get(f"we_nummern_{left_n}")
                        or expose_data.get(f"we_beispiel_{right_n}") or expose_data.get(f"we_nummern_{right_n}")
                        or typ == 1):  # Typ 1 immer
                    active_we_n.add(left_n)
                    active_we_n.add(right_n)
            for k in PLATZHALTER:
                if not k.startswith("bild_") or k in already_filled:
                    continue
                m = re.match(r'^bild_we_(\d+)$', k)
                if m and int(m.group(1)) not in active_we_n:
                    continue
                slot_list.append({
                    "slot":  k,
                    "label": _slot_label(k),
                    "slide": 0,
                })
        slot_list.sort(key=lambda x: (x["slide"] or 99, x["slot"]))
        print(f"[{job_id}] Slot-Liste: {len(slot_list)} Upload-Kandidaten "
              f"({sum(1 for s in slot_list if s['slide'])} mit Slide-Index)")

        # Preview wird IMMER gezeigt – auch wenn Slide-Bilder nicht gerendert
        # werden konnten (CloudConvert / pdftoppm Fehler). Dann zeigt das UI
        # nur die Slot-Liste (Karten-Grid) ohne Slide-Vorschau-Bilder.
        if not slide_jpgs:
            print(f"[{job_id}] ⚠️  Keine Slide-Bilder verfügbar – Preview zeigt nur Slot-Liste")
        _set(
            status="preview",
            phase="Bereit für Bilder-Upload",
            name=projekt_name,
            slide_jpgs=slide_jpgs,
            bbox_map=bbox_map,
            slot_list=slot_list,
            already_filled=already_filled,
        )
        gc.collect()
        print(f"[{job_id}] ✓ Preview bereit – {len(slide_jpgs)} JPGs, "
              f"{len(slot_list)} Upload-Slots")

    except Exception as e:
        import traceback as tb
        err = f"{e}\n{tb.format_exc()}"
        print(f"[{job_id}] ✗ Fehler: {err[:500]}")
        _set(status="error", phase="Fehler", error=str(e))
    finally:
        # Input-Verzeichnis aufräumen (falls noch nicht geschehen)
        shutil.rmtree(input_dir, ignore_errors=True)


@app.route("/generate-expose", methods=["POST"])
def generate_expose():
    """Empfängt ZIPs, speichert sie auf Disk und startet sofort den Hintergrund-Job.
    KEINE schwere Verarbeitung hier – alles im Thread – damit kein 30s-Proxy-Timeout."""
    if request.headers.get("X-API-Token") != API_TOKEN:
        return jsonify({"error": "Unauthorized"}), 401
    try:
        _cleanup_old_jobs()

        job_id = str(uuid.uuid4())
        input_dir = os.path.join(_JOB_DIR, f"input_{job_id}")
        os.makedirs(input_dir, exist_ok=True)

        zip_paths = []

        session_ids = request.form.getlist("session_ids")
        if session_ids:
            for i, sid in enumerate(session_ids):
                zip_bytes = assemble_session(sid)
                zip_path = os.path.join(input_dir, f"zip_{i}.zip")
                with open(zip_path, "wb") as fh:
                    fh.write(zip_bytes)
                zip_paths.append(zip_path)
                print(f"[{job_id}] Session {sid}: {len(zip_bytes)//1024} KB gespeichert")

        elif request.content_type and "multipart" in request.content_type:
            uploaded = request.files.getlist("files") or request.files.getlist("file")
            if not uploaded:
                shutil.rmtree(input_dir, ignore_errors=True)
                return jsonify({"error": "Keine Dateien im Request"}), 400
            for i, f in enumerate(uploaded):
                raw = f.read()
                zip_path = os.path.join(input_dir, f"upload_{i}.zip")
                with open(zip_path, "wb") as fh:
                    fh.write(raw)
                zip_paths.append(zip_path)

        else:
            body = request.get_json(force=True) or {}
            b64_list = body.get("zip_base64_list") or (
                [body["zip_base64"]] if "zip_base64" in body else None
            )
            if not b64_list:
                shutil.rmtree(input_dir, ignore_errors=True)
                return jsonify({"error": "zip_base64 oder zip_base64_list fehlt"}), 400
            for i, b64 in enumerate(b64_list):
                raw = base64.b64decode(b64)
                zip_path = os.path.join(input_dir, f"b64_{i}.zip")
                with open(zip_path, "wb") as fh:
                    fh.write(raw)
                zip_paths.append(zip_path)

        if not zip_paths:
            shutil.rmtree(input_dir, ignore_errors=True)
            return jsonify({"error": "Keine ZIP-Dateien empfangen"}), 400

        # ── Direkt hochgeladene Projektfotos speichern ────────────────────────
        direct_images = request.files.getlist("images") or request.files.getlist("images[]")
        saved_direct = 0
        for i, img_file in enumerate(direct_images):
            if not img_file or not img_file.filename:
                continue
            ext = os.path.splitext(img_file.filename)[1].lower()
            if ext not in (".jpg", ".jpeg", ".png", ".webp"):
                continue
            img_path = os.path.join(input_dir, f"_direct_img_{i}{ext}")
            img_file.save(img_path)
            saved_direct += 1
        if saved_direct:
            print(f"[{job_id}] {saved_direct} direkte Projektfotos gespeichert")

        # ── Optionale Kundennamen (überschreiben Claude-Output) ──────────────
        user_entwicklername = (request.form.get("entwicklername") or "").strip()
        user_projekttitel   = (request.form.get("projekttitel")   or "").strip()
        # Kompatibilität: alter 'projektname' Parameter mappt auf Projekttitel
        if not user_projekttitel:
            user_projekttitel = (request.form.get("projektname") or "").strip()
        if user_entwicklername or user_projekttitel:
            print(f"[{job_id}] Kunden-Namen: entwickler={user_entwicklername!r}, "
                  f"titel={user_projekttitel!r}")

        _write_job(job_id,
                   status="processing",
                   phase="Wird gestartet …",
                   created=_time.time(),
                   pdf_path=None,
                   name=None,
                   error=None,
                   user_entwicklername=user_entwicklername,
                   user_projekttitel=user_projekttitel)

        t = _threading.Thread(
            target=_run_expose_job,
            args=(job_id, zip_paths),
            daemon=True
        )
        t.start()

        print(f"[{job_id}] Job gestartet mit {len(zip_paths)} ZIP(s)")
        return jsonify({"job_id": job_id})

    except Exception as e:
        import traceback
        return jsonify({"error": str(e), "trace": traceback.format_exc()}), 500


@app.route("/job/<job_id>", methods=["GET", "OPTIONS"])
def job_status(job_id):
    if request.method == "OPTIONS":
        return make_response("", 204)
    if not _authorize_job(job_id):
        return jsonify({"error": "Unauthorized"}), 401
    job = _read_job(job_id)
    if not job:
        return jsonify({"error": "Job nicht gefunden"}), 404
    status = job.get("status", "processing")
    if status == "processing":
        return jsonify({"status": "processing", "phase": job.get("phase", "")})
    if status == "error":
        return jsonify({"status": "error", "error": job.get("error", "Unbekannter Fehler")}), 500
    if status == "preview":
        # Preview bereit – noch kein Download
        return jsonify({
            "status":         "preview",
            "phase":          job.get("phase", "Bereit für Upload"),
            "name":           job.get("name", "Expose"),
            "slide_count":    len(job.get("slide_jpgs", []) or []),
            "bbox_map":       job.get("bbox_map", {}),
            "slot_list":      job.get("slot_list", []),
            "already_filled": job.get("already_filled", []),
        })
    if status == "done":
        pdf_path = job.get("pdf_path")
        if not pdf_path or not os.path.exists(pdf_path):
            return jsonify({"error": "Datei nicht mehr verfügbar"}), 410
        ext = os.path.splitext(pdf_path)[1].lower()
        if ext == ".pdf":
            mimetype = "application/pdf"
            dl_name  = f"{job.get('name', 'Expose')}_Expose.pdf"
        else:
            mimetype = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            dl_name  = f"{job.get('name', 'Expose')}_Expose.pptx"
        return send_file(pdf_path, mimetype=mimetype,
                         as_attachment=True, download_name=dl_name)
    return jsonify({"error": f"Unbekannter Job-Status: {status}"}), 500


@app.route("/job/<job_id>/slide/<int:n>", methods=["GET", "OPTIONS"])
def job_slide_image(job_id, n):
    """Liefert die Slide-JPG-Datei für die Preview-Ansicht.
    ?typ=marketing|kurz|rechtlich (Default marketing) waehlt den Unterordner."""
    if request.method == "OPTIONS":
        return make_response("", 204)
    if not _authorize_job(job_id):
        return jsonify({"error": "Unauthorized"}), 401
    typ = (request.args.get("typ") or "marketing").lower()
    if typ not in ("marketing", "kurz", "rechtlich"):
        typ = "marketing"
    path = os.path.join(_job_slides_dir(job_id, typ), f"slide_{n}.jpg")
    if not os.path.exists(path):
        return jsonify({"error": "Slide nicht gefunden"}), 404
    return send_file(path, mimetype="image/jpeg")


def _run_finalize_job(job_id):
    """Second-Pass: liest gespeicherten State + User-Uploads, refillt Template,
    konvertiert zu PDF (falls möglich) und setzt Status=done."""
    def _set(**kw):
        _write_job(job_id, **kw)

    try:
        _set(status="processing", phase="Finale Datei wird erstellt …")
        state_path = _job_state_path(job_id)
        with open(state_path) as fh:
            state = json.load(fh)
        expose_data         = state["expose_data"]
        projekt_name        = state.get("projekt_name", "Expose")
        cust_files          = state.get("customer_images_files", {}) or {}

        # Auto-zugeordnete Bilder zuerst laden
        customer_images = {}
        for slot, fpath in cust_files.items():
            try:
                with open(fpath, "rb") as fh:
                    customer_images[slot] = fh.read()
            except Exception:
                pass

        # User-Uploads OVERRIDEN auto-Zuweisungen
        ud = _job_uploads_dir(job_id)
        if os.path.isdir(ud):
            for fname in os.listdir(ud):
                # Format: <slot>.<ext>
                base, ext = os.path.splitext(fname)
                slot = base.lower()
                if not slot.startswith("bild_"):
                    continue
                with open(os.path.join(ud, fname), "rb") as fh:
                    customer_images[slot] = fh.read()
                # Wenn User ein Bild für einen Slot hochlädt, sollte die alte URL nicht
                # konkurrieren – auf "" setzen, damit fill_pptx das Kunden-Bild nimmt.
                expose_data[slot] = ""

        # Template laden + neu füllen
        _set(phase="Template wird befüllt …")
        tmpl_bytes = requests.get(TEMPLATE_URL, timeout=30).content
        pptx_bytes = fill_pptx(tmpl_bytes, expose_data, customer_images=customer_images)

        # Konvertierung
        _set(phase="PDF wird erstellt …")
        out_path = None
        if _can_convert_to_pdf():
            try:
                pdf_bytes = convert_to_pdf(pptx_bytes, f"{projekt_name}.pptx")
                pdf_bytes = _add_hyperlinks_to_pdf(pdf_bytes)
                out_path = _job_pdf_path(job_id)
                with open(out_path, "wb") as fh:
                    fh.write(pdf_bytes)
            except Exception as pe:
                print(f"[{job_id}] Final PDF-Fehler: {pe}")
                out_path = None
        if not out_path:
            out_path = os.path.join(_JOB_DIR, f"{job_id}.pptx")
            with open(out_path, "wb") as fh:
                fh.write(pptx_bytes)

        _set(status="done", phase="Fertig", pdf_path=out_path, name=projekt_name)
        print(f"[{job_id}] ✓ Final fertig: {out_path}")
    except Exception as e:
        import traceback as tb
        err = f"{e}\n{tb.format_exc()}"
        print(f"[{job_id}] ✗ Finalize Fehler: {err[:500]}")
        _set(status="error", phase="Fehler", error=str(e))


@app.route("/job/<job_id>/upload", methods=["POST", "OPTIONS"])
def job_upload_image(job_id):
    """Speichert ein vom User hochgeladenes Bild für einen bestimmten Slot.
    Multipart: field 'image' (Datei) + 'slot' (z.B. 'bild_titel')."""
    if request.method == "OPTIONS":
        return make_response("", 204)
    if not _authorize_job(job_id):
        return jsonify({"error": "Unauthorized"}), 401
    job = _read_job(job_id)
    if not job:
        return jsonify({"error": "Job nicht gefunden"}), 404
    if job.get("status") not in ("preview", "done"):
        return jsonify({"error": f"Job nicht im Preview-Status (ist: {job.get('status')})"}), 400

    slot = (request.form.get("slot") or "").strip().lower()
    if not slot.startswith("bild_"):
        return jsonify({"error": "Ungültiger Slot"}), 400
    img = request.files.get("image")
    if not img or not img.filename:
        return jsonify({"error": "Keine Bilddatei"}), 400
    ext = os.path.splitext(img.filename)[1].lower()
    if ext not in (".jpg", ".jpeg", ".png", ".webp"):
        return jsonify({"error": "Nur JPG/PNG/WEBP erlaubt"}), 400

    ud = _job_uploads_dir(job_id)
    # Vorherige Datei für diesen Slot löschen (egal welche Extension)
    for fname in os.listdir(ud):
        if os.path.splitext(fname)[0].lower() == slot:
            try: os.remove(os.path.join(ud, fname))
            except OSError: pass
    target = os.path.join(ud, f"{slot}{ext}")
    img.save(target)
    print(f"[{job_id}] Upload: {slot} → {os.path.basename(target)} ({os.path.getsize(target)//1024} KB)")
    return jsonify({"ok": True, "slot": slot, "size": os.path.getsize(target)})


@app.route("/job/<job_id>/upload/<slot>", methods=["DELETE", "OPTIONS"])
def job_remove_upload(job_id, slot):
    """Entfernt das vom User hochgeladene Bild für einen Slot."""
    if request.method == "OPTIONS":
        return make_response("", 204)
    if not _authorize_job(job_id):
        return jsonify({"error": "Unauthorized"}), 401
    slot = slot.lower()
    if not slot.startswith("bild_"):
        return jsonify({"error": "Ungültiger Slot"}), 400
    ud = _job_uploads_dir(job_id)
    removed = 0
    if os.path.isdir(ud):
        for fname in os.listdir(ud):
            if os.path.splitext(fname)[0].lower() == slot:
                try:
                    os.remove(os.path.join(ud, fname))
                    removed += 1
                except OSError:
                    pass
    return jsonify({"ok": True, "removed": removed})


@app.route("/job/<job_id>/finalize", methods=["POST", "OPTIONS"])
def job_finalize(job_id):
    """Startet Second-Pass im Hintergrund. Gibt sofort 202 zurück."""
    if request.method == "OPTIONS":
        return make_response("", 204)
    if not _authorize_job(job_id):
        return jsonify({"error": "Unauthorized"}), 401
    job = _read_job(job_id)
    if not job:
        return jsonify({"error": "Job nicht gefunden"}), 404
    if not os.path.exists(_job_state_path(job_id)):
        return jsonify({"error": "Kein Job-State zum Finalisieren"}), 400

    _write_job(job_id, status="processing", phase="Finale Datei wird vorbereitet …")
    t = _threading.Thread(target=_run_finalize_job, args=(job_id,), daemon=True)
    t.start()
    return jsonify({"ok": True, "job_id": job_id})


@app.route("/job/<job_id>/uploaded", methods=["GET", "OPTIONS"])
def job_list_uploads(job_id):
    """Listet aktuell hochgeladene User-Bilder pro Slot."""
    if request.method == "OPTIONS":
        return make_response("", 204)
    if not _authorize_job(job_id):
        return jsonify({"error": "Unauthorized"}), 401
    ud = _job_uploads_dir(job_id)
    items = {}
    if os.path.isdir(ud):
        for fname in os.listdir(ud):
            slot = os.path.splitext(fname)[0].lower()
            items[slot] = {
                "filename": fname,
                "size":     os.path.getsize(os.path.join(ud, fname)),
            }
    return jsonify({"uploads": items})


@app.route("/job/<job_id>/uploaded/<slot>/preview", methods=["GET", "OPTIONS"])
def job_uploaded_preview(job_id, slot):
    """Liefert das hochgeladene Bild zur Anzeige im Frontend."""
    if request.method == "OPTIONS":
        return make_response("", 204)
    if not _authorize_job(job_id):
        return jsonify({"error": "Unauthorized"}), 401
    ud = _job_uploads_dir(job_id)
    slot = slot.lower()
    if not os.path.isdir(ud):
        return jsonify({"error": "Keine Uploads"}), 404
    for fname in os.listdir(ud):
        if os.path.splitext(fname)[0].lower() == slot:
            mt = "image/jpeg"
            ext = os.path.splitext(fname)[1].lower()
            if ext == ".png":
                mt = "image/png"
            elif ext == ".webp":
                mt = "image/webp"
            return send_file(os.path.join(ud, fname), mimetype=mt)
    return jsonify({"error": "Slot nicht hochgeladen"}), 404


# ── Rechtlicher Teil B ────────────────────────────────────────────────────────
# Pfad lokal: erst aus Repo (urbanunits_Rechtlich_v1.pptx), Fallback URL.
_RECHTLICH_TPL_LOCAL = os.path.join(os.path.dirname(__file__), "urbanunits_Rechtlich_v1.pptx")
# PDF-Seitenanzahl im PPTX-Template, die VOR der Teilungserklaerung kommen
# (Cover, Angebot, Vor-/Nachteile, Konzept, Steuern, Vertrag = S. 1-102 im DQN-Original).
# Danach (Index ≥ _RECHTLICH_SPLIT_AFTER) kommen Verwaltungs-Muster (S. 308-322).
_RECHTLICH_SPLIT_AFTER = 102


def _load_rechtlich_template_bytes():
    if os.path.exists(_RECHTLICH_TPL_LOCAL):
        with open(_RECHTLICH_TPL_LOCAL, "rb") as fh:
            return fh.read()
    return requests.get(RECHTLICH_TEMPLATE_URL, timeout=60).content


def _merge_rechtlich_pdf(dynamic_pdf_bytes, teil_pdf_bytes):
    """Fuegt Teilungserklaerung-PDF zwischen Seite 102 (Vertrag) und Seite 103
    (Verwaltungs-Muster) des dynamisch erzeugten Rechtlich-PDFs ein."""
    import io as _io
    from pypdf import PdfReader, PdfWriter
    dyn = PdfReader(_io.BytesIO(dynamic_pdf_bytes))
    n   = len(dyn.pages)
    w   = PdfWriter()
    cut = min(_RECHTLICH_SPLIT_AFTER, n)
    for i in range(cut):
        w.add_page(dyn.pages[i])
    if teil_pdf_bytes:
        teil = PdfReader(_io.BytesIO(teil_pdf_bytes))
        for p in teil.pages:
            w.add_page(p)
    for i in range(cut, n):
        w.add_page(dyn.pages[i])
    buf = _io.BytesIO()
    w.write(buf)
    return buf.getvalue()


@app.route("/job/<job_id>/teilungserklaerung", methods=["GET", "POST", "DELETE", "OPTIONS"])
def job_teilungserklaerung(job_id):
    """PDF-Upload der Teilungserklärung (vom Notar geliefert). Wird beim
    Render des rechtlichen Teils B zwischen Vertrag und Verwaltungs-Muster
    eingefuegt. Multipart-Field: 'file' (PDF)."""
    if request.method == "OPTIONS":
        return make_response("", 204)
    if not _authorize_job(job_id):
        return jsonify({"error": "Unauthorized"}), 401
    job = _read_job(job_id)
    if not job:
        return jsonify({"error": "Job nicht gefunden"}), 404
    path = os.path.join(_job_dir(job_id), "teilungserklaerung.pdf")

    if request.method == "GET":
        if os.path.exists(path):
            return jsonify({"uploaded": True, "size": os.path.getsize(path)})
        return jsonify({"uploaded": False})

    if request.method == "DELETE":
        if os.path.exists(path):
            try: os.remove(path)
            except OSError: pass
        return jsonify({"ok": True})

    # POST: Upload
    f = request.files.get("file")
    if not f or not f.filename:
        return jsonify({"error": "Keine Datei"}), 400
    if not f.filename.lower().endswith(".pdf"):
        return jsonify({"error": "Nur PDF erlaubt"}), 400
    f.save(path)
    sz = os.path.getsize(path)
    print(f"[{job_id}] Teilungserklaerung hochgeladen: {sz//1024} KB")
    return jsonify({"ok": True, "size": sz})


def _run_render_rechtlich(job_id):
    """Background-Worker: Rechtlich-PPTX befuellen, PDF konvertieren, mit
    Teilungserklaerung mergen, Ergebnis unter rechtlich.pdf ablegen."""
    try:
        _write_job(job_id, rechtlich_status="processing",
                   rechtlich_phase="Template wird geladen …")
        # Daten aus State holen
        state_path = _job_state_path(job_id)
        if not os.path.exists(state_path):
            raise RuntimeError("Job-State nicht gefunden")
        with open(state_path) as fh:
            state = json.load(fh)
        expose_data  = dict(state.get("expose_data") or {})
        projekt_name = state.get("projekt_name", "Expose")

        # Derived: ENTWICKLER_ADRESSE_INVERS = "{plz_ort}, {strasse}"
        if not expose_data.get("entwickler_adresse_invers"):
            plz_ort = (expose_data.get("entwickler_plz_ort") or "").strip()
            strasse = (expose_data.get("entwickler_strasse") or "").strip()
            if plz_ort and strasse:
                expose_data["entwickler_adresse_invers"] = f"{plz_ort}, {strasse}"

        _write_job(job_id, rechtlich_phase="PPTX wird befuellt …")
        tmpl_bytes = _load_rechtlich_template_bytes()
        filled     = fill_pptx(tmpl_bytes, expose_data)

        _write_job(job_id, rechtlich_phase="PDF wird konvertiert (CloudConvert) …")
        pdf_bytes  = convert_to_pdf(filled, f"{projekt_name}_Rechtlich.pptx")

        teil_path  = os.path.join(_job_dir(job_id), "teilungserklaerung.pdf")
        if os.path.exists(teil_path):
            _write_job(job_id, rechtlich_phase="Teilungserklaerung wird angehaengt …")
            with open(teil_path, "rb") as fh:
                teil_bytes = fh.read()
            pdf_bytes = _merge_rechtlich_pdf(pdf_bytes, teil_bytes)
            print(f"[{job_id}] Teilungserklaerung gemerged ({len(teil_bytes)//1024} KB)")

        out_path = os.path.join(_job_dir(job_id), "rechtlich.pdf")
        with open(out_path, "wb") as fh:
            fh.write(pdf_bytes)
        _write_job(job_id,
                   rechtlich_status="done",
                   rechtlich_phase="Fertig",
                   rechtlich_size=len(pdf_bytes))
        print(f"[{job_id}] Rechtlich-PDF fertig: {len(pdf_bytes)//1024} KB → {out_path}")
    except Exception as e:
        import traceback as _tb
        err = f"{e}\n{_tb.format_exc()[:600]}"
        print(f"[{job_id}] Rechtlich-Render fehlgeschlagen: {err}")
        _write_job(job_id, rechtlich_status="error", rechtlich_error=str(e))


@app.route("/job/<job_id>/render-rechtlich", methods=["POST", "OPTIONS"])
def job_render_rechtlich(job_id):
    """Startet Hintergrund-Render des rechtlichen Teils B. Gibt sofort 202 zurueck.
    Status kann via GET /job/<id> abgefragt werden (Felder rechtlich_status/_phase)."""
    if request.method == "OPTIONS":
        return make_response("", 204)
    if not _authorize_job(job_id):
        return jsonify({"error": "Unauthorized"}), 401
    job = _read_job(job_id)
    if not job:
        return jsonify({"error": "Job nicht gefunden"}), 404
    if not os.path.exists(_job_state_path(job_id)):
        return jsonify({"error": "Job-State nicht vorhanden — erst Marketing finalisieren"}), 400

    _write_job(job_id, rechtlich_status="processing",
               rechtlich_phase="Wird gestartet …", rechtlich_error=None)
    t = _threading.Thread(target=_run_render_rechtlich, args=(job_id,), daemon=True)
    t.start()
    return jsonify({"ok": True, "job_id": job_id})


@app.route("/job/<job_id>/download-rechtlich", methods=["GET", "OPTIONS"])
def job_download_rechtlich(job_id):
    """Liefert das generierte Rechtlich-PDF als Download."""
    if request.method == "OPTIONS":
        return make_response("", 204)
    if not _authorize_job(job_id):
        return jsonify({"error": "Unauthorized"}), 401
    path = os.path.join(_job_dir(job_id), "rechtlich.pdf")
    if not os.path.exists(path):
        return jsonify({"error": "Noch nicht generiert"}), 404
    name = (_read_job(job_id) or {}).get("name", "Expose")
    return send_file(path, mimetype="application/pdf",
                     as_attachment=True,
                     download_name=f"{name}_Rechtlicher-Teil-B.pdf")


@app.route("/debug-images", methods=["GET"])
def debug_images():
    """Testet den kompletten Bild-Pipeline ohne Upload: Unsplash → Download → PPTX.
    Gibt JSON-Bericht zurück. Nur mit korrektem API-Token aufrufbar."""
    if request.headers.get("X-API-Token") != API_TOKEN:
        return jsonify({"error": "Unauthorized"}), 401

    report = {
        "unsplash_key_set": bool(UNSPLASH_ACCESS_KEY),
        "unsplash_key_len": len(UNSPLASH_ACCESS_KEY) if UNSPLASH_ACCESS_KEY else 0,
        "unsplash_test": None,
        "picsum_test": None,
        "image_urls_count": 0,
        "image_downloads_ok": 0,
        "image_downloads_fail": 0,
    }

    # Unsplash Test
    try:
        r = requests.get(
            "https://api.unsplash.com/photos/random",
            params={"query": "modern apartment", "orientation": "landscape"},
            headers={"Authorization": f"Client-ID {UNSPLASH_ACCESS_KEY}"},
            timeout=10
        )
        report["unsplash_test"] = {"status": r.status_code, "body_preview": r.text[:200]}
        if r.status_code == 200:
            report["unsplash_sample_url"] = r.json()["urls"]["regular"]
    except Exception as e:
        report["unsplash_test"] = {"error": str(e)}

    # Picsum Test
    try:
        r = requests.get("https://picsum.photos/seed/42/200/150", timeout=10)
        report["picsum_test"] = {"status": r.status_code, "size_kb": len(r.content) // 1024}
    except Exception as e:
        report["picsum_test"] = {"error": str(e)}

    # Voller Image-Flow mit Dummy-Daten
    data = DUMMY_EXPOSE_DATA.copy()
    data = fill_image_placeholders(data)
    urls = {k: v for k, v in data.items() if k.startswith("bild_") and isinstance(v, str) and v.startswith("http")}
    report["image_urls_count"] = len(urls)
    report["sample_urls"] = dict(list(urls.items())[:3])

    for key, url in urls.items():
        try:
            r = requests.get(url, timeout=10)
            if r.status_code == 200:
                report["image_downloads_ok"] += 1
            else:
                report["image_downloads_fail"] += 1
        except Exception:
            report["image_downloads_fail"] += 1

    return jsonify(report)


@app.route("/fill-pptx", methods=["POST"])
def fill_pptx_endpoint():
    """Debug-Endpunkt: Gibt das rohe PPTX ohne PDF-Konvertierung zurück.
    Body: JSON mit optionalem 'data'-Objekt. Ohne 'data' → DUMMY_EXPOSE_DATA + Unsplash."""
    if request.headers.get("X-API-Token") != API_TOKEN:
        return jsonify({"error": "Unauthorized"}), 401
    try:
        body = request.get_json(force=True) or {}
        data = body.get("data") or DUMMY_EXPOSE_DATA.copy()
        data = {**PLATZHALTER, **data}
        data = fill_image_placeholders(data)

        bild_count = sum(1 for k, v in data.items()
                         if k.startswith("bild_") and isinstance(v, str) and v.startswith("http"))
        print(f"fill-pptx endpoint: {bild_count} bild_* URLs")

        tmpl_bytes = requests.get(TEMPLATE_URL, timeout=30).content
        pptx_bytes = fill_pptx(tmpl_bytes, data)

        projekt_name = data.get("projekt_name", "Debug").replace(" ", "_")
        return send_file(io.BytesIO(pptx_bytes),
                         mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                         as_attachment=True,
                         download_name=f"{projekt_name}_debug.pptx")
    except Exception as e:
        import traceback
        return jsonify({"error": str(e), "trace": traceback.format_exc()}), 500


if __name__ == "__main__":
    app.run(host="0.0.0.0", port=int(os.environ.get("PORT", 5000)))
