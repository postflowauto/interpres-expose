"""Generiert urbanunits_Rechtlich_v1.pptx aus dem DQN-Rechtlich-PDF.

Strategie: A4-Hochformat-PPTX, eine Folie pro Quell-Seite. Projektspezifische
Strings werden durch Platzhalter (z.B. {{ENTWICKLER_NAME}}) ersetzt, sodass das
Template per fill_pptx-Mechanismus pro Projekt befüllt werden kann.

Aufruf:
    python3 tools/build_rechtlich_template.py /path/to/DQN_Rechtlich.pdf

Output: urbanunits_Rechtlich_v1.pptx im Repo-Root.

Hinweis: Die Teilungserklärung (PDF-Seiten 103-307 im DQN-Original) wird
ABSICHTLICH NICHT ins Template uebernommen — sie kommt als User-Upload (jeder
Notar liefert sein eigenes PDF) und wird per PyPDF-Merge angehaengt.
"""
import sys, os, re
from pathlib import Path
from pypdf import PdfReader
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ROOT          = Path(__file__).resolve().parent.parent
DEFAULT_OUT   = ROOT / "urbanunits_Rechtlich_v1.pptx"

# Quell-Bereiche: alles AUSSER Teilungserklärung (103-307)
DYNAMIC_RANGES = [(1, 102), (308, 322)]

# Reihenfolge ist wichtig — längere/spezifischere Patterns ZUERST.
# Jeder Eintrag: (regex, replacement). Regex matched case-sensitive ausser
# wenn re.IGNORECASE im pattern selbst genutzt wird.
PLACEHOLDER_REPLACEMENTS = [
    # Notar / Abwicklung
    (r"Notar Peter Krolopp in 39112 Magdeburg,?\s*Humboldtstr\.\s*2",
        "Notar {{NOTAR_NAME}} in {{NOTAR_ADRESSE}}"),
    (r"Notar Peter Krolopp in Magdeburg", "Notar {{NOTAR_NAME}} in {{NOTAR_STADT}}"),
    (r"Peter Krolopp",        "{{NOTAR_NAME}}"),
    # Entwickler — voller Adress-Pattern zuerst
    (r"Diamant Quartier GmbH,?\s*Heyrothsberger Straße 8,?\s*39175 Biederitz",
        "{{ENTWICKLER_NAME}}, {{ENTWICKLER_ADRESSE}}"),
    (r"39175 Biederitz,?\s*Heyrothsberger Str\.?\s*8",
        "{{ENTWICKLER_ADRESSE_INVERS}}"),
    (r"Heyrothsberger Straße 8,?\s*39175 Biederitz",   "{{ENTWICKLER_ADRESSE}}"),
    (r"Heyrothsberger Str\.?\s*8",                     "{{ENTWICKLER_STRASSE}}"),
    (r"39175 Biederitz",                               "{{ENTWICKLER_PLZ_ORT}}"),
    (r"Diamant Quartier GmbH",                         "{{ENTWICKLER_NAME}}"),
    # WEG-Bezeichnung "Diamant Quartier" (ohne GmbH) — gemeinschaftliche Wohnanlage
    (r"Diamant Quartier(?!\s*GmbH)",                   "{{WEG_BEZEICHNUNG}}"),
    (r"Amtsgericht Stendal unter HRB 29572",
        "{{ENTWICKLER_HANDELSREGISTER}}"),
    # Verwalter
    (r"Büschel Immobilien GmbH",                       "{{VERWALTER_NAME}}"),
    (r"Listemannstraße 10",                            "{{VERWALTER_STRASSE}}"),
    (r"Herr Henry Berg und Herr Vincent Kühn",         "{{VERWALTER_GESCHAEFTSFUEHRER}}"),
    (r"Henry Berg und Vincent Kühn",                   "{{VERWALTER_GESCHAEFTSFUEHRER}}"),
    # Projekt-Adresse / Lage
    (r"Lübecker Str\.?/?\s*Gröperstr\.?\s*in 39124 Magdeburg",
        "{{PROJEKT_ADRESSE}}, {{PROJEKT_PLZ_STADT}}"),
    (r"Lübecker Str\.?/?\s*Gröperstr\.?",              "{{PROJEKT_ADRESSE}}"),
    (r"Lübecker Str\.?/?\s*Gröper Str\.?",             "{{PROJEKT_ADRESSE}}"),
    (r"Lübecker Straße",                               "{{PROJEKT_ADRESSE}}"),
    (r"Lübecker Str\.?",                               "{{PROJEKT_ADRESSE}}"),
    (r"Alten Diamant Brauerei",                        "{{PROJEKT_AREAL}}"),
    (r"Diamant Brauerei",                              "{{PROJEKT_AREAL}}"),
    (r"39124 Magdeburg",                               "{{PROJEKT_PLZ_STADT}}"),
    # Projekt-Daten
    (r"Neubau Mikroapartments",                        "{{PROJEKT_TYP}}"),
    (r"\b104 Eigentumswohnungen\b",                    "{{ANZAHL_WE}} Eigentumswohnungen"),
    (r"\b104 Wohneinheiten\b",                         "{{ANZAHL_WE}} Wohneinheiten"),
    (r"\b104 Wohnungen\b",                             "{{ANZAHL_WE}} Wohnungen"),
    # Termine (projektspezifisch)
    (r"30\.06\.2027",                                  "{{FERTIG_GEPLANT}}"),
    (r"31\.12\.2027",                                  "{{FERTIG_VERBINDLICH}}"),
    (r"04\.07\.2025",                                  "{{VERWEISURKUNDE_DATUM}}"),
    # Stand-Daten (Prospekt + Steuer)
    (r"Stand Oktober 2025",                            "Stand {{PROSPEKT_DATUM}}"),
    (r"Oktober 2025",                                  "{{PROSPEKT_DATUM}}"),
    (r"im Juli 2025",                                  "im {{STAND_STEUER}}"),
    (r"\bJuli 2025\b",                                 "{{STAND_STEUER}}"),
    # Preise + Größen
    (r"184\.092,00 €.*?bis\s*325\.612,00 €",
        "{{KAUFPREIS_VON}} € bis {{KAUFPREIS_BIS}} €"),
    (r"184\.092,00 €",                                 "{{KAUFPREIS_VON}} €"),
    (r"325\.612,00 €",                                 "{{KAUFPREIS_BIS}} €"),
    (r"31,74 m ?² bis 56,14 m²",
        "{{GROESSE_VON}} m² bis {{GROESSE_BIS}} m²"),
    # KfW + Baujahr
    (r'KfW 40 [„"]Klimafreundlicher Neubau – Wohngebäude["”]?\s*mit\s*QNG',
        "{{KFW_STUFE}}"),
    (r'KfW 40 [„"]Klimafreundlicher Neubau\s*–\s*Wohngebäude["”]?\s*mit QNG \(QNG-Zertifizierung wurde seitens\s*der Verkäuferin beauftragt\)',
        "{{KFW_STUFE}}"),
    (r"\b2025-2027\b",                                 "{{BAUJAHR_RANGE}}"),
]


def clean_pdf_text(t: str) -> str:
    """Glättet PDF-Extraktions-Artefakte (mehrfache Spaces, Zerstückelung)."""
    # Soft-Hyphens / weiche Trenner
    t = t.replace("­", "")
    # Trennstriche am Zeilenende: "Mikro-\napartments" → "Mikroapartments"
    t = re.sub(r"-\s*\n\s*", "", t)
    # Erst Doppel-Newlines (echte Absatztrenner) als Sentinel sichern
    t = t.replace("\n\n", "<<PARA>>")
    # Alle restlichen Single-Newlines → Space (PDF-Wortbruch-Artefakte)
    t = t.replace("\n", " ")
    # Sentinel zurück zu echten Absätzen
    t = t.replace("<<PARA>>", "\n\n")
    # Mehrere Spaces → 1 Space
    t = re.sub(r"[ \t]{2,}", " ", t)
    # Mehrere Leerabsätze → max 2
    t = re.sub(r"\n{3,}", "\n\n", t)
    return t.strip()


def apply_placeholders(text: str) -> tuple[str, int]:
    """Wendet alle Replacement-Regeln an. Gibt (neuer Text, Anzahl Ersetzungen) zurück."""
    total = 0
    for pat, rep in PLACEHOLDER_REPLACEMENTS:
        new, n = re.subn(pat, rep, text, flags=re.IGNORECASE)
        if n:
            total += n
            text = new
    return text, total


def build(src_pdf: str, out_pptx: str):
    reader = PdfReader(src_pdf)
    total_src = len(reader.pages)
    print(f"Quelle: {src_pdf} ({total_src} Seiten)")

    prs = Presentation()
    # A4 Hochformat: 21 cm × 29.7 cm
    prs.slide_width  = Cm(21)
    prs.slide_height = Cm(29.7)

    blank = prs.slide_layouts[6]  # blank-layout

    total_replacements = 0
    slides_built       = 0
    truncated_pages    = []

    for start, end in DYNAMIC_RANGES:
        for src_p in range(start, min(end, total_src) + 1):
            raw   = reader.pages[src_p - 1].extract_text() or ""
            clean = clean_pdf_text(raw)
            text, n = apply_placeholders(clean)
            total_replacements += n

            slide = prs.slides.add_slide(blank)
            tb = slide.shapes.add_textbox(
                left   = Cm(1.5),
                top    = Cm(1.5),
                width  = Cm(18),
                height = Cm(26.7),
            )
            tf = tb.text_frame
            tf.word_wrap = True
            try:
                tf.auto_size = None  # NO autofit — kann sonst Text shrink en
            except Exception:
                pass

            paragraphs = text.split("\n")
            # Erster Absatz: dem default-paragraph zuweisen
            p0 = tf.paragraphs[0]
            r0 = p0.add_run()
            r0.text = paragraphs[0] if paragraphs else ""
            r0.font.size = Pt(9)
            r0.font.name = "Arial"

            for line in paragraphs[1:]:
                p = tf.add_paragraph()
                r = p.add_run()
                r.text = line
                r.font.size = Pt(9)
                r.font.name = "Arial"

            # Wenn der Text sehr lang ist, koennte er ueberlaufen. Logge.
            if len(text) > 3500:
                truncated_pages.append(src_p)

            # Footer mit Original-Seitenzahl (Debug-Aid)
            ftr = slide.shapes.add_textbox(
                left=Cm(1.5), top=Cm(28.5), width=Cm(18), height=Cm(0.8),
            )
            fp = ftr.text_frame.paragraphs[0]
            fr = fp.add_run()
            fr.text = f"— DQN-Original S. {src_p} —"
            fr.font.size = Pt(7)
            fr.font.name = "Arial"
            fr.font.color.rgb = RGBColor(0x80, 0x80, 0x80)
            fp.alignment = PP_ALIGN.CENTER

            slides_built += 1

    prs.save(out_pptx)
    print(f"\n✓ {slides_built} Folien generiert")
    print(f"  Platzhalter-Ersetzungen: {total_replacements}")
    if truncated_pages:
        print(f"  ⚠️  {len(truncated_pages)} Seiten >3500 Zeichen (evtl. Überlauf): {truncated_pages}")
    print(f"  Output: {out_pptx} ({os.path.getsize(out_pptx)//1024} KB)")


if __name__ == "__main__":
    src = sys.argv[1] if len(sys.argv) > 1 else \
          "/Users/anakinrosner/Downloads/DQN_RECHTLICHER TEIL B urbanunits - The Central_21.10.2025 final.pdf"
    out = sys.argv[2] if len(sys.argv) > 2 else str(DEFAULT_OUT)
    build(src, out)
