"""V2 Editor — schlanke Editor-Schicht über V1.

Konzept:
  - V1 macht die ganze schwere Arbeit (PDF-Analyse, Claude, Bilder, PPTX,
    PDF-Konvertierung, Slide-JPGs). Das Template kommt vom Designer.
  - V2 bietet eine schöne Editor-UI: Folien-Liste links, große Vorschau in
    der Mitte, Properties-Panel rechts mit allen editierbaren Texten und
    Bild-Slots der gerade ausgewählten Folie.
  - Edit → V2-API patcht state.json (expose_data) → triggert V1-Re-Render
    (PPTX neu füllen + PDF + Slide-JPGs) → Editor zeigt neue Vorschau.
  - "PDF herunterladen" → V1-Finalize (gleicher Pfad wie bisher).

Kein eigener Look mehr. V1-Pipeline + dein Template = Quelle der Wahrheit.
"""
from __future__ import annotations
import io
import os
import json
import threading
import time
from pathlib import Path
from flask import (request, send_file, jsonify, send_from_directory,
                   Response, redirect)

V2_DIR      = Path(__file__).parent
STATIC_DIR  = V2_DIR / "static"
EDITOR_HTML = V2_DIR / "editor.html"
JOB_DIR     = "/tmp/interpres_jobs"

# ── Per-Slide-Platzhalter-Scan (gecached, pro Typ) ──────────────────────────
# Liest das Template einmal pro Typ und merkt sich pro Slide alle {{KEY}}-
# Platzhalter (lowercase). Wird im Editor genutzt um pro Folie nur die
# relevanten Edit-Felder anzuzeigen. Cache pro typ ('marketing','kurz','rechtlich').
_PER_SLIDE_PLACEHOLDERS_BY_TYP: dict[str, list[list[str]]] = {}


def _scan_template_placeholders(pptx_bytes: bytes) -> list:
    """Liefert pro Slide die Liste der Platzhalter-Keys (lowercase, dedupliziert).
    Robust gegen Whitespace-Splits, Line-Breaks innerhalb von {{...}} und |Xpt-Hints.
    """
    import io as _io
    import re as _re
    from pptx import Presentation
    PH = _re.compile(r"\{\{(.*?)\}\}", _re.DOTALL)

    def _extract(text):
        out = set()
        for m in PH.finditer(text):
            inner = m.group(1)
            if "|" in inner:
                inner = inner.split("|")[0]
            key = _re.sub(r"\s+", "", inner).lower().replace("-", "")
            if key and _re.match(r"^[a-z][a-z0-9_]*$", key):
                out.add(key)
        return out

    prs = Presentation(_io.BytesIO(pptx_bytes))
    result = []
    for slide in prs.slides:
        keys = set()
        def _scan_tf(tf):
            keys.update(_extract(tf.text))
        for shape in slide.shapes:
            try:
                if shape.has_text_frame:
                    _scan_tf(shape.text_frame)
                if shape.shape_type == 6:
                    for child in shape.shapes:
                        if child.has_text_frame:
                            _scan_tf(child.text_frame)
            except Exception:
                continue
        result.append(sorted(keys))
    return result


_LOCAL_TEMPLATE_FILES = {
    "marketing": "urbanunits_Marketing_Expose_v3.pdf-24.pptx",
    "kurz":      "urbanunits_Kurzexpose-4.pptx",
    "rechtlich": "urbanunits_Rechtlich_v1.pptx",
}


def _load_template_bytes(typ: str = "marketing") -> bytes | None:
    """Laed das PPTX-Template fuer den angegebenen Typ. Erst lokal aus Repo,
    Fallback URL aus TEMPLATE_URLS. Gibt None zurueck wenn beides nicht klappt."""
    repo_root = os.path.dirname(os.path.dirname(__file__))
    fname = _LOCAL_TEMPLATE_FILES.get(typ, "")
    if fname:
        path = os.path.join(repo_root, fname)
        if os.path.exists(path):
            with open(path, "rb") as fh:
                return fh.read()
    try:
        import importlib, requests
        appmod = importlib.import_module("app")
        url = appmod.TEMPLATE_URLS.get(typ) or appmod.TEMPLATE_URL
        if not url:
            return None
        resp = requests.get(url, timeout=30)
        if resp.status_code != 200:
            print(f"[v2] Template-URL ({typ}): HTTP {resp.status_code} — {url[:80]}")
            return None
        return resp.content
    except Exception as e:
        print(f"[v2] Template-Load ({typ}) Fehler: {e}")
        return None


def _get_template_placeholders(typ: str = "marketing") -> list[list[str]]:
    """Lazy-loads + cached: pro Slide alle Template-Platzhalter fuer den
    angegebenen Typ. Wenn das Template fuer den Typ nicht erreichbar ist
    (z.B. KURZ-Template noch nicht gepusht), wird [] zurueckgegeben — der
    Editor zeigt dann eine Hinweis-Meldung statt zu crashen."""
    if typ in _PER_SLIDE_PLACEHOLDERS_BY_TYP:
        return _PER_SLIDE_PLACEHOLDERS_BY_TYP[typ]
    tmpl_bytes = _load_template_bytes(typ)
    if not tmpl_bytes:
        _PER_SLIDE_PLACEHOLDERS_BY_TYP[typ] = []
        return []
    try:
        scanned = _scan_template_placeholders(tmpl_bytes)
        _PER_SLIDE_PLACEHOLDERS_BY_TYP[typ] = scanned
        print(f"[v2] Template-Scan ({typ}): {len(scanned)} Slides, "
              f"{sum(len(s) for s in scanned)} Platzhalter total")
    except Exception as e:
        print(f"[v2] Template-Scan ({typ}) Fehler: {e}")
        _PER_SLIDE_PLACEHOLDERS_BY_TYP[typ] = []
    return _PER_SLIDE_PLACEHOLDERS_BY_TYP[typ]


def _v1_state_path(job_id: str) -> str:
    return os.path.join(JOB_DIR, f"work_{job_id}", "state.json")


def _v1_meta_path(job_id: str) -> str:
    return os.path.join(JOB_DIR, f"{job_id}.json")


def _v1_slides_dir(job_id: str, typ: str = "marketing") -> str:
    """Slide-JPGs pro Expose-Typ in eigenem Unterordner.
    marketing → ../slides   (Default, Backwards-Compat)
    kurz/rechtlich → ../slides_<typ>
    """
    base = os.path.join(JOB_DIR, f"work_{job_id}")
    if typ == "marketing":
        return os.path.join(base, "slides")
    return os.path.join(base, f"slides_{typ}")


def _output_path(job_id: str, typ: str, ext: str = "pdf") -> str:
    """PDF/PPTX-Cache pro Typ. marketing behaelt den Original-Pfad fuer Backwards-Compat."""
    if typ == "marketing":
        return os.path.join(JOB_DIR, f"{job_id}.{ext}")
    return os.path.join(JOB_DIR, f"{job_id}.{typ}.{ext}")


VALID_EXPOSE_TYPS = ("marketing", "kurz", "rechtlich")


def _v1_uploads_dir(job_id: str) -> str:
    return os.path.join(JOB_DIR, f"work_{job_id}", "uploads")


def _read_state(job_id: str) -> dict:
    """Liest state.json (expose_data + customer_images_files + projekt_name)."""
    with open(_v1_state_path(job_id)) as f:
        return json.load(f)


def _write_state(job_id: str, state: dict):
    path = _v1_state_path(job_id)
    tmp = path + ".tmp"
    with open(tmp, "w") as f:
        json.dump(state, f, ensure_ascii=False)
    os.replace(tmp, path)


def _read_meta(job_id: str) -> dict:
    try:
        with open(_v1_meta_path(job_id)) as f:
            return json.load(f)
    except Exception:
        return {}


def _write_meta(job_id: str, **fields):
    path = _v1_meta_path(job_id)
    try:
        with open(path) as f:
            data = json.load(f)
    except Exception:
        data = {}
    data.update(fields)
    tmp = path + ".tmp"
    with open(tmp, "w") as f:
        json.dump(data, f)
    os.replace(tmp, path)


# ── Slot-Labels (übernommen aus V1, hier im V2-Scope für Editor-Anzeige) ──
SLOT_LABELS = {
    "bild_titel":              "Titelbild (Außenansicht)",
    "bild_projekt_aussen":     "Projekt – Außenansicht",
    "bild_projekt":            "Projekt – Bild",
    "bild_quartier":           "Quartier / Umgebung",
    "bild_greenliving_1":      "Green Living – Bild 1",
    "bild_greenliving_2":      "Green Living – Bild 2",
    "bild_interior":           "Innenraum / Interior",
    "bild_ausstattung_1":      "Ausstattung 1",
    "bild_ausstattung_2":      "Ausstattung 2",
    "bild_ausstattung_3":      "Ausstattung 3",
    "bild_ausstattung_4":      "Ausstattung 4",
    "bild_ausstattung_5":      "Ausstattung 5",
    "bild_ausstattung_6":      "Ausstattung 6",
    "bild_ansicht_1":          "Außenansicht 1",
    "bild_ansicht_2":          "Außenansicht 2",
    "bild_standort_innen":     "Standort innen",
    "bild_standort_aussen":    "Standort außen",
    "bild_lageplan":           "Lageplan",
    "bild_stadt_gross":        "Stadtbild groß",
    "bild_stadt_klein":        "Stadtbild klein",
    "bild_grundriss_intro_1":  "Grundriss-Intro 1",
    "bild_grundriss_intro_2":  "Grundriss-Intro 2",
    "bild_grundriss_intro_3":  "Grundriss-Intro 3",
    "bild_grundriss_1":        "Grundriss 1",
    "bild_grundriss_2":        "Grundriss 2",
    "bild_grundriss_3":        "Grundriss 3",
    "bild_grundriss_4":        "Grundriss 4",
    "bild_collage_1":          "Collage 1",
    "bild_collage_2":          "Collage 2",
    "bild_collage_3":          "Collage 3",
    "bild_collage_4":          "Collage 4",
    "bild_collage_5":          "Collage 5",
    "bild_hotel_1":            "Hotel-Feeling 1",
    "bild_hotel_2":            "Hotel-Feeling 2",
    "bild_rechtlich_1":        "Rechtliches – Bild 1",
    "bild_rechtlich_2":        "Rechtliches – Bild 2",
    "bild_stadt_presse":       "Stadtbild – Presse",
    "bild_stadt_branche":      "Stadtbild – Branche",
    # Kurz-Exposé-Bildslots
    "bild_titel_1":            "Kurz-Cover Bild 1",
    "bild_titel_2":            "Kurz-Cover Bild 2",
    "bild_titel_3":            "Kurz-Cover Bild 3",
    "bild_titel_4":            "Kurz-Cover Bild 4 (Hauptbild)",
    "bild_titel_5":            "Kurz-Cover Bild 5",
    "bild_titel_6":            "Kurz-Cover Bild 6",
    "bild_kurz_1":             "Kurz-Seite 2 – Bild 1",
    "bild_kurz_2":             "Kurz-Seite 2 – Bild 2",
    "bild_kurz_3":             "Kurz-Seite 2 – Bild 3",
    "bild_kurz_4":             "Kurz-Seite 2 – Bild 4",
}


def _slot_label(key: str) -> str:
    if key in SLOT_LABELS:
        return SLOT_LABELS[key]
    import re
    m = re.match(r"^bild_amenity_(\d+)$", key)
    if m:
        return f"Amenity {m.group(1)}"
    m = re.match(r"^bild_we_(\d+)$", key)
    if m:
        return f"WE-Bild {m.group(1)}"
    return key.replace("_", " ").title()


# Field-Gruppierung — welche expose_data-Keys gehören thematisch zusammen
FIELD_GROUPS = [
    {"name": "Projekt", "keys": [
        ("projekt_titel",         "Projekttitel"),
        ("entwickler_name",       "Entwicklername"),
        ("anzahl_we",             "Anzahl Wohneinheiten"),
        ("groesse_von",           "Größe von (m²)"),
        ("groesse_bis",           "Größe bis (m²)"),
        ("produkt_beschreibung",  "Produkt-Beschreibung"),
        ("kaufpreis_ab",          "Kaufpreis ab (€)"),
        ("kfw_standard",          "KfW-Standard"),
        ("kfw_darlehen",          "KfW-Darlehen (€)"),
        ("energieversorgung",     "Energieversorgung"),
        ("stellplaetze",          "Stellplätze"),
        ("zitat_intro",           "Intro-Zitat"),
    ]},
    {"name": "Stadt & Lage", "keys": [
        ("stadt",                 "Stadt"),
        ("stadtteil",             "Stadtteil"),
        ("adresse_lang",          "Adresse"),
        ("plz",                   "PLZ"),
        ("bundesland",            "Bundesland"),
        ("stadt_einwohner",       "Einwohner"),
        ("bundesland_bip",        "BIP Bundesland"),
        ("stadt_mietsteigerung",  "Mietsteigerung"),
        ("stadt_studierende",     "Studierende"),
    ]},
    {"name": "Texte – Investment", "keys": [
        ("text_intro",                "Intro-Text"),
        ("text_investment_pitch",     "Investment-Pitch"),
        ("text_kapitel_invest_1",     "Kapitel Invest – Lead 1"),
        ("text_kapitel_invest_2",     "Kapitel Invest – Lead 2"),
    ]},
    {"name": "Texte – Standort", "keys": [
        ("text_kapitel_live_1",       "Kapitel Live – Lead 1"),
        ("text_kapitel_live_2",       "Kapitel Live – Lead 2"),
        ("text_standort_1",           "Standort-Text 1"),
        ("text_standort_2",           "Standort-Text 2"),
        ("text_stadt_intro",          "Stadt-Intro"),
        ("text_stadt_wachstum_1",     "Stadt-Wachstum 1"),
        ("text_stadt_wachstum_2",     "Stadt-Wachstum 2"),
    ]},
    {"name": "Texte – Projekt & Stay", "keys": [
        ("text_kapitel_stay_1",       "Kapitel Stay – Lead 1"),
        ("text_kapitel_stay_2",       "Kapitel Stay – Lead 2"),
        ("text_greenliving_1",        "Greenliving 1"),
        ("text_greenliving_2",        "Greenliving 2"),
        ("text_ausstattung_kurz",     "Ausstattung – kurz"),
        ("text_ausstattung_detail",   "Ausstattung – Detail"),
        ("text_architektur",          "Architektur"),
    ]},
    {"name": "Min zu …", "keys": [
        ("min_uni",        "Min – Uni"),
        ("label_min_uni",  "Label – Uni"),
        ("min_bahnhof",    "Min – Bahnhof"),
        ("label_min_bahnhof", "Label – Bahnhof"),
        ("min_altstadt",   "Min – Altstadt"),
        ("label_min_altstadt", "Label – Altstadt"),
    ]},
    {"name": "Kurz-Exposé", "keys": [
        ("projekt_untertitel",        "Untertitel / Tagline"),
        ("projekt_beschreibung",      "Beschreibung (Pitch-Text Seite 2)"),
        ("projekt_beschreibung_kurz", "Beschreibung (Backup-Alias)"),
        ("text_relevanz",             "USP – Relevanz"),
        ("text_design",               "USP – Design"),
        ("text_foerderung",           "USP – Förderung"),
        ("text_tech",                 "USP – Tech"),
        ("besonderheiten_liste",      "Besonderheiten (Liste)"),
        ("gesamtwohnflaeche",         "Gesamtwohnfläche"),
        ("zimmer_anzahl_min",         "Zimmer min"),
        ("zimmer_anzahl_max",         "Zimmer max"),
    ]},
]


def register_v2(app):
    """Registriert alle /v2/*-Routen an der Flask-App."""

    @app.route("/v2/static/<path:filename>")
    def v2_static(filename):
        return send_from_directory(STATIC_DIR, filename)

    @app.route("/v2/health")
    def v2_health():
        return jsonify({"v2": True, "mode": "editor-on-v1"})

    @app.route("/v2/from-job/<job_id>")
    def v2_from_job(job_id):
        """Öffnet den V2-Editor für einen V1-Job."""
        if not os.path.exists(_v1_state_path(job_id)):
            return Response("""<!doctype html>
<html><head><meta charset="utf-8"><title>Job nicht gefunden</title>
<style>
body{background:#0a1220;color:#e8d9b3;font-family:-apple-system,sans-serif;
     display:flex;align-items:center;justify-content:center;min-height:100vh;margin:0;}
.box{max-width:520px;text-align:center;padding:40px;}
h1{font-family:'Playfair Display',serif;color:#C8A96E;font-size:24px;margin-bottom:16px;}
p{color:#c8d1de;line-height:1.6;margin-bottom:24px;}
a.btn{display:inline-block;background:#C8A96E;color:#0a1220;text-decoration:none;
      padding:12px 28px;border-radius:6px;font-weight:600;letter-spacing:0.05em;}
a.btn:hover{background:#d4ba84;}
small{color:#6b7d96;font-size:11px;display:block;margin-top:30px;}
</style></head>
<body><div class="box">
<h1>Dieser Job ist nicht mehr verfügbar</h1>
<p>Der Server hat zwischenzeitlich neu gestartet oder der Job ist abgelaufen.
Erstelle einfach ein neues Exposé — du wirst automatisch in den Editor geleitet.</p>
<a class="btn" href="/">Neues Exposé generieren</a>
<small>Job-ID: """ + job_id + """</small>
</div></body></html>""",
                status=404, mimetype="text/html; charset=utf-8"
            )
        return redirect(f"/v2/editor/{job_id}")

    @app.route("/v2/editor/<job_id>")
    def v2_editor(job_id):
        if not os.path.exists(_v1_state_path(job_id)):
            return Response("Job nicht gefunden.", status=404)
        with open(EDITOR_HTML, encoding="utf-8") as f:
            html = f.read()
        # Job-ID + Token ins Frontend injizieren
        api_token = os.environ.get("API_TOKEN", "interpres-secret-2026")
        inject = (
            f"<script>"
            f"window.JOB_ID = {json.dumps(job_id)};"
            f"window.API_TOKEN = {json.dumps(api_token)};"
            f"</script>"
        )
        html = html.replace("<head>", "<head>\n" + inject)
        return Response(html, mimetype="text/html")

    # ── API: Job-State ───────────────────────────────────────────────────
    @app.route("/v2/api/job/<job_id>")
    def v2_api_job_get(job_id):
        if not os.path.exists(_v1_state_path(job_id)):
            return jsonify({"error": "not found"}), 404
        # Optional ?typ=marketing|kurz|rechtlich → liefert Slide-Count fuer den Typ.
        # Default = last_render_typ aus meta, sonst 'marketing'.
        state = _read_state(job_id)
        meta  = _read_meta(job_id)
        expose = state.get("expose_data", {})
        typ = (request.args.get("typ") or meta.get("last_render_typ") or "marketing").lower()
        if typ not in VALID_EXPOSE_TYPS:
            typ = "marketing"
        slides_dir = _v1_slides_dir(job_id, typ)
        slide_count = 0
        if os.path.isdir(slides_dir):
            slide_count = len([n for n in os.listdir(slides_dir)
                              if n.startswith("slide_") and n.endswith(".jpg")])

        # Per-Typ Render-Status (welche PDFs sind schon gecacht)
        rendered_typs = [
            t for t in VALID_EXPOSE_TYPS
            if os.path.exists(_output_path(job_id, t, "pdf"))
            or os.path.exists(_output_path(job_id, t, "pptx"))
        ]

        # Welche Templates sind ueberhaupt konfiguriert/erreichbar?
        # marketing ist immer da (Pflicht), kurz/rechtlich nur wenn URL erreichbar.
        import importlib as _il
        _appmod = _il.import_module("app")
        available_typs = ["marketing"]
        for t in ("kurz", "rechtlich"):
            url = _appmod.TEMPLATE_URLS.get(t)
            if not url:
                continue
            # Template-Verfuegbarkeit ueber den Scan-Cache pruefen — wenn
            # _get_template_placeholders eine nicht-leere Liste liefert,
            # ist das Template erreichbar.
            phs = _get_template_placeholders(t)
            if phs:
                available_typs.append(t)

        # Bild-Slots aus expose_data extrahieren
        bild_slots = []
        for k, v in expose.items():
            if k.startswith("bild_"):
                bild_slots.append({
                    "key":      k,
                    "label":    _slot_label(k),
                    "value":    v if isinstance(v, str) else "",
                    "has_url":  bool(v and isinstance(v, str) and v.startswith("http")),
                })
        bild_slots.sort(key=lambda x: x["key"])

        # Hochgeladene Customer-Images
        uploads_dir = _v1_uploads_dir(job_id)
        uploaded = {}
        if os.path.isdir(uploads_dir):
            for fname in os.listdir(uploads_dir):
                slot = os.path.splitext(fname)[0].lower()
                uploaded[slot] = fname

        # Pro-Slide-Platzhalter aus dem Template des AKTUELLEN Typs
        # (Editor zeigt damit pro Folie nur die relevanten Edit-Felder).
        slide_placeholders = _get_template_placeholders(typ)
        # Flache Liste ALLER Platzhalter dieses Typs (ueber alle Folien) —
        # dient als typ-sicherer Fallback im Editor: wenn das Pro-Folie-
        # Scanning mal leer ist, werden trotzdem NUR Felder dieses Typs
        # gezeigt (nie marketing/kurz/rechtlich gemischt).
        typ_placeholders = sorted({k for slide in slide_placeholders for k in slide})

        return jsonify({
            "job_id":      job_id,
            "name":        meta.get("name") or expose.get("projekt_titel", "Expose"),
            "status":      meta.get("status", "unknown"),
            "phase":       meta.get("phase", ""),
            "current_typ": typ,
            "slide_count": slide_count,
            "rendered_typs": rendered_typs,
            "available_typs": available_typs,
            "expose":      expose,
            "field_groups": FIELD_GROUPS,
            "bild_slots":  bild_slots,
            "uploaded":    uploaded,
            "slide_placeholders": slide_placeholders,
            "typ_placeholders":   typ_placeholders,
        })

    @app.route("/v2/api/job/<job_id>/text", methods=["PUT", "OPTIONS"])
    def v2_api_text_put(job_id):
        if request.method == "OPTIONS":
            return ("", 204)
        if not os.path.exists(_v1_state_path(job_id)):
            return jsonify({"error": "not found"}), 404
        body = request.get_json(force=True) or {}
        # body = { "key1": "value1", "key2": "value2", ... }
        state = _read_state(job_id)
        expose = state.get("expose_data", {})
        for k, v in body.items():
            # Akzeptiere nur Strings, keine Bild-URLs (die werden separat gehandhabt)
            if k.startswith("bild_"):
                continue
            expose[k] = v if isinstance(v, str) else str(v)
        state["expose_data"] = expose
        _write_state(job_id, state)
        return jsonify({"ok": True, "updated": list(body.keys())})

    @app.route("/v2/api/job/<job_id>/render", methods=["POST", "OPTIONS"])
    def v2_api_render(job_id):
        """Triggert V1-Re-Render: PPTX neu füllen mit aktueller expose_data
        + Customer-Uploads, dann PDF + Slide-JPGs erzeugen.
        Setzt status="processing" + phase, läuft im Background-Thread.

        Param: ?typ=marketing|kurz|rechtlich (Default: marketing)
        """
        if request.method == "OPTIONS":
            return ("", 204)
        if not os.path.exists(_v1_state_path(job_id)):
            return jsonify({"error": "not found"}), 404

        typ = (request.args.get("typ") or request.form.get("typ") or "marketing").lower()
        if typ not in VALID_EXPOSE_TYPS:
            return jsonify({"error": f"unknown typ: {typ}"}), 400

        _write_meta(job_id, status="processing", phase=f"V2-Render läuft … ({typ})",
                    last_render_typ=typ)
        t = threading.Thread(target=_v2_render_worker,
                             args=(job_id, typ),
                             daemon=True)
        t.start()
        return jsonify({"ok": True, "typ": typ})

    @app.route("/v2/api/job/<job_id>/teilungserklaerung",
               methods=["GET", "POST", "DELETE", "OPTIONS"])
    def v2_api_teilungserklaerung(job_id):
        """PDF-Upload der Teilungserklaerung fuer den rechtlichen Teil B.

        Speicherung doppelt:
        - job-spezifisch:  work_<id>/teilungserklaerung.pdf (Override pro Job)
        - projekt-spezifisch: appmod._project_teilung_path (gilt fuer alle
          Jobs desselben Projekts).
        Beim POST geht der Default-Scope auf 'project' (Vertriebs-Komfort);
        ?scope=job nur fuer Sonderfaelle.
        """
        if request.method == "OPTIONS":
            return ("", 204)
        import importlib
        appmod = importlib.import_module("app")
        state = _read_state(job_id)
        projekt_name = (state.get("projekt_name") or "default") if state else "default"

        work_dir  = os.path.join(JOB_DIR, f"work_{job_id}")
        os.makedirs(work_dir, exist_ok=True)
        job_path  = os.path.join(work_dir, "teilungserklaerung.pdf")
        proj_path = appmod._project_teilung_path(projekt_name)

        if request.method == "GET":
            if os.path.exists(job_path):
                return jsonify({"uploaded": True, "size": os.path.getsize(job_path),
                                "scope": "job"})
            if os.path.exists(proj_path):
                return jsonify({"uploaded": True, "size": os.path.getsize(proj_path),
                                "scope": "project", "projekt_name": projekt_name})
            return jsonify({"uploaded": False})

        if request.method == "DELETE":
            scope = (request.args.get("scope") or "job").lower()
            removed = []
            if os.path.exists(job_path):
                try: os.remove(job_path); removed.append("job")
                except OSError: pass
            if scope == "project" and os.path.exists(proj_path):
                try: os.remove(proj_path); removed.append("project")
                except OSError: pass
            return jsonify({"ok": True, "removed": removed})

        # POST
        f = request.files.get("file")
        if not f or not f.filename:
            return jsonify({"error": "Keine Datei"}), 400
        if not f.filename.lower().endswith(".pdf"):
            return jsonify({"error": "Nur PDF erlaubt"}), 400
        scope = (request.args.get("scope") or "project").lower()
        target = job_path if scope == "job" else proj_path
        f.save(target)
        sz = os.path.getsize(target)
        print(f"[v2] [{job_id}] Teilungserklaerung gespeichert ({scope}, {sz//1024} KB): {target}")
        return jsonify({"ok": True, "size": sz, "scope": scope,
                        "projekt_name": projekt_name})

    @app.route("/v2/api/job/<job_id>/render-status")
    def v2_api_render_status(job_id):
        meta = _read_meta(job_id)
        return jsonify({
            "status": meta.get("status", "unknown"),
            "phase":  meta.get("phase", ""),
            "error":  meta.get("error", None),
            "last_render_typ": meta.get("last_render_typ", "marketing"),
            # Per-Typ-Status: zeigt welche PDFs schon gerendert wurden (lazy-Cache)
            "rendered_typs": [
                t for t in VALID_EXPOSE_TYPS
                if os.path.exists(_output_path(job_id, t, "pdf"))
                or os.path.exists(_output_path(job_id, t, "pptx"))
            ],
        })

    @app.route("/v2/api/job/<job_id>/download")
    def v2_api_download(job_id):
        """Liefert PDF zurück. ?typ=marketing|kurz|rechtlich (Default: marketing)."""
        typ = (request.args.get("typ") or "marketing").lower()
        if typ not in VALID_EXPOSE_TYPS:
            return jsonify({"error": f"unknown typ: {typ}"}), 400

        meta = _read_meta(job_id)
        # Pro-Typ-Pfad pruefen (PDF bevorzugt, PPTX als Fallback wenn keine
        # Konvertierung verfuegbar war)
        pdf_path = _output_path(job_id, typ, "pdf")
        pptx_path = _output_path(job_id, typ, "pptx")
        if os.path.exists(pdf_path):
            out_path, mt, ext_label = pdf_path, "application/pdf", "pdf"
        elif os.path.exists(pptx_path):
            out_path = pptx_path
            mt = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            ext_label = "pptx"
        else:
            return jsonify({
                "error": f"Kein {typ}-PDF vorhanden – erst 'Aktualisieren' klicken."
            }), 409

        name_base = meta.get("name", "Expose")
        suffix = "" if typ == "marketing" else f"_{typ}"
        # ENTWURF-Markierung im Dateinamen wenn rechtlich ohne Teilung
        if typ == "rechtlich" and meta.get("rechtlich_is_entwurf"):
            suffix += "_ENTWURF"
        return send_file(out_path, mimetype=mt, as_attachment=True,
                         download_name=f"{name_base}{suffix}.{ext_label}")


# ── Kurz-Expose Image Auto-Mapping ───────────────────────────────────────
# Welche Marketing-Slots (in Prio-Reihenfolge) sollen welchen Kurz-Slot
# befuellen. Wird durchsucht: erstes verfuegbares Marketing-Bild gewinnt.
_KURZ_IMAGE_MAPPING = {
    # Cover Slide 1 — 6er Collage (bild_titel_4 ist das Hauptbild)
    "bild_titel_1": ["bild_quartier", "bild_amenity_1", "bild_ansicht_1"],
    "bild_titel_2": ["bild_interior", "bild_ausstattung_1", "bild_hotel_1"],
    "bild_titel_3": ["bild_greenliving_1", "bild_amenity_2", "bild_quartier"],
    "bild_titel_4": ["bild_titel", "bild_projekt_aussen", "bild_ansicht_1"],
    "bild_titel_5": ["bild_projekt_aussen", "bild_ansicht_2", "bild_ansicht_1"],
    "bild_titel_6": ["bild_ausstattung_2", "bild_hotel_2", "bild_amenity_3"],
    # Seite 2 — 4 Bilder
    "bild_kurz_1": ["bild_projekt_aussen", "bild_titel", "bild_ansicht_1"],
    "bild_kurz_2": ["bild_greenliving_2", "bild_ansicht_2", "bild_quartier"],
    "bild_kurz_3": ["bild_hotel_1", "bild_interior", "bild_ausstattung_1"],
    "bild_kurz_4": ["bild_hotel_2", "bild_amenity_1", "bild_ausstattung_2"],
}
_KURZ_UNSPLASH_FALLBACKS = {
    "bild_titel_1": "city lifestyle people street",
    "bild_titel_2": "modern apartment bedroom interior",
    "bild_titel_3": "modern apartment building exterior",
    "bild_titel_4": "modern residential building facade",
    "bild_titel_5": "residential apartment building street",
    "bild_titel_6": "young woman lifestyle outdoor",
    "bild_kurz_1": "modern apartment building exterior",
    "bild_kurz_2": "apartment building aerial green",
    "bild_kurz_3": "coffee shop friends modern",
    "bild_kurz_4": "young people lifestyle group friends",
}


def _auto_map_kurz_images(expose: dict, customer_images: dict, appmod) -> None:
    """Befuellt leere Kurz-Bildslots mit (1) passenden Marketing-Bildern des
    Kunden oder (2) Unsplash-Fallbacks. Mutiert expose + customer_images in-place."""
    for kurz_slot, sources in _KURZ_IMAGE_MAPPING.items():
        # Kunde hat selbst hochgeladen → unverändert lassen
        if customer_images.get(kurz_slot):
            continue
        # Marketing-Bild als Bytes uebertragen
        for src in sources:
            if customer_images.get(src):
                customer_images[kurz_slot] = customer_images[src]
                expose[kurz_slot] = ""
                print(f"[v2 kurz] auto-map: {kurz_slot} <- {src} (customer)")
                break
        else:
            # Marketing hat eine Unsplash-URL fuer diesen Marketing-Slot
            for src in sources:
                val = expose.get(src)
                if isinstance(val, str) and val.startswith("http"):
                    expose[kurz_slot] = val
                    print(f"[v2 kurz] auto-map: {kurz_slot} <- {src} (url)")
                    break
            else:
                # Letzter Schritt: frische Unsplash-Suche
                query = _KURZ_UNSPLASH_FALLBACKS.get(kurz_slot, "modern apartment")
                try:
                    url = appmod.fetch_unsplash_image(query)
                    if url:
                        expose[kurz_slot] = url
                        print(f"[v2 kurz] auto-map: {kurz_slot} <- unsplash({query!r})")
                except Exception as e:
                    print(f"[v2 kurz] unsplash {kurz_slot}: {e}")


def _v2_render_worker(job_id: str, typ: str = "marketing"):
    """Background-Worker: lädt expose_data + Customer-Images, ruft V1-fill_pptx,
    konvertiert zu PDF, rendert Slide-JPGs.

    typ: 'marketing' (Default) | 'kurz' | 'rechtlich'
      Bestimmt welches Template geladen + welche Output-Pfade benutzt werden.

    RAM-Strategie für 512 MB Render Starter:
      - Nach jedem grossen Schritt: del + gc.collect()
      - tmpl_bytes/customer_images sofort nach fill_pptx wegwerfen
      - pptx_bytes wegwerfen sobald pdf_bytes existiert
      - Slide-JPGs in temp-Dir rendern, dann atomic swap → kein Datenverlust bei Render-Fail
    """
    import gc
    try:
        # V1-Funktionen importieren (Lazy, um Zirkular-Import zu vermeiden)
        import importlib
        appmod = importlib.import_module("app")

        _write_meta(job_id, status="processing", phase="State + Bilder werden geladen …")
        state = _read_state(job_id)
        expose = state.get("expose_data", {})
        cust_files = state.get("customer_images_files", {}) or {}

        customer_images = {}
        for slot, fpath in cust_files.items():
            try:
                with open(fpath, "rb") as f:
                    customer_images[slot] = f.read()
            except Exception:
                pass

        # User-Uploads übersteuern Auto-Zuweisungen
        uploads_dir = _v1_uploads_dir(job_id)
        if os.path.isdir(uploads_dir):
            for fname in os.listdir(uploads_dir):
                base, ext = os.path.splitext(fname)
                slot = base.lower()
                if not slot.startswith("bild_"):
                    continue
                with open(os.path.join(uploads_dir, fname), "rb") as f:
                    customer_images[slot] = f.read()
                expose[slot] = ""

        # ── Auto-Mapping fuer Kurz-Exposé ──────────────────────────────────
        # Wenn der Kunde im Marketing schon passende Bilder hochgeladen hat,
        # werden sie hier auf die Kurz-Slots uebertragen. Was leer bleibt,
        # bekommt einen Unsplash-Fallback.
        if typ == "kurz":
            _auto_map_kurz_images(expose, customer_images, appmod)

        projekt_name = expose.get("projekt_titel", "Expose")

        _write_meta(job_id, phase=f"Template wird gefüllt … ({typ})")
        tmpl_bytes = _load_template_bytes(typ)
        if not tmpl_bytes:
            _write_meta(job_id, status="error",
                        phase=f"Template '{typ}' nicht ladbar",
                        error=f"Kein lokales File und URL nicht erreichbar fuer typ={typ}")
            return

        # Derived field fuer Rechtlich: ENTWICKLER_ADRESSE_INVERS
        if typ == "rechtlich" and not expose.get("entwickler_adresse_invers"):
            plz_ort = (expose.get("entwickler_plz_ort") or "").strip()
            strasse = (expose.get("entwickler_strasse") or "").strip()
            if plz_ort and strasse:
                expose["entwickler_adresse_invers"] = f"{plz_ort}, {strasse}"

        pptx_bytes = appmod.fill_pptx(tmpl_bytes, expose, customer_images=customer_images)
        # Nach fill_pptx: tmpl + customer_images sind nicht mehr gebraucht.
        del tmpl_bytes, customer_images, state, cust_files
        gc.collect()

        _write_meta(job_id, phase="PDF wird konvertiert …")
        pdf_bytes = None
        pdf_error = None
        if appmod._can_convert_to_pdf():
            try:
                projekt_safe = projekt_name.replace(" ", "_")
                pdf_bytes = appmod.convert_to_pdf(pptx_bytes, f"{projekt_safe}.pptx")
                # Quellen-URLs als clickable Hyperlinks anreichern
                pdf_bytes = appmod._add_hyperlinks_to_pdf(pdf_bytes)
                # Rechtlich: Teilungserklaerung (oder Platzhalter) einfuegen.
                # Bei fehlender Teilung: ENTWURF-Wasserzeichen auf jeder Seite.
                if typ == "rechtlich":
                    teil_bytes = appmod._resolve_teilung_bytes(job_id, projekt_name)
                    is_entwurf = teil_bytes is None
                    if teil_bytes:
                        print(f"[v2] Teilungserklaerung gefunden ({len(teil_bytes)//1024} KB)")
                    else:
                        print(f"[v2] Keine Teilungserklaerung — Entwurf-Modus mit Wasserzeichen")
                    pdf_bytes = appmod._merge_rechtlich_pdf(pdf_bytes, teil_bytes)
                    if is_entwurf:
                        pdf_bytes = appmod._add_entwurf_watermark(pdf_bytes)
                    # Meta-Flag fuer Download-Filename setzen
                    _write_meta(job_id, rechtlich_is_entwurf=is_entwurf)
            except Exception as e:
                import traceback as _tb
                pdf_error = str(e)
                print(f"[v2] PDF-Konvertierung Fehler: {e}\n{_tb.format_exc()[:600]}")
        else:
            pdf_error = "Keine PDF-Konvertierung verfügbar (CloudConvert/LibreOffice fehlen)"

        # Output-Pfad (PDF wenn möglich, sonst PPTX)
        slide_render_error = None
        if pdf_bytes:
            out_path = _output_path(job_id, typ, "pdf")
            with open(out_path, "wb") as f:
                f.write(pdf_bytes)
            # PPTX-Bytes können nun freigegeben werden — PDF ist persistent.
            del pptx_bytes
            gc.collect()

            _write_meta(job_id, phase=f"Slide-Vorschau wird erstellt … ({typ})")
            slides_dir = _v1_slides_dir(job_id, typ)
            # Atomic swap: erst in tmp-Dir rendern, bei Erfolg alten Dir-Inhalt swappen.
            # Vorteil: bei Render-Fail (OOM/Exception) bleiben alte Slides erhalten.
            tmp_slides_dir = slides_dir + ".tmp"
            try:
                # Alten tmp-Dir aufräumen falls Reste vom letzten Crash
                if os.path.isdir(tmp_slides_dir):
                    import shutil as _sh
                    _sh.rmtree(tmp_slides_dir, ignore_errors=True)
                os.makedirs(tmp_slides_dir, exist_ok=True)
                # dpi=110 statt 150: pdftoppm spawnt subprocess, parent muss Memory haben
                # Bei 512 MB Render-Plan ist 110 dpi der robuste Sweet-Spot.
                appmod.render_pdf_to_jpgs(pdf_bytes, tmp_slides_dir, dpi=110)
                # Erfolg → alte Slides löschen, neue an die Stelle verschieben
                if os.path.isdir(slides_dir):
                    for fname in os.listdir(slides_dir):
                        if fname.endswith(".jpg"):
                            try: os.remove(os.path.join(slides_dir, fname))
                            except OSError: pass
                else:
                    os.makedirs(slides_dir, exist_ok=True)
                for fname in os.listdir(tmp_slides_dir):
                    os.replace(os.path.join(tmp_slides_dir, fname),
                               os.path.join(slides_dir, fname))
                import shutil as _sh
                _sh.rmtree(tmp_slides_dir, ignore_errors=True)
                # Slide-Render durch → pdf_bytes nicht mehr nötig
                del pdf_bytes
                gc.collect()
            except Exception as e:
                import traceback as _tb
                slide_render_error = str(e)
                print(f"[v2] Slide-JPG-Render Fehler: {e}\n{_tb.format_exc()[:600]}")
                # Tmp-Dir aufräumen
                import shutil as _sh
                _sh.rmtree(tmp_slides_dir, ignore_errors=True)
        else:
            out_path = _output_path(job_id, typ, "pptx")
            with open(out_path, "wb") as f:
                f.write(pptx_bytes)
            del pptx_bytes
            gc.collect()

        # Status: error wenn render gescheitert ist (sonst sieht User "✓" trotz 0 Slides)
        if slide_render_error or (pdf_error and not pdf_bytes):
            err_msg = slide_render_error or pdf_error
            _write_meta(job_id,
                        status="error",
                        phase=f"Vorschau-Render fehlgeschlagen ({typ})",
                        pdf_path=out_path,
                        name=projekt_name,
                        last_render_typ=typ,
                        error=err_msg)
            print(f"[v2] ✗ Render mit Fehler abgeschlossen ({typ}): {err_msg[:200]}")
        else:
            _write_meta(job_id,
                        status="preview",
                        phase=f"Vorschau aktualisiert ({typ})",
                        pdf_path=out_path,
                        name=projekt_name,
                        last_render_typ=typ,
                        error=None)
            print(f"[v2] ✓ Render fertig für Job {job_id} (typ={typ})")
    except Exception as e:
        import traceback as tb
        err = f"{e}\n{tb.format_exc()}"
        print(f"[v2] ✗ Render Fehler: {err[:500]}")
        _write_meta(job_id, status="error", phase="Fehler", error=str(e))